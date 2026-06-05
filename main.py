import os
import json
import uuid
import random
from datetime import datetime, timedelta, timezone
from dateutil.relativedelta import relativedelta
from typing import Optional, List
from urllib.parse import unquote

from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Request, Depends, status, WebSocket, WebSocketDisconnect
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.trustedhost import TrustedHostMiddleware
from starlette.middleware.base import BaseHTTPMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from passlib.context import CryptContext
from jose import JWTError, jwt

# استيراد مكتبة السحاب (Supabase)
from supabase import create_client, Client

# Rate Limiting — منطق يدوي في الذاكرة (لا يحتاج مكتبة خارجية)
import time as _time
from collections import defaultdict as _defaultdict

_rate_store: dict = _defaultdict(list)  # {ip: [timestamps]}

def _is_rate_limited(ip: str, max_calls: int, window_seconds: int, key_prefix: str = "") -> bool:
    """يتحقق إذا تجاوز الـ IP الحد المسموح — يُرجع True إذا محظور
    
    key_prefix: لفصل الحدود حسب نوع الطلب (login vs check_answer vs parent_search)
    مثال: _is_rate_limited(ip, 10, 60, "login") لا يتداخل مع check_answer
    """
    key = f"{key_prefix}:{ip}" if key_prefix else ip
    now = _time.time()
    calls = _rate_store[key]
    # احتفظ فقط بالطلبات داخل النافذة الزمنية
    _rate_store[key] = [t for t in calls if now - t < window_seconds]
    if len(_rate_store[key]) >= max_calls:
        return True
    _rate_store[key].append(now)
    return False


# ══════════════════════════════════════════════════
# 🛡️ حماية رفع الملفات
# ══════════════════════════════════════════════════
ALLOWED_FILE_EXTENSIONS = {
    # مستندات
    ".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".txt",
    # صور
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg",
    # فيديو/صوت
    ".mp4", ".webm", ".mp3", ".wav", ".m4a",
}

MAX_FILE_SIZE_MB = 30  # الحد الأقصى لحجم الملف بالميجابايت
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024


def _validate_upload(file_content: bytes, filename: str, allowed_exts: set = None) -> str:
    """
    يتحقق من ملف مرفوع — يرفع HTTPException إذا فشل
    يُرجع الـ extension الآمن (مع نقطة)
    
    التحققات:
    1. الملف ليس فارغاً
    2. الحجم ≤ MAX_FILE_SIZE_BYTES
    3. الـ extension في القائمة المسموحة
    4. لا يحتوي path traversal (../, ./, /)
    """
    if not file_content:
        raise HTTPException(status_code=400, detail="الملف فارغ")
    
    if len(file_content) > MAX_FILE_SIZE_BYTES:
        size_mb = len(file_content) / (1024 * 1024)
        raise HTTPException(
            status_code=413,
            detail=f"الملف كبير جداً ({size_mb:.1f} MB). الحد الأقصى {MAX_FILE_SIZE_MB} MB"
        )
    
    if not filename:
        raise HTTPException(status_code=400, detail="اسم الملف مفقود")
    
    # منع path traversal
    if "/" in filename or "\\" in filename or ".." in filename:
        raise HTTPException(status_code=400, detail="اسم الملف يحتوي على رموز ممنوعة")
    
    # استخراج وفحص الـ extension
    ext = os.path.splitext(filename)[1].lower().strip()
    if not ext:
        raise HTTPException(status_code=400, detail="الملف بدون امتداد")
    
    allowed = allowed_exts if allowed_exts else ALLOWED_FILE_EXTENSIONS
    if ext not in allowed:
        raise HTTPException(
            status_code=415,
            detail=f"امتداد '{ext}' غير مسموح. المسموح: {', '.join(sorted(allowed))}"
        )
    
    return ext


def _is_safe_url(url: str) -> bool:
    """يتحقق من سلامة URL خارجي — يمنع javascript:, data:, file:, إلخ"""
    if not url:
        return False
    url_lower = url.strip().lower()
    # نقبل فقط http/https
    if not (url_lower.startswith("https://") or url_lower.startswith("http://")):
        return False
    # نمنع localhost / private IPs (يتطلب extra parsing لكن الأساس)
    blocked_hosts = ["localhost", "127.0.0.1", "0.0.0.0", "169.254.", "::1"]
    for host in blocked_hosts:
        if host in url_lower:
            return False
    return True


# ==========================================
# --- 1. الإعدادات الأمنية والاتصال ---
# ==========================================
# ══════════════════════════════════════════════════
# الأسرار تُقرأ من متغيرات البيئة (.env) — لا تكتب
# أي قيمة حرفية هنا أبداً
# ══════════════════════════════════════════════════
# python-dotenv اختياري — يمكن ضبط المتغيرات مباشرة في البيئة
try:
    import importlib.util as _ilu
    if _ilu.find_spec("dotenv") is not None:
        from dotenv import load_dotenv  # type: ignore[import]
        load_dotenv()
except Exception:
    pass

SECRET_KEY = os.getenv("JWT_SECRET_KEY")
if not SECRET_KEY:
    import secrets
    SECRET_KEY = secrets.token_hex(32)
    print("⚠️ JWT_SECRET_KEY not set — using random key (tokens reset on restart)")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 480

# بيانات الاتصال بـ Supabase — من متغيرات البيئة (إلزامي)
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
if not SUPABASE_URL or not SUPABASE_KEY:
    raise RuntimeError(
        "❌ SUPABASE_URL و SUPABASE_KEY يجب تعيينهما في متغيرات البيئة. "
        "اذهب إلى Render Dashboard → Environment → أضفهما."
    )

# كلمة مرور الأدمن — من متغيرات البيئة (إلزامي)
ADMIN_PASSWORD = os.getenv("ADMIN_PASSWORD")
if not ADMIN_PASSWORD:
    raise RuntimeError(
        "❌ ADMIN_PASSWORD يجب تعيينه في متغيرات البيئة. "
        "اذهب إلى Render Dashboard → Environment → أضفها."
    )
if len(ADMIN_PASSWORD) < 8:
    raise RuntimeError("❌ ADMIN_PASSWORD يجب أن تكون 8 أحرف على الأقل.")

supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
pwd_context = CryptContext(schemes=["pbkdf2_sha256"], deprecated="auto")

def hash_password(password: str): return pwd_context.hash(password)
def verify_password(plain, hashed): return pwd_context.verify(plain, hashed)

def create_access_token(data: dict):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

# ==========================================
# --- 2. تهيئة التطبيق ومعالجة الأخطاء ---
# ==========================================
app = FastAPI(title="Math Empire API")

@app.middleware("http")
async def security_headers(request: Request, call_next):
    """إضافة security headers لكل response"""
    response = await call_next(request)
    response.headers["X-Content-Type-Options"]    = "nosniff"
    response.headers["X-Frame-Options"]           = "SAMEORIGIN"
    response.headers["Referrer-Policy"]           = "strict-origin-when-cross-origin"
    response.headers["Permissions-Policy"]        = "geolocation=(), camera=(), microphone=(self)"

    # ═══ Content Security Policy ═══
    # نسمح بالـ CDNs المستخدمة فعلاً + Supabase + fonts.googleapis.com + unsafe-inline/eval
    # (unsafe-inline ضروري بسبب كثرة inline scripts/styles في الملفات الحالية)
    csp_directives = [
        "default-src 'self'",
        "script-src 'self' 'unsafe-inline' 'unsafe-eval' "
            "https://cdnjs.cloudflare.com "
            "https://cdn.jsdelivr.net "
            "https://unpkg.com "
            "https://generativelanguage.googleapis.com "
            "https://api.x.ai",
        # 🆕 السماح بـ Web Workers (لـ pdf.js والمكتبات الأخرى)
        "worker-src 'self' blob: "
            "https://cdnjs.cloudflare.com "
            "https://cdn.jsdelivr.net "
            "https://unpkg.com",
        "child-src 'self' blob:",
        "style-src 'self' 'unsafe-inline' "
            "https://fonts.googleapis.com "
            "https://cdnjs.cloudflare.com",
        "font-src 'self' data: "
            "https://fonts.gstatic.com "
            "https://cdnjs.cloudflare.com",
        "img-src 'self' data: blob: https:",
        "media-src 'self' data: blob: https:",
        "connect-src 'self' "
            "https://*.supabase.co "
            "https://generativelanguage.googleapis.com "
            "https://api.x.ai "
            "https://fonts.googleapis.com "
            "https://fonts.gstatic.com "
            "https://cdnjs.cloudflare.com "
            "https://cdn.jsdelivr.net "
            "https://unpkg.com "
            "https://www.gstatic.com "
            "wss: ws:",
        "frame-ancestors 'self'",
        "base-uri 'self'",
        "form-action 'self'",
    ]
    response.headers["Content-Security-Policy"] = "; ".join(csp_directives)

    # HSTS — فقط للـ HTTPS (Render يستخدم HTTPS)
    if request.url.scheme == "https":
        response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"

    return response

# Rate Limiter — يدوي في الذاكرة، لا يحتاج تسجيل

@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    print(f"Error occurred: {str(exc)}")
    return JSONResponse(
        status_code=500,
        content={"status": "error", "message": "عطل في الديوان الملكي"},
    )

# ضمان ترميز UTF-8 لجميع الاستجابات
from fastapi.responses import Response
import json as _json

# ═══ Web Push Notifications (VAPID) ═══
VAPID_PUBLIC_KEY  = os.environ.get("VAPID_PUBLIC_KEY", "")
VAPID_PRIVATE_KEY = os.environ.get("VAPID_PRIVATE_KEY", "")
VAPID_CLAIMS_SUB  = os.environ.get("VAPID_CLAIMS_SUB", "mailto:rashdy.sayed@example.com")

# pywebpush قد لا يكون مثبتاً — نستورده بحذر
try:
    from pywebpush import webpush, WebPushException  # type: ignore
    PUSH_ENABLED = bool(VAPID_PUBLIC_KEY and VAPID_PRIVATE_KEY)
except ImportError:
    PUSH_ENABLED = False
    print("⚠️ pywebpush غير مثبت — سيتم تخطي push notifications")


def _send_push_to_endpoint(subscription_info: dict, payload: dict) -> bool:
    """يُرسل push لاشتراك واحد. يُرجع True عند النجاح."""
    if not PUSH_ENABLED:
        return False
    try:
        webpush(
            subscription_info=subscription_info,
            data=json.dumps(payload, ensure_ascii=False),
            vapid_private_key=VAPID_PRIVATE_KEY,
            vapid_claims={"sub": VAPID_CLAIMS_SUB}
        )
        return True
    except WebPushException as e:
        # لو 410/404 → الاشتراك منتهي، احذفه
        if e.response and e.response.status_code in (410, 404):
            try:
                supabase.table("push_subscriptions").delete().eq(
                    "endpoint", subscription_info.get("endpoint", "")
                ).execute()
            except Exception:
                pass
        return False
    except Exception:
        return False


def _push_to_student(student_id: int, title: str, body: str,
                     url: str = "/student", tag: str = "general",
                     icon: str = "/static/icon-192.png", require_interaction: bool = False) -> int:
    """يُرسل push لكل اشتراكات طالب معيّن. يُرجع عدد الإرسالات الناجحة."""
    if not PUSH_ENABLED:
        return 0
    try:
        res = supabase.table("push_subscriptions").select(
            "endpoint, p256dh, auth"
        ).eq("student_id", student_id).execute()
        subs = res.data or []
    except Exception:
        return 0
    
    sent = 0
    payload = {
        "title": title,
        "body": body,
        "url": url,
        "tag": tag,
        "icon": icon,
        "requireInteraction": require_interaction,
    }
    for s in subs:
        info = {
            "endpoint": s["endpoint"],
            "keys": {"p256dh": s["p256dh"], "auth": s["auth"]}
        }
        if _send_push_to_endpoint(info, payload):
            sent += 1
    return sent


def _push_to_grade(grade: str, title: str, body: str, url: str = "/student", tag: str = "general") -> int:
    """يُرسل push لكل طلاب صف معيّن"""
    if not PUSH_ENABLED:
        return 0
    try:
        # اجلب IDs الطلاب في الصف
        st_res = supabase.table("students").select("id").eq("grade", grade).execute()
        student_ids = [s["id"] for s in (st_res.data or [])]
        if not student_ids:
            return 0
        # اجلب اشتراكاتهم
        sub_res = supabase.table("push_subscriptions").select(
            "endpoint, p256dh, auth, student_id"
        ).in_("student_id", student_ids).execute()
        subs = sub_res.data or []
    except Exception:
        return 0
    
    sent = 0
    payload = {"title": title, "body": body, "url": url, "tag": tag, "icon": "/static/icon-192.png"}
    for s in subs:
        info = {"endpoint": s["endpoint"], "keys": {"p256dh": s["p256dh"], "auth": s["auth"]}}
        if _send_push_to_endpoint(info, payload):
            sent += 1
    return sent



class UTF8JSONResponse(JSONResponse):
    def render(self, content) -> bytes:
        return _json.dumps(content, ensure_ascii=False, allow_nan=False, indent=None, separators=(",", ":")).encode("utf-8")

app.router.default_response_class = UTF8JSONResponse

# النطاقات المسموح بها — أضف نطاق إنتاجك هنا
_ALLOWED_ORIGINS = [o.strip() for o in os.getenv(
    "ALLOWED_ORIGINS",
    "http://localhost:8001,http://127.0.0.1:8001"
).split(",") if o.strip()]


# ════════════════════════════════════════════════════════════
# 🛡️ SECURITY MIDDLEWARE — حماية متعددة الطبقات
# ════════════════════════════════════════════════════════════
class SecurityHeadersMiddleware(BaseHTTPMiddleware):
    """يُضيف security headers لكل response"""
    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        # HSTS: يُجبر HTTPS لمدة سنة كاملة (production فقط)
        if os.getenv("ENV", "production").lower() == "production":
            response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains; preload"
        # منع التضمين في iframes (clickjacking)
        response.headers["X-Frame-Options"] = "SAMEORIGIN"
        # منع MIME sniffing
        response.headers["X-Content-Type-Options"] = "nosniff"
        # سياسة المُحيل (لا نُسرّب URLs لمواقع خارجية)
        response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
        # منع تتبع المتصفحات
        response.headers["Permissions-Policy"] = "geolocation=(), microphone=(self), camera=()"
        # XSS Protection للمتصفحات القديمة
        response.headers["X-XSS-Protection"] = "1; mode=block"
        return response


class HTTPSRedirectInProduction(BaseHTTPMiddleware):
    """يُحوّل HTTP → HTTPS في الإنتاج فقط"""
    async def dispatch(self, request: Request, call_next):
        if os.getenv("ENV", "production").lower() == "production":
            # نتحقق من X-Forwarded-Proto (لأن Render يُمرّر HTTPS عبر proxy)
            proto = request.headers.get("x-forwarded-proto", "").lower()
            if proto == "http":
                # نُعيد التوجيه لـ HTTPS
                url = str(request.url).replace("http://", "https://", 1)
                from fastapi.responses import RedirectResponse
                return RedirectResponse(url=url, status_code=301)
        return await call_next(request)

# نُضيف الـ middlewares (الترتيب مهم: من الخارج للداخل)
app.add_middleware(SecurityHeadersMiddleware)
app.add_middleware(HTTPSRedirectInProduction)

# Trusted Hosts — يمنع Host Header Injection
_TRUSTED_HOSTS = [
    "math-empire22.onrender.com",
    "*.onrender.com",
    "localhost",
    "127.0.0.1",
]
# نُضيف custom domains من env إن وُجدت
_extra_hosts = os.getenv("TRUSTED_HOSTS", "").split(",")
_TRUSTED_HOSTS.extend([h.strip() for h in _extra_hosts if h.strip()])

app.add_middleware(
    TrustedHostMiddleware,
    allowed_hosts=_TRUSTED_HOSTS
)


app.add_middleware(
    CORSMiddleware,
    allow_origins=_ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["Authorization", "Content-Type"],
)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
for folder in ["static", "templates"]:
    path = os.path.join(BASE_DIR, folder)
    if not os.path.exists(path): os.makedirs(path)

templates = Jinja2Templates(directory=os.path.join(BASE_DIR, "templates"))
app.mount("/static", StaticFiles(directory=os.path.join(BASE_DIR, "static")), name="static")


# ════════════════════════════════════════════════════════════
# 🔐 PROGRESSIVE LOCKOUT — قفل متدرج للحسابات
# ════════════════════════════════════════════════════════════
# يحفظ في الذاكرة (للسرعة) - يُمكن نقله لـ Redis لاحقاً
_login_attempts = {}  # {key: {"count": int, "first": ts, "locked_until": ts}}

def _get_lockout_duration(attempts: int) -> int:
    """يُرجع مدة القفل بالثواني حسب عدد المحاولات الفاشلة"""
    if attempts < 3:   return 0          # لا قفل
    if attempts < 5:   return 60         # دقيقة
    if attempts < 10:  return 300        # 5 دقائق
    if attempts < 15:  return 3600       # ساعة
    if attempts < 20:  return 21600      # 6 ساعات
    return 86400                          # 24 ساعة

def check_progressive_lockout(key: str) -> tuple:
    """يتحقق من حالة القفل. يُرجع (مقفول؟, ثواني متبقية)"""
    import time as _time
    now = _time.time()
    rec = _login_attempts.get(key)
    if not rec:
        return False, 0
    # نُزيل السجلات القديمة (أكثر من 24 ساعة)
    if now - rec.get("first", 0) > 86400:
        _login_attempts.pop(key, None)
        return False, 0
    # هل ما زال مقفولاً؟
    locked_until = rec.get("locked_until", 0)
    if locked_until > now:
        return True, int(locked_until - now)
    return False, 0

def record_login_failure(key: str):
    """يُسجّل محاولة فاشلة + يحسب القفل"""
    import time as _time
    now = _time.time()
    rec = _login_attempts.get(key, {"count": 0, "first": now})
    rec["count"] += 1
    rec["last"] = now
    duration = _get_lockout_duration(rec["count"])
    if duration > 0:
        rec["locked_until"] = now + duration
    _login_attempts[key] = rec
    return rec["count"], duration

def record_login_success(key: str):
    """يُزيل السجل عند نجاح الدخول"""
    _login_attempts.pop(key, None)

def cleanup_old_attempts():
    """تنظيف دوري - يحذف السجلات الأقدم من 24 ساعة"""
    import time as _time
    now = _time.time()
    to_remove = [k for k, v in _login_attempts.items() if now - v.get("first", 0) > 86400]
    for k in to_remove:
        _login_attempts.pop(k, None)


# ════════════════════════════════════════════════════════════
# 📝 AUDIT LOG — سجل التدقيق للعمليات الأمنية
# ════════════════════════════════════════════════════════════
def security_log(event: str, ip: str, details: dict = None):
    """يُسجّل عمليات أمنية حساسة"""
    try:
        from datetime import datetime, timezone
        log_entry = {
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "event": event,
            "ip": ip,
            "details": details or {}
        }
        # طباعة في console (يلتقطها Render Logs)
        print(f"[SECURITY] {log_entry}", flush=True)
        # محاولة الحفظ في DB (لو الجدول موجود)
        try:
            supabase.table("security_log").insert(log_entry).execute()
        except Exception:
            pass  # الجدول غير موجود - نكتفي بـ console
    except Exception as e:
        print(f"[SECURITY LOG ERROR] {e}", flush=True)



async def get_current_admin(request: Request):
    auth_header = request.headers.get("Authorization")
    if not auth_header or not auth_header.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="يرجى تسجيل دخول المعلم")
    token = auth_header.split(" ")[1]
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("sub") != "admin": raise Exception()
        return payload.get("sub")
    except:
        raise HTTPException(status_code=401, detail="جلسة العمل غير صالحة")


# ════════════════════════════════════════════════════════════
# 🛡️ SUPERVISOR AUTH HELPER
# ════════════════════════════════════════════════════════════
async def get_current_supervisor(request: Request):
    """يتحقق من JWT المشرف ويُرجع supervisor record كاملاً"""
    auth_header = request.headers.get("Authorization")
    if not auth_header or not auth_header.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="يرجى تسجيل دخول المشرف")
    token = auth_header.split(" ")[1]
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("role") != "supervisor":
            raise HTTPException(status_code=401, detail="ليس لديك صلاحية مشرف")
        sup_id = payload.get("sub")
        if not sup_id:
            raise HTTPException(status_code=401, detail="توكن غير صالح")
        # نجلب المشرف من DB
        res = supabase.table("supervisors").select("*").eq("id", sup_id).eq("is_active", True).execute()
        if not res.data:
            raise HTTPException(status_code=401, detail="الحساب غير موجود أو معطّل")
        return res.data[0]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=401, detail="جلسة العمل غير صالحة")


def _get_supervisor_student_ids(supervisor_id: int) -> list:
    """يُرجع قائمة IDs الطلاب التابعين لمشرف معيّن"""
    try:
        res = supabase.table("supervisor_students").select("student_id").eq("supervisor_id", supervisor_id).execute()
        return [r["student_id"] for r in (res.data or [])]
    except Exception:
        return []

# ==========================================
# --- 3. مسارات العرض (HTML) والملفات التقنية ---
# ==========================================
@app.get("/")
async def read_root(request: Request): return templates.TemplateResponse(request=request, name="index.html")

@app.get("/admin")
async def read_admin(request: Request): return templates.TemplateResponse(request=request, name="admin.html")

@app.get("/student")
async def read_student(request: Request): return templates.TemplateResponse(request=request, name="student.html")

@app.get("/parent")
async def read_parent(request: Request): return templates.TemplateResponse(request=request, name="parent.html")

@app.get("/teachers")
async def read_teachers(request: Request): return templates.TemplateResponse(request=request, name="teachers.html")

@app.get("/supervisor")
async def read_supervisor(request: Request): return templates.TemplateResponse(request=request, name="supervisor.html")

@app.get("/manifest.json")
async def get_manifest(): return FileResponse(os.path.join(BASE_DIR, "manifest.json"))

# ════════════════════════════════════════════════════════════
# 🔍 SEO Routes — robots.txt + sitemap.xml لإدراج جوجل
# ════════════════════════════════════════════════════════════
@app.get("/robots.txt", include_in_schema=False)
async def get_robots():
    """🤖 ملف توجيه محركات البحث — ديناميكي"""
    from fastapi.responses import PlainTextResponse
    content = (
        "User-agent: *\n"
        "Allow: /\n"
        "Allow: /student\n"
        "Allow: /parent\n"
        "Allow: /teachers\n"
        "Allow: /static/\n"
        "\n"
        "Disallow: /admin\n"
        "Disallow: /api/\n"
        "Disallow: /static/avatars/\n"
        "\n"
        "User-agent: Googlebot\n"
        "Allow: /\n"
        "Crawl-delay: 1\n"
        "\n"
        "User-agent: Bingbot\n"
        "Allow: /\n"
        "Crawl-delay: 1\n"
        "\n"
        "User-agent: GPTBot\n"
        "Disallow: /\n"
        "\n"
        "User-agent: SemrushBot\n"
        "Disallow: /\n"
        "\n"
        "Sitemap: https://math-empire22.onrender.com/sitemap.xml\n"
    )
    return PlainTextResponse(
        content=content,
        headers={
            "Content-Type": "text/plain; charset=utf-8",
            "Cache-Control": "public, max-age=3600"
        }
    )


@app.get("/sitemap.xml", include_in_schema=False)
async def get_sitemap():
    """خريطة الموقع — ديناميكية مع headers صحيحة"""
    from fastapi.responses import Response
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    base_url = "https://math-empire22.onrender.com"
    
    xml_content = (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<urlset xmlns="http://www.sitemaps.org/schemas/sitemap/0.9" '
        'xmlns:image="http://www.google.com/schemas/sitemap-image/0.9">\n'
        f'  <url>\n'
        f'    <loc>{base_url}/</loc>\n'
        f'    <lastmod>{today}</lastmod>\n'
        f'    <changefreq>weekly</changefreq>\n'
        f'    <priority>1.0</priority>\n'
        f'    <image:image>\n'
        f'      <image:loc>{base_url}/static/logo.jpg</image:loc>\n'
        f'      <image:title>إمبراطورية الرياضيات الملكية</image:title>\n'
        f'    </image:image>\n'
        f'  </url>\n'
        f'  <url>\n'
        f'    <loc>{base_url}/student</loc>\n'
        f'    <lastmod>{today}</lastmod>\n'
        f'    <changefreq>daily</changefreq>\n'
        f'    <priority>0.9</priority>\n'
        f'  </url>\n'
        f'  <url>\n'
        f'    <loc>{base_url}/parent</loc>\n'
        f'    <lastmod>{today}</lastmod>\n'
        f'    <changefreq>weekly</changefreq>\n'
        f'    <priority>0.8</priority>\n'
        f'  </url>\n'
        f'  <url>\n'
        f'    <loc>{base_url}/teachers</loc>\n'
        f'    <lastmod>{today}</lastmod>\n'
        f'    <changefreq>weekly</changefreq>\n'
        f'    <priority>0.8</priority>\n'
        f'  </url>\n'
        '</urlset>'
    )
    return Response(
        content=xml_content,
        media_type="application/xml",
        headers={
            "Content-Type": "application/xml; charset=utf-8",
            "Cache-Control": "public, max-age=3600"
        }
    )


@app.get("/favicon.ico", include_in_schema=False)
async def get_favicon():
    """أيقونة المتصفح"""
    ico = os.path.join(BASE_DIR, "static", "favicon.ico")
    if os.path.exists(ico):
        return FileResponse(ico)
    logo = os.path.join(BASE_DIR, "static", "logo.jpg")
    if os.path.exists(logo):
        return FileResponse(logo)
    raise HTTPException(status_code=404)

@app.get("/sw.js")
async def get_sw(): return FileResponse(os.path.join(BASE_DIR, "static", "sw.js"))

# ==========================================
# --- 4. نظام الدخول (إمبراطور / طالب / معلم) ---
# ==========================================
@app.post("/api/admin/login")
async def admin_login(request: Request, username: str = Form(...), password: str = Form(...)):
    """🔐 تسجيل دخول الأدمن مع حماية شاملة"""
    client_ip = request.client.host if request.client else "unknown"
    lockout_key = f"admin:{client_ip}"
    
    # 1️⃣ تحقق من القفل المتدرج
    is_locked, seconds_left = check_progressive_lockout(lockout_key)
    if is_locked:
        minutes = max(1, seconds_left // 60)
        security_log("admin_login_blocked", client_ip, {"reason": "lockout", "seconds_left": seconds_left})
        raise HTTPException(
            status_code=429,
            detail=f"🔒 الحساب مقفول مؤقتاً — حاول بعد {minutes} دقيقة"
        )
    
    # 2️⃣ rate limit بسيط للهجمات السريعة
    if _is_rate_limited(client_ip, max_calls=10, window_seconds=60):
        security_log("admin_login_rate_limited", client_ip)
        raise HTTPException(status_code=429, detail="⏳ تجاوزت عدد المحاولات — انتظر دقيقة")
    
    # 3️⃣ التحقق من البيانات
    expected_user = os.getenv("ADMIN_USERNAME", "admin")
    if username == expected_user and password == ADMIN_PASSWORD:
        # نجح الدخول
        record_login_success(lockout_key)
        token = create_access_token(data={"sub": username})
        security_log("admin_login_success", client_ip, {"username": username})
        return {"access_token": token, "token_type": "bearer"}
    
    # 4️⃣ فشل الدخول — نُسجّل المحاولة
    count, duration = record_login_failure(lockout_key)
    security_log("admin_login_failed", client_ip, {
        "username": username,
        "attempts": count,
        "lockout_seconds": duration
    })
    
    if duration > 0:
        minutes = max(1, duration // 60)
        raise HTTPException(
            status_code=401,
            detail=f"❌ بيانات خاطئة — تم قفل المحاولات لمدة {minutes} دقيقة"
        )
    
    remaining = max(0, 3 - count)
    if remaining > 0:
        raise HTTPException(status_code=401, detail=f"❌ بيانات خاطئة — تبقّى {remaining} محاولات قبل القفل")
    raise HTTPException(status_code=401, detail="❌ بيانات دخول المعلم خاطئة")

@app.post("/api/teacher/register")
async def register_teacher(full_name: str=Form(...), username: str=Form(...), password: str=Form(...)):
    existing = supabase.table("teachers").select("username").eq("username", username).execute()
    if existing.data: raise HTTPException(status_code=400, detail="المستخدم موجود مسبقاً")
    supabase.table("teachers").insert({
        "full_name": full_name, 
        "username": username, 
        "password": hash_password(password)
    }).execute()
    return {"status": "success"}

@app.post("/api/teacher/login")
async def teacher_login(request: Request, username: str = Form(...), password: str = Form(...)):
    """🔐 تسجيل دخول معلم مع حماية متدرجة"""
    ip = request.client.host if request.client else "unknown"
    lockout_key = f"teacher:{username}:{ip}"
    
    # القفل المتدرج
    is_locked, seconds_left = check_progressive_lockout(lockout_key)
    if is_locked:
        minutes = max(1, seconds_left // 60)
        security_log("teacher_login_blocked", ip, {"username": username})
        raise HTTPException(status_code=429, detail=f"🔒 الحساب مقفول مؤقتاً — حاول بعد {minutes} دقيقة")
    
    if _is_rate_limited(f"teacher_login:{ip}", max_calls=10, window_seconds=60):
        raise HTTPException(status_code=429, detail="محاولات كثيرة — انتظر دقيقة")

    res = supabase.table("teachers").select("*").eq("username", username).execute()
    if res.data and verify_password(password, res.data[0]['password']):
        # ═══ تحديث آخر دخول (last_login) — قد لا يوجد العمود في قواعد قديمة ═══
        try:
            supabase.table("teachers").update({
                "last_login": datetime.now(timezone.utc).isoformat()
            }).eq("id", res.data[0]['id']).execute()
        except Exception:
            pass  # last_login column may not exist yet
        # ═══ نسخة نظيفة من بيانات المعلم بدون كلمة المرور ═══
        user_clean = {k: v for k, v in res.data[0].items() if k != 'password'}
        return {"status": "success", "user": user_clean}

    raise HTTPException(status_code=401, detail="بيانات الدخول خاطئة")


# ════════════════════════════════════════════════════════════════════════════
# ☁️ TEACHER CLOUD DATA — حفظ/جلب بيانات المعلم في السحابة (تتبعه عبر الأجهزة)
# ════════════════════════════════════════════════════════════════════════════
# المفاتيح المسموحة (whitelist للأمان)
_ALLOWED_TD_KEYS = {
    "exam_library", "teacher_stats", "recent_activity", "favorites",
    "teacher_theme", "teacher_theme_v2", "saved_shapes", "lesson_drafts",
    "gradebook_data", "plan_drafts", "worksheet_drafts",
    # ── سجلات وبيانات إضافية ──
    "gradebook_v1", "weak_students_v1", "gifted_students_v1",
    "math_empire_games_v1", "exam_draft", "default_teacher_name",
    "default_school_name", "teacherProfilePhoto", "sidebar_compact",
    "theme_mode", "math_empire_games_sound"
}

@app.get("/api/teacher/data/get")
async def teacher_data_get(teacher_id: int, data_key: str = ""):
    """☁️ جلب بيانات المعلم من السحابة. بدون data_key يجلب كل البيانات."""
    if not teacher_id:
        raise HTTPException(status_code=400, detail="teacher_id مطلوب")
    try:
        q = supabase.table("teacher_data").select("data_key, data_value, updated_at").eq("teacher_id", teacher_id)
        if data_key:
            q = q.eq("data_key", data_key)
        res = q.execute()
        out = {}
        for row in (res.data or []):
            out[row["data_key"]] = row["data_value"]
        return {"status": "ok", "data": out}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل الجلب: {str(e)[:160]}")


@app.post("/api/teacher/data/save")
async def teacher_data_save(
    teacher_id: int = Form(...),
    data_key: str = Form(...),
    data_value: str = Form(...),
):
    """☁️ حفظ بيانات المعلم في السحابة (upsert)."""
    if not teacher_id:
        raise HTTPException(status_code=400, detail="teacher_id مطلوب")
    if data_key not in _ALLOWED_TD_KEYS:
        raise HTTPException(status_code=400, detail=f"مفتاح غير مسموح: {data_key}")
    if len(data_value) > 4 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="البيانات أكبر من 4 ميجابايت")
    import json as _jl
    try:
        parsed = _jl.loads(data_value)  # تحقّق أنها JSON صحيح
    except Exception:
        # بعض القيم نصوص عادية (اسم المعلم، الصورة، الثيم) — نخزّنها كنص داخل JSON
        parsed = data_value
    try:
        # تأكد أن المعلم موجود
        t = supabase.table("teachers").select("id").eq("id", teacher_id).limit(1).execute()
        if not t.data:
            raise HTTPException(status_code=404, detail="المعلم غير موجود")
        row = {
            "teacher_id": teacher_id,
            "data_key": data_key,
            "data_value": parsed,
            "updated_at": datetime.now(timezone.utc).isoformat()
        }
        supabase.table("teacher_data").upsert(row, on_conflict="teacher_id,data_key").execute()
        return {"status": "ok", "message": "تم الحفظ السحابي"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل الحفظ: {str(e)[:160]}")


def _normalize_phone(phone: str) -> str:
    """تطبيع رقم الهاتف للمقارنة: إزالة كل ما عدا الأرقام"""
    if not phone:
        return ""
    import re as _re
    # احتفظ بالأرقام فقط
    digits = _re.sub(r'\D', '', str(phone))
    # احذف 00 أو + من البداية
    if digits.startswith("00"):
        digits = digits[2:]
    # احذف 968 إن كان رقماً عمانياً (لتوحيد المقارنة)
    return digits


def _phone_matches(p1: str, p2: str) -> bool:
    """مقارنة ذكية بين رقمين — يدعم اختلاف بادئة الدولة"""
    if not p1 or not p2:
        return False
    n1 = _normalize_phone(p1)
    n2 = _normalize_phone(p2)
    if n1 == n2:
        return True
    # تطابق الأرقام الـ 8 الأخيرة (للأرقام المحلية)
    if len(n1) >= 8 and len(n2) >= 8 and n1[-8:] == n2[-8:]:
        return True
    return False


def _generate_unique_parent_code() -> str:
    """توليد كود ولي أمر فريد بصيغة RS-XXXXXX (50 محاولة)"""
    import secrets
    import string
    # حروف بدون 0,O,1,I لتجنّب الالتباس
    chars = "ABCDEFGHJKLMNPQRSTUVWXYZ23456789"
    for _ in range(50):
        code = "RS-" + "".join(secrets.choice(chars) for _ in range(6))
        try:
            existing = supabase.table("students").select("id").eq("parent_code", code).limit(1).execute()
            if not existing.data:
                return code
        except Exception:
            pass
    # fallback نادر — كود أطول
    return "RS-" + "".join(secrets.choice(chars) for _ in range(8))


@app.post("/api/student/register")
async def register_student(
    full_name: str=Form(...), username: str=Form(...),
    password: str=Form(...), grade: str=Form(...),
    parent_code: str=Form(default="")
):
    existing = supabase.table("students").select("username").eq("username", username).execute()
    if existing.data: raise HTTPException(status_code=400, detail="المستخدم موجود مسبقاً")
    
    # 🔧 توليد كود فريد دائماً (سواء أدخل المستخدم أم لا)
    pc = (parent_code or "").strip().upper()
    if pc:
        # إن أدخل كود يدوياً — تأكد من فرادته
        if not pc.startswith("RS-"):
            pc = f"RS-{pc}"
        existing_code = supabase.table("students").select("id").eq("parent_code", pc).limit(1).execute()
        if existing_code.data:
            # الكود مكرر — ولّد كود جديد بدلاً منه
            pc = _generate_unique_parent_code()
    else:
        # لم يُدخل كود — ولّد واحداً
        pc = _generate_unique_parent_code()
    
    res = supabase.table("students").insert({
        "full_name":   full_name,
        "username":    username,
        "password":    hash_password(password),
        "grade":       grade,
        "parent_code": pc,
    }).execute()
    
    return {"status": "success", "parent_code": pc}

def get_current_student(request: Request):
    """التحقق من JWT الطالب"""
    auth = request.headers.get("Authorization", "")
    token = auth.replace("Bearer ", "").strip()
    if not token:
        # fallback: قبول student_id في الـ form (للتوافق مع الكود القديم)
        return None
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("role") != "student":
            return None
        return payload
    except JWTError:
        return None


@app.post("/api/student/login")
async def login_student(request: Request, username: str = Form(...), password: str = Form(...)):
    """🔐 تسجيل دخول الطالب مع حماية متدرجة"""
    client_ip = request.client.host if request.client else "unknown"
    lockout_key = f"student:{username}:{client_ip}"
    
    # القفل المتدرج
    is_locked, seconds_left = check_progressive_lockout(lockout_key)
    if is_locked:
        minutes = max(1, seconds_left // 60)
        security_log("student_login_blocked", client_ip, {"username": username})
        raise HTTPException(status_code=429, detail=f"🔒 الحساب مقفول مؤقتاً — حاول بعد {minutes} دقيقة")
    
    # Rate limit
    if _is_rate_limited(client_ip, max_calls=10, window_seconds=60):
        raise HTTPException(status_code=429, detail="⏳ تجاوزت عدد المحاولات المسموحة — انتظر دقيقة")
    
    res = supabase.table("students").select("*").eq("username", username).execute()
    if res.data and verify_password(password, res.data[0]['password']):
        user = res.data[0]

        # ═══ فحص is_active: الحسابات المعطلة تُمنع من الدخول ═══
        # نعتبر الحساب نشطاً افتراضياً إذا ما كان الحقل موجوداً (للتوافق مع القديم)
        if user.get('is_active') is False:
            raise HTTPException(
                status_code=403,
                detail="🚫 حسابك معطّل حالياً. تواصل مع الأستاذ رشدي لإعادة التفعيل."
            )

        user.pop('password', None)
        
        # ═══ ضمان وجود حقل xp (نأخذ الأعلى بين total_points و xp إن وُجد) ═══
        server_xp = max(int(user.get('total_points', 0) or 0), int(user.get('xp', 0) or 0))
        user['xp'] = server_xp
        user['total_points'] = server_xp

        # ═══ تحديث last_active للطالب ═══
        try:
            try:

                supabase.table("students").update({
                "last_active": datetime.now(timezone.utc).isoformat()
            }).eq("id", user["id"]).execute()

            except Exception:

                pass  # last_active column may not exist
        except Exception:
            pass  # الحقل قد لا يكون موجوداً في جداول قديمة

        # إنشاء JWT للطالب
        token = create_access_token({
            "sub":   str(user["id"]),
            "role":  "student",
            "grade": user.get("grade", "")
        })
        return {"status": "success", "access_token": token, "user": user}
    raise HTTPException(status_code=401, detail="بيانات الدخول خاطئة")


# ════════════════════════════════════════════════════════════
# 🌐 مزامنة XP من العميل للسيرفر
# ════════════════════════════════════════════════════════════
@app.get("/api/student/{student_id}/xp")
async def get_student_xp(student_id: int):
    """🔍 جلب XP الحالي للطالب من السيرفر (مصدر الحقيقة)"""
    try:
        res = supabase.table("students").select(
            "id, full_name, total_points, grade"
        ).eq("id", student_id).limit(1).execute()
        
        if not res.data:
            raise HTTPException(status_code=404, detail="الطالب غير موجود")
        
        student = res.data[0]
        return {
            "status": "ok",
            "id": student["id"],
            "full_name": student.get("full_name", ""),
            "xp": int(student.get("total_points", 0) or 0),
            "grade": student.get("grade", "")
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/student/sync_xp")
async def sync_student_xp(
    student_id: int = Form(...),
    xp: int = Form(...),
    full_name: str = Form(default="")
):
    """🌐 مزامنة XP — يأخذ الأعلى بين السيرفر والعميل + يُرجع القيمة الصحيحة"""
    try:
        if xp < 0: xp = 0
        if xp > 9999999: xp = 9999999  # سقف منطقي
        
        # نأخذ القيمة الحالية من السيرفر
        res = supabase.table("students").select("id, total_points, full_name").eq("id", student_id).limit(1).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="الطالب غير موجود")
        
        student = res.data[0]
        server_xp = int(student.get('total_points', 0) or 0)
        
        # 🛡️ نُحدّث فقط لو القيمة الجديدة أعلى (حماية ضد التراجع)
        if xp > server_xp:
            try:

                supabase.table("students").update({
                "total_points": xp,
                "last_active": datetime.now(timezone.utc).isoformat()
            }).eq("id", student_id).execute()

            except Exception:

                pass  # last_active column may not exist
            return {"status": "ok", "synced": True, "new_xp": xp, "old_xp": server_xp}
        
        # السيرفر أعلى أو مساوٍ — نُرجع قيمة السيرفر (للتصحيح في العميل)
        return {"status": "ok", "synced": False, "new_xp": server_xp, "old_xp": server_xp}
        
    except HTTPException:
        raise
    except Exception as e:
        return {"status": "error", "message": str(e)[:200]}


# ==========================================
# --- 5. مسار المنحة الملكية (XP اليدوي) ---
# ==========================================
@app.post("/api/admin/grant_xp")
async def grant_xp(student_name: str = Form(...), points: int = Form(...), admin=Depends(get_current_admin)):
    supabase.table("results").insert({
        "student_name": student_name,
        "lesson": "💎 منحة ملكية تقديرية من الأستاذ رشدي",
        "score": points,
        "total": points
    }).execute()
    return {"status": "success"}

# ==========================================
# --- 6. إدارة المنهج الدراسي ---
# ==========================================
@app.post("/api/admin/curriculum/grades")
async def add_grade(name: str = Form(...), admin=Depends(get_current_admin)):
    return supabase.table("grades").insert({"name": name}).execute()

@app.post("/api/admin/curriculum/semesters")
async def add_semester(grade_id: int = Form(...), name: str = Form(...), admin=Depends(get_current_admin)):
    return supabase.table("semesters").insert({"grade_id": grade_id, "name": name}).execute()

@app.post("/api/admin/curriculum/units")
async def add_unit(semester_id: int = Form(...), name: str = Form(...), admin=Depends(get_current_admin)):
    return supabase.table("units").insert({"semester_id": semester_id, "name": name}).execute()

@app.post("/api/admin/curriculum/lessons")
async def add_lesson(unit_id: int = Form(...), name: str = Form(...), admin=Depends(get_current_admin)):
    return supabase.table("lessons").insert({"unit_id": unit_id, "name": name}).execute()

@app.put("/api/admin/curriculum/grades/{item_id}")
async def update_grade(item_id: int, name: str = Form(...), admin=Depends(get_current_admin)):
    supabase.table("grades").update({"name": name}).eq("id", item_id).execute()
    return {"status": "success"}

@app.put("/api/admin/curriculum/semesters/{item_id}")
async def update_semester(item_id: int, name: str = Form(...), admin=Depends(get_current_admin)):
    supabase.table("semesters").update({"name": name}).eq("id", item_id).execute()
    return {"status": "success"}

@app.put("/api/admin/curriculum/units/{item_id}")
async def update_unit(item_id: int, name: str = Form(...), admin=Depends(get_current_admin)):
    supabase.table("units").update({"name": name}).eq("id", item_id).execute()
    return {"status": "success"}

@app.put("/api/admin/curriculum/lessons/{item_id}")
async def update_lesson(item_id: int, name: str = Form(...), admin=Depends(get_current_admin)):
    supabase.table("lessons").update({"name": name}).eq("id", item_id).execute()
    return {"status": "success"}

@app.delete("/api/admin/curriculum/grades/{item_id}")
async def delete_grade(item_id: int, admin=Depends(get_current_admin)):
    supabase.table("grades").delete().eq("id", item_id).execute()
    return {"status": "success"}

@app.delete("/api/admin/curriculum/semesters/{item_id}")
async def delete_semester(item_id: int, admin=Depends(get_current_admin)):
    supabase.table("semesters").delete().eq("id", item_id).execute()
    return {"status": "success"}

@app.delete("/api/admin/curriculum/units/{item_id}")
async def delete_unit(item_id: int, admin=Depends(get_current_admin)):
    supabase.table("units").delete().eq("id", item_id).execute()
    return {"status": "success"}

@app.delete("/api/admin/curriculum/lessons/{item_id}")
async def delete_lesson(item_id: int, admin=Depends(get_current_admin)):
    supabase.table("lessons").delete().eq("id", item_id).execute()
    return {"status": "success"}

@app.get("/api/curriculum/structure")
async def get_full_structure():
    try:
        res = supabase.table("grades").select("*, semesters(*, units(*, lessons(*)))").execute()
        data = res.data or []
    except Exception as e:
        print(f"curriculum/structure error: {e}")
        # fallback: جلب الصفوف فقط بدون nested
        try:
            res = supabase.table("grades").select("id, name").execute()
            data = [{"id": g["id"], "name": g["name"], "semesters": []} for g in (res.data or [])]
        except:
            return []

    # ترتيب احتياطي إذا لم يكن sort_order موجوداً
    grade_order = [
        'الصف الخامس','الصف السادس','الصف السابع','الصف الثامن',
        'الصف التاسع','الصف العاشر',
        'الصف الحادي عشر (متقدم)','الصف الحادي عشر(اساسي)',
        'الصف الثاني عشر (متقدم)','الصف الثاني عشر (أساسي)',
    ]

    def grade_sort_key(g):
        name = (g.get("name") or "").strip()
        try:
            return grade_order.index(name)
        except ValueError:
            return 999

    data.sort(key=grade_sort_key)

    # trim الفراغات الزائدة من جميع الأسماء
    for g in data:
        g["name"] = (g.get("name") or "").strip()
        for s in (g.get("semesters") or []):
            s["name"] = (s.get("name") or "").strip()
            for u in (s.get("units") or []):
                u["name"] = (u.get("name") or "").strip()
                for l in (u.get("lessons") or []):
                    l["name"] = (l.get("name") or "").strip()
    return data

# ==========================================
# --- 7. بنك الأسئلة والامتحانات ---
# ==========================================
@app.post("/api/admin/questions")
async def add_question(
    grade:      str         = Form(...),
    semester:   str         = Form(default=""),
    unit:       str         = Form(default=""),
    lesson:     str         = Form(...),
    subject:    str         = Form(...),
    q_type:     str         = Form(...),
    question:   str         = Form(...),
    options:    str         = Form(""),
    answer:     str         = Form(...),
    image:      UploadFile  = File(None),
    is_elite:   str         = Form(default="false"),
    difficulty: str         = Form(default="hard"),
    admin=Depends(get_current_admin)
):
    img_url = ""
    if image and image.filename:
        img_name = f"q_img_{uuid.uuid4().hex}{os.path.splitext(image.filename)[1]}"
        content = await image.read()
        supabase.storage.from_("resources").upload(path=img_name, file=content, file_options={"content-type": image.content_type})
        img_url = supabase.storage.from_("resources").get_public_url(img_name)

    row = {
        "grade":    grade,
        "semester": semester,
        "unit":     unit,
        "lesson":   lesson,
        "subject":  subject,
        "q_type":   q_type,
        "question": question,
        "options":  options,
        "answer":   answer,
        "image_url": img_url,
    }
    # حقول النخبة — اختيارية
    elite_val = is_elite.lower().strip() not in ('false', '0', 'no', '')
    if elite_val:
        row["is_elite"]   = True
        row["difficulty"] = difficulty.strip() or "hard"

    try:
        supabase.table("questions").insert(row).execute()
    except Exception as e:
        # fallback: لو الأعمدة الجديدة ما موجودة بعد، نحذفها ونعيد المحاولة
        if "semester" in str(e) or "unit" in str(e) or "column" in str(e).lower():
            row.pop("semester", None)
            row.pop("unit", None)
            supabase.table("questions").insert(row).execute()
        else:
            raise
    # 🧹 invalidate cache
    try:
        cache_invalidate("questions:", "stats:")
    except Exception: pass
    
    return {"status": "success"}

@app.put("/api/admin/questions/{q_id}")
async def update_question(
    q_id: int,
    grade: str=Form(...),
    semester: str=Form(default=""),
    unit: str=Form(default=""),
    lesson: str=Form(...),
    subject: str=Form(...),
    q_type: str=Form(...), question: str=Form(...), options: str=Form(""), 
    answer: str=Form(...), image: UploadFile=File(None), admin=Depends(get_current_admin)
):
    update_data = {
        "grade": grade, "semester": semester, "unit": unit,
        "lesson": lesson, "subject": subject, "q_type": q_type, 
        "question": question, "options": options, "answer": answer
    }
    
    if image and image.filename:
        img_name = f"q_img_{uuid.uuid4().hex}{os.path.splitext(image.filename)[1]}"
        content = await image.read()
        supabase.storage.from_("resources").upload(path=img_name, file=content, file_options={"content-type": image.content_type})
        update_data["image_url"] = supabase.storage.from_("resources").get_public_url(img_name)

    try:
        supabase.table("questions").update(update_data).eq("id", q_id).execute()
    except Exception as e:
        # fallback: لو الأعمدة الجديدة ما موجودة
        if "semester" in str(e) or "unit" in str(e) or "column" in str(e).lower():
            update_data.pop("semester", None)
            update_data.pop("unit", None)
            supabase.table("questions").update(update_data).eq("id", q_id).execute()
        else:
            raise
    # 🧹 invalidate cache
    try:
        cache_invalidate("questions:", "stats:")
    except Exception: pass
    
    return {"status": "success"}

@app.delete("/api/admin/questions/{q_id}")
async def delete_question(q_id: int, admin=Depends(get_current_admin)):
    """🗑️ حذف سؤال — مع تحقق فعلي + invalidate cache"""
    try:
        # نتحقق أولاً أن السؤال موجود
        check = supabase.table("questions").select("id").eq("id", q_id).limit(1).execute()
        if not check.data:
            raise HTTPException(status_code=404, detail="السؤال غير موجود")
        
        # نحذف فعلياً
        result = supabase.table("questions").delete().eq("id", q_id).execute()
        
        # 🧹 invalidate cache (مهم جداً - بدونه السؤال يظهر بعد الحذف)
        try:
            cache_invalidate("questions:", "stats:")
        except Exception:
            pass
        
        return {"status": "success", "deleted_id": q_id}
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل الحذف: {str(e)[:200]}")


@app.post("/api/admin/questions/bulk_delete")
async def bulk_delete_questions(
    request: Request,
    admin=Depends(get_current_admin)
):
    """
    حذف أسئلة بالجملة حسب الفلتر الملكي.
    يتطلب كلمة مرور الأدمن كحماية إضافية لأن العملية مدمّرة.

    Body (form-urlencoded):
      - admin_password (إلزامي): كلمة مرور الأدمن للتأكيد
      - grade          (إلزامي): اسم الصف — نرفض الحذف بدونه لمنع مسح كامل
      - semester       (اختياري): اسم الفصل
      - unit           (اختياري): اسم الوحدة
      - lesson         (اختياري): اسم الدرس
      - dry_run        (اختياري): "true" لعدّ الأسئلة فقط دون حذف (معاينة)
    """
    body = await request.form()
    admin_pass = body.get("admin_password", "")
    grade      = (body.get("grade") or "").strip()
    semester   = (body.get("semester") or "").strip()
    unit       = (body.get("unit") or "").strip()
    lesson     = (body.get("lesson") or "").strip()
    dry_run    = str(body.get("dry_run", "")).lower() in ("true", "1", "yes")

    # حماية 1: كلمة مرور الأدمن مطلوبة
    if not admin_pass or admin_pass != ADMIN_PASSWORD:
        raise HTTPException(status_code=403, detail="كلمة مرور الأدمن خاطئة")

    # حماية 2: الصف إلزامي — لا يُسمح بحذف كل الأسئلة دفعة واحدة عن طريق الخطأ
    if not grade:
        raise HTTPException(status_code=400, detail="يجب تحديد الصف على الأقل")

    # بناء الاستعلام مع كل صيغ الصف (السادس / 6 / الصف السادس …)
    variants = _grade_variants(grade)
    if not variants:
        variants = [grade]

    # 🔍 منطق ذكي مرن — يجلب كل أسئلة الصف ثم يفلتر بـ Python
    # هذا يحل مشكلة:
    # 1. الفراغات الزائدة في القيم المخزنة
    # 2. اختلاف بسيط في النص (مع/بدون أرقام الدرس)
    # 3. أعمدة semester/unit مفقودة في بعض الصفوف

    def _normalize_curr_text(s):
        """ينظّف النص للمقارنة المرنة"""
        if not s:
            return ""
        s = str(s).strip()
        # توحيد الفراغات
        import re as _re
        s = _re.sub(r'\s+', ' ', s)
        # حذف الأقواس والترقيم
        s = s.replace('  ', ' ')
        return s.lower()

    matched_ids = set()
    matched_preview = []
    diagnostics = {
        "total_grade_questions": 0,
        "filtered_out": {"semester": 0, "unit": 0, "lesson": 0},
        "sample_db_values": {"semesters": set(), "units": set(), "lessons": set()},
    }

    norm_semester = _normalize_curr_text(semester)
    norm_unit     = _normalize_curr_text(unit)
    norm_lesson   = _normalize_curr_text(lesson)

    for v in variants:
        # 🛡️ نجلب كل أسئلة الصف بدون فلترة على semester/unit/lesson من قاعدة البيانات
        try:
            res = supabase.table("questions").select(
                "id, grade, semester, unit, lesson, question"
            ).eq("grade", v.strip()).execute()
            rows = res.data or []
        except Exception as e:
            # fallback: لو semester/unit مش موجودة في الجدول
            print(f"[bulk_delete] جلب بأعمدة كاملة فشل: {e}")
            try:
                res = supabase.table("questions").select(
                    "id, grade, lesson, question"
                ).eq("grade", v.strip()).execute()
                rows = res.data or []
                # نُكمّل السطور المفقودة
                for r in rows:
                    r.setdefault("semester", "")
                    r.setdefault("unit", "")
            except Exception as e2:
                print(f"[bulk_delete] جلب بسيط فشل أيضاً: {e2}")
                continue

        diagnostics["total_grade_questions"] += len(rows)

        # 🎯 الفلترة المرنة في Python
        for row in rows:
            row_sem    = _normalize_curr_text(row.get("semester", ""))
            row_unit   = _normalize_curr_text(row.get("unit", ""))
            row_lesson = _normalize_curr_text(row.get("lesson", ""))

            # تجميع عيّنة من القيم الفعلية في قاعدة البيانات (للتشخيص)
            if row.get("semester"):
                diagnostics["sample_db_values"]["semesters"].add(row["semester"])
            if row.get("unit"):
                diagnostics["sample_db_values"]["units"].add(row["unit"])
            if row.get("lesson"):
                diagnostics["sample_db_values"]["lessons"].add(row["lesson"])

            # فلتر الفصل
            if norm_semester:
                if row_sem != norm_semester and norm_semester not in row_sem and row_sem not in norm_semester:
                    diagnostics["filtered_out"]["semester"] += 1
                    continue

            # فلتر الوحدة
            if norm_unit:
                if row_unit != norm_unit and norm_unit not in row_unit and row_unit not in norm_unit:
                    diagnostics["filtered_out"]["unit"] += 1
                    continue

            # فلتر الدرس
            if norm_lesson:
                if row_lesson != norm_lesson and norm_lesson not in row_lesson and row_lesson not in norm_lesson:
                    diagnostics["filtered_out"]["lesson"] += 1
                    continue

            if row["id"] not in matched_ids:
                matched_ids.add(row["id"])
                if len(matched_preview) < 5:
                    matched_preview.append({
                        "id": row["id"],
                        "lesson": row.get("lesson", ""),
                        "question": (row.get("question") or "")[:80]
                    })

    count = len(matched_ids)

    # تحويل sets إلى lists للـ JSON
    diagnostics["sample_db_values"]["semesters"] = list(diagnostics["sample_db_values"]["semesters"])[:5]
    diagnostics["sample_db_values"]["units"]     = list(diagnostics["sample_db_values"]["units"])[:5]
    diagnostics["sample_db_values"]["lessons"]   = list(diagnostics["sample_db_values"]["lessons"])[:10]

    if dry_run:
        return {
            "status": "preview",
            "count": count,
            "preview": matched_preview,
            "filter": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson},
            "diagnostics": diagnostics  # 🩺 تفاصيل للتشخيص
        }

    if count == 0:
        return {"status": "empty", "deleted": 0, "message": "لا توجد أسئلة مطابقة"}

    # الحذف الفعلي — على دفعات من 50 لتجنّب تجاوز حدود Supabase
    ids_list = list(matched_ids)
    deleted = 0
    for i in range(0, len(ids_list), 50):
        batch = ids_list[i:i+50]
        try:
            supabase.table("questions").delete().in_("id", batch).execute()
            deleted += len(batch)
        except Exception as e:
            # نكمل حتى لو فشلت دفعة
            print(f"[bulk_delete] فشلت دفعة: {e}")

    # 🧹 invalidate cache
    try:
        cache_invalidate("questions:", "stats:")
    except Exception: pass
    
    return {
        "status": "success",
        "deleted": deleted,
        "requested": count,
        "filter": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson}
    }


@app.post("/api/admin/questions/bulk")
async def bulk_add_questions(request: Request, admin=Depends(get_current_admin)):
    """ضخ دفعة أسئلة مولّدة بالذكاء الاصطناعي دفعة واحدة — أسرع من الإرسال الفردي"""
    body = await request.json()
    questions = body.get("questions", [])
    if not questions:
        raise HTTPException(status_code=400, detail="لا توجد أسئلة في الطلب")
    inserted = 0
    errors   = 0
    for q in questions:
        try:
            row = {
                "grade":    str(q.get("grade",    "") or "").strip(),
                "semester": str(q.get("semester", "") or "").strip(),
                "unit":     str(q.get("unit",    "") or "").strip(),
                "lesson":   str(q.get("lesson",   "") or "").strip(),
                "subject":  str(q.get("subject",  "رياضيات") or "رياضيات").strip(),
                "q_type":   str(q.get("q_type",   "choice") or "choice").strip(),
                "question": str(q.get("question", "") or "").strip(),
                "options":  str(q.get("options",  "") or "").strip(),
                "answer":   str(q.get("answer",   "") or "").strip(),
                "image_url": "",
            }
            if not row["question"] or not row["answer"]:
                errors += 1
                continue
            is_elite_val = str(q.get("is_elite", "false")).lower().strip()
            if is_elite_val not in ("false", "0", "no", ""):
                row["is_elite"]   = True
                row["difficulty"] = str(q.get("difficulty", "hard") or "hard").strip()
            else:
                diff = str(q.get("difficulty", "") or "").strip()
                if diff in ("easy", "medium", "hard"):
                    row["difficulty"] = diff
            try:
                supabase.table("questions").insert(row).execute()
                inserted += 1
            except Exception as e:
                # fallback: الأعمدة الجديدة ما موجودة
                if "semester" in str(e) or "unit" in str(e) or "column" in str(e).lower():
                    row.pop("semester", None)
                    row.pop("unit", None)
                    supabase.table("questions").insert(row).execute()
                    inserted += 1
                else:
                    errors += 1
        except Exception:
            errors += 1
    return {"inserted": inserted, "errors": errors, "total": len(questions)}


# ═══════════════════════════════════════════════════════════════
# 💾 CACHE SYSTEM — تخزين مؤقت للبيانات الثقيلة
# يُقلّل ضغط Supabase بنسبة 80% على البيانات المتكررة
# ═══════════════════════════════════════════════════════════════
import time as _time

class _SimpleCache:
    """Cache في الذاكرة مع TTL تلقائي + max size"""
    def __init__(self, max_size=200):
        self._store = {}
        self._max = max_size
    
    def get(self, key):
        item = self._store.get(key)
        if not item:
            return None
        value, expires_at = item
        if _time.time() > expires_at:
            del self._store[key]
            return None
        return value
    
    def set(self, key, value, ttl_seconds=300):
        if len(self._store) >= self._max:
            # نزيل أقدم العناصر
            oldest = sorted(self._store.items(), key=lambda x: x[1][1])[:50]
            for k, _ in oldest:
                self._store.pop(k, None)
        self._store[key] = (value, _time.time() + ttl_seconds)
    
    def delete(self, key):
        self._store.pop(key, None)
    
    def clear_pattern(self, pattern):
        """يحذف كل المفاتيح التي تحتوي على pattern (مفيد عند تعديل البيانات)"""
        keys_to_del = [k for k in self._store.keys() if pattern in k]
        for k in keys_to_del:
            del self._store[k]
    
    def stats(self):
        now = _time.time()
        valid = sum(1 for _, exp in self._store.values() if now <= exp)
        return {"total": len(self._store), "valid": valid, "max": self._max}

_cache = _SimpleCache(max_size=300)


def cache_invalidate(*patterns):
    """يلغي cache لأنماط محددة - يُستدعى عند تعديل البيانات"""
    for p in patterns:
        _cache.clear_pattern(p)


# ═══════════════════════════════════════════════════════════════
# 📊 PAGINATION HELPERS
# ═══════════════════════════════════════════════════════════════
def _paginate_params(page: int = 1, page_size: int = 50, max_size: int = 200):
    """يضمن أن الـ pagination params آمنة وضمن الحدود"""
    page = max(1, int(page))
    page_size = max(1, min(int(page_size), max_size))
    offset = (page - 1) * page_size
    return page, page_size, offset


def _paginate_response(data: list, total: int, page: int, page_size: int):
    """يبني response موحّد للـ pagination"""
    total_pages = (total + page_size - 1) // page_size if total else 1
    return {
        "data": data,
        "pagination": {
            "page": page,
            "page_size": page_size,
            "total_items": total,
            "total_pages": total_pages,
            "has_next": page < total_pages,
            "has_prev": page > 1,
        }
    }


@app.get("/api/admin/questions")
async def get_all_questions(
    admin = Depends(get_current_admin),
    page: int = 1,
    page_size: int = 100,
    grade: str = "",
    lesson: str = "",
    search: str = "",
):
    """
    🚀 جلب الأسئلة مع pagination + cache + فلترة
    
    Parameters:
    - page: رقم الصفحة (افتراضي 1)
    - page_size: عدد العناصر (افتراضي 100، أقصى 200)
    - grade: فلتر الصف
    - lesson: فلتر الدرس
    - search: بحث في نص السؤال
    """
    page, page_size, offset = _paginate_params(page, page_size, max_size=200)
    
    # cache key يعتمد على كل الفلاتر
    cache_key = f"questions:p{page}:s{page_size}:g{grade}:l{lesson}:q{search}"
    cached = _cache.get(cache_key)
    if cached:
        return cached
    
    try:
        # نبني الاستعلام مع الفلاتر
        query = supabase.table("questions").select("*", count="exact")
        if grade:
            query = query.eq("grade", grade.strip())
        if lesson:
            query = query.eq("lesson", lesson.strip())
        if search:
            query = query.ilike("question", f"%{search.strip()}%")
        
        # نطبّق الـ pagination + الترتيب
        res = query.order("id", desc=True).range(offset, offset + page_size - 1).execute()
        
        result = _paginate_response(
            data=res.data or [],
            total=res.count or 0,
            page=page,
            page_size=page_size
        )
        
        # نخزّن في cache لمدة 60 ثانية فقط (الأسئلة قد تتغيّر)
        _cache.set(cache_key, result, ttl_seconds=60)
        return result
    except Exception as e:
        print(f"[questions] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/admin/questions/all")
async def get_all_questions_legacy(admin = Depends(get_current_admin)):
    """
    🔄 endpoint قديم متوافق — يجلب كل الأسئلة (استخدمه فقط للتصدير الكامل)
    ⚠️ لا تستخدمه في الواجهة العادية - استخدم /api/admin/questions مع pagination
    """
    cache_key = "questions:all"
    cached = _cache.get(cache_key)
    if cached:
        return cached
    
    all_questions = []
    page_size = 1000
    offset = 0
    
    for _ in range(50):
        try:
            res = supabase.table("questions").select("*").order("id", desc=True).range(offset, offset + page_size - 1).execute()
            batch = res.data or []
            if not batch:
                break
            all_questions.extend(batch)
            if len(batch) < page_size:
                break
            offset += page_size
        except Exception as e:
            print(f"questions pagination error at offset {offset}: {e}")
            break
    
    _cache.set(cache_key, all_questions, ttl_seconds=120)
    return all_questions

def _grade_variants(grade: str) -> list:
    """يولّد جميع الصيغ الممكنة لاسم الصف لضمان تطابق الأسئلة المخزونة بأي صيغة"""
    if not grade:
        return []
    s = grade.strip()
    # نضيف الصيغة الأصلية + مع مسافات (لأن Supabase قد يخزنها بمسافات زائدة)
    variants = {s, f" {s} ", f" {s}", f"{s} "}

    # خريطة الأرقام العربية ↔ الأرقام
    ar_to_num = {
        "الأول":"1","الثاني":"2","الثالث":"3","الرابع":"4","الخامس":"5","السادس":"6",
        "السابع":"7","الثامن":"8","التاسع":"9","العاشر":"10","الحادي عشر":"11","الثاني عشر":"12",
    }
    num_to_ar = {v: k for k, v in ar_to_num.items()}

    import re
    # حالة 1: "الصف السادس" → نستخرج "السادس" ثم "6"
    m = re.match(r"الصف\s+(.+)", s)
    if m:
        word = m.group(1).strip()
        num  = ar_to_num.get(word, word)
        variants.update([num, f"الصف {word}", f"الصف {num}", word])

    # حالة 2: "السادس" فقط (بدون "الصف") → نضيف "6" و"الصف السادس"
    elif s in ar_to_num:
        num  = ar_to_num[s]
        variants.update([num, f"الصف {s}", f"الصف {num}"])

    # حالة 3: رقم مجرد "6" → نضيف "السادس" و"الصف السادس"
    elif re.match(r"^\d+$", s):
        word = num_to_ar.get(s, s)
        variants.update([f"الصف {word}", f"الصف {s}", word])

    return list(variants)


@app.get("/api/student/questions")
async def get_questions_for_student(grade: str, lesson: str = ""):
    """
    جلب الأسئلة للطالب — بدون حقل answer
    - يدعم جميع صيغ اسم الصف (الصف السادس / 6 / السادس)
    - إذا أُرسل lesson يُفلتر به، وإلا يُرجع كل أسئلة الصف
    - يُرجع الأسئلة حتى لو لم يكن هناك منهج مبني (للأسئلة القديمة)
    - الإجابات تُفحص عبر /api/student/check_answer فقط (لا تُرسل للعميل)
    """
    variants = _grade_variants(grade)

    all_questions = []
    seen_ids: set = set()

    # إذا لم تتوفر أي صيغة — أرجع كل الأسئلة (fallback للأسئلة القديمة)
    search_variants = variants if variants else [grade] if grade else []

    # ═══ حقول آمنة فقط — بدون answer ═══
    SAFE_FIELDS = "id, grade, lesson, subject, q_type, question, options, image_url"

    for v in search_variants:
        v_stripped = v.strip()
        query = supabase.table("questions").select(SAFE_FIELDS).eq("grade", v_stripped)
        if lesson:
            query = query.ilike("lesson", lesson.strip())
        res = query.execute()
        for q in (res.data or []):
            if q["id"] not in seen_ids:
                seen_ids.add(q["id"])
                all_questions.append(q)

    # إذا لم نجد شيئاً وكان lesson محدداً — جرّب بحث جزئي في lesson
    if not all_questions and lesson:
        for v in search_variants:
            res = supabase.table("questions").select(SAFE_FIELDS)\
                .eq("grade", v).ilike("lesson", f"%{lesson.strip()}%").execute()
            for q in (res.data or []):
                if q["id"] not in seen_ids:
                    seen_ids.add(q["id"])
                    all_questions.append(q)

    return all_questions


@app.post("/api/student/check_answer")
async def check_answer(request: Request):
    """
    التحقق من إجابة الطالب على السيرفر
    يُرجع is_correct + correct_answer عند الخطأ (ليُعرض في مستشفى الأرقام)
    Body JSON: { "question_id": int, "student_answer": str }
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=120, window_seconds=60):
        raise HTTPException(status_code=429, detail="طلبات كثيرة جداً — انتظر لحظة")

    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="صيغة الطلب غير صحيحة — يُتوقع JSON")

    question_id    = body.get("question_id")
    student_answer = body.get("student_answer", "")

    if not question_id or not isinstance(question_id, int):
        raise HTTPException(status_code=400, detail="question_id مطلوب وصحيح")

    # جلب الإجابة الصحيحة من قاعدة البيانات
    res = supabase.table("questions").select("id, answer, q_type").eq("id", question_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="السؤال غير موجود")

    correct_answer  = str(res.data[0].get("answer", "")).strip()
    q_type          = res.data[0].get("q_type", "")
    student_cleaned = str(student_answer).strip()

    # مقارنة غير حساسة لحالة الأحرف + تجاهل المسافات الزائدة
    normalize = lambda s: " ".join((s or "").strip().split())
    is_correct = normalize(correct_answer).lower() == normalize(student_cleaned).lower()

    response_data = {
        "is_correct": is_correct,
        "q_type":     q_type,
    }
    # نُرجع الإجابة الصحيحة فقط عند الخطأ — لعرضها في مستشفى الأرقام
    if not is_correct:
        response_data["correct_answer"] = correct_answer

    return response_data


@app.get("/api/admin/debug/questions")
async def debug_questions(admin=Depends(get_current_admin)):
    """
    نقطة تشخيصية — تُظهر قيم grade و lesson المخزونة فعلياً في Supabase
    مفيدة لمعرفة لماذا لا تظهر الأسئلة عند الطلاب
    """
    res = supabase.table("questions").select("id, grade, lesson, q_type").order("id", desc=True).limit(200).execute()
    if not res.data:
        return {"count": 0, "grades": [], "lessons": [], "samples": []}
    
    grades  = sorted(set(q["grade"]  or "" for q in res.data))
    lessons = sorted(set(q["lesson"] or "" for q in res.data))
    return {
        "count":   len(res.data),
        "grades":  grades,
        "lessons": lessons,
        "samples": res.data[:10],
    }

@app.post("/api/admin/exams")
async def create_exam(
    title: str=Form(...), exam_date: str=Form(...), 
    exam_time: str=Form(...), target_lesson: str=Form(...), duration: int=Form(...),
    exam_type: str=Form(default="ملحمة أسبوعية"), num_questions: int=Form(default=10), 
    points_per_q: int=Form(default=10), target_q_type: str=Form(default="all"),
    notif_lifetime_hours: int=Form(default=72),
    admin=Depends(get_current_admin)
):
    # حساب تاريخ انتهاء صلاحية الإشعار (0 = دائم)
    expires_at = None
    if notif_lifetime_hours and notif_lifetime_hours > 0:
        expires_at = (datetime.now(timezone.utc) + timedelta(hours=notif_lifetime_hours)).isoformat()

    payload = {
        "title": title, "exam_type": exam_type, "exam_date": exam_date, "exam_time": exam_time, 
        "target_lesson": target_lesson, "duration": duration, "num_questions": num_questions, 
        "points_per_q": points_per_q, "target_q_type": target_q_type,
        "notif_lifetime_hours": notif_lifetime_hours,
        "expires_at": expires_at
    }
    try:
        supabase.table("exams").insert(payload).execute()
    except Exception as e:
        # لو الأعمدة الجديدة غير موجودة في DB بعد، نُحاول بدون
        err_msg = str(e).lower()
        if "notif_lifetime_hours" in err_msg or "expires_at" in err_msg or "column" in err_msg:
            payload.pop("notif_lifetime_hours", None)
            payload.pop("expires_at", None)
            supabase.table("exams").insert(payload).execute()
        else:
            raise
    return {"status": "success"}

@app.get("/api/exams/upcoming")
async def get_upcoming_exams(student_id: Optional[int] = None, username: Optional[str] = None):
    res = supabase.table("exams").select("*").order("id", desc=True).execute()
    exams = res.data if res.data else []

    if not exams:
        return []

    now = datetime.now(timezone.utc)

    # 1) فلترة الاختبارات منتهية الصلاحية (expires_at سابق للوقت الحالي)
    active_exams = []
    for e in exams:
        exp = e.get("expires_at")
        if exp:
            try:
                # Supabase يُرجع ISO string — نُحوّل لـ datetime
                exp_dt = datetime.fromisoformat(exp.replace("Z", "+00:00"))
                if exp_dt < now:
                    continue  # منتهي
            except Exception:
                pass
        active_exams.append(e)

    # 2) لو فيه student_id أو username، نستبعد الاختبارات المكتملة
    sid = student_id
    if not sid and username:
        st = supabase.table("students").select("id").eq("username", username).execute()
        if st.data:
            sid = st.data[0]["id"]

    if sid:
        try:
            done = supabase.table("exam_completions").select("exam_id").eq("student_id", sid).execute()
            done_ids = {r["exam_id"] for r in (done.data or [])}
            active_exams = [e for e in active_exams if e["id"] not in done_ids]
        except Exception:
            # لو الجدول لسه ما اتعمل — ما نكسر الـ endpoint
            pass

    return active_exams

@app.delete("/api/admin/exams/{exam_id}")
async def delete_exam(exam_id: int, admin=Depends(get_current_admin)):
    supabase.table("exams").delete().eq("id", exam_id).execute()
    return {"status": "success"}

# ─── تسجيل إكمال الاختبار (لإخفائه من إشعارات الطالب) ───
@app.post("/api/student/exam_completed")
async def mark_exam_completed(
    exam_id: int = Form(...),
    student_id: int = Form(...)
):
    # تحقق من وجود الطالب والاختبار
    st = supabase.table("students").select("id").eq("id", student_id).execute()
    if not st.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    ex = supabase.table("exams").select("id").eq("id", exam_id).execute()
    if not ex.data:
        raise HTTPException(status_code=404, detail="الاختبار غير موجود")

    try:
        # upsert: لو موجود مسبقاً ما نكرّره
        existing = supabase.table("exam_completions").select("id")\
            .eq("exam_id", exam_id).eq("student_id", student_id).execute()
        if not existing.data:
            supabase.table("exam_completions").insert({
                "exam_id": exam_id,
                "student_id": student_id,
                "completed_at": datetime.now(timezone.utc).isoformat()
            }).execute()
    except Exception as e:
        # لو الجدول غير موجود، نُعيد رسالة واضحة
        raise HTTPException(
            status_code=500,
            detail=f"جدول exam_completions غير متوفر — شغّل migration أولاً: {e}"
        )
    return {"status": "success"}

# ==========================================
# --- 8. ديوان الموارد ولفائف المعرفة ---
# ==========================================
@app.get("/api/resources")
async def get_resources(grade: str, semester: str, category: str = "all"):
    query = supabase.table("teacher_resources").select("*").eq("grade", grade).eq("semester", semester)
    if category != "all":
        query = query.eq("category", category)
    res = query.execute()
    return res.data

@app.get("/api/admin/all_resources")
async def get_all_resources(admin=Depends(get_current_admin)):
    res = supabase.table("teacher_resources").select("*").order("id", desc=True).execute()
    return res.data if res.data else []

@app.post("/api/admin/resources")
async def add_resource(
    title: str=Form(...), grade: str=Form(...), semester: str=Form(...), 
    category: str=Form(...), description: str=Form(""), 
    file: UploadFile=File(...), admin=Depends(get_current_admin)
):
    # ═══ تحقق آمن من الملف ═══
    content = await file.read()
    file_extension = _validate_upload(content, file.filename or "")

    file_name = f"res_{uuid.uuid4().hex}{file_extension}"

    supabase.storage.from_("resources").upload(
        path=file_name,
        file=content,
        file_options={"content-type": file.content_type or "application/octet-stream"}
    )
    file_url = supabase.storage.from_("resources").get_public_url(file_name)

    supabase.table("teacher_resources").insert({
        "title": title, "grade": grade, "semester": semester,
        "category": category, "description": description, "file_url": file_url
    }).execute()
    return {"status": "success"}

@app.delete("/api/admin/resources/{res_id}")
async def delete_resource(res_id: int, admin=Depends(get_current_admin)):
    supabase.table("teacher_resources").delete().eq("id", res_id).execute()
    return {"status": "success"}

@app.post("/api/admin/summaries")
async def upload_summary(
    lesson:         str         = Form(...),
    resource_type:  str         = Form(default="pdf"),
    resource_label: str         = Form(default=""),
    external_url:   str         = Form(default=""),
    pdf:            UploadFile  = File(default=None),
    admin=Depends(get_current_admin)
):
    """رفع مصدر تعليمي — PDF أو فيديو أو رابط خارجي"""
    try:
        is_file_type = resource_type in ("pdf", "worksheet")

        if is_file_type:
            if not pdf or not pdf.filename:
                raise HTTPException(status_code=400, detail="يرجى إرفاق ملف")
            # ═══ تحقق آمن من الملف ═══
            content        = await pdf.read()
            file_extension = _validate_upload(content, pdf.filename)
            file_name      = f"res_{uuid.uuid4().hex}{file_extension}"
            content_type   = pdf.content_type or "application/pdf"
            supabase.storage.from_("resources").upload(
                path=file_name, file=content,
                file_options={"content-type": content_type}
            )
            resource_url = supabase.storage.from_("resources").get_public_url(file_name)
        else:
            if not external_url:
                raise HTTPException(status_code=400, detail="يرجى إدخال الرابط الخارجي")
            # ═══ تحقق من سلامة الرابط الخارجي (منع javascript:, file:, إلخ) ═══
            if not _is_safe_url(external_url):
                raise HTTPException(status_code=400, detail="الرابط غير صالح — يجب أن يبدأ بـ https:// أو http://")
            resource_url = external_url

        row = {
            "lesson":         lesson,
            "pdf_url":        resource_url,
            "resource_type":  resource_type,
            "resource_label": resource_label,
        }
        supabase.table("summaries").insert(row).execute()
        return {"status": "success", "url": resource_url}

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/admin/summaries_list")
async def get_summaries():
    """قائمة المصادر التعليمية — متاحة لكل المستخدمين (لا تكشف بيانات حساسة)
    تُرجع: lesson, resource_type, resource_label, external_url, file_url
    """
    res = supabase.table("summaries").select("*").order("id", desc=True).execute()
    return res.data if res.data else []


@app.delete("/api/admin/summaries/{resource_id}")
async def delete_summary(resource_id: str, admin=Depends(get_current_admin)):
    """حذف مصدر بالمعرّف أو باسم الدرس (للتوافق مع الكود القديم)"""
    if resource_id.isdigit():
        supabase.table("summaries").delete().eq("id", int(resource_id)).execute()
    else:
        clean = unquote(resource_id)
        supabase.table("summaries").delete().eq("lesson", clean).execute()
    return {"status": "success"}

# ==========================================
# --- 9. النتائج ولوحة الشرف والبحث ---
# ==========================================
@app.post("/api/student/results")
async def save_result(
    request: Request,
    student_id: int=Form(...), student_name: str=Form(...),
    lesson: str=Form(...), score: int=Form(...), total: int=Form(...)
):
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=30, window_seconds=60):
        raise HTTPException(status_code=429, detail="طلبات كثيرة جداً")
    if score < 0 or total <= 0 or score > total:
        raise HTTPException(status_code=400, detail="بيانات غير صحيحة")
    if total > 200:
        raise HTTPException(status_code=400, detail="عدد أسئلة غير منطقي")
    if len(lesson.strip()) < 2 or len(lesson) > 300:
        raise HTTPException(status_code=400, detail="اسم درس غير صحيح")
    # التحقق أن الطالب موجود ونتيجة واحدة لكل درس/جلسة
    st = supabase.table("students").select("id, username").eq("id", student_id).execute()
    if not st.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    supabase.table("results").insert({
        "student_id": student_id,
        "student_name": student_name[:100],
        "lesson": lesson[:300],
        "score": score,
        "total": total
    }).execute()

    # ═══ تحديث last_active للطالب عند حفظ النتيجة ═══
    try:
        try:

            supabase.table("students").update({
            "last_active": datetime.now(timezone.utc).isoformat()
        }).eq("id", student_id).execute()

        except Exception:

            pass  # last_active column may not exist
    except Exception:
        pass
    
    # 🏅 فحص الشارات المستحقة تلقائياً
    new_badges = []
    try:
        new_badges = _check_and_grant_achievements(student_id, {
            "score": score,
            "total": total,
        })
    except Exception:
        pass

    # 🧹 invalidate cache (leaderboard + stats يجب تحديثها)
    try:
        cache_invalidate("leaderboard:", "stats:")
    except Exception:
        pass

    return {"status": "success", "new_badges": new_badges}

@app.post("/api/student/heartbeat")
async def student_heartbeat(student_id: int = Form(...)):
    """
    تحديث last_active للطالب + حفظ session bucket
    يُستدعى كل دقيقة من العميل أثناء النشاط
    
    Bucket = 5 دقائق → الطالب يُسجَّل مرة واحدة في كل bucket
    وقت الطالب اليومي = عدد buckets فريدة × 5 دقائق
    """
    now = datetime.now(timezone.utc)
    try:
        # 1. تحديث last_active في students (للتوافق)
        try:

            supabase.table("students").update({
            "last_active": now.isoformat()
        }).eq("id", student_id).execute()

        except Exception:

            pass  # last_active column may not exist
        
        # 2. حفظ session bucket (مدة 5 دقائق)
        # نقرّب الوقت لأقرب 5 دقائق (12:00, 12:05, 12:10, ...)
        bucket_minute = (now.minute // 5) * 5
        bucket = now.replace(minute=bucket_minute, second=0, microsecond=0)
        
        # جلب بيانات الطالب لتسجيلها مع الجلسة
        try:
            stu_res = supabase.table("students").select("full_name, grade").eq("id", student_id).limit(1).execute()
            stu_data = stu_res.data[0] if stu_res.data else {}
        except Exception:
            stu_data = {}
        
        # upsert: إن كان الـ bucket موجود، نحدّث last_seen + counter
        try:
            existing = supabase.table("student_sessions").select("id, heartbeat_count")\
                .eq("student_id", student_id).eq("session_bucket", bucket.isoformat()).limit(1).execute()
            if existing.data:
                supabase.table("student_sessions").update({
                    "last_seen": now.isoformat(),
                    "heartbeat_count": (existing.data[0].get("heartbeat_count", 1) or 1) + 1
                }).eq("id", existing.data[0]["id"]).execute()
            else:
                supabase.table("student_sessions").insert({
                    "student_id": student_id,
                    "student_name": stu_data.get("full_name", ""),
                    "grade": stu_data.get("grade", ""),
                    "session_bucket": bucket.isoformat(),
                    "last_seen": now.isoformat(),
                    "heartbeat_count": 1
                }).execute()
        except Exception as e:
            # الجدول قد لا يكون موجوداً بعد
            print(f"[heartbeat session] {str(e)[:100]}")
        
        return {"status": "ok"}
    except Exception:
        return {"status": "skipped"}


@app.get("/api/admin/stats/daily_activity")
async def get_daily_activity_stats(
    days: int = 7,
    admin = Depends(get_current_admin)
):
    """
    📊 إحصائيات الحضور اليومي — حقيقية ومن الـ buckets
    
    Returns:
        - daily_stats: قائمة بآخر N أيام (اليوم + N-1 يوم سابقاً)
            * date: التاريخ
            * unique_students: عدد الطلاب الفريدين
            * total_minutes: إجمالي الدقائق المُحتسبة (بكل الطلاب)
            * avg_minutes_per_student: متوسط الدقائق لكل طالب
        - today_live: الطلاب النشطون الآن (آخر 5 دقائق)
        - today_total_unique: إجمالي طلاب اليوم الفريدين
    """
    from datetime import datetime, timezone, timedelta
    
    days = max(1, min(days, 30))  # حد بين 1 و 30
    now = datetime.now(timezone.utc)
    today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
    range_start = today_start - timedelta(days=days - 1)
    
    # جلب كل sessions في النطاق المطلوب (مع pagination)
    all_sessions = []
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("student_sessions").select(
                "student_id, student_name, grade, session_bucket, last_seen, heartbeat_count"
            ).gte("session_bucket", range_start.isoformat()).order("session_bucket", desc=True).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_sessions.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        except Exception as e:
            print(f"[daily_activity] {e}")
            break
    
    # تجميع حسب التاريخ
    from collections import defaultdict
    daily = defaultdict(lambda: {"buckets": 0, "students": set(), "by_grade": defaultdict(set)})
    
    for s in all_sessions:
        bucket_str = s.get("session_bucket", "")
        if not bucket_str:
            continue
        try:
            dt = datetime.fromisoformat(bucket_str.replace("Z", "+00:00"))
            if dt.tzinfo is None:
                dt = dt.replace(tzinfo=timezone.utc)
            date_key = dt.date().isoformat()
            daily[date_key]["buckets"] += 1
            sid = s.get("student_id")
            if sid:
                daily[date_key]["students"].add(sid)
                grade = s.get("grade") or "غير محدد"
                daily[date_key]["by_grade"][grade].add(sid)
        except Exception:
            continue
    
    # بناء قائمة آخر N أيام (حتى لو لا توجد بيانات)
    daily_stats = []
    for i in range(days):
        d = (today_start - timedelta(days=days - 1 - i)).date()
        key = d.isoformat()
        info = daily.get(key, {"buckets": 0, "students": set(), "by_grade": defaultdict(set)})
        unique = len(info["students"])
        total_min = info["buckets"] * 5
        daily_stats.append({
            "date": key,
            "day_name": d.strftime("%A"),  # سيُترجم في الواجهة
            "unique_students": unique,
            "total_minutes": total_min,
            "total_hours": round(total_min / 60, 1),
            "avg_minutes_per_student": round(total_min / unique, 1) if unique > 0 else 0,
            "grades_count": len(info["by_grade"]),
        })
    
    # طلاب اليوم
    today_key = today_start.date().isoformat()
    today_info = daily.get(today_key, {"buckets": 0, "students": set()})
    
    # نشطون الآن (آخر 5 دقائق)
    five_min_ago = now - timedelta(minutes=5)
    live_students = set()
    for s in all_sessions:
        try:
            seen = datetime.fromisoformat(s.get("last_seen", "").replace("Z", "+00:00"))
            if seen.tzinfo is None:
                seen = seen.replace(tzinfo=timezone.utc)
            if seen >= five_min_ago:
                sid = s.get("student_id")
                if sid:
                    live_students.add(sid)
        except Exception:
            continue
    
    return {
        "daily_stats": daily_stats,
        "today_live": len(live_students),
        "today_total_unique": len(today_info["students"]),
        "today_total_minutes": today_info["buckets"] * 5,
        "now": now.isoformat(),
    }


@app.get("/api/admin/stats/student_time/{student_id}")
async def get_student_time_breakdown(
    student_id: int,
    days: int = 30,
    admin = Depends(get_current_admin)
):
    """⏱️ تفاصيل وقت طالب معيّن خلال آخر N يوم"""
    from datetime import datetime, timezone, timedelta
    days = max(1, min(days, 90))
    range_start = datetime.now(timezone.utc) - timedelta(days=days)
    
    sessions = []
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("student_sessions").select(
                "session_bucket, heartbeat_count"
            ).eq("student_id", student_id).gte("session_bucket", range_start.isoformat()).order("session_bucket", desc=True).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            sessions.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        except Exception:
            break
    
    # تجميع يومي
    from collections import defaultdict
    daily = defaultdict(int)
    for s in sessions:
        try:
            dt = datetime.fromisoformat(s["session_bucket"].replace("Z", "+00:00"))
            daily[dt.date().isoformat()] += 5  # كل bucket = 5 دقائق
        except Exception:
            continue
    
    total_minutes = sum(daily.values())
    days_active   = len(daily)
    
    return {
        "student_id":      student_id,
        "total_minutes":   total_minutes,
        "total_hours":     round(total_minutes / 60, 1),
        "days_active":     days_active,
        "avg_per_day":     round(total_minutes / days_active, 1) if days_active else 0,
        "daily_breakdown": [{"date": k, "minutes": v} for k, v in sorted(daily.items())],
    }





@app.post("/api/student/update_profile")
async def student_update_profile(
    student_id: int = Form(...),
    full_name: str  = Form(default=""),
    school_name: str = Form(default=""),
    avatar_url: str = Form(default="")
):
    """
    تحديث بيانات ملف الطالب الشخصي (الاسم، المدرسة، الأفاتار).
    لا يمكن تغيير اسم المستخدم أو كلمة المرور أو الصف من هنا.
    """
    update_data = {}
    if full_name.strip():
        update_data["full_name"] = full_name.strip()[:120]
    if school_name.strip():
        update_data["school_name"] = school_name.strip()[:120]
    if avatar_url.strip():
        update_data["avatar_url"] = avatar_url.strip()[:500]

    if not update_data:
        raise HTTPException(status_code=400, detail="لا توجد بيانات للتحديث")

    try:
        res = supabase.table("students").update(update_data).eq("id", student_id).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="الطالب غير موجود")
        # نُرجع البيانات المحدّثة (بدون كلمة المرور)
        student = res.data[0]
        student.pop("password", None)
        return {"status": "success", "user": student}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل تحديث البيانات: {str(e)}")


@app.get("/api/config/public")
async def get_public_config():
    """يُعيد إعدادات عامة آمنة للواجهة (مثل Supabase URL و anon key للـ Realtime)"""
    return {
        "supabase_url": SUPABASE_URL,
        "supabase_key": SUPABASE_KEY,  # anon key — آمن للعرض العام
    }


@app.get("/api/leaderboard")
async def get_lb(top: int = 100, grade: str = ""):
    """
    🚀 لوحة الصدارة المحسّنة — أسرع 10x
    
    التحسينات:
    - يستخدم total_points من جدول students مباشرةً (بدلاً من جمع كل results)
    - cache لمدة 60 ثانية
    - يدعم top + فلتر الصف
    """
    cache_key = f"leaderboard:top{top}:g{grade}"
    cached = _cache.get(cache_key)
    if cached:
        return cached
    
    try:
        # 🚀 طريقة محسّنة: جلب من جدول students مباشرةً
        # (total_points مُحدّث تلقائياً عند كل نتيجة)
        query = supabase.table("students").select(
            "full_name, grade, total_points"
        ).gt("total_points", 0)
        
        if grade.strip():
            query = query.eq("grade", grade.strip())
        
        # ترتيب تنازلي + حد أقصى
        top = max(1, min(int(top), 500))
        res = query.order("total_points", desc=True).limit(top).execute()
        students = res.data or []
        
        result = [
            {
                "student_name": s.get("full_name", ""),
                "total_points": s.get("total_points", 0),
                "grade": s.get("grade", "")
            }
            for s in students if s.get("full_name")
        ]
        
        # cache 60 ثانية
        _cache.set(cache_key, result, ttl_seconds=60)
        return result
    except Exception as e:
        print(f"[leaderboard] error: {e}")
        # fallback للطريقة القديمة
        try:
            all_results = []
            offset = 0
            for _ in range(20):
                res = supabase.table("results").select("student_name, score, grade").range(offset, offset + 999).execute()
                batch = res.data or []
                if not batch: break
                all_results.extend(batch)
                if len(batch) < 1000: break
                offset += 1000
            
            lb = {}
            grades = {}
            for r in all_results:
                name = r.get("student_name") or ""
                if not name: continue
                lb[name] = lb.get(name, 0) + (r.get("score") or 0)
                if r.get("grade"): grades[name] = r["grade"]
            
            sorted_lb = sorted(lb.items(), key=lambda x: x[1], reverse=True)[:top]
            return [{"student_name": k, "total_points": v, "grade": grades.get(k, "")} for k, v in sorted_lb]
        except Exception:
            return []

@app.get("/api/parent/search/{query:path}")
async def parent_search(query: str, request: Request):
    """
    بحث ولي الأمر — يقبل **كود ولي الأمر فقط** (RS-XXXXX) للحفاظ على خصوصية الطلاب
    تم تقييد البحث بالـ ID/username/الاسم لمنع كشف بيانات الطلاب لأي شخص.
    Rate limit: 20 محاولة/دقيقة لكل IP لمنع brute force على الكودات
    """
    # ═══ Rate limiting لمنع تخمين الأكواد ═══
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=20, window_seconds=60, key_prefix="parent_search"):
        raise HTTPException(status_code=429, detail="⏳ محاولات كثيرة — انتظر دقيقة")

    clean = unquote(query).strip()
    if not clean:
        return {"found": False, "message": "يرجى إدخال كود ولي الأمر"}

    # ═══ نقبل فقط: parent_code (RS-XXXXX) أو الكود بدون البادئة ═══
    pc = clean.upper() if clean.upper().startswith("RS-") else f"RS-{clean.upper()}"

    # تحقق من الصيغة: RS- + 4-12 أحرف/أرقام
    import re as _re
    if not _re.match(r'^RS-[A-Z0-9]{4,12}$', pc):
        return {"found": False, "message": "صيغة الكود غير صحيحة. مثال: RS-AB12CD"}

    # 🔧 يُرجع parent_code أيضاً (حتى تستخدمه parent.html بدل التخمين)
    st = supabase.table("students").select(
        "id, full_name, grade, username, parent_code, created_at"
    ).eq("parent_code", pc).limit(1).execute()

    if not st.data:
        return {"found": False, "message": "لم يعثر على طالب بهذا الكود"}

    student = st.data[0]
    student.pop("password", None)

    # ═══ التاريخ الكامل ═══
    history = supabase.table("results").select(
        "id, lesson, score, total, timestamp"
    ).eq("student_id", student["id"]).order("timestamp", desc=True).limit(200).execute().data or []

    # ═══ إحصائيات الجلسات (وقت فعلي في المنصة) ═══
    from datetime import datetime, timezone, timedelta
    session_minutes_30d = 0
    session_minutes_7d = 0
    session_minutes_total = 0
    days_active_30d = set()
    last_seen_iso = None
    
    try:
        now = datetime.now(timezone.utc)
        cutoff_30d = now - timedelta(days=30)
        cutoff_7d = now - timedelta(days=7)
        
        # جلب كل الجلسات (pagination)
        all_sessions = []
        offset = 0
        for _ in range(20):
            sess = supabase.table("student_sessions").select(
                "session_bucket, last_seen"
            ).eq("student_id", student["id"]).order("session_bucket", desc=True).range(offset, offset + 999).execute()
            batch = sess.data or []
            if not batch:
                break
            all_sessions.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        for s in all_sessions:
            session_minutes_total += 5
            try:
                bucket_dt = datetime.fromisoformat(s["session_bucket"].replace("Z", "+00:00"))
                if bucket_dt.tzinfo is None:
                    bucket_dt = bucket_dt.replace(tzinfo=timezone.utc)
                if bucket_dt >= cutoff_30d:
                    session_minutes_30d += 5
                    days_active_30d.add(bucket_dt.date().isoformat())
                if bucket_dt >= cutoff_7d:
                    session_minutes_7d += 5
                # last seen
                if s.get("last_seen"):
                    if last_seen_iso is None or s["last_seen"] > last_seen_iso:
                        last_seen_iso = s["last_seen"]
            except Exception:
                continue
    except Exception:
        pass
    
    # ═══ ترتيب الطالب في الـ leaderboard ═══
    rank_in_grade = None
    rank_total = None
    grade_total_students = 0
    try:
        # جلب كل النتائج للحساب الدقيق
        all_results = []
        offset = 0
        for _ in range(50):
            rb = supabase.table("results").select(
                "student_id, student_name, score, grade"
            ).range(offset, offset + 999).execute()
            batch = rb.data or []
            if not batch:
                break
            all_results.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        # XP لكل طالب (إجمالي) + XP لكل طالب في نفس الصف
        from collections import defaultdict
        xp_total = defaultdict(int)
        xp_grade = defaultdict(int)  # student_id -> XP, للطلاب في نفس الصف
        
        for r in all_results:
            sid = r.get("student_id")
            if not sid: continue
            s_xp = r.get("score") or 0
            xp_total[sid] += s_xp
            if r.get("grade") == student.get("grade"):
                xp_grade[sid] += s_xp
        
        # الترتيب الإجمالي
        sorted_total = sorted(xp_total.items(), key=lambda x: x[1], reverse=True)
        for idx, (sid, _) in enumerate(sorted_total, 1):
            if sid == student["id"]:
                rank_total = idx
                break
        
        # الترتيب داخل نفس الصف
        sorted_grade = sorted(xp_grade.items(), key=lambda x: x[1], reverse=True)
        grade_total_students = len(sorted_grade)
        for idx, (sid, _) in enumerate(sorted_grade, 1):
            if sid == student["id"]:
                rank_in_grade = idx
                break
    except Exception:
        pass
    
    # ═══ إحصاء الإشعارات الحديثة ═══
    notifications_count = 0
    try:
        cutoff = (datetime.now(timezone.utc) - timedelta(days=30)).isoformat()
        nt_res = supabase.table("notifications").select("id", count="exact").gte(
            "created_at", cutoff
        ).limit(1).execute()
        notifications_count = nt_res.count if hasattr(nt_res, "count") else 0
    except Exception:
        pass
    
    # ═══ تحليل الأداء ═══
    total_tests = len(history)
    total_score = sum((r.get("score") or 0) for r in history)
    total_max   = sum((r.get("total") or 0) for r in history)
    avg_pct     = round((total_score / total_max * 100), 1) if total_max > 0 else 0
    
    # تجميع حسب الدرس (نقاط الضعف)
    lesson_stats = {}
    for r in history:
        lesson = r.get("lesson") or "غير محدد"
        if lesson not in lesson_stats:
            lesson_stats[lesson] = {"score": 0, "total": 0, "count": 0}
        lesson_stats[lesson]["score"] += (r.get("score") or 0)
        lesson_stats[lesson]["total"] += (r.get("total") or 0)
        lesson_stats[lesson]["count"] += 1
    
    weak_lessons = []
    strong_lessons = []
    for lesson, st_data in lesson_stats.items():
        if st_data["total"] > 0 and st_data["count"] >= 2:
            pct = (st_data["score"] / st_data["total"]) * 100
            entry = {"lesson": lesson, "pct": round(pct, 1), "count": st_data["count"]}
            if pct < 65:
                weak_lessons.append(entry)
            elif pct >= 85:
                strong_lessons.append(entry)
    
    weak_lessons.sort(key=lambda x: x["pct"])
    strong_lessons.sort(key=lambda x: -x["pct"])

    return {
        "found":   True,
        "student": student,
        "history": history,
        "stats": {
            "total_xp":             total_score,
            "total_tests":          total_tests,
            "avg_score_pct":        avg_pct,
            "minutes_30d":          session_minutes_30d,
            "hours_30d":            round(session_minutes_30d / 60, 1),
            "minutes_7d":           session_minutes_7d,
            "minutes_total":        session_minutes_total,
            "days_active_30d":      len(days_active_30d),
            "last_seen":            last_seen_iso,
            "rank_total":           rank_total,
            "rank_in_grade":        rank_in_grade,
            "grade_total_students": grade_total_students,
            "notifications_recent": notifications_count,
            "weak_lessons":         weak_lessons[:5],
            "strong_lessons":       strong_lessons[:5],
        },
    }


# ==========================================
# --- 10. نظام ساحة المبارزة المباشرة (WebSockets) ---
# ==========================================
class ArenaConnectionManager:
    def __init__(self):
        # طابور الانتظار مخصص حسب الصف الدراسي: {"الصف السادس": [{"ws": socket, "name": "أحمد"}], ...}
        self.waiting_players = {}
        # حفظ الغرف النشطة لتبادل النقاط: {"room_id": {"p1": p1, "p2": p2}}
        self.active_rooms = {}

    async def connect(self, websocket: WebSocket, student_name: str, grade: str):
        await websocket.accept()
        
        # إنشاء الطابور الخاص بالصف إذا لم يكن موجوداً
        if grade not in self.waiting_players:
            self.waiting_players[grade] = []
            
        self.waiting_players[grade].append({"ws": websocket, "name": student_name})
        await self.matchmake(grade)

    async def matchmake(self, grade: str):
        queue = self.waiting_players[grade]
        # إذا توفر طالبان من نفس الصف، يتم إنشاء الغرفة وبدء المعركة
        if len(queue) >= 2:
            p1 = queue.pop(0)
            p2 = queue.pop(0)
            
            room_id = f"room_{uuid.uuid4().hex[:8]}"
            
            # جلب أسئلة مخصصة لصف الطالبين من السحاب
            res = supabase.table("questions").select("*").eq("grade", grade).execute()
            all_qs = res.data if res.data else []
            
            # اختيار 5 أسئلة عشوائياً (أو أقل إذا لم يتوفر)
            if len(all_qs) >= 5:
                match_qs = random.sample(all_qs, 5)
            else:
                match_qs = all_qs 
                
            self.active_rooms[room_id] = {"p1": p1, "p2": p2}
            
            # إرسال إشارة بدء المعركة مع نفس الأسئلة لكلا الطالبين في نفس اللحظة
            try:
                await p1["ws"].send_json({
                    "type": "match_found", 
                    "opponent": p2["name"], 
                    "room_id": room_id,
                    "questions": match_qs
                })
                await p2["ws"].send_json({
                    "type": "match_found", 
                    "opponent": p1["name"], 
                    "room_id": room_id,
                    "questions": match_qs
                })
            except Exception as e:
                print(f"Error starting match: {e}")

    async def broadcast_score(self, room_id: str, sender_name: str, new_score: int):
        room = self.active_rooms.get(room_id)
        if room:
            target = room["p2"] if room["p1"]["name"] == sender_name else room["p1"]
            try:
                await target["ws"].send_json({
                    "type": "score_update", 
                    "opponent_score": new_score
                })
            except Exception as e:
                print(f"Error broadcasting score: {e}")

    async def disconnect(self, websocket: WebSocket, grade: str):
        # إزالة الطالب من طابور الانتظار إذا انسحب
        if grade in self.waiting_players:
            self.waiting_players[grade] = [p for p in self.waiting_players[grade] if p["ws"] != websocket]
        
        # إذا انسحب الطالب أثناء المعركة النشطة
        for room_id, room in list(self.active_rooms.items()):
            if room["p1"]["ws"] == websocket or room["p2"]["ws"] == websocket:
                target = room["p2"] if room["p1"]["ws"] == websocket else room["p1"]
                try:
                    await target["ws"].send_json({"type": "opponent_disconnected"})
                except:
                    pass
                del self.active_rooms[room_id]
                break

arena_manager = ArenaConnectionManager()

@app.websocket("/api/arena/ws/{student_name}/{grade}")
async def arena_websocket(websocket: WebSocket, student_name: str, grade: str, token: Optional[str] = None):
    """
    WebSocket Arena — يقبل توكن JWT اختياري كـ query parameter (?token=...)
    - لو التوكن موجود وصحيح: نستخدم بيانات الـ JWT (أكثر أماناً)
    - لو مش موجود: نقبل الاسم من URL للتوافق مع العملاء القدامى (سنُهمل هذا المسار لاحقاً)
    """
    # نقوم بفك تشفير الأسماء والصفوف التي قد تحتوي على مسافات
    clean_name = unquote(student_name)
    clean_grade = unquote(grade)
    verified = False

    # محاولة التحقق من التوكن إن وُجد
    if token:
        try:
            payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
            if payload.get("role") == "student":
                sid = payload.get("sub")
                if sid:
                    # جلب البيانات الحقيقية من DB لاستبدال الاسم/الصف القادم من URL
                    try:
                        st = supabase.table("students").select("full_name, grade, is_active")\
                             .eq("id", int(sid)).execute()
                        if st.data:
                            # منع الحسابات المعطلة من الدخول للساحة
                            if st.data[0].get("is_active") is False:
                                await websocket.close(code=1008, reason="account_disabled")
                                return
                            clean_name = st.data[0].get("full_name", clean_name)
                            clean_grade = st.data[0].get("grade", clean_grade)
                            verified = True
                    except Exception:
                        pass
        except JWTError:
            # توكن غير صالح — نرفض الاتصال
            await websocket.close(code=1008, reason="invalid_token")
            return

    await arena_manager.connect(websocket, clean_name, clean_grade)
    try:
        while True:
            data = await websocket.receive_json()
            if data.get("type") == "score_update":
                room_id = data.get("room_id")
                new_score = data.get("score")
                await arena_manager.broadcast_score(room_id, clean_name, new_score)
    except WebSocketDisconnect:
        await arena_manager.disconnect(websocket, clean_grade)


# ==========================================
# --- 12. نظام أكواد الاشتراك الملكي ---
# ==========================================
# ⚠️ يجب إنشاء جدول subscription_codes في Supabase بهذا SQL:
# CREATE TABLE subscription_codes (
#   id              BIGSERIAL PRIMARY KEY,
#   code            TEXT UNIQUE NOT NULL,
#   months          INTEGER NOT NULL DEFAULT 1,  -- -1 = دائم
#   note            TEXT DEFAULT '',
#   is_used         BOOLEAN DEFAULT FALSE,
#   used_at         TIMESTAMPTZ,
#   student_id      BIGINT REFERENCES students(id) ON DELETE SET NULL,
#   activated_by_student BIGINT REFERENCES students(id) ON DELETE SET NULL,
#   created_at      TIMESTAMPTZ DEFAULT NOW()
# );
# CREATE INDEX idx_sub_codes_code ON subscription_codes(code);
# ──────────────────────────────────────────────────────────────────

@app.post("/api/admin/subscription/generate")
async def generate_subscription_codes(
    months: int = Form(...),
    count: int = Form(default=10),
    student_id: Optional[int] = Form(default=None),
    note: str = Form(default=""),
    admin=Depends(get_current_admin)
):
    """توليد أكواد اشتراك آمنة وحفظها في قاعدة البيانات"""
    import hmac, hashlib
    codes = []
    for _ in range(min(count, 50)):  # حد أقصى 50 كوداً في المرة الواحدة
        raw_uuid = uuid.uuid4().hex.upper()
        prefix = "LIFE" if months == -1 else ("YEAR" if months == 12 else ("HALF" if months == 6 else ("QRTR" if months == 3 else "MNTH")))
        # HMAC-SHA256 للتحقق من صحة الكود لاحقاً
        sig = hmac.new(SECRET_KEY.encode(), raw_uuid.encode(), hashlib.sha256).hexdigest()[:8].upper()
        code = f"ME-{prefix}-{raw_uuid[:8]}-{raw_uuid[8:16]}-{sig}"
        
        data = {
            "code": code,
            "months": months,
            "note": note,
            "is_used": False,
            "student_id": student_id
        }
        result = supabase.table("subscription_codes").insert(data).execute()
        codes.append({"code": code, "id": result.data[0]["id"] if result.data else None})
    
    return {"status": "success", "codes": codes}


@app.post("/api/subscription/activate")
async def activate_subscription_code(
    code: str = Form(...),
    student_id: Optional[int] = Form(default=None)
):
    """تفعيل كود اشتراك من قِبَل الطالب"""
    code_upper = code.strip().upper()
    
    # البحث عن الكود في قاعدة البيانات
    res = supabase.table("subscription_codes").select("*").eq("code", code_upper).execute()
    
    if not res.data:
        raise HTTPException(status_code=404, detail="الكود غير موجود")
    
    entry = res.data[0]
    
    if entry.get("is_used"):
        raise HTTPException(status_code=400, detail="هذا الكود مستخدَم مسبقاً")
    
    # التحقق من أن الكود مربوط بطالب معين (إن وُجد)
    if entry.get("student_id") and student_id and entry["student_id"] != student_id:
        raise HTTPException(status_code=403, detail="هذا الكود مخصص لطالب آخر")
    
    # حساب تاريخ الانتهاء — باستخدام relativedelta لتجنب مشكلة ديسمبر
    months = entry.get("months", 1)
    if months == -1:
        expiry = None
    else:
        now = datetime.now(timezone.utc)
        expiry = (now + relativedelta(months=months)).isoformat()
    
    # تحديث الكود كمستخدَم
    update_data = {
        "is_used": True,
        "used_at": datetime.now(timezone.utc).isoformat(),
        "activated_by_student": student_id
    }
    supabase.table("subscription_codes").update(update_data).eq("id", entry["id"]).execute()
    
    return {
        "status": "success",
        "months": months,
        "expiry": expiry,
        "note": entry.get("note", "")
    }


@app.get("/api/admin/subscription/codes")
async def get_all_subscription_codes(admin=Depends(get_current_admin)):
    """جلب جميع أكواد الاشتراك للأدمن"""
    res = supabase.table("subscription_codes").select("*, students(full_name, grade)").order("id", desc=True).execute()
    return res.data if res.data else []


@app.delete("/api/admin/subscription/codes/{code_id}")
async def delete_subscription_code(code_id: int, admin=Depends(get_current_admin)):
    """حذف كود اشتراك"""
    supabase.table("subscription_codes").delete().eq("id", code_id).execute()
    return {"status": "success"}


@app.post("/api/admin/sub_codes/batch")
async def batch_save_sub_codes(request: Request, admin=Depends(get_current_admin)):
    """حفظ دفعة من الأكواد المولَّدة من واجهة الأدمن دفعة واحدة في Supabase"""
    body = await request.json()
    codes_list = body.get("codes", [])
    if not codes_list:
        raise HTTPException(status_code=400, detail="لا توجد أكواد للحفظ")

    rows = []
    for c in codes_list:
        rows.append({
            "code":       c.get("code", ""),
            "months":     int(c.get("months", 1)),
            "note":       c.get("note", ""),
            "student_id": int(c["studentId"]) if c.get("studentId") else None,
            "is_used":    False,
        })

    try:
        supabase.table("subscription_codes").insert(rows).execute()
    except Exception:
        saved = 0
        for row in rows:
            try:
                supabase.table("subscription_codes").insert(row).execute()
                saved += 1
            except Exception:
                pass
        return {"status": "partial", "saved": saved, "total": len(rows)}

    return {"status": "success", "saved": len(rows)}


@app.get("/api/admin/students")
async def get_all_students_admin(
    admin = Depends(get_current_admin),
    page: int = 0,           # 0 = جلب الكل (للتوافق مع الكود القديم)
    page_size: int = 100,
    grade: str = "",
    search: str = "",
):
    """
    🚀 قائمة الطلاب للأدمن — مع pagination + cache + بحث
    
    - page=0 → يجلب الكل (للتوافق العكسي)
    - page>=1 → يستخدم pagination
    """
    # cache key
    cache_key = f"students:p{page}:s{page_size}:g{grade}:q{search}"
    cached = _cache.get(cache_key)
    if cached:
        return cached
    
    try:
        if page == 0:
            # 🔄 الوضع القديم: جلب الكل (للتوافق)
            res = supabase.table("students").select(
                "id, full_name, grade, username, parent_code"
            ).order("full_name").execute()
            result = res.data or []
            _cache.set(cache_key, result, ttl_seconds=120)
            return result
        
        # 🚀 الوضع الجديد: pagination
        page, page_size, offset = _paginate_params(page, page_size, max_size=200)
        
        query = supabase.table("students").select(
            "id, full_name, grade, username, parent_code, total_points, created_at",
            count="exact"
        )
        
        if grade.strip():
            query = query.eq("grade", grade.strip())
        if search.strip():
            query = query.ilike("full_name", f"%{search.strip()}%")
        
        res = query.order("full_name").range(offset, offset + page_size - 1).execute()
        
        result = _paginate_response(
            data=res.data or [],
            total=res.count or 0,
            page=page,
            page_size=page_size
        )
        
        _cache.set(cache_key, result, ttl_seconds=120)
        return result
    except Exception as e:
        print(f"[students] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


def _count_active_recent(results: list, days: int = 7) -> int:
    """يحسب عدد الطلاب الفريدين الذين شاركوا في آخر N أيام"""
    from datetime import datetime, timezone, timedelta
    cutoff = datetime.now(timezone.utc) - timedelta(days=days)
    active_ids = set()
    for r in results:
        ts = r.get("timestamp", "")
        if not ts:
            continue
        try:
            if isinstance(ts, str):
                ts = ts.replace("Z", "+00:00")
                dt = datetime.fromisoformat(ts)
                if dt.tzinfo is None:
                    dt = dt.replace(tzinfo=timezone.utc)
                if dt >= cutoff:
                    sid = r.get("student_id")
                    if sid:
                        active_ids.add(sid)
        except Exception:
            continue
    return len(active_ids)


@app.post("/api/student/forgot_password")
async def student_forgot_password(
    request: Request,
    username: str       = Form(...),
    parent_code: str    = Form(...),
    new_password: str   = Form(...),
):
    """
    استعادة كلمة مرور طالب — يتطلب كود ولي الأمر للأمان
    Rate-limit: 5 محاولات / دقيقة لكل IP
    """
    import re
    
    # 🛡️ rate limiting يدوي
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=5, window_seconds=60, key_prefix="forgot_pwd"):
        raise HTTPException(status_code=429, detail="محاولات كثيرة — انتظر دقيقة وحاول مرة أخرى")
    
    # تحقق من شكل الكود
    if not re.match(r'^RS-[A-Z0-9]{4,12}$', parent_code.strip().upper()):
        raise HTTPException(status_code=400, detail="كود ولي الأمر غير صالح")
    if len(new_password) < 6:
        raise HTTPException(status_code=400, detail="كلمة المرور قصيرة (6+ أحرف)")
    if len(username) < 3:
        raise HTTPException(status_code=400, detail="اسم المستخدم غير صالح")

    try:
        # ابحث عن الطالب باسم المستخدم + كود ولي الأمر
        res = supabase.table("students").select("id, username, parent_code").eq(
            "username", username.strip()
        ).eq("parent_code", parent_code.strip().upper()).limit(1).execute()
        
        if not res.data:
            raise HTTPException(
                status_code=404,
                detail="لم نعثر على حساب بهذا الاسم وكود ولي الأمر — تحقق من البيانات"
            )

        student_id = res.data[0]["id"]
        # حدّث كلمة المرور (مع hashing)
        hashed = hash_password(new_password)
        supabase.table("students").update({
            "password": hashed
        }).eq("id", student_id).execute()
        
        return {"status": "success", "message": "تم تحديث كلمة المرور"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



# ═══════════════════════════════════════════════════════════════
# 👨‍👩‍👦 PARENT SESSIONS — تسجيل دخول دائم لأولياء الأمور
# ═══════════════════════════════════════════════════════════════

def _generate_parent_token() -> str:
    """توليد token عشوائي قوي للجلسة"""
    import secrets
    return secrets.token_urlsafe(48)


@app.post("/api/parent/login")
async def parent_login(
    request: Request,
    parent_code: str  = Form(...),
    parent_name: str  = Form(default=""),
    parent_email: str = Form(default=""),
    parent_phone: str = Form(default=""),
):
    """
    تسجيل دخول ولي الأمر — يُنشئ session token دائم (180 يوم)
    + يحفظ بيانات ولي الأمر للتواصل لاحقاً
    """
    # rate limit
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=20, window_seconds=60, key_prefix="parent_login"):
        raise HTTPException(status_code=429, detail="محاولات كثيرة — انتظر دقيقة")
    
    pc = parent_code.strip().upper()
    if not pc.startswith("RS-"):
        pc = f"RS-{pc}"
    
    import re as _re
    if not _re.match(r'^RS-[A-Z0-9]{4,12}$', pc):
        raise HTTPException(status_code=400, detail="صيغة الكود غير صحيحة")
    
    # ابحث عن الطالب
    try:
        st = supabase.table("students").select(
            "id, full_name, grade, parent_email, parent_name"
        ).eq("parent_code", pc).limit(1).execute()
        if not st.data:
            raise HTTPException(status_code=404, detail="لم نعثر على طالب بهذا الكود")
        student = st.data[0]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")
    
    # حدّث بيانات ولي الأمر إن قُدّمت (للتواصل لاحقاً)
    update_data = {}
    if parent_name.strip() and not student.get("parent_name"):
        update_data["parent_name"] = parent_name.strip()[:200]
    if parent_email.strip():
        # تحقق صيغة الإيميل
        if _re.match(r'^[^@\s]+@[^@\s]+\.[^@\s]+$', parent_email.strip()):
            update_data["parent_email"] = parent_email.strip().lower()[:200]
    if parent_phone.strip():
        # رقم بسيط (أي طول معقول)
        clean_phone = _re.sub(r'[^\d+]', '', parent_phone.strip())
        if 6 <= len(clean_phone) <= 20:
            update_data["parent_phone"] = clean_phone
    
    if update_data:
        try:
            supabase.table("students").update(update_data).eq("id", student["id"]).execute()
        except Exception:
            pass  # نُكمل حتى لو فشل التحديث
    
    # أنشئ session token
    token = _generate_parent_token()
    try:
        supabase.table("parent_sessions").insert({
            "student_id":    student["id"],
            "parent_code":   pc,
            "session_token": token,
            "parent_name":   parent_name.strip()[:200] or None,
            "parent_email":  (parent_email.strip().lower() or None) if parent_email else None,
            "parent_phone":  update_data.get("parent_phone"),
            "user_agent":    str(request.headers.get("user-agent", ""))[:300],
        }).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل إنشاء الجلسة: {str(e)[:150]}")
    
    return {
        "status":  "success",
        "token":   token,
        "expires_in_days": 180,
        "student": {
            "id":         student["id"],
            "full_name":  student.get("full_name"),
            "grade":      student.get("grade"),
        }
    }


@app.get("/api/parent/auto_login/{token}")
async def parent_auto_login(token: str, request: Request):
    """
    دخول تلقائي بـ token — يُستخدم عند تحميل الصفحة
    """
    if not token or len(token) < 30:
        raise HTTPException(status_code=400, detail="token غير صالح")
    
    try:
        sess = supabase.table("parent_sessions").select(
            "id, student_id, parent_code, parent_name, parent_email, parent_phone, expires_at"
        ).eq("session_token", token).limit(1).execute()
        
        if not sess.data:
            raise HTTPException(status_code=401, detail="جلسة غير موجودة — سجّل الدخول مرة أخرى")
        
        session = sess.data[0]
        
        # تحقق من الصلاحية
        from datetime import datetime, timezone
        if session.get("expires_at"):
            try:
                exp = datetime.fromisoformat(session["expires_at"].replace("Z", "+00:00"))
                if exp.tzinfo is None:
                    exp = exp.replace(tzinfo=timezone.utc)
                if exp < datetime.now(timezone.utc):
                    # احذف الجلسة المنتهية
                    supabase.table("parent_sessions").delete().eq("id", session["id"]).execute()
                    raise HTTPException(status_code=401, detail="انتهت صلاحية الجلسة")
            except HTTPException:
                raise
            except Exception:
                pass
        
        # حدّث last_used_at
        try:
            supabase.table("parent_sessions").update({
                "last_used_at": datetime.now(timezone.utc).isoformat()
            }).eq("id", session["id"]).execute()
        except Exception:
            pass
        
        # 🆕 جلب username الطالب أيضاً (للتوافق مع نظام v2)
        student_username = None
        try:
            stud_res = supabase.table("students").select("username, full_name").eq(
                "id", session["student_id"]
            ).limit(1).execute()
            if stud_res.data:
                student_username = stud_res.data[0].get("username")
        except Exception:
            pass
        
        return {
            "valid": True,
            "parent_code": session["parent_code"],  # للتوافق العكسي
            "username": student_username,  # 🆕 للنظام الجديد
            "parent_name": session.get("parent_name"),
            "parent_email": session.get("parent_email"),
            "parent_phone": session.get("parent_phone"),
            "student_id": session["student_id"],
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



# ═══════════════════════════════════════════════════════════════
# 👨‍👩 PARENT ACCESS V2 — دخول مبسّط: username + phone
# ═══════════════════════════════════════════════════════════════
@app.post("/api/parent/login_v2")
async def parent_login_v2(
    request: Request,
    username: str     = Form(...),
    parent_phone: str = Form(...),
    parent_name: str  = Form(default=""),
    parent_email: str = Form(default=""),
):
    """
    🔐 دخول ولي الأمر بـ:
       • username الطالب (يعرفه الطالب وأسرته)
       • parent_phone رقم هاتف ولي الأمر (سرّي)
    
    منطق الذكاء:
    - أول دخول: يُسجّل رقم الهاتف للطالب (يصبح "كلمة السر")
    - الدخول التالي: يجب أن يطابق الرقم المُسجّل
    """
    # rate limit
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=15, window_seconds=60, key_prefix="parent_login_v2"):
        raise HTTPException(status_code=429, detail="محاولات كثيرة — انتظر دقيقة")
    
    uname = username.strip()
    if not uname:
        raise HTTPException(status_code=400, detail="اسم المستخدم مطلوب")
    
    phone_clean = _normalize_phone(parent_phone)
    if len(phone_clean) < 6 or len(phone_clean) > 20:
        raise HTTPException(status_code=400, detail="رقم الهاتف غير صحيح")
    
    # ابحث عن الطالب
    try:
        st_res = supabase.table("students").select(
            "id, full_name, grade, username, parent_phone, parent_name, parent_email"
        ).eq("username", uname).limit(1).execute()
        
        if not st_res.data:
            raise HTTPException(status_code=404, detail="لم نعثر على طالب باسم المستخدم هذا")
        
        student = st_res.data[0]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")
    
    saved_phone = (student.get("parent_phone") or "").strip()
    
    # ─── السيناريو 1: أول دخول لهذا الطالب — لا يوجد هاتف مُسجّل ───
    if not saved_phone:
        # سجّل الهاتف + الاسم (تأمين الحساب لأول مرة)
        update_data = {
            "parent_phone": phone_clean,
        }
        if parent_name.strip():
            update_data["parent_name"] = parent_name.strip()[:200]
        if parent_email.strip():
            import re as _re
            em = parent_email.strip().lower()
            if _re.match(r'^[^@\s]+@[^@\s]+\.[^@\s]+$', em):
                update_data["parent_email"] = em[:200]
        
        try:
            supabase.table("students").update(update_data).eq("id", student["id"]).execute()
            student.update(update_data)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"فشل التسجيل: {str(e)[:150]}")
        
        is_first_time = True
    else:
        # ─── السيناريو 2: دخول لاحق — يجب أن يطابق الهاتف ───
        if not _phone_matches(saved_phone, phone_clean):
            raise HTTPException(
                status_code=403,
                detail="رقم الهاتف لا يطابق المُسجّل لهذا الطالب — تواصل مع الأستاذ إن نسيت"
            )
        
        # حدّث الاسم/الإيميل إن قُدّما
        update_data = {}
        if parent_name.strip() and not student.get("parent_name"):
            update_data["parent_name"] = parent_name.strip()[:200]
        if parent_email.strip():
            import re as _re
            em = parent_email.strip().lower()
            if _re.match(r'^[^@\s]+@[^@\s]+\.[^@\s]+$', em):
                update_data["parent_email"] = em[:200]
        
        if update_data:
            try:
                supabase.table("students").update(update_data).eq("id", student["id"]).execute()
            except Exception:
                pass
        
        is_first_time = False
    
    # ─── إنشاء session token ───
    token = _generate_parent_token()
    try:
        supabase.table("parent_sessions").insert({
            "student_id":    student["id"],
            "parent_code":   "USERNAME_AUTH",  # مَعلَم: لا يستخدم parent_code
            "session_token": token,
            "parent_name":   (parent_name.strip() or student.get("parent_name") or "")[:200] or None,
            "parent_email":  ((parent_email.strip().lower() if parent_email else None) or student.get("parent_email")),
            "parent_phone":  phone_clean,
            "user_agent":    str(request.headers.get("user-agent", ""))[:300],
        }).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل إنشاء الجلسة: {str(e)[:150]}")
    
    return {
        "status": "success",
        "token": token,
        "expires_in_days": 180,
        "is_first_time": is_first_time,
        "student": {
            "id":         student["id"],
            "full_name":  student.get("full_name"),
            "grade":      student.get("grade"),
            "username":   student.get("username"),
        }
    }


@app.get("/api/parent/search_by_username/{username}")
async def parent_search_by_username(username: str, request: Request):
    """
    📖 جلب بيانات الطالب الكاملة + الإحصاءات (بعد تسجيل الدخول الناجح)
    تعتمد على token في المتصفح للمصادقة (auto_login يتم قبلها)
    """
    # rate limit
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=60, window_seconds=60, key_prefix="parent_search_uname"):
        raise HTTPException(status_code=429, detail="محاولات كثيرة")
    
    uname = username.strip()
    if not uname:
        raise HTTPException(status_code=400, detail="اسم المستخدم مطلوب")
    
    # ابحث
    st = supabase.table("students").select(
        "id, full_name, grade, username, parent_code, created_at"
    ).eq("username", uname).limit(1).execute()
    
    if not st.data:
        return {"found": False, "message": "لم يعثر على طالب باسم المستخدم هذا"}
    
    student = st.data[0]
    student.pop("password", None)
    
    # ═══ نفس منطق parent_search القديم (الإحصاءات الكاملة) ═══
    history = supabase.table("results").select(
        "id, lesson, score, total, timestamp"
    ).eq("student_id", student["id"]).order("timestamp", desc=True).limit(200).execute().data or []

    from datetime import datetime, timezone, timedelta
    session_minutes_30d = 0
    session_minutes_7d = 0
    session_minutes_total = 0
    days_active_30d = set()
    last_seen_iso = None
    
    try:
        now = datetime.now(timezone.utc)
        cutoff_30d = now - timedelta(days=30)
        cutoff_7d = now - timedelta(days=7)
        all_sessions = []
        offset = 0
        for _ in range(20):
            sess = supabase.table("student_sessions").select(
                "session_bucket, last_seen"
            ).eq("student_id", student["id"]).order("session_bucket", desc=True).range(offset, offset + 999).execute()
            batch = sess.data or []
            if not batch: break
            all_sessions.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        
        for s in all_sessions:
            session_minutes_total += 5
            try:
                bucket_dt = datetime.fromisoformat(s["session_bucket"].replace("Z", "+00:00"))
                if bucket_dt.tzinfo is None:
                    bucket_dt = bucket_dt.replace(tzinfo=timezone.utc)
                if bucket_dt >= cutoff_30d:
                    session_minutes_30d += 5
                    days_active_30d.add(bucket_dt.date().isoformat())
                if bucket_dt >= cutoff_7d:
                    session_minutes_7d += 5
                if s.get("last_seen"):
                    if last_seen_iso is None or s["last_seen"] > last_seen_iso:
                        last_seen_iso = s["last_seen"]
            except Exception:
                continue
    except Exception:
        pass
    
    # الترتيب
    rank_in_grade = None
    grade_total_students = 0
    try:
        all_results = []
        offset = 0
        for _ in range(50):
            rb = supabase.table("results").select(
                "student_id, score, grade"
            ).range(offset, offset + 999).execute()
            batch = rb.data or []
            if not batch: break
            all_results.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        
        from collections import defaultdict
        xp_grade = defaultdict(int)
        for r in all_results:
            sid = r.get("student_id")
            if not sid: continue
            if r.get("grade") == student.get("grade"):
                xp_grade[sid] += (r.get("score") or 0)
        
        sorted_grade = sorted(xp_grade.items(), key=lambda x: x[1], reverse=True)
        grade_total_students = len(sorted_grade)
        for idx, (sid, _) in enumerate(sorted_grade, 1):
            if sid == student["id"]:
                rank_in_grade = idx
                break
    except Exception:
        pass
    
    total_tests = len(history)
    total_score = sum((r.get("score") or 0) for r in history)
    total_max = sum((r.get("total") or 0) for r in history)
    avg_pct = round((total_score / total_max * 100), 1) if total_max > 0 else 0
    
    lesson_stats = {}
    for r in history:
        lesson = r.get("lesson") or "غير محدد"
        if lesson not in lesson_stats:
            lesson_stats[lesson] = {"score": 0, "total": 0, "count": 0}
        lesson_stats[lesson]["score"] += (r.get("score") or 0)
        lesson_stats[lesson]["total"] += (r.get("total") or 0)
        lesson_stats[lesson]["count"] += 1
    
    weak_lessons = []
    strong_lessons = []
    for lesson, st_data in lesson_stats.items():
        if st_data["total"] > 0 and st_data["count"] >= 2:
            pct = (st_data["score"] / st_data["total"]) * 100
            entry = {"lesson": lesson, "pct": round(pct, 1), "count": st_data["count"]}
            if pct < 65: weak_lessons.append(entry)
            elif pct >= 85: strong_lessons.append(entry)
    weak_lessons.sort(key=lambda x: x["pct"])
    strong_lessons.sort(key=lambda x: -x["pct"])
    
    return {
        "found":   True,
        "student": student,
        "history": history,
        "stats": {
            "total_xp":             total_score,
            "total_tests":          total_tests,
            "avg_score_pct":        avg_pct,
            "minutes_30d":          session_minutes_30d,
            "hours_30d":            round(session_minutes_30d / 60, 1),
            "minutes_7d":           session_minutes_7d,
            "minutes_total":        session_minutes_total,
            "days_active_30d":      len(days_active_30d),
            "last_seen":            last_seen_iso,
            "rank_in_grade":        rank_in_grade,
            "grade_total_students": grade_total_students,
            "weak_lessons":         weak_lessons[:5],
            "strong_lessons":       strong_lessons[:5],
        },
    }


@app.post("/api/parent/logout")
async def parent_logout(token: str = Form(...)):
    """تسجيل خروج وحذف الجلسة"""
    try:
        supabase.table("parent_sessions").delete().eq("session_token", token).execute()
    except Exception:
        pass
    return {"status": "ok"}


@app.post("/api/parent/update_profile")
async def parent_update_profile(
    token: str        = Form(...),
    parent_name: str  = Form(default=""),
    parent_email: str = Form(default=""),
    parent_phone: str = Form(default=""),
):
    """تحديث بيانات ولي الأمر"""
    import re as _re
    
    # تحقق من الـ token
    sess = supabase.table("parent_sessions").select(
        "student_id, id"
    ).eq("session_token", token).limit(1).execute()
    
    if not sess.data:
        raise HTTPException(status_code=401, detail="جلسة غير موجودة")
    
    student_id = sess.data[0]["student_id"]
    session_id = sess.data[0]["id"]
    
    update_student = {}
    update_session = {}
    
    if parent_name.strip():
        v = parent_name.strip()[:200]
        update_student["parent_name"] = v
        update_session["parent_name"] = v
    
    if parent_email.strip():
        em = parent_email.strip().lower()
        if _re.match(r'^[^@\s]+@[^@\s]+\.[^@\s]+$', em):
            update_student["parent_email"] = em[:200]
            update_session["parent_email"] = em[:200]
        else:
            raise HTTPException(status_code=400, detail="صيغة البريد الإلكتروني غير صحيحة")
    
    if parent_phone.strip():
        clean_phone = _re.sub(r'[^\d+]', '', parent_phone.strip())
        if 6 <= len(clean_phone) <= 20:
            update_student["parent_phone"] = clean_phone
            update_session["parent_phone"] = clean_phone
    
    try:
        if update_student:
            supabase.table("students").update(update_student).eq("id", student_id).execute()
        if update_session:
            supabase.table("parent_sessions").update(update_session).eq("id", session_id).execute()
        return {"status": "updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


# ═══════════════════════════════════════════════════════════════
# 📧 ADMIN — إدارة بريد أولياء الأمور
# ═══════════════════════════════════════════════════════════════

@app.post("/api/admin/students/{student_id}/regenerate_parent_code")
async def regenerate_parent_code(student_id: int, admin = Depends(get_current_admin)):
    """إعادة توليد كود ولي الأمر لطالب معيّن"""
    new_code = _generate_unique_parent_code()
    try:
        # احذف جلسات ولي الأمر السابقة (الكود القديم لم يعد صالحاً)
        supabase.table("parent_sessions").delete().eq("student_id", student_id).execute()
        # حدّث الكود
        supabase.table("students").update({"parent_code": new_code}).eq("id", student_id).execute()
        return {"status": "success", "parent_code": new_code}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])



# ═══════════════════════════════════════════════════════════════
# 🛡️ ADMIN — إصلاح أكواد ولي الأمر جماعياً (للحالات الطارئة)
# ═══════════════════════════════════════════════════════════════
@app.post("/api/admin/parent_codes/fix_all")
async def admin_fix_all_parent_codes(admin = Depends(get_current_admin)):
    """
    🔧 إصلاح أكواد ولي الأمر لكل الطلاب بدون كود أو بكود مكرر.
    يُولّد كوداً فريداً لكل طالب بحاجة، ويترك الأكواد الصحيحة الفريدة كما هي.
    """
    try:
        # 1. جلب كل الطلاب (pagination)
        all_students = []
        offset = 0
        for _ in range(20):
            res = supabase.table("students").select(
                "id, full_name, parent_code"
            ).order("id").range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_students.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        # 2. تشخيص: من بحاجة لكود جديد؟
        import re as _re
        code_to_students = {}  # parent_code -> [student_ids]
        for s in all_students:
            pc = (s.get("parent_code") or "").strip().upper()
            if pc:
                code_to_students.setdefault(pc, []).append(s["id"])
        
        # الطلاب الذين يحتاجون كوداً جديداً:
        need_new_code = []
        for s in all_students:
            pc = (s.get("parent_code") or "").strip()
            sid = s["id"]
            
            # شرط 1: لا كود أصلاً
            if not pc or pc.lower() == "null":
                need_new_code.append(sid)
                continue
            
            # شرط 2: صيغة غير صالحة
            if not _re.match(r'^RS-[A-Z0-9]{4,12}$', pc.upper()):
                need_new_code.append(sid)
                continue
            
            # شرط 3: كود مكرر — نُبقي على أقدم طالب فقط
            siblings = code_to_students.get(pc.upper(), [])
            if len(siblings) > 1:
                # أقدم طالب (أصغر id) يحتفظ بالكود
                if sid != min(siblings):
                    need_new_code.append(sid)
        
        # 3. توليد كود فريد لكل طالب يحتاج
        used_codes = set()
        for s in all_students:
            pc = (s.get("parent_code") or "").strip().upper()
            if pc and s["id"] not in need_new_code:
                used_codes.add(pc)
        
        import secrets
        chars = "ABCDEFGHJKLMNPQRSTUVWXYZ23456789"
        
        def gen_code():
            for _ in range(100):
                candidate = "RS-" + "".join(secrets.choice(chars) for _ in range(6))
                if candidate not in used_codes:
                    used_codes.add(candidate)
                    return candidate
            # fallback
            return "RS-" + "".join(secrets.choice(chars) for _ in range(8))
        
        # 4. تطبيق التحديثات
        updated_count = 0
        failed_count = 0
        
        for sid in need_new_code:
            new_code = gen_code()
            try:
                supabase.table("students").update({
                    "parent_code": new_code
                }).eq("id", sid).execute()
                
                # احذف جلسات parent_sessions القديمة (الكود تغيّر)
                try:
                    supabase.table("parent_sessions").delete().eq("student_id", sid).execute()
                except Exception:
                    pass
                
                updated_count += 1
            except Exception:
                failed_count += 1
        
        return {
            "status": "success",
            "total_students": len(all_students),
            "needed_fix": len(need_new_code),
            "updated": updated_count,
            "failed": failed_count,
            "diagnosis": {
                "no_code_or_invalid": sum(1 for s in all_students 
                    if not (s.get("parent_code") or "").strip() 
                    or not _re.match(r'^RS-[A-Z0-9]{4,12}$', (s.get("parent_code") or "").strip().upper())),
                "duplicates": sum(1 for codes in code_to_students.values() if len(codes) > 1),
            }
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:300]}")


@app.get("/api/admin/parent_codes/diagnose")
async def admin_diagnose_parent_codes(admin = Depends(get_current_admin)):
    """
    🔍 تشخيص حالة أكواد ولي الأمر — قبل تشغيل الإصلاح
    """
    try:
        all_students = []
        offset = 0
        for _ in range(20):
            res = supabase.table("students").select(
                "id, full_name, grade, parent_code"
            ).order("id").range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_students.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        import re as _re
        code_to_count = {}
        no_code = []
        invalid_format = []
        
        for s in all_students:
            pc = (s.get("parent_code") or "").strip()
            if not pc or pc.lower() == "null":
                no_code.append({"id": s["id"], "name": s.get("full_name"), "grade": s.get("grade")})
            elif not _re.match(r'^RS-[A-Z0-9]{4,12}$', pc.upper()):
                invalid_format.append({"id": s["id"], "name": s.get("full_name"), "code": pc})
            else:
                code_to_count[pc.upper()] = code_to_count.get(pc.upper(), 0) + 1
        
        duplicates = []
        for code, cnt in code_to_count.items():
            if cnt > 1:
                # احصل على كل الطلاب بهذا الكود
                students_with_code = [
                    {"id": s["id"], "name": s.get("full_name"), "grade": s.get("grade")}
                    for s in all_students
                    if (s.get("parent_code") or "").strip().upper() == code
                ]
                duplicates.append({"code": code, "count": cnt, "students": students_with_code})
        
        return {
            "total_students": len(all_students),
            "no_code_count": len(no_code),
            "invalid_format_count": len(invalid_format),
            "duplicate_groups_count": len(duplicates),
            "students_in_duplicates": sum(d["count"] for d in duplicates),
            "no_code_sample": no_code[:10],
            "invalid_format_sample": invalid_format[:10],
            "duplicates_sample": duplicates[:5],
            "is_healthy": len(no_code) == 0 and len(invalid_format) == 0 and len(duplicates) == 0,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



@app.post("/api/admin/students/{student_id}/reset_parent_login")
async def reset_parent_login(student_id: int, admin = Depends(get_current_admin)):
    """
    🔄 إعادة تعيين دخول ولي الأمر لطالب معيّن.
    يحذف parent_phone (يصبح أول دخول جديد) + يحذف الجلسات القديمة.
    """
    try:
        # احذف الجلسات
        supabase.table("parent_sessions").delete().eq("student_id", student_id).execute()
        # امسح parent_phone (سيُسجَّل من جديد عند أول دخول)
        supabase.table("students").update({
            "parent_phone": None,
        }).eq("id", student_id).execute()
        return {"status": "success", "message": "تم — يستطيع ولي الأمر التسجيل من جديد"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.get("/api/admin/parent_emails")
async def admin_list_parent_emails(admin = Depends(get_current_admin)):
    """
    قائمة كل أولياء الأمور الذين أدخلوا إيميلاتهم — لإرسال نشرات
    """
    try:
        # جلب من students مع pagination
        all_students = []
        offset = 0
        for _ in range(20):
            res = supabase.table("students").select(
                "id, full_name, grade, parent_email, parent_name, parent_phone"
            ).neq("parent_email", "").not_.is_("parent_email", "null").range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_students.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        # رتّب وحدد الفريد
        unique_emails = {}
        for s in all_students:
            em = (s.get("parent_email") or "").strip().lower()
            if not em:
                continue
            if em not in unique_emails:
                unique_emails[em] = {
                    "email": em,
                    "parent_name": s.get("parent_name") or "",
                    "parent_phone": s.get("parent_phone") or "",
                    "students": [],
                }
            unique_emails[em]["students"].append({
                "id": s["id"],
                "name": s.get("full_name"),
                "grade": s.get("grade"),
            })
        
        return {
            "total_emails": len(unique_emails),
            "total_students_with_email": len(all_students),
            "parents": list(unique_emails.values()),
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.get("/api/admin/parent_emails/export")
async def admin_export_parent_emails(admin = Depends(get_current_admin)):
    """
    تصدير قائمة الإيميلات كـ CSV
    """
    from fastapi.responses import PlainTextResponse
    try:
        all_students = []
        offset = 0
        for _ in range(20):
            res = supabase.table("students").select(
                "id, full_name, grade, parent_email, parent_name, parent_phone"
            ).neq("parent_email", "").not_.is_("parent_email", "null").range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_students.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        # CSV
        lines = ["البريد,اسم ولي الأمر,الهاتف,اسم الطالب,الصف"]
        for s in all_students:
            em = (s.get("parent_email") or "").strip()
            if not em:
                continue
            row = [
                em,
                (s.get("parent_name") or "").replace(",", " "),
                (s.get("parent_phone") or "").replace(",", " "),
                (s.get("full_name") or "").replace(",", " "),
                (s.get("grade") or "").replace(",", " "),
            ]
            lines.append(",".join(row))
        
        csv_content = "\n".join(lines)
        # BOM لـ Excel ليقرأ العربية
        csv_with_bom = "\ufeff" + csv_content
        
        return PlainTextResponse(
            content=csv_with_bom,
            headers={
                "Content-Type": "text/csv; charset=utf-8",
                "Content-Disposition": 'attachment; filename="parent_emails.csv"'
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])



# ═══════════════════════════════════════════════════════════════
# 🏅 ACHIEVEMENTS / BADGES SYSTEM
# ═══════════════════════════════════════════════════════════════

# تعريف كل الشارات المتاحة
ACHIEVEMENT_DEFINITIONS = {
    # ─── سلاسل الإنجاز ───
    "streak_3":  {"name": "🔥 سلسلة 3 أيام",  "icon": "🔥", "color": "#e67e22", "tier": "bronze", "xp": 30,  "description": "دخول 3 أيام متتالية"},
    "streak_7":  {"name": "⚡ أسبوع متواصل",   "icon": "⚡", "color": "#f39c12", "tier": "silver", "xp": 75,  "description": "دخول 7 أيام متتالية"},
    "streak_30": {"name": "👑 شهر كامل",       "icon": "👑", "color": "#d4af37", "tier": "gold",   "xp": 300, "description": "دخول 30 يوم متتالي"},
    
    # ─── الإتقان ───
    "perfect_first":   {"name": "💯 أول علامة كاملة",     "icon": "💯", "color": "#2ecc71", "tier": "bronze", "xp": 50,  "description": "أول تحدي بـ 100%"},
    "perfect_5":       {"name": "🎯 خمس مرات مثالية",     "icon": "🎯", "color": "#27ae60", "tier": "silver", "xp": 100, "description": "5 تحديات بـ 100%"},
    "perfect_20":      {"name": "🏆 سيد الإتقان",         "icon": "🏆", "color": "#16a085", "tier": "gold",   "xp": 250, "description": "20 تحدي بـ 100%"},
    
    # ─── السرعة ───
    "speed_demon":     {"name": "💨 السريع",              "icon": "💨", "color": "#3498db", "tier": "silver", "xp": 80,  "description": "تحدي بكل الإجابات صحيحة في أقل من دقيقتين"},
    
    # ─── الكمية ───
    "tests_10":        {"name": "📚 المثابر",            "icon": "📚", "color": "#9b59b6", "tier": "bronze", "xp": 50,  "description": "إنجاز 10 تحديات"},
    "tests_50":        {"name": "🎓 المجتهد",            "icon": "🎓", "color": "#8e44ad", "tier": "silver", "xp": 150, "description": "إنجاز 50 تحدي"},
    "tests_100":       {"name": "⚔️ المحارب",            "icon": "⚔️", "color": "#6c3483", "tier": "gold",   "xp": 400, "description": "إنجاز 100 تحدي"},
    "tests_250":       {"name": "🏛️ الأسطورة",          "icon": "🏛️", "color": "#bf953f", "tier": "legendary", "xp": 1000, "description": "إنجاز 250 تحدي"},
    
    # ─── الوقت ───
    "early_bird":      {"name": "🌅 البكور",              "icon": "🌅", "color": "#f1c40f", "tier": "bronze", "xp": 40, "description": "دخول قبل الساعة 7 صباحاً"},
    "night_owl":       {"name": "🦉 ساهر الليل",          "icon": "🦉", "color": "#34495e", "tier": "bronze", "xp": 40, "description": "دخول بعد منتصف الليل"},
    
    # ─── XP الكلي ───
    "xp_500":          {"name": "🥉 برونزي",              "icon": "🥉", "color": "#cd7f32", "tier": "bronze", "xp": 25, "description": "جمع 500 XP"},
    "xp_1500":         {"name": "🥈 فضي",                 "icon": "🥈", "color": "#c0c0c0", "tier": "silver", "xp": 75, "description": "جمع 1500 XP"},
    "xp_5000":         {"name": "🥇 ذهبي",                "icon": "🥇", "color": "#ffd700", "tier": "gold",   "xp": 200, "description": "جمع 5000 XP"},
    "xp_10000":        {"name": "💎 ماسي",                "icon": "💎", "color": "#b9f2ff", "tier": "legendary", "xp": 500, "description": "جمع 10000 XP"},
    
    # ─── متنوع ───
    "first_test":      {"name": "🎉 البداية",            "icon": "🎉", "color": "#e74c3c", "tier": "bronze", "xp": 20, "description": "إنجاز أول تحدي"},
    "comeback":        {"name": "💪 العائد القوي",       "icon": "💪", "color": "#16a085", "tier": "silver", "xp": 60, "description": "العودة بعد غياب أسبوع+"},
    "weekend_warrior": {"name": "🛡️ محارب الإجازة",       "icon": "🛡️", "color": "#27ae60", "tier": "bronze", "xp": 35, "description": "دخول في عطلة نهاية الأسبوع"},
}


def _grant_badge(student_id: int, badge_id: str) -> dict:
    """يمنح شارة لطالب — يتجاهل لو ممنوحة سابقاً"""
    if badge_id not in ACHIEVEMENT_DEFINITIONS:
        return {"granted": False, "reason": "غير معروفة"}
    
    badge = ACHIEVEMENT_DEFINITIONS[badge_id]
    try:
        # تحقق من عدم التكرار
        existing = supabase.table("achievements").select("id").eq(
            "student_id", student_id
        ).eq("badge_id", badge_id).limit(1).execute()
        
        if existing.data:
            return {"granted": False, "reason": "ممنوحة سابقاً"}
        
        # امنح الشارة
        supabase.table("achievements").insert({
            "student_id":  student_id,
            "badge_id":    badge_id,
            "badge_name":  badge["name"],
            "badge_icon":  badge["icon"],
            "badge_color": badge["color"],
            "badge_tier":  badge["tier"],
            "description": badge["description"],
            "xp_reward":   badge["xp"],
        }).execute()
        
        # امنح XP إضافي مكافأة
        try:
            st_res = supabase.table("students").select("total_points").eq("id", student_id).limit(1).execute()
            if st_res.data:
                cur_xp = st_res.data[0].get("total_xp") or 0
                supabase.table("students").update({
                    "total_xp": cur_xp + badge["xp"]
                }).eq("id", student_id).execute()
        except Exception:
            pass
        
        # أرسل push notification
        try:
            _push_to_student(
                student_id,
                f"{badge['icon']} شارة جديدة!",
                f"حصلت على شارة \"{badge['name']}\" — +{badge['xp']} XP",
                url="/student",
                tag=f"badge_{badge_id}"
            )
        except Exception:
            pass
        
        return {"granted": True, "badge": badge}
    except Exception as e:
        return {"granted": False, "reason": str(e)[:100]}


def _check_and_grant_achievements(student_id: int, context: dict = None) -> list:
    """
    يفحص كل الشروط الممكنة ويمنح الشارات المستحقة.
    context: {"score": 10, "total": 10, "duration_seconds": 90, ...}
    """
    granted = []
    context = context or {}
    
    try:
        # 1. اجلب بيانات الطالب
        st_res = supabase.table("students").select(
            "id, total_points, created_at"
        ).eq("id", student_id).limit(1).execute()
        
        if not st_res.data:
            return []
        
        student = st_res.data[0]
        total_xp = student.get("total_points") or 0
        
        # 2. اجلب نتائج الطالب
        results_res = supabase.table("results").select(
            "score, total, timestamp"
        ).eq("student_id", student_id).execute()
        results = results_res.data or []
        total_tests = len(results)
        perfect_count = sum(1 for r in results if r.get("total") and r.get("score") == r.get("total"))
        
        # 3. اجلب الجلسات للسلاسل
        from datetime import datetime, timezone, timedelta
        sessions_res = supabase.table("student_sessions").select(
            "session_bucket"
        ).eq("student_id", student_id).order("session_bucket", desc=True).limit(500).execute()
        
        active_dates = set()
        for s in (sessions_res.data or []):
            try:
                bucket_dt = datetime.fromisoformat(s["session_bucket"].replace("Z", "+00:00"))
                if bucket_dt.tzinfo is None:
                    bucket_dt = bucket_dt.replace(tzinfo=timezone.utc)
                active_dates.add(bucket_dt.date())
            except Exception:
                continue
        
        # حساب أطول سلسلة
        sorted_dates = sorted(active_dates, reverse=True)
        max_streak = 0
        if sorted_dates:
            current_streak = 1
            for i in range(1, len(sorted_dates)):
                if (sorted_dates[i-1] - sorted_dates[i]).days == 1:
                    current_streak += 1
                    max_streak = max(max_streak, current_streak)
                else:
                    current_streak = 1
            max_streak = max(max_streak, current_streak)
        
        # ─── الفحوصات ───
        candidates = []
        
        # السلاسل
        if max_streak >= 3:  candidates.append("streak_3")
        if max_streak >= 7:  candidates.append("streak_7")
        if max_streak >= 30: candidates.append("streak_30")
        
        # عدد التحديات
        if total_tests >= 1:   candidates.append("first_test")
        if total_tests >= 10:  candidates.append("tests_10")
        if total_tests >= 50:  candidates.append("tests_50")
        if total_tests >= 100: candidates.append("tests_100")
        if total_tests >= 250: candidates.append("tests_250")
        
        # العلامات الكاملة
        if perfect_count >= 1:  candidates.append("perfect_first")
        if perfect_count >= 5:  candidates.append("perfect_5")
        if perfect_count >= 20: candidates.append("perfect_20")
        
        # XP
        if total_xp >= 500:   candidates.append("xp_500")
        if total_xp >= 1500:  candidates.append("xp_1500")
        if total_xp >= 5000:  candidates.append("xp_5000")
        if total_xp >= 10000: candidates.append("xp_10000")
        
        # السرعة (من السياق)
        if (context.get("duration_seconds") and context.get("duration_seconds") < 120
            and context.get("score") and context.get("total")
            and context.get("score") == context.get("total")):
            candidates.append("speed_demon")
        
        # الوقت الحالي (للبكور والساهر)
        now = datetime.now(timezone.utc)
        # نُحوّل للتوقيت المحلي (عُمان UTC+4)
        local_h = (now.hour + 4) % 24
        if 4 <= local_h < 7:
            candidates.append("early_bird")
        elif local_h >= 0 and local_h < 4:
            candidates.append("night_owl")
        
        # عطلة نهاية الأسبوع (الجمعة/السبت في عُمان)
        if now.weekday() in (4, 5):  # Friday=4, Saturday=5
            candidates.append("weekend_warrior")
        
        # العائد القوي (أكثر من أسبوع غياب ثم عودة)
        if len(sorted_dates) >= 2:
            gap = (sorted_dates[0] - sorted_dates[1]).days
            if gap >= 7:
                candidates.append("comeback")
        
        # امنح كل المرشحات (التكرار محظور بالـ UNIQUE)
        for badge_id in candidates:
            result = _grant_badge(student_id, badge_id)
            if result.get("granted"):
                granted.append(result["badge"])
    
    except Exception as e:
        print(f"[achievements] error: {e}")
    
    return granted


@app.get("/api/student/{student_id}/achievements")
async def get_student_achievements(student_id: int, request: Request):
    """قائمة الشارات الممنوحة لطالب"""
    try:
        # rate limit
        ip = request.client.host if request.client else "unknown"
        if _is_rate_limited(ip, max_calls=30, window_seconds=60, key_prefix="ach_get"):
            raise HTTPException(status_code=429, detail="محاولات كثيرة")
        
        res = supabase.table("achievements").select("*").eq(
            "student_id", student_id
        ).order("earned_at", desc=True).execute()
        
        # احسب الإحصائيات
        earned = res.data or []
        total_xp_from_badges = sum(b.get("xp_reward", 0) for b in earned)
        
        by_tier = {"bronze": 0, "silver": 0, "gold": 0, "legendary": 0}
        for b in earned:
            tier = b.get("badge_tier", "bronze")
            if tier in by_tier:
                by_tier[tier] += 1
        
        return {
            "earned": earned,
            "total_count": len(earned),
            "total_available": len(ACHIEVEMENT_DEFINITIONS),
            "total_xp_from_badges": total_xp_from_badges,
            "by_tier": by_tier,
            "all_badges": [
                {"id": k, **v, "earned": any(e["badge_id"] == k for e in earned)}
                for k, v in ACHIEVEMENT_DEFINITIONS.items()
            ]
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.post("/api/student/{student_id}/check_achievements")
async def trigger_check_achievements(
    student_id: int,
    request: Request,
    score: int = Form(default=0),
    total: int = Form(default=0),
    duration_seconds: int = Form(default=0),
):
    """يُستدعى بعد كل تحدي — يفحص ويمنح الشارات الجديدة"""
    granted = _check_and_grant_achievements(student_id, {
        "score": score,
        "total": total,
        "duration_seconds": duration_seconds,
    })
    return {"new_badges": granted}


# ═══════════════════════════════════════════════════════════════
# 🚨 ISSUE REPORTS — بلاغات الأسئلة
# ═══════════════════════════════════════════════════════════════
@app.post("/api/student/report_issue")
async def report_issue(
    request: Request,
    student_id: int   = Form(...),
    student_name: str = Form(default=""),
    question_id: int  = Form(default=0),
    issue_type: str   = Form(...),
    description: str  = Form(default=""),
):
    """طالب يبلّغ عن مشكلة في سؤال"""
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=10, window_seconds=300, key_prefix="report_issue"):
        raise HTTPException(status_code=429, detail="محاولات كثيرة")
    
    if issue_type not in ("wrong_answer", "typo", "unclear", "other"):
        raise HTTPException(status_code=400, detail="نوع غير صالح")
    
    try:
        supabase.table("issue_reports").insert({
            "student_id":   student_id,
            "student_name": student_name[:200],
            "question_id":  question_id if question_id > 0 else None,
            "issue_type":   issue_type,
            "description":  description.strip()[:1000],
        }).execute()
        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.get("/api/admin/issue_reports")
async def list_issue_reports(
    status: str = "open",
    admin = Depends(get_current_admin)
):
    """قائمة البلاغات"""
    try:
        q = supabase.table("issue_reports").select("*")
        if status and status != "all":
            q = q.eq("status", status)
        res = q.order("created_at", desc=True).limit(200).execute()
        return res.data or []
    except Exception as e:
        return {"error": str(e)[:200]}


@app.put("/api/admin/issue_reports/{report_id}")
async def update_issue_report(
    report_id: int,
    status: str       = Form(...),
    admin_notes: str  = Form(default=""),
    admin = Depends(get_current_admin)
):
    """تحديث حالة البلاغ"""
    if status not in ("open", "reviewing", "resolved", "dismissed"):
        raise HTTPException(status_code=400, detail="حالة غير صالحة")
    
    update = {"status": status, "admin_notes": admin_notes[:500]}
    if status in ("resolved", "dismissed"):
        from datetime import datetime, timezone
        update["resolved_at"] = datetime.now(timezone.utc).isoformat()
    
    try:
        supabase.table("issue_reports").update(update).eq("id", report_id).execute()
        return {"status": "updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


# ═══════════════════════════════════════════════════════════════
# 🎓 TUTORING BOOKINGS — حجز الدروس الخصوصية
# ═══════════════════════════════════════════════════════════════
@app.post("/api/tutoring/book")
async def book_tutoring(
    request: Request,
    student_id: int   = Form(...),
    student_name: str = Form(default=""),
    parent_phone: str = Form(default=""),
    booking_date: str = Form(...),
    booking_time: str = Form(...),
    topic: str        = Form(default=""),
    notes: str        = Form(default=""),
    method: str       = Form(default="whatsapp"),
):
    """طالب يحجز جلسة استشارة"""
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=5, window_seconds=300, key_prefix="tutoring_book"):
        raise HTTPException(status_code=429, detail="حجوزات كثيرة")
    
    if method not in ("whatsapp", "google_meet", "in_person"):
        method = "whatsapp"
    
    try:
        res = supabase.table("tutoring_bookings").insert({
            "student_id":   student_id,
            "student_name": student_name[:200],
            "parent_phone": _normalize_phone(parent_phone) if parent_phone else None,
            "booking_date": booking_date,
            "booking_time": booking_time,
            "topic":        topic.strip()[:200],
            "notes":        notes.strip()[:1000],
            "method":       method,
        }).execute()
        
        # أرسل إشعار للأدمن (push للطالب نفسه)
        try:
            _push_to_student(
                student_id,
                "📅 تم استلام طلب الحجز",
                f"موعدك: {booking_date} {booking_time} — في انتظار تأكيد الأستاذ",
                tag="tutoring_pending"
            )
        except Exception:
            pass
        
        return {"status": "success", "booking_id": res.data[0]["id"] if res.data else None}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.get("/api/student/{student_id}/bookings")
async def get_student_bookings(student_id: int):
    """حجوزات طالب معيّن"""
    try:
        res = supabase.table("tutoring_bookings").select("*").eq(
            "student_id", student_id
        ).order("booking_date", desc=True).limit(50).execute()
        return res.data or []
    except Exception as e:
        return {"error": str(e)[:200]}


@app.get("/api/admin/tutoring_bookings")
async def admin_list_bookings(
    status: str = "all",
    admin = Depends(get_current_admin)
):
    """كل الحجوزات للأدمن"""
    try:
        q = supabase.table("tutoring_bookings").select("*")
        if status and status != "all":
            q = q.eq("status", status)
        res = q.order("booking_date", desc=False).limit(200).execute()
        return res.data or []
    except Exception as e:
        return {"error": str(e)[:200]}


@app.put("/api/admin/tutoring_bookings/{booking_id}")
async def update_booking(
    booking_id: int,
    status: str           = Form(...),
    meeting_link: str     = Form(default=""),
    admin = Depends(get_current_admin)
):
    """تحديث حالة الحجز"""
    if status not in ("pending", "confirmed", "completed", "cancelled"):
        raise HTTPException(status_code=400, detail="حالة غير صالحة")
    
    update = {"status": status}
    if meeting_link:
        update["meeting_link"] = meeting_link.strip()[:500]
    
    try:
        # اجلب الحجز للحصول على student_id
        bk = supabase.table("tutoring_bookings").select("student_id, booking_date, booking_time").eq("id", booking_id).limit(1).execute()
        
        supabase.table("tutoring_bookings").update(update).eq("id", booking_id).execute()
        
        # أرسل إشعار للطالب
        if bk.data:
            sid = bk.data[0]["student_id"]
            d = bk.data[0].get("booking_date", "")
            t = bk.data[0].get("booking_time", "")
            if status == "confirmed":
                msg = f"📅 موعدك مؤكد: {d} {t}"
                if meeting_link:
                    msg += f"\nالرابط: {meeting_link}"
                _push_to_student(sid, "✅ تم تأكيد الحجز", msg, tag="tutoring_confirmed", url="/student")
            elif status == "cancelled":
                _push_to_student(sid, "❌ تم إلغاء الحجز", f"تواصل مع الأستاذ للتفاصيل", tag="tutoring_cancelled")
        
        return {"status": "updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


# ═══════════════════════════════════════════════════════════════
# 🎬 SHORT VIDEOS — الفيديوهات القصيرة
# ═══════════════════════════════════════════════════════════════
@app.post("/api/admin/videos")
async def upload_short_video(
    title: str         = Form(...),
    description: str   = Form(default=""),
    grade: str         = Form(default=""),
    lesson: str        = Form(default=""),
    video_url: str     = Form(...),
    thumbnail_url: str = Form(default=""),
    duration_sec: int  = Form(default=0),
    admin = Depends(get_current_admin)
):
    """إضافة فيديو قصير (يقبل YouTube/Vimeo URL أو رابط فيديو مباشر)"""
    if not video_url.startswith(("http://", "https://")):
        raise HTTPException(status_code=400, detail="رابط الفيديو غير صالح")
    
    try:
        res = supabase.table("short_videos").insert({
            "title":         title.strip()[:200],
            "description":   description.strip()[:1000],
            "grade":         grade[:100],
            "lesson":        lesson[:200],
            "video_url":     video_url[:1000],
            "thumbnail_url": thumbnail_url[:1000] or None,
            "duration_sec":  max(0, duration_sec),
        }).execute()
        return {"status": "success", "id": res.data[0]["id"] if res.data else None}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.get("/api/videos")
async def list_short_videos(grade: str = "", lesson: str = ""):
    """قائمة الفيديوهات (للطلاب)"""
    try:
        q = supabase.table("short_videos").select("*").eq("is_published", True)
        if grade:  q = q.eq("grade", grade)
        if lesson: q = q.eq("lesson", lesson)
        res = q.order("uploaded_at", desc=True).limit(100).execute()
        return res.data or []
    except Exception as e:
        return {"error": str(e)[:200]}


@app.post("/api/videos/{video_id}/view")
async def increment_video_view(video_id: int):
    """زيادة عداد المشاهدات"""
    try:
        cur = supabase.table("short_videos").select("views_count").eq("id", video_id).limit(1).execute()
        if cur.data:
            new_count = (cur.data[0].get("views_count") or 0) + 1
            supabase.table("short_videos").update({"views_count": new_count}).eq("id", video_id).execute()
        return {"status": "ok"}
    except Exception:
        return {"status": "error"}


@app.post("/api/videos/{video_id}/like")
async def toggle_video_like(
    video_id: int,
    student_id: int = Form(...),
):
    """إعجاب/إلغاء إعجاب فيديو"""
    try:
        # تحقق من وجود إعجاب
        existing = supabase.table("video_likes").select("id").eq(
            "video_id", video_id
        ).eq("student_id", student_id).limit(1).execute()
        
        if existing.data:
            # احذف
            supabase.table("video_likes").delete().eq("id", existing.data[0]["id"]).execute()
            # نقص العداد
            cur = supabase.table("short_videos").select("likes_count").eq("id", video_id).limit(1).execute()
            if cur.data:
                new_count = max(0, (cur.data[0].get("likes_count") or 0) - 1)
                supabase.table("short_videos").update({"likes_count": new_count}).eq("id", video_id).execute()
            return {"liked": False}
        else:
            # أضف
            supabase.table("video_likes").insert({
                "video_id": video_id,
                "student_id": student_id,
            }).execute()
            cur = supabase.table("short_videos").select("likes_count").eq("id", video_id).limit(1).execute()
            if cur.data:
                new_count = (cur.data[0].get("likes_count") or 0) + 1
                supabase.table("short_videos").update({"likes_count": new_count}).eq("id", video_id).execute()
            return {"liked": True}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.delete("/api/admin/videos/{video_id}")
async def delete_video(video_id: int, admin = Depends(get_current_admin)):
    try:
        supabase.table("video_likes").delete().eq("video_id", video_id).execute()
        supabase.table("short_videos").delete().eq("id", video_id).execute()
        return {"status": "deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


# ═══════════════════════════════════════════════════════════════
# 💾 BACKUP & EXPORT
# ═══════════════════════════════════════════════════════════════
@app.get("/api/admin/export/{table_name}")
async def export_table_csv(
    table_name: str,
    admin = Depends(get_current_admin)
):
    """تصدير جدول كـ CSV"""
    from fastapi.responses import PlainTextResponse
    
    # قائمة الجداول المسموح بها
    allowed = ["students", "results", "questions", "exams", "achievements", 
               "tutoring_bookings", "issue_reports", "admin_tasks"]
    if table_name not in allowed:
        raise HTTPException(status_code=400, detail="جدول غير مسموح")
    
    try:
        # جلب كل البيانات (pagination)
        all_rows = []
        offset = 0
        for _ in range(100):
            res = supabase.table(table_name).select("*").range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_rows.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        
        if not all_rows:
            csv_content = "(لا توجد بيانات)"
        else:
            # CSV header من المفاتيح
            keys = list(all_rows[0].keys())
            lines = [",".join(keys)]
            for row in all_rows:
                row_values = []
                for k in keys:
                    v = row.get(k, "")
                    if v is None:
                        v = ""
                    s = str(v).replace('"', '""').replace(",", "،").replace("\n", " ")
                    row_values.append(f'"{s}"')
                lines.append(",".join(row_values))
            csv_content = "\n".join(lines)
        
        # BOM للعربية في Excel
        csv_with_bom = "\ufeff" + csv_content
        
        return PlainTextResponse(
            content=csv_with_bom,
            headers={
                "Content-Type": "text/csv; charset=utf-8",
                "Content-Disposition": f'attachment; filename="{table_name}_export.csv"'
            }
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])



# ═══════════════════════════════════════════════════════════════
# 🤖 AI ASSISTANT — المساعد الذكي للطالب
# يستخدم Google Gemini (مجاني) — يحتاج GEMINI_API_KEY في env vars
# ═══════════════════════════════════════════════════════════════

# قواعد للحماية وتحسين الجودة
AI_SYSTEM_PROMPT = """أنت مساعد رياضيات ذكي للطلاب في إمبراطورية الرياضيات.

⚜️ قواعدك:
1. تتكلم بالعربية الفصحى البسيطة المناسبة لطلاب المدارس
2. تشرح خطوة بخطوة بصبر وحب
3. لا تعطي الإجابة النهائية مباشرة — اشرح الطريقة وادفع الطالب للتفكير
4. استخدم أمثلة بسيطة من حياة الطلاب
5. إذا كان السؤال خارج الرياضيات، اعتذر بلطف
6. كن مشجّعاً ومحفّزاً — استخدم إيموجي مناسبة بشكل مدروس
7. لا تتجاوز 200 كلمة في الإجابة
8. لا تتحدث عن نفسك إلا لو سُئلت
"""



# ════════════════════════════════════════════════════════════
# 🤖 AI VISUAL GENERATOR — توليد أشكال بصرية بالذكاء الاصطناعي
# ════════════════════════════════════════════════════════════
@app.post("/api/admin/ai/generate_visual")
async def ai_generate_visual(
    description: str = Form(...),
    admin = Depends(get_current_admin)
):
    """🤖 يحوّل وصفاً عربياً لإعدادات شكل بصري (JSON)"""
    description = (description or "").strip()
    if not description:
        raise HTTPException(status_code=400, detail="الوصف فارغ")
    if len(description) > 500:
        raise HTTPException(status_code=400, detail="الوصف طويل جداً")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="الذكاء الاصطناعي غير مفعّل")
    
    # System prompt - شامل لكل الأشكال المتاحة
    system_prompt = """أنت مساعد ذكي لمعلم رياضيات. مهمتك تحويل وصف عربي لشكل رياضي إلى JSON بإعدادات الشكل.

# الأشكال المتاحة (TOOL):
## ✏️ أشكال ثنائية الأبعاد (2D):
- triangle (مثلث)
- rectangle (مستطيل)
- circle (دائرة)
- circle-sector (قطاع دائري)
- polygon (مضلع منتظم)
- fraction-bar (شريط كسر)
- grid (شبكة)
- number-line (محور أعداد)
- angle (زاوية)

## 🧊 أشكال ثلاثية الأبعاد (3D):
- cube (مكعب)
- cuboid (متوازي مستطيلات)
- pyramid-3 (هرم ثلاثي)
- pyramid-4 (هرم رباعي)
- pyramid-5 (هرم خماسي)
- pyramid-6 (هرم سداسي)
- pyramid-8 (هرم ثماني)
- cylinder (أسطوانة)
- cone (مخروط)
- prism-tri (موشور/منشور ثلاثي)
- prism-quad (موشور/منشور رباعي)
- prism-penta (موشور/منشور خماسي)
- prism-hexa (موشور/منشور سداسي)
- prism-octa (موشور/منشور ثماني)

## 📊 رسوم بيانية وإحصائية:
- bar-chart (رسم بياني بالأعمدة)
- pie-chart (رسم دائري)
- box-plot (مخطط الصندوق)
- scatter (مخطط الانتشار)

## 🧮 مفاهيم متقدمة:
- venn (مخطط فن للمجموعات)
- vector (متجه)

# الحقول التفصيلية لكل شكل:

## 2D Shapes
- triangle: {tri_type: "right" أو "equilateral" أو "isosceles" أو "scalene", tri_a, tri_b, tri_c}
- rectangle: {rect_w, rect_h, rect_color}
- circle: {circle_r, circle_color}
- circle-sector: {sec_angle (درجة), sec_r, sec_color}
- polygon: {poly_sides (3-12), poly_side, poly_color}
- fraction-bar: {frac_denom (المقام), frac_num (البسط), frac_color}
- grid: {grid_cols, grid_rows}
- number-line: {nl_min, nl_max, nl_points (نص مفصول بفاصلة "1,3,5"), nl_color}
- angle: {ang_deg (الزاوية بالدرجات), ang_label (رمز مثل θ أو α), ang_color}

## 3D Shapes
- cube: {cube_side, cube_color}
- cuboid: {cb_l (طول), cb_w (عرض), cb_h (ارتفاع), cb_color}
- pyramid-3, pyramid-4, pyramid-5, pyramid-6, pyramid-8: {py_base (طول ضلع القاعدة), py_height, py_color}
- cylinder: {cyl_r (نصف القطر), cyl_h (الارتفاع), cyl_color}
- cone: {cone_r, cone_h, cone_color}
- prism-tri, prism-quad, prism-penta, prism-hexa, prism-octa: {pr_base (طول ضلع القاعدة), pr_len (طول الموشور), pr_color}

## Charts
- bar-chart: {bar_data ("اسم1:قيمة1,اسم2:قيمة2"), bar_color, bar_title}
- pie-chart: {pie_data ("اسم1:نسبة1,اسم2:نسبة2"), pie_title}
- box-plot: {box_data ("3,7,8,12,14"), box_title, box_color}
- scatter: {sc_data ("1,2; 3,5; 4,4"), sc_title, sc_trend (true/false), sc_color}

## Advanced
- venn: {venn_count (2 أو 3), venn_a, venn_b, venn_c, venn_shade (true/false)}
- vector: {vec_name (مثل "v"), vec_x, vec_y, vec_color}

# الألوان (hex):
- #3498db (أزرق)، #e74c3c (أحمر)، #2ecc71 (أخضر)، #f39c12 (برتقالي)
- #9b59b6 (بنفسجي)، #16a085 (تركواز)، #e67e22 (برتقالي داكن)
- #a855f7 (بنفسجي فاتح)، #ec4899 (وردي)، #00bfff (سماوي)

# قواعد مهمة:
1. أرجِع JSON فقط - لا تشرح أو تُضف نص قبل/بعد JSON
2. التنسيق: {"tool": "TOOL_NAME", "settings": {...}}
3. لو طلب المستخدم شكلاً غير موجود، اختر الأقرب
4. الكلمات المرادفة: "موشور" = "منشور" = "prism"، "مثلث قائم" → tri_type="right"
5. إذا لم يحدد لون، اختر لوناً مناسباً للسياق
6. إذا لم يحدد أبعاد، استخدم قيم افتراضية معقولة

# أمثلة شاملة لكل الأشكال:

الوصف: "موشور خماسي ضلع قاعدته 4 وطوله 8"
الرد: {"tool":"prism-penta","settings":{"pr_base":4,"pr_len":8,"pr_color":"#2ecc71"}}

الوصف: "منشور سداسي"
الرد: {"tool":"prism-hexa","settings":{"pr_base":5,"pr_len":10,"pr_color":"#16a085"}}

الوصف: "هرم خماسي قاعدته 6 وارتفاعه 9"
الرد: {"tool":"pyramid-5","settings":{"py_base":6,"py_height":9,"py_color":"#f59e0b"}}

الوصف: "مكعب طول ضلعه 5"
الرد: {"tool":"cube","settings":{"cube_side":5,"cube_color":"#3498db"}}

الوصف: "مثلث قائم الأضلاع 3 و 4 و 5"
الرد: {"tool":"triangle","settings":{"tri_type":"right","tri_a":3,"tri_b":4,"tri_c":5}}

الوصف: "مثلث متساوي الأضلاع"
الرد: {"tool":"triangle","settings":{"tri_type":"equilateral","tri_a":6,"tri_b":6,"tri_c":6}}

الوصف: "أسطوانة نصف قطرها 4 وارتفاعها 12"
الرد: {"tool":"cylinder","settings":{"cyl_r":4,"cyl_h":12,"cyl_color":"#16a085"}}

الوصف: "مخروط ارتفاعه 10"
الرد: {"tool":"cone","settings":{"cone_r":4,"cone_h":10,"cone_color":"#f39c12"}}

الوصف: "متوازي مستطيلات 8×5×3"
الرد: {"tool":"cuboid","settings":{"cb_l":8,"cb_w":5,"cb_h":3,"cb_color":"#e74c3c"}}

الوصف: "زاوية 45 درجة"
الرد: {"tool":"angle","settings":{"ang_deg":45,"ang_label":"θ","ang_color":"#a855f7"}}

الوصف: "زاوية قائمة"
الرد: {"tool":"angle","settings":{"ang_deg":90,"ang_label":"∟","ang_color":"#a855f7"}}

الوصف: "قطاع دائري بزاوية 90 درجة"
الرد: {"tool":"circle-sector","settings":{"sec_angle":90,"sec_r":5,"sec_color":"#e74c3c"}}

الوصف: "دائرة نصف قطرها 7"
الرد: {"tool":"circle","settings":{"circle_r":7,"circle_color":"#3498db"}}

الوصف: "كسر 3/4"
الرد: {"tool":"fraction-bar","settings":{"frac_denom":4,"frac_num":3,"frac_color":"#9b59b6"}}

الوصف: "خط الأعداد من -10 إلى 10 مع علامات عند -5, 0, 5"
الرد: {"tool":"number-line","settings":{"nl_min":-10,"nl_max":10,"nl_points":"-5,0,5","nl_color":"#e74c3c"}}

الوصف: "سداسي منتظم"
الرد: {"tool":"polygon","settings":{"poly_sides":6,"poly_side":5,"poly_color":"#00bfff"}}

الوصف: "ثماني منتظم"
الرد: {"tool":"polygon","settings":{"poly_sides":8,"poly_side":4,"poly_color":"#00bfff"}}

الوصف: "رسم بياني لدرجات أحمد 90 وسارة 85 ومحمد 70"
الرد: {"tool":"bar-chart","settings":{"bar_data":"أحمد:90,سارة:85,محمد:70","bar_color":"#3498db","bar_title":"الدرجات"}}

الوصف: "رسم دائري للفواكه: تفاح 40%، موز 30%، عنب 30%"
الرد: {"tool":"pie-chart","settings":{"pie_data":"تفاح:40,موز:30,عنب:30","pie_title":"مبيعات الفواكه"}}

الوصف: "مخطط صندوق للدرجات 3,7,8,12,13,14,18,21,23,27"
الرد: {"tool":"box-plot","settings":{"box_data":"3,7,8,12,13,14,18,21,23,27","box_title":"الدرجات","box_color":"#ec4899"}}

الوصف: "مخطط انتشار للنقاط (1,2),(3,5),(4,4),(6,7)"
الرد: {"tool":"scatter","settings":{"sc_data":"1,2; 3,5; 4,4; 6,7","sc_title":"البيانات","sc_trend":true,"sc_color":"#ec4899"}}

الوصف: "مخطط فن لمجموعتين A و B متقاطعتين"
الرد: {"tool":"venn","settings":{"venn_count":2,"venn_a":"A","venn_b":"B","venn_shade":true}}

الوصف: "مخطط فن لثلاث مجموعات"
الرد: {"tool":"venn","settings":{"venn_count":3,"venn_a":"A","venn_b":"B","venn_c":"C","venn_shade":false}}

الوصف: "متجه (5, 3)"
الرد: {"tool":"vector","settings":{"vec_name":"v","vec_x":5,"vec_y":3,"vec_color":"#e91e63"}}
"""
    
    full_prompt = system_prompt + f"\n\nالوصف: \"{description}\"\nالرد:"
    
    import httpx, json as _json
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
        payload = {
            "contents": [{"parts": [{"text": full_prompt}]}],
            "generationConfig": {
                "temperature": 0.3,  # دقة عالية للJSON
                "maxOutputTokens": 400,
                "topP": 0.9,
                "responseMimeType": "application/json"
            }
        }
        
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.post(url, json=payload)
            
            # fallback لو الأساسي فشل
            if r.status_code in (404, 429):
                fallback_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-lite:generateContent?key={api_key}"
                r = await client.post(fallback_url, json=payload)
            
            r.raise_for_status()
            data = r.json()
            
            # استخراج النص
            candidates = data.get("candidates", [])
            if not candidates:
                raise HTTPException(status_code=500, detail="لم يُرجع الذكاء أي نتيجة")
            
            parts = candidates[0].get("content", {}).get("parts", [])
            if not parts:
                raise HTTPException(status_code=500, detail="استجابة فارغة")
            
            ai_text = parts[0].get("text", "").strip()
            
            # نُحاول parse JSON (قد يحوي markdown blocks)
            ai_text = ai_text.strip()
            if ai_text.startswith("```"):
                # نُزيل ```json ... ```
                ai_text = ai_text.split("```")[1] if "```" in ai_text else ai_text
                if ai_text.startswith("json"):
                    ai_text = ai_text[4:].strip()
            
            try:
                parsed = _json.loads(ai_text)
            except Exception as parse_err:
                # نحاول استخراج JSON من النص
                import re
                match = re.search(r'\{[\s\S]*\}', ai_text)
                if match:
                    parsed = _json.loads(match.group(0))
                else:
                    raise HTTPException(status_code=500, detail=f"فشل تحليل JSON: {ai_text[:100]}")
            
            # نتحقق من البنية
            if "tool" not in parsed or "settings" not in parsed:
                raise HTTPException(status_code=500, detail="JSON غير صالح من AI")
            
            return parsed
            
    except HTTPException:
        raise
    except httpx.TimeoutException:
        raise HTTPException(status_code=504, detail="انتهت مهلة الذكاء — حاول مرة أخرى")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/ai/ask")
async def ai_ask(
    request: Request,
    student_id: int = Form(...),
    question: str   = Form(...),
    context: str    = Form(default=""),  # سياق اختياري (الدرس الحالي/سؤال معيّن)
):
    """
    🤖 يسأل الطالب المساعد الذكي
    """
    # rate limit صارم
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=15, window_seconds=300, key_prefix="ai_ask"):
        raise HTTPException(status_code=429, detail="استخدمت أسئلة كثيرة — انتظر 5 دقائق")
    
    # تحقق من المدخلات
    q = question.strip()
    if len(q) < 3:
        raise HTTPException(status_code=400, detail="السؤال قصير جداً")
    if len(q) > 1000:
        raise HTTPException(status_code=400, detail="السؤال طويل — اختصره")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="المساعد الذكي غير مفعّل — تواصل مع الأستاذ"
        )
    
    # بناء الرسالة الكاملة
    full_prompt = AI_SYSTEM_PROMPT
    if context.strip():
        full_prompt += f"\n\n📚 سياق الدرس: {context.strip()[:500]}"
    full_prompt += f"\n\n❓ سؤال الطالب: {q}\n\n💡 إجابتك:"
    
    # استدعاء Gemini API مباشرة عبر REST (لا حاجة لمكتبة)
    import httpx
    try:
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
        payload = {
            "contents": [{"parts": [{"text": full_prompt}]}],
            "generationConfig": {
                "temperature": 0.7,
                "maxOutputTokens": 500,
                "topP": 0.9,
            },
            "safetySettings": [
                {"category": "HARM_CATEGORY_HARASSMENT",        "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_HATE_SPEECH",       "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_LOW_AND_ABOVE"},
                {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
            ]
        }
        
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.post(url, json=payload)
            
            # 🔄 لو الموديل الأساسي فشل (مثلاً deprecated/quota)، جرّب flash-lite كاحتياط
            if r.status_code == 404 or r.status_code == 429:
                fallback_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-lite:generateContent?key={api_key}"
                try:
                    r = await client.post(fallback_url, json=payload)
                except Exception:
                    pass
            
            if r.status_code != 200:
                err_msg = "فشل الاتصال بالمساعد"
                try:
                    err_data = r.json()
                    if "error" in err_data:
                        err_msg = err_data["error"].get("message", err_msg)[:200]
                except:
                    pass
                # إن كانت رسالة "model not found" — اقترح حل
                if "not found" in err_msg.lower() or "not supported" in err_msg.lower():
                    err_msg = "الموديل غير معروف — تحديث المنصة مطلوب. تواصل مع الأستاذ."
                raise HTTPException(status_code=500, detail=err_msg)
            
            data = r.json()
            
            # استخراج النص
            answer = ""
            try:
                candidates = data.get("candidates", [])
                if candidates:
                    parts = candidates[0].get("content", {}).get("parts", [])
                    if parts:
                        answer = parts[0].get("text", "").strip()
            except Exception:
                pass
            
            if not answer:
                # ربما المحتوى مرفوض من safety
                finish_reason = (candidates[0].get("finishReason", "") if candidates else "")
                if finish_reason == "SAFETY":
                    raise HTTPException(status_code=400, detail="السؤال غير مناسب")
                raise HTTPException(status_code=500, detail="لم نتلقَّ إجابة")
            
            # سجّل الاستخدام (اختياري - للإحصاء)
            try:
                supabase.table("ai_questions_log").insert({
                    "student_id": student_id,
                    "question": q[:500],
                    "answer_length": len(answer),
                }).execute()
            except Exception:
                pass  # الجدول قد لا يكون موجوداً
            
            return {
                "answer": answer,
                "model": "gemini-2.5-flash",
            }
    
    except HTTPException:
        raise
    except httpx.TimeoutException:
        raise HTTPException(status_code=504, detail="استغرق المساعد وقتاً طويلاً — حاول مرة أخرى")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/ai/status")
async def ai_status():
    """فحص حالة المساعد الذكي"""
    has_key = bool(os.getenv("GEMINI_API_KEY", "").strip())
    return {"enabled": has_key}


# ═══════════════════════════════════════════════════════════════
# 📧 MONTHLY EMAIL REPORTS — التقارير الشهرية لأولياء الأمور
# ═══════════════════════════════════════════════════════════════

def _generate_monthly_report_html(student: dict, stats: dict, parent_name: str = "") -> str:
    """يولّد HTML لتقرير شهري احترافي"""
    
    name = student.get("full_name", "البطل")
    grade = student.get("grade", "")
    
    rank = stats.get("rank_in_grade") or "—"
    rank_total = stats.get("grade_total_students") or "—"
    
    # الإنجازات
    perfect = stats.get("perfect_count", 0)
    avg_pct = stats.get("avg_score_pct", 0)
    minutes = stats.get("minutes_30d", 0)
    days_active = stats.get("days_active_30d", 0)
    total_tests = stats.get("total_tests_30d", 0)
    
    # التقييم
    if avg_pct >= 85:
        rating_emoji = "🏆"
        rating_text = "ممتاز"
        rating_color = "#2ecc71"
    elif avg_pct >= 70:
        rating_emoji = "⭐"
        rating_text = "جيد جداً"
        rating_color = "#3498db"
    elif avg_pct >= 50:
        rating_emoji = "📈"
        rating_text = "في تحسّن"
        rating_color = "#f39c12"
    else:
        rating_emoji = "💪"
        rating_text = "يحتاج دعم"
        rating_color = "#e74c3c"
    
    greeting = f"السلام عليكم {parent_name}" if parent_name else "السلام عليكم ولي الأمر الكريم"
    
    return f"""<!DOCTYPE html>
<html dir="rtl" lang="ar">
<head>
<meta charset="UTF-8">
<style>
    body {{ font-family: 'Segoe UI', Tahoma, Arial; background: #f5f5f0; margin: 0; padding: 20px; color: #333; }}
    .container {{ max-width: 600px; margin: 0 auto; background: white; border-radius: 16px; overflow: hidden; box-shadow: 0 8px 24px rgba(0,0,0,0.1); }}
    .header {{ background: linear-gradient(135deg, #d4af37, #b8860b); color: white; padding: 30px 20px; text-align: center; }}
    .header h1 {{ margin: 0 0 8px; font-size: 24px; }}
    .header p {{ margin: 0; opacity: 0.9; font-size: 14px; }}
    .greeting {{ padding: 24px; background: #fdfaf0; }}
    .greeting p {{ margin: 0; line-height: 1.7; font-size: 15px; }}
    .student-card {{ background: linear-gradient(135deg, #1a1a2e, #0d1b3e); color: white; padding: 24px; text-align: center; }}
    .student-card h2 {{ margin: 0 0 6px; font-size: 22px; color: #fcf6ba; }}
    .student-card p {{ margin: 0; opacity: 0.85; font-size: 14px; }}
    .rating-badge {{ display: inline-block; background: {rating_color}; color: white; padding: 8px 20px; border-radius: 20px; margin-top: 12px; font-weight: bold; }}
    .stats-grid {{ display: table; width: 100%; padding: 24px; }}
    .stat-row {{ display: table-row; }}
    .stat-cell {{ display: table-cell; padding: 14px; text-align: center; border-bottom: 1px solid #eee; }}
    .stat-num {{ font-size: 28px; font-weight: bold; color: #d4af37; }}
    .stat-label {{ font-size: 13px; color: #666; margin-top: 4px; }}
    .insights {{ background: #fdfaf0; padding: 24px; border-top: 3px solid #d4af37; }}
    .insights h3 {{ color: #b8860b; margin: 0 0 14px; font-size: 18px; }}
    .insight-item {{ background: white; padding: 12px 16px; border-radius: 8px; margin-bottom: 8px; border-right: 4px solid #d4af37; }}
    .footer {{ background: #1a1a2e; color: white; padding: 20px; text-align: center; font-size: 12px; }}
    .footer a {{ color: #d4af37; text-decoration: none; }}
</style>
</head>
<body>
<div class="container">
    <div class="header">
        <h1>👑 التقرير الشهري</h1>
        <p>إمبراطورية الرياضيات — أ. رشدي سيد</p>
    </div>
    
    <div class="greeting">
        <p>{greeting}،</p>
        <p style="margin-top:10px;">يسرّنا أن نُرسل لك تقرير الشهر الماضي عن أداء ابنك في المنصة:</p>
    </div>
    
    <div class="student-card">
        <h2>{name}</h2>
        <p>📚 {grade}</p>
        <div class="rating-badge">{rating_emoji} {rating_text}</div>
    </div>
    
    <div class="stats-grid">
        <div class="stat-row">
            <div class="stat-cell">
                <div class="stat-num">{total_tests}</div>
                <div class="stat-label">📊 تحديات منجزة</div>
            </div>
            <div class="stat-cell">
                <div class="stat-num">{avg_pct}%</div>
                <div class="stat-label">🎯 متوسط الدقة</div>
            </div>
        </div>
        <div class="stat-row">
            <div class="stat-cell">
                <div class="stat-num">{minutes}</div>
                <div class="stat-label">⏱️ دقيقة دراسة</div>
            </div>
            <div class="stat-cell">
                <div class="stat-num">{days_active}</div>
                <div class="stat-label">📅 يوم نشاط</div>
            </div>
        </div>
        <div class="stat-row">
            <div class="stat-cell">
                <div class="stat-num">{perfect}</div>
                <div class="stat-label">💯 علامة كاملة</div>
            </div>
            <div class="stat-cell">
                <div class="stat-num">#{rank}<span style="font-size:14px;color:#999;">/{rank_total}</span></div>
                <div class="stat-label">🏆 الترتيب في الصف</div>
            </div>
        </div>
    </div>
    
    <div class="insights">
        <h3>💡 ملاحظات الأستاذ</h3>
        <div class="insight-item">
            <strong style="color:{rating_color};">{rating_emoji} التقييم العام:</strong>
            <span>{rating_text} — {('استمر على هذا المستوى الرائع!' if avg_pct >= 70 else 'يحتاج لمزيد من المتابعة والتشجيع')}</span>
        </div>
        {"<div class='insight-item'><strong>📈 نشاط جيد:</strong> دخل المنصة " + str(days_active) + " يوم في الشهر الماضي.</div>" if days_active >= 10 else "<div class='insight-item'><strong>⚠️ نشاط قليل:</strong> ندعو ولي الأمر لتشجيع الطالب على دخول المنصة بشكل أكثر انتظاماً.</div>"}
        {"<div class='insight-item'><strong>🏆 إنجاز رائع:</strong> حقق " + str(perfect) + " علامة كاملة هذا الشهر!</div>" if perfect >= 3 else ""}
    </div>
    
    <div class="footer">
        <p>هذا تقرير تلقائي شهري من <strong>إمبراطورية الرياضيات</strong></p>
        <p>للتفاصيل الكاملة: <a href="https://your-domain.com/parent">ادخل بوابة ولي الأمر</a></p>
        <p style="margin-top:14px;opacity:0.6;">⚜️ الأستاذ رشدي سيد ⚜️</p>
    </div>
</div>
</body>
</html>"""


@app.get("/api/admin/email/status")
async def email_status_diagnose(admin = Depends(get_current_admin)):
    """🩺 تشخيص حالة إعدادات البريد"""
    smtp_host = os.getenv("SMTP_HOST", "").strip()
    smtp_user = os.getenv("SMTP_USER", "").strip()
    smtp_pass = os.getenv("SMTP_PASS", "").strip()
    smtp_from = os.getenv("SMTP_FROM", smtp_user).strip()
    smtp_port = os.getenv("SMTP_PORT", "587").strip()
    
    # عدّ الطلاب الذين لديهم parent_email
    parents_with_email = 0
    try:
        res = supabase.table("students").select("id", count="exact").not_.is_("parent_email", "null").execute()
        parents_with_email = res.count or 0
    except Exception:
        pass
    
    masked_user = (smtp_user[:3] + "***" + smtp_user[-10:]) if len(smtp_user) > 13 else (smtp_user[:3] + "***" if smtp_user else "(غير مُعرّف)")
    
    return {
        "smtp_configured": bool(smtp_host and smtp_user and smtp_pass),
        "smtp_host": smtp_host or "(غير مُعرّف)",
        "smtp_port": smtp_port,
        "smtp_user": masked_user,
        "smtp_pass_set": bool(smtp_pass),
        "parents_with_email_count": parents_with_email,
        "missing_vars": [
            v for v in ["SMTP_HOST", "SMTP_USER", "SMTP_PASS"]
            if not os.getenv(v, "").strip()
        ]
    }


@app.post("/api/admin/send_monthly_reports")
async def send_monthly_reports(
    test_email: str = Form(default=""),  # لو مُدخل: يُرسل لهذا فقط للاختبار
    admin = Depends(get_current_admin)
):
    """
    📧 يُرسل تقارير شهرية لكل أولياء الأمور بإيميل
    يستخدم SendGrid أو Mailgun (تحتاج SMTP credentials)
    """
    smtp_host = os.getenv("SMTP_HOST", "").strip()
    smtp_user = os.getenv("SMTP_USER", "").strip()
    smtp_pass = os.getenv("SMTP_PASS", "").strip()
    smtp_from = os.getenv("SMTP_FROM", smtp_user).strip()
    smtp_port = int(os.getenv("SMTP_PORT", "587"))
    
    if not smtp_host or not smtp_user or not smtp_pass:
        raise HTTPException(
            status_code=503,
            detail="إعدادات البريد غير مكتملة — أضف SMTP_HOST/USER/PASS في env"
        )
    
    # اجلب كل الطلاب الذين لديهم parent_email (مع معالجة لو العمود غير موجود)
    try:
        all_students = []
        offset = 0
        for _ in range(20):
            try:
                # نحاول الفلتر بـ parent_email
                res = supabase.table("students").select(
                    "id, full_name, grade, parent_email, parent_name"
                ).neq("parent_email", "").not_.is_("parent_email", "null").range(offset, offset + 999).execute()
            except Exception as e_filter:
                # لو العمود غير موجود، نجلب الكل ونفلتر بـ Python
                print(f"[reports] العمود parent_email غير قابل للفلترة، نفلتر يدوياً: {e_filter}")
                try:
                    res = supabase.table("students").select(
                        "id, full_name, grade, parent_email, parent_name"
                    ).range(offset, offset + 999).execute()
                except Exception as e2:
                    # العمود غير موجود نهائياً
                    raise HTTPException(
                        status_code=503,
                        detail="❌ عمود parent_email غير موجود في قاعدة البيانات. شغّل ملف parent_email_migration.sql أولاً"
                    )
            
            batch = res.data or []
            if not batch: break
            # فلترة يدوية للأمان
            valid = [s for s in batch if s.get("parent_email") and "@" in str(s.get("parent_email", ""))]
            all_students.extend(valid)
            if len(batch) < 1000: break
            offset += 1000
        
        if not all_students:
            return {
                "status": "no_recipients",
                "sent": 0,
                "failed": 0,
                "total": 0,
                "errors": [],
                "warning": "⚠️ لا يوجد أي طالب لديه parent_email مسجّل. يجب أولاً إضافة بريد ولي الأمر للطلاب."
            }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل جلب الطلاب: {str(e)[:200]}")
    
    if test_email:
        # وضع اختبار — فقط أول طالب لإيميل اختبار
        if all_students:
            all_students = [{**all_students[0], "parent_email": test_email}]
        else:
            raise HTTPException(status_code=404, detail="لا يوجد طلاب")
    
    sent = 0
    failed = 0
    errors = []
    
    import smtplib
    from email.mime.text import MIMEText
    from email.mime.multipart import MIMEMultipart
    from datetime import datetime, timezone, timedelta
    
    try:
        smtp = smtplib.SMTP(smtp_host, smtp_port)
        smtp.starttls()
        smtp.login(smtp_user, smtp_pass)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل اتصال SMTP: {str(e)[:200]}")
    
    now = datetime.now(timezone.utc)
    cutoff = now - timedelta(days=30)
    
    for student in all_students:
        try:
            # احسب الإحصائيات للشهر الماضي
            results_res = supabase.table("results").select(
                "score, total, timestamp"
            ).eq("student_id", student["id"]).gte("timestamp", cutoff.isoformat()).execute()
            results = results_res.data or []
            
            total_tests = len(results)
            total_score = sum((r.get("score") or 0) for r in results)
            total_max = sum((r.get("total") or 0) for r in results)
            avg_pct = round((total_score / total_max * 100), 1) if total_max > 0 else 0
            perfect_count = sum(1 for r in results if r.get("total") and r.get("score") == r.get("total"))
            
            # الجلسات
            sess_res = supabase.table("student_sessions").select(
                "session_bucket"
            ).eq("student_id", student["id"]).gte("session_bucket", cutoff.isoformat()).execute()
            sessions = sess_res.data or []
            minutes_30d = len(sessions) * 5
            days_active = len(set(s["session_bucket"][:10] for s in sessions if s.get("session_bucket")))
            
            stats = {
                "total_tests_30d": total_tests,
                "avg_score_pct": avg_pct,
                "minutes_30d": minutes_30d,
                "days_active_30d": days_active,
                "perfect_count": perfect_count,
            }
            
            # بناء HTML
            html_body = _generate_monthly_report_html(
                student, stats, student.get("parent_name", "")
            )
            
            # إرسال
            msg = MIMEMultipart("alternative")
            msg["Subject"] = f"📊 تقرير الشهر الماضي للطالب {student['full_name']}"
            msg["From"] = smtp_from
            msg["To"] = student["parent_email"]
            msg.attach(MIMEText(html_body, "html", "utf-8"))
            
            smtp.send_message(msg)
            sent += 1
            
            # سجّل في email_log
            try:
                supabase.table("email_log").insert({
                    "parent_email": student["parent_email"],
                    "student_id": student["id"],
                    "notif_type": "monthly_report",
                    "subject": msg["Subject"],
                }).execute()
            except Exception:
                pass
            
        except Exception as e:
            failed += 1
            errors.append({"student": student.get("full_name"), "error": str(e)[:100]})
    
    try:
        smtp.quit()
    except Exception:
        pass
    
    return {
        "status": "done",
        "sent": sent,
        "failed": failed,
        "total": len(all_students),
        "errors": errors[:10],
    }


# ═══════════════════════════════════════════════════════════════
# 🏛️ VIRTUAL CLASSES — الفصول الافتراضية
# ═══════════════════════════════════════════════════════════════

@app.get("/api/virtual_classes/student/{grade}")
async def get_student_virtual_classes(grade: str):
    """
    🎓 يجلب الفصول الافتراضية الخاصة بصف الطالب
    """
    try:
        # نُجرّب كل صيغ الصف الممكنة
        all_results = []
        seen_ids = set()
        for v in _grade_variants(grade):
            try:
                res = supabase.table("virtual_classes").select("*").eq(
                    "grade", v.strip()
                ).eq("is_active", True).order("starts_at", desc=False).execute()
                for row in (res.data or []):
                    if row["id"] not in seen_ids:
                        seen_ids.add(row["id"])
                        all_results.append(row)
            except Exception as e:
                print(f"[vc/student] failed for grade '{v}': {e}")
        
        return {"classes": all_results}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/admin/virtual_classes")
async def admin_list_virtual_classes(admin = Depends(get_current_admin)):
    """📋 قائمة كل الفصول الافتراضية للأدمن"""
    try:
        res = supabase.table("virtual_classes").select("*").order(
            "created_at", desc=True
        ).execute()
        return {"classes": res.data or []}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/admin/virtual_classes/create")
async def admin_create_virtual_class(
    title: str         = Form(...),
    grade: str         = Form(...),
    meeting_url: str   = Form(...),
    semester: str      = Form(default=""),
    unit: str          = Form(default=""),
    lesson: str        = Form(default=""),
    description: str   = Form(default=""),
    starts_at: str     = Form(default=""),  # ISO string
    duration_min: int  = Form(default=60),
    admin = Depends(get_current_admin)
):
    """➕ إنشاء فصل افتراضي جديد"""
    try:
        if not title.strip() or not grade.strip() or not meeting_url.strip():
            raise HTTPException(status_code=400, detail="العنوان والصف والرابط مطلوبون")
        
        # تحقق من صيغة الرابط
        if not (meeting_url.startswith("http://") or meeting_url.startswith("https://")):
            raise HTTPException(status_code=400, detail="الرابط يجب أن يبدأ بـ http:// أو https://")
        
        data = {
            "title": title.strip(),
            "grade": grade.strip(),
            "meeting_url": meeting_url.strip(),
            "semester": semester.strip() or None,
            "unit": unit.strip() or None,
            "lesson": lesson.strip() or None,
            "description": description.strip() or None,
            "duration_min": duration_min,
            "is_active": True,
        }
        if starts_at and starts_at.strip():
            data["starts_at"] = starts_at.strip()
        
        res = supabase.table("virtual_classes").insert(data).execute()
        return {"status": "created", "class": (res.data or [None])[0]}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/admin/virtual_classes/{class_id}/update")
async def admin_update_virtual_class(
    class_id: int,
    title: str        = Form(default=""),
    meeting_url: str  = Form(default=""),
    description: str  = Form(default=""),
    starts_at: str    = Form(default=""),
    is_active: bool   = Form(default=True),
    admin = Depends(get_current_admin)
):
    """✏️ تعديل فصل افتراضي"""
    try:
        updates = {"is_active": is_active}
        if title.strip():       updates["title"] = title.strip()
        if meeting_url.strip(): updates["meeting_url"] = meeting_url.strip()
        if description:         updates["description"] = description.strip() or None
        if starts_at:           updates["starts_at"] = starts_at.strip() or None
        
        supabase.table("virtual_classes").update(updates).eq("id", class_id).execute()
        return {"status": "updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/admin/virtual_classes/{class_id}")
async def admin_delete_virtual_class(
    class_id: int,
    admin = Depends(get_current_admin)
):
    """🗑️ حذف فصل افتراضي"""
    try:
        supabase.table("virtual_classes").delete().eq("id", class_id).execute()
        return {"status": "deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



# ═══════════════════════════════════════════════════════════════
# 👑 EMPIRE STATS — إحصائيات العرش الكاملة الدقيقة
# ═══════════════════════════════════════════════════════════════

@app.get("/api/admin/empire/ping")
async def empire_ping(admin = Depends(get_current_admin)):
    """🏓 يُختبر استجابة endpoint البسيط (للتشخيص)"""
    try:
        res = supabase.table("students").select("id", count="exact").limit(1).execute()
        return {"ok": True, "students_count": res.count or 0}
    except Exception as e:
        return {"ok": False, "error": str(e)[:200]}


@app.get("/api/admin/empire/stats")
async def empire_full_stats(admin = Depends(get_current_admin)):
    """
    📊 إحصائيات كاملة دقيقة لكل المنصة:
    - عدد الطلاب الإجمالي + حسب الصف
    - الطلاب الجدد (آخر 7 أيام، 30 يوم)
    - الطلاب النشطين (آخر 24 ساعة، 7 أيام، 30 يوم)
    - الإجابات: إجمالي + صحيحة + خاطئة
    - النقاط: إجمالي + متوسط + قمة
    - بنك الأسئلة: إجمالي + حسب الصف
    - أولياء الأمور: مسجّلين + نشطين
    - الإنجازات والشارات
    """
    from datetime import datetime, timedelta, timezone
    
    now = datetime.now(timezone.utc)
    last_24h = (now - timedelta(hours=24)).isoformat()
    last_7d  = (now - timedelta(days=7)).isoformat()
    last_30d = (now - timedelta(days=30)).isoformat()
    
    stats = {
        "students": {},
        "results": {},
        "questions": {},
        "parents": {},
        "achievements": {},
        "activity": {},
        "lists": {
            "new_students_7d": [],
            "new_students_30d": [],
            "active_today": [],
            "active_7d": [],
            "top_xp": [],
        }
    }
    
    try:
        # 🎓 الطلاب — جلب كامل
        all_students = []
        offset = 0
        # نُجرّب أعمدة محتملة (last_login أو last_active)
        select_cols_options = [
            "id, full_name, grade, total_points, created_at",
            "id, full_name, grade, total_points, created_at, last_login",
            "id, full_name, grade, total_points, created_at",
            "id, full_name, grade, total_points",
            "id, full_name, grade",
        ]
        select_cols = None
        for cols in select_cols_options:
            try:
                # نحاول مع صفحة اختبار
                test_res = supabase.table("students").select(cols).limit(1).execute()
                select_cols = cols
                break
            except Exception:
                continue
        
        if not select_cols:
            select_cols = "id, full_name, grade"
        
        print(f"[empire] using columns: {select_cols}")
        
        for _ in range(20):
            try:
                res = supabase.table("students").select(select_cols).range(offset, offset + 999).execute()
                batch = res.data or []
                if not batch: break
                all_students.extend(batch)
                if len(batch) < 1000: break
                offset += 1000
            except Exception as e:
                print(f"[empire] students fetch error: {e}")
                break
        
        stats["students"]["total"] = len(all_students)
        
        # طلاب جدد
        new_7d  = [s for s in all_students if s.get("created_at") and s["created_at"] >= last_7d]
        new_30d = [s for s in all_students if s.get("created_at") and s["created_at"] >= last_30d]
        stats["students"]["new_7d"]  = len(new_7d)
        stats["students"]["new_30d"] = len(new_30d)
        
        # حسب الصف
        by_grade = {}
        for s in all_students:
            g = s.get("grade", "غير محدد")
            by_grade[g] = by_grade.get(g, 0) + 1
        stats["students"]["by_grade"] = by_grade
        
        # قائمة الطلاب الجدد (مرتبين بالأحدث)
        new_30d_sorted = sorted(new_30d, key=lambda x: x.get("created_at", ""), reverse=True)
        stats["lists"]["new_students_7d"] = [
            {
                "id": s["id"],
                "name": s.get("full_name", ""),
                "grade": s.get("grade", ""),
                "joined": s.get("created_at", "")[:10],
            }
            for s in new_30d_sorted[:30] if s.get("created_at", "") >= last_7d
        ]
        stats["lists"]["new_students_30d"] = [
            {
                "id": s["id"],
                "name": s.get("full_name", ""),
                "grade": s.get("grade", ""),
                "joined": s.get("created_at", "")[:10],
            }
            for s in new_30d_sorted[:50]
        ]
        
        # 📝 النتائج — إجمالي + النشاط
        try:
            res_count = supabase.table("results").select("id", count="exact").execute()
            stats["results"]["total"] = res_count.count or 0
        except:
            stats["results"]["total"] = 0
        
        # نتائج آخر 24 ساعة
        try:
            res_24h = supabase.table("results").select(
                "id, student_id, score, total, timestamp"
            ).gte("timestamp", last_24h).execute()
            results_24h_data = res_24h.data or []
            stats["results"]["last_24h"] = len(results_24h_data)
            
            # طلاب نشطين اليوم (لديهم نتيجة في آخر 24 ساعة)
            active_today_ids = set(r.get("student_id") for r in results_24h_data if r.get("student_id"))
            stats["activity"]["active_today_count"] = len(active_today_ids)
            
            students_dict = {s["id"]: s for s in all_students}
            stats["lists"]["active_today"] = [
                {
                    "id": sid,
                    "name": students_dict[sid].get("full_name", ""),
                    "grade": students_dict[sid].get("grade", ""),
                    "xp": students_dict[sid].get("total_points", 0),
                }
                for sid in active_today_ids if sid in students_dict
            ][:30]
        except Exception as e:
            print(f"[stats] active today error: {e}")
            stats["results"]["last_24h"] = 0
            stats["activity"]["active_today_count"] = 0
        
        # نتائج آخر 7 أيام
        try:
            res_7d = supabase.table("results").select(
                "id, student_id, score, total, timestamp"
            ).gte("timestamp", last_7d).execute()
            results_7d_data = res_7d.data or []
            stats["results"]["last_7d"] = len(results_7d_data)
            
            # طلاب نشطين الأسبوع
            active_7d_ids = set(r.get("student_id") for r in results_7d_data if r.get("student_id"))
            stats["activity"]["active_7d_count"] = len(active_7d_ids)
            
            stats["lists"]["active_7d"] = [
                {
                    "id": sid,
                    "name": students_dict[sid].get("full_name", ""),
                    "grade": students_dict[sid].get("grade", ""),
                    "xp": students_dict[sid].get("total_points", 0),
                }
                for sid in active_7d_ids if sid in students_dict
            ][:50]
            
            # حساب الدقة
            total_correct = sum((r.get("score") or 0) for r in results_7d_data)
            total_questions = sum((r.get("total") or 0) for r in results_7d_data)
            stats["results"]["accuracy_7d"] = round(
                (total_correct / total_questions * 100), 1
            ) if total_questions > 0 else 0
            stats["results"]["correct_7d"] = total_correct
            stats["results"]["wrong_7d"] = total_questions - total_correct
        except Exception as e:
            print(f"[stats] 7d error: {e}")
        
        # 🏆 أعلى XP
        top_students = sorted(
            all_students,
            key=lambda x: (x.get("total_points") or 0),
            reverse=True
        )[:10]
        stats["lists"]["top_xp"] = [
            {
                "id": s["id"],
                "name": s.get("full_name", ""),
                "grade": s.get("grade", ""),
                "xp": s.get("total_points", 0),
            }
            for s in top_students
        ]
        
        # إجمالي XP
        total_xp = sum((s.get("total_points") or 0) for s in all_students)
        stats["students"]["total_xp"] = total_xp
        stats["students"]["avg_xp"] = round(total_xp / len(all_students), 1) if all_students else 0
        
        # ❓ بنك الأسئلة
        try:
            q_count = supabase.table("questions").select("id", count="exact").execute()
            stats["questions"]["total"] = q_count.count or 0
        except:
            stats["questions"]["total"] = 0
        
        # 👨‍👩‍👧 أولياء الأمور
        try:
            parents_with_phone = [s for s in all_students if s.get("parent_phone") if hasattr(s, 'get')]
            # نُعيد جلب مع parent_phone
            res_parents = supabase.table("students").select("id, parent_phone, parent_email").execute()
            parents_data = res_parents.data or []
            with_phone = sum(1 for s in parents_data if s.get("parent_phone"))
            with_email = sum(1 for s in parents_data if s.get("parent_email"))
            stats["parents"]["total_with_phone"] = with_phone
            stats["parents"]["total_with_email"] = with_email
        except Exception as e:
            print(f"[stats] parents error: {e}")
            stats["parents"]["total_with_phone"] = 0
            stats["parents"]["total_with_email"] = 0
        
        # 🏅 الإنجازات
        try:
            ach_count = supabase.table("achievements").select("id", count="exact").execute()
            stats["achievements"]["total_granted"] = ach_count.count or 0
        except:
            stats["achievements"]["total_granted"] = 0
        
        return stats
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ في الإحصائيات: {str(e)[:200]}")



# ═══════════════════════════════════════════════════════════════
# ⚡ QUICK STATS — إحصاءات سريعة (count فقط، بدون جلب البيانات)
# ═══════════════════════════════════════════════════════════════

@app.get("/api/admin/stats/counts")
async def quick_counts(admin = Depends(get_current_admin)):
    """
    ⚡ يُرجع counts بسرعة فائقة (Supabase count exact)
    مفيد للوحات الإحصائيات بدون تحميل آلاف الصفوف
    """
    cache_key = "stats:counts"
    cached = _cache.get(cache_key)
    if cached:
        return cached
    
    counts = {}
    
    # تخصيص الجداول لجلب counts فقط
    tables_to_count = [
        "students", "questions", "results", "achievements",
        "issue_reports", "tutoring_bookings", "html_lessons",
        "virtual_classes", "short_videos"
    ]
    
    for tbl in tables_to_count:
        try:
            res = supabase.table(tbl).select("id", count="exact").limit(1).execute()
            counts[tbl] = res.count or 0
        except Exception as e:
            print(f"[quick_counts] {tbl}: {e}")
            counts[tbl] = 0
    
    # cache لـ 5 دقائق (الـ counts لا تتغيّر بسرعة)
    _cache.set(cache_key, counts, ttl_seconds=300)
    return counts


# ═══════════════════════════════════════════════════════════════
# 💾 CACHE MANAGEMENT — للأدمن
# ═══════════════════════════════════════════════════════════════

@app.get("/api/admin/cache/stats")
async def cache_stats(admin = Depends(get_current_admin)):
    """📊 إحصاءات الـ cache"""
    return _cache.stats()


@app.post("/api/admin/cache/clear")
async def cache_clear(
    pattern: str = Form(default=""),
    admin = Depends(get_current_admin)
):
    """🧹 مسح الـ cache (الكل أو نمط محدد)"""
    if pattern:
        _cache.clear_pattern(pattern)
        return {"status": "cleared", "pattern": pattern}
    else:
        _cache._store.clear()
        return {"status": "cleared_all"}



# ═══════════════════════════════════════════════════════════════
# 📝 EXAM GENERATOR — مولّد الاختبارات الفوري للمعلمين
# ═══════════════════════════════════════════════════════════════

@app.get("/api/teacher/exam_generator/preview")
async def teacher_exam_preview(
    grade: str,
    semester: str = "",
    unit: str = "",
    lesson: str = "",
    lessons: str = "",      # 📖 جديد: دروس متعددة مفصولة بـ |
    count: int = 10,
    difficulty: str = "mixed",
    types: str = "all"
):
    """📝 معاينة أسئلة لمولّد الاختبار"""
    if not grade.strip():
        raise HTTPException(status_code=400, detail="grade مطلوب")
    
    count = max(1, min(int(count), 50))
    
    try:
        # 🔄 استخدام _grade_variants لدعم كل صيغ اسم الصف
        variants = _grade_variants(grade) or [grade.strip()]
        
        all_questions = []
        seen_ids = set()
        
        # 📖 الدروس: دعم لقائمة دروس متعددة
        lessons_list = []
        if lessons.strip():
            lessons_list = [l.strip() for l in lessons.split("|") if l.strip()]
        elif lesson.strip():
            lessons_list = [lesson.strip()]
        
        for v in variants:
            if lessons_list:
                # نجلب لكل درس على حدة
                for L in lessons_list:
                    q = supabase.table("questions").select("*").eq("grade", v.strip()).ilike("lesson", L)
                    res = q.execute()
                    for row in (res.data or []):
                        if row.get("id") not in seen_ids:
                            seen_ids.add(row.get("id"))
                            all_questions.append(row)
            else:
                query = supabase.table("questions").select("*").eq("grade", v.strip())
                res = query.execute()
                for row in (res.data or []):
                    if row.get("id") not in seen_ids:
                        seen_ids.add(row.get("id"))
                        all_questions.append(row)
        
        def _norm(s):
            return str(s or "").strip().lower()
        
        if semester.strip():
            ns = _norm(semester)
            all_questions = [q for q in all_questions if 
                _norm(q.get("semester")) == ns or ns in _norm(q.get("semester")) or _norm(q.get("semester")) in ns]
        
        if unit.strip():
            nu = _norm(unit)
            all_questions = [q for q in all_questions if 
                _norm(q.get("unit")) == nu or nu in _norm(q.get("unit")) or _norm(q.get("unit")) in nu]
        
        # فلترة بنوع السؤال — ندعم q_type (الصيغة الصحيحة في DB)
        if types != "all":
            type_map = {
                "mcq":   ["mcq", "اختياري", "متعدد"],
                "tf":    ["true_false", "tf", "صواب", "خطأ"],
                "short": ["short", "essay", "قصير", "مقالي"],
            }
            keywords = type_map.get(types, [])
            if keywords:
                filtered = []
                for q in all_questions:
                    qtype = _norm(q.get("q_type") or q.get("type") or q.get("question_type") or "")
                    if any(k in qtype for k in keywords):
                        filtered.append(q)
                if filtered:
                    all_questions = filtered
        
        import random
        random.shuffle(all_questions)
        selected = all_questions[:count]
        
        return {
            "status": "ok",
            "total_available": len(all_questions),
            "selected_count": len(selected),
            "questions": selected,
            "filter": {
                "grade": grade, "semester": semester, "unit": unit,
                "lesson": lesson, "count": count,
                "difficulty": difficulty, "types": types,
            }
        }
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/teacher/exam_generator/build_pdf")
async def teacher_exam_build_pdf(
    title: str           = Form(default="اختبار رياضيات"),
    school_name: str     = Form(default=""),
    teacher_name: str    = Form(default=""),
    class_name: str      = Form(default=""),
    grade: str           = Form(...),
    semester: str        = Form(default=""),
    unit: str            = Form(default=""),
    lesson: str          = Form(default=""),
    lessons: str         = Form(default=""),    # 📖 جديد: دروس متعددة بـ |
    count: int           = Form(default=10),
    types: str           = Form(default="all"),
    include_answers: bool = Form(default=True),
    show_marks: bool     = Form(default=True),
    total_marks: int     = Form(default=20),
    mode: str            = Form(default="exam"),  # 🆕 exam | activity
):
    """📄 بناء ورقة اختبار/نشاط جاهزة كـ HTML قابل للطباعة"""
    try:
        # 🔄 _grade_variants لدعم كل الصياغات
        variants = _grade_variants(grade) or [grade.strip()]
        
        all_questions = []
        seen_ids = set()
        
        # 📖 الدروس: دعم لقائمة دروس متعددة
        lessons_list = []
        if lessons.strip():
            lessons_list = [l.strip() for l in lessons.split("|") if l.strip()]
        elif lesson.strip():
            lessons_list = [lesson.strip()]
        
        for v in variants:
            if lessons_list:
                # نجلب لكل درس على حدة
                for L in lessons_list:
                    q = supabase.table("questions").select("*").eq("grade", v.strip()).ilike("lesson", L)
                    res = q.execute()
                    for row in (res.data or []):
                        if row.get("id") not in seen_ids:
                            seen_ids.add(row.get("id"))
                            all_questions.append(row)
            else:
                query = supabase.table("questions").select("*").eq("grade", v.strip())
                res = query.execute()
                for row in (res.data or []):
                    if row.get("id") not in seen_ids:
                        seen_ids.add(row.get("id"))
                        all_questions.append(row)
        
        def _norm(s):
            return str(s or "").strip().lower()
        
        if semester.strip():
            ns = _norm(semester)
            all_questions = [q for q in all_questions if _norm(q.get("semester")) == ns or ns in _norm(q.get("semester")) or _norm(q.get("semester")) in ns]
        if unit.strip():
            nu = _norm(unit)
            all_questions = [q for q in all_questions if _norm(q.get("unit")) == nu or nu in _norm(q.get("unit")) or _norm(q.get("unit")) in nu]
        
        import random
        random.shuffle(all_questions)
        questions = all_questions[:max(1, min(count, 50))]
        
        if not questions:
            raise HTTPException(status_code=404, detail="لا توجد أسئلة مطابقة للفلتر")
        
        marks_per_q = round(total_marks / len(questions), 2) if show_marks else 0
        
        from datetime import datetime
        date_str = datetime.now().strftime("%Y/%m/%d")
        
        # نبني محتوى الأسئلة
        questions_html = ""
        answers_html = ""
        arabic_letters = ["أ", "ب", "ج", "د", "هـ", "و"]
        
        for i, q in enumerate(questions, 1):
            q_text = (q.get("question") or "").replace("<", "&lt;").replace(">", "&gt;")
            options = q.get("options") or q.get("choices") or ""
            # نقرأ من كل الأعمدة المحتملة
            correct = (q.get("correct_answer") or q.get("answer") or 
                       q.get("correct") or q.get("right_answer") or "")
            
            opts_list = []
            if options:
                if isinstance(options, list):
                    opts_list = options
                elif isinstance(options, str):
                    opts_list = [o.strip() for o in options.replace("،", ",").split(",") if o.strip()]
            
            mark_label = ('<span class="qmark">[' + str(marks_per_q) + ' درجة]</span>') if show_marks else ""
            
            q_block = '<div class="question">'
            q_block += '<div class="q-header"><span class="q-num">' + str(i) + '.</span>'
            q_block += '<span class="q-text">' + q_text + '</span>' + mark_label + '</div>'
            
            if opts_list:
                q_block += '<div class="q-options">'
                for idx, opt in enumerate(opts_list[:6]):
                    letter = arabic_letters[idx] if idx < len(arabic_letters) else str(idx+1)
                    opt_clean = str(opt).replace("<", "&lt;").replace(">", "&gt;")
                    q_block += '<div class="q-option"><span class="q-letter">(' + letter + ')</span> ' + opt_clean + '</div>'
                q_block += '</div>'
            else:
                q_block += '<div class="q-answer-lines"><div class="line"></div><div class="line"></div><div class="line"></div></div>'
            
            q_block += '</div>'
            questions_html += q_block
            
            ans_clean = str(correct).replace("<", "&lt;").replace(">", "&gt;")
            answers_html += '<div class="ans-item"><strong>' + str(i) + '.</strong> ' + ans_clean + '</div>'
        
        # نبني الـ HTML الكامل
        css_styles = """
            * { box-sizing: border-box; margin: 0; padding: 0; }
            body {
                font-family: 'Cairo', sans-serif;
                background: #fff; color: #000;
                padding: 30px 25px;
                max-width: 800px; margin: 0 auto;
                line-height: 1.8;
            }
            .header { text-align: center; border-bottom: 3px double #000; padding-bottom: 18px; margin-bottom: 24px; }
            .school { font-size: 18px; font-weight: 700; margin-bottom: 4px; }
            .title-main { font-family: 'Amiri', serif; font-size: 32px; font-weight: 700; margin: 8px 0; }
            .meta-row { display: flex; justify-content: center; flex-wrap: wrap; gap: 10px; margin-top: 16px; font-size: 14px; }
            .meta-item { background: #f5f5f5; padding: 6px 14px; border-radius: 4px; border: 1px solid #ddd; }
            .student-info { display: grid; grid-template-columns: 1fr 1fr; gap: 14px; margin: 20px 0; padding: 14px; background: #fafafa; border: 1px solid #ddd; border-radius: 6px; }
            .student-info > div { border-bottom: 1px dashed #999; padding-bottom: 4px; }
            .student-info b { color: #444; }
            .instructions { background: #fff8e1; border-right: 4px solid #f9a825; padding: 12px 16px; margin: 16px 0; font-size: 14px; border-radius: 4px; }
            .question { margin: 22px 0; padding: 14px; border-radius: 6px; border: 1px solid #eee; page-break-inside: avoid; }
            .q-header { font-size: 16px; font-weight: 600; margin-bottom: 10px; }
            .q-num { color: #1976d2; font-weight: 700; margin-left: 6px; }
            .qmark { color: #d32f2f; font-size: 13px; font-weight: 600; background: #ffebee; padding: 2px 8px; border-radius: 3px; margin-right: 8px; }
            .q-options { padding-right: 24px; margin-top: 8px; }
            .q-option { margin: 6px 0; padding: 4px 0; font-size: 15px; }
            .q-letter { color: #1976d2; font-weight: 700; margin-left: 6px; }
            .q-answer-lines { margin-top: 8px; padding-right: 16px; }
            .line { border-bottom: 1px solid #999; height: 28px; margin: 6px 0; }
            body.is-activity .line { height: 38px; }
            body.is-activity .q-answer-lines { padding-right: 8px; }
            body.is-activity .question { background: #f0fdf4; border-color: #86efac; }
            body.is-activity .qmark { background: #dcfce7; color: #15803d; }
            body.is-activity .title-main { color: #15803d; }
            body.is-activity .answers-title { background: #15803d !important; }
            body.is-activity .instructions { background: #ecfdf5; border-right-color: #10b981; }
            .footer { margin-top: 40px; padding-top: 16px; border-top: 2px solid #000; text-align: center; font-size: 13px; color: #555; }
            .signature { display: flex; justify-content: space-between; margin-top: 30px; font-size: 14px; }
            .sign-block { text-align: center; min-width: 180px; }
            .sign-line { border-top: 1px solid #000; margin-top: 36px; padding-top: 4px; }
            .answers-page { page-break-before: always; padding: 20px; background: #f9f9f9; }
            .answers-title { text-align: center; font-size: 24px; font-weight: 700; background: #1976d2; color: #fff; padding: 14px; border-radius: 6px; margin-bottom: 20px; }
            .ans-item { padding: 8px 14px; background: #fff; margin: 6px 0; border-right: 4px solid #4caf50; border-radius: 4px; font-size: 15px; }
            .controls { position: fixed; top: 12px; left: 12px; background: #fff; border: 1px solid #ccc; border-radius: 8px; padding: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.15); z-index: 9999; }
            .controls button { background: #1976d2; color: #fff; border: none; padding: 8px 16px; margin: 0 4px; border-radius: 6px; cursor: pointer; font-family: 'Cairo', sans-serif; font-weight: 600; font-size: 14px; }
            .controls button:hover { background: #1565c0; }
            .controls button.green { background: #43a047; }
            .controls button.green:hover { background: #2e7d32; }
            @media print { .controls { display: none !important; } body { padding: 15px; } .question { border: none; } }
        """
        
        # بناء meta items
        meta_items = '<div class="meta-item">📚 ' + grade + '</div>'
        if semester:
            meta_items += '<div class="meta-item">📅 ' + semester + '</div>'
        if unit:
            meta_items += '<div class="meta-item">📦 ' + unit + '</div>'
        if lesson:
            meta_items += '<div class="meta-item">📖 ' + lesson + '</div>'
        meta_items += '<div class="meta-item">📅 ' + date_str + '</div>'
        if show_marks:
            meta_items += '<div class="meta-item">⭐ الدرجة الكلية: ' + str(total_marks) + '</div>'
        
        school_block = ('<div class="school">' + school_name + '</div>') if school_name else ''
        class_text = class_name if class_name else "........................"
        
        is_activity = (mode == "activity")
        
        if is_activity:
            instructions_text = "🎯 نشاط صفي: ناقش مع زملائك وحلّ الأسئلة في الوقت المحدد. "
            instructions_text += "خصّص مساحة كافية للتفكير قبل الإجابة."
        else:
            instructions_text = "اقرأ الأسئلة بعناية قبل الإجابة. "
            if show_marks:
                instructions_text += "كل سؤال بـ " + str(marks_per_q) + " درجة. "
            instructions_text += "اكتب الإجابة في المكان المخصص."
        
        teacher_sign = "المعلم " + (teacher_name if teacher_name else "")
        
        answers_block = ""
        if include_answers:
            answers_block = '<div class="answers-page"><div class="answers-title">📋 ورقة الإجابات النموذجية</div>' + answers_html + '</div>'
        
        full_html = (
            '<!DOCTYPE html>\n'
            '<html lang="ar" dir="rtl">\n'
            '<head>\n'
            '<meta charset="UTF-8">\n'
            '<title>' + title + '</title>\n'
            '<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
            '<style>' + css_styles + '</style>\n'
            '</head>\n'
            '<body>\n'
            '<div class="controls">'
            '<button onclick="window.print()">🖨️ طباعة / حفظ PDF</button>'
            '<button class="green" onclick="window.close()">✕ إغلاق</button>'
            '</div>\n'
            '<div class="header">'
            + school_block +
            '<div class="title-main">' + title + '</div>'
            '<div class="meta-row">' + meta_items + '</div>'
            '</div>\n'
            '<div class="student-info">'
            '<div><b>اسم الطالب:</b> ........................................</div>'
            '<div><b>الفصل:</b> ' + class_text + '</div>'
            '<div><b>رقم الجلوس:</b> ............................</div>'
            '<div><b>الدرجة:</b> ............ من ' + str(total_marks) + '</div>'
            '</div>\n'
            '<div class="instructions"><b>📌 تعليمات:</b> ' + instructions_text + '</div>\n'
            '<div class="questions">' + questions_html + '</div>\n'
            '<div class="signature">'
            '<div class="sign-block"><div class="sign-line">' + teacher_sign + '</div></div>'
            '<div class="sign-block"><div class="sign-line">المراجع</div></div>'
            '</div>\n'
            '<div class="footer">🎓 إمبراطورية الرياضيات &nbsp;·&nbsp; مع تمنياتي بالتوفيق</div>\n'
            + answers_block +
            '\n</body>\n</html>'
        )
        
        from fastapi.responses import HTMLResponse
        return HTMLResponse(content=full_html)
    
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ═══════════════════════════════════════════════════════════════
# 🏆 CERTIFICATES & TEMPLATES — قوالب وشهادات للطباعة
# ═══════════════════════════════════════════════════════════════

@app.post("/api/teacher/exam_generator/build_pdf_v2")
async def teacher_exam_build_pdf_v2(
    title: str           = Form(default="اختبار رياضيات"),
    school_name: str     = Form(default=""),
    teacher_name: str    = Form(default=""),
    class_name: str      = Form(default=""),
    grade: str           = Form(...),
    semester: str        = Form(default=""),
    unit: str            = Form(default=""),
    lesson: str          = Form(default=""),
    lessons: str         = Form(default=""),
    # 🆕 عدد كل نوع من الأسئلة
    n_essay: int         = Form(default=3),   # ✍️ مقالي
    n_tf: int            = Form(default=5),    # ✅ صواب/خطأ
    n_match: int         = Form(default=1),    # 🔗 وصّل
    n_mcq: int           = Form(default=4),    # 🔵 اختياري
    include_answers: bool = Form(default=True),
    show_marks: bool     = Form(default=True),
    total_marks: int     = Form(default=20),
    use_ai: bool         = Form(default=True),  # 🆕 توليد بالـ AI إذا لم يكف البنك
):
    """📄 مولّد اختبارات متنوع v2: مقالي + صواب/خطأ + وصّل + اختياري"""
    try:
        variants = _grade_variants(grade) or [grade.strip()]
        
        # 📖 الدروس
        lessons_list = []
        if lessons.strip():
            lessons_list = [l.strip() for l in lessons.split("|") if l.strip()]
        elif lesson.strip():
            lessons_list = [lesson.strip()]
        
        # نجلب أسئلة البنك المتاحة
        bank_questions = []
        seen_ids = set()
        for v in variants:
            if lessons_list:
                for L in lessons_list:
                    res = supabase.table("questions").select("*").eq("grade", v.strip()).ilike("lesson", L).execute()
                    for row in (res.data or []):
                        if row.get("id") not in seen_ids:
                            seen_ids.add(row.get("id"))
                            bank_questions.append(row)
            else:
                res = supabase.table("questions").select("*").eq("grade", v.strip()).execute()
                for row in (res.data or []):
                    if row.get("id") not in seen_ids:
                        seen_ids.add(row.get("id"))
                        bank_questions.append(row)
        
        def _norm(s):
            return str(s or "").strip().lower()
        if semester.strip():
            ns = _norm(semester)
            bank_questions = [q for q in bank_questions if ns in _norm(q.get("semester")) or _norm(q.get("semester")) in ns]
        if unit.strip():
            nu = _norm(unit)
            bank_questions = [q for q in bank_questions if nu in _norm(q.get("unit")) or _norm(q.get("unit")) in nu]
        
        import random
        random.shuffle(bank_questions)
        
        # نصنّف أسئلة البنك حسب النوع
        def detect_type(q):
            t = _norm(q.get("q_type") or q.get("type"))
            opts = q.get("options") or ""
            ans = _norm(q.get("answer") or q.get("correct_answer") or q.get("correct"))
            if "true_false" in t or "صواب" in t or ans in ("صواب", "خطأ", "true", "false"):
                return "tf"
            if "multiple" in t or "اختياري" in t or (opts and len(str(opts)) > 3):
                return "mcq"
            if "short" in t or "essay" in t or "مقالي" in t or "قصير" in t:
                return "essay"
            return "essay"  # افتراضي
        
        pools = {"essay": [], "tf": [], "mcq": [], "match": []}
        for q in bank_questions:
            pools[detect_type(q)].append(q)
        
        # نجمع المطلوب من كل نوع
        selected = {"essay": [], "tf": [], "mcq": [], "match": []}
        needs_ai = {}
        for typ, want in [("essay", n_essay), ("tf", n_tf), ("mcq", n_mcq), ("match", n_match)]:
            available = pools[typ][:want]
            selected[typ] = available
            shortage = want - len(available)
            if shortage > 0:
                needs_ai[typ] = shortage
        
        # 🤖 نولّد الناقص بالـ AI
        ai_generated = {"essay": [], "tf": [], "mcq": [], "match": []}
        if use_ai and needs_ai:
            api_key = os.getenv("GEMINI_API_KEY", "").strip()
            if api_key:
                ctx = f"الصف: {grade}"
                if unit: ctx += f" | الوحدة: {unit}"
                if lessons_list: ctx += f" | الدروس: {'، '.join(lessons_list)}"
                
                type_names = {"essay": "مقالي (يتطلب شرح/حل)", "tf": "صواب أو خطأ", "mcq": "اختيار من متعدد (4 خيارات)", "match": "وصّل من عمودين"}
                
                for typ, cnt in needs_ai.items():
                    prompt = f"""أنت معلم رياضيات خبير في سلطنة عُمان. أنشئ {cnt} سؤال من نوع "{type_names[typ]}".

السياق: {ctx}

أخرج JSON فقط:"""
                    if typ == "match":
                        prompt += """
{"questions":[{"type":"match","q":"وصّل كل عبارة بما يناسبها","pairs":[{"left":"العمود أ","right":"العمود ب"}],"answer":"أ-1، ب-2"}]}"""
                    elif typ == "tf":
                        prompt += """
{"questions":[{"type":"tf","q":"العبارة","answer":"صواب"}]}"""
                    elif typ == "mcq":
                        prompt += """
{"questions":[{"type":"mcq","q":"السؤال","options":["أ","ب","ج","د"],"answer":"أ"}]}"""
                    else:
                        prompt += """
{"questions":[{"type":"essay","q":"السؤال المقالي","answer":"نموذج الإجابة"}]}"""
                    
                    try:
                        import httpx
                        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
                        with httpx.Client(timeout=45.0) as client:
                            r = client.post(url, json={"contents":[{"parts":[{"text":prompt}]}],"generationConfig":{"temperature":0.8,"maxOutputTokens":3000,"responseMimeType":"application/json"}})
                            if r.status_code == 200:
                                raw = r.json().get("candidates",[{}])[0].get("content",{}).get("parts",[{}])[0].get("text","").strip()
                                if raw.startswith("```"):
                                    raw = raw.split("\n",1)[1] if "\n" in raw else raw
                                    raw = raw.rsplit("```",1)[0] if "```" in raw else raw
                                import json as jl
                                parsed = jl.loads(raw)
                                ai_generated[typ] = parsed.get("questions", [])
                    except Exception as e:
                        print(f"[exam_v2 AI] {typ}: {str(e)[:100]}")
        
        # ندمج البنك + AI لكل نوع
        def normalize_bank_q(q, typ):
            opts = q.get("options") or ""
            if isinstance(opts, str):
                opts = [o.strip() for o in opts.replace("،", ",").replace("|", ",").split(",") if o.strip()]
            return {
                "type": typ,
                "q": q.get("question") or "",
                "options": opts,
                "answer": q.get("answer") or q.get("correct_answer") or q.get("correct") or "",
                "pairs": q.get("pairs") or []
            }
        
        final = {"essay": [], "tf": [], "mcq": [], "match": []}
        for typ in ["essay", "tf", "mcq", "match"]:
            for q in selected[typ]:
                final[typ].append(normalize_bank_q(q, typ))
            for q in ai_generated[typ]:
                final[typ].append(q)
        
        total_q = sum(len(final[t]) for t in final)
        if total_q == 0:
            raise HTTPException(status_code=404, detail="لا توجد أسئلة. تأكد من وجود أسئلة في البنك أو فعّل التوليد بالـ AI.")
        
        marks_per_q = round(total_marks / total_q, 2) if (show_marks and total_q > 0) else 0
        
        from datetime import datetime
        date_str = datetime.now().strftime("%Y/%m/%d")
        arabic_letters = ["أ", "ب", "ج", "د", "هـ", "و"]
        
        # نبني الأقسام
        sections_html = ""
        answers_html = ""
        q_counter = [0]  # mutable
        
        def esc(s):
            return str(s or "").replace("<", "&lt;").replace(">", "&gt;")
        
        section_defs = [
            ("essay", "✍️ السؤال الأول: أجب عن الأسئلة التالية", "essay"),
            ("tf", "✅ السؤال الثاني: ضع علامة (✓) أو (✗)", "tf"),
            ("match", "🔗 السؤال الثالث: صِل من العمود (أ) ما يناسبه من العمود (ب)", "match"),
            ("mcq", "🔵 السؤال الرابع: اختر الإجابة الصحيحة", "mcq"),
        ]
        
        section_num = 0
        for typ, sec_title, _ in section_defs:
            qs = final[typ]
            if not qs:
                continue
            section_num += 1
            sections_html += '<div class="exam-section"><div class="section-title">' + sec_title + '</div>'
            
            for q in qs:
                q_counter[0] += 1
                n = q_counter[0]
                qtext = esc(q.get("q") or q.get("question"))
                mark_label = ('<span class="qmark">[' + str(marks_per_q) + ' د]</span>') if show_marks else ""
                
                blk = '<div class="question"><div class="q-header"><span class="q-num">' + str(n) + '.</span> <span class="q-text">' + qtext + '</span>' + mark_label + '</div>'
                
                if typ == "mcq":
                    opts = q.get("options") or []
                    blk += '<div class="q-options">'
                    for idx, opt in enumerate(opts[:6]):
                        letter = arabic_letters[idx] if idx < len(arabic_letters) else str(idx+1)
                        blk += '<div class="q-option"><span class="q-circle">⃝</span> <span class="q-letter">(' + letter + ')</span> ' + esc(opt) + '</div>'
                    blk += '</div>'
                elif typ == "tf":
                    blk += '<div class="tf-box"><span class="tf-opt">صواب ( )</span> <span class="tf-opt">خطأ ( )</span></div>'
                elif typ == "match":
                    pairs = q.get("pairs") or []
                    if pairs:
                        import random as rnd
                        rights = [p.get("right","") for p in pairs]
                        rnd.shuffle(rights)
                        blk += '<table class="match-table"><tr><th>العمود (أ)</th><th>العمود (ب)</th></tr>'
                        for pi, p in enumerate(pairs):
                            r = rights[pi] if pi < len(rights) else ""
                            blk += '<tr><td>' + str(pi+1) + '. ' + esc(p.get("left","")) + ' (....)</td><td>' + arabic_letters[pi] + '. ' + esc(r) + '</td></tr>'
                        blk += '</table>'
                else:  # essay
                    blk += '<div class="q-answer-lines"><div class="line"></div><div class="line"></div><div class="line"></div></div>'
                
                blk += '</div>'
                sections_html += blk
                
                # الإجابات
                ans = esc(q.get("answer") or q.get("correct"))
                answers_html += '<div class="ans-item"><strong>' + str(n) + '.</strong> ' + ans + '</div>'
            
            sections_html += '</div>'
        
        css_styles = """
            * { box-sizing: border-box; margin: 0; padding: 0; }
            body { font-family: 'Cairo', sans-serif; background: #fff; color: #000; padding: 30px 25px; max-width: 800px; margin: 0 auto; line-height: 1.9; }
            .header { text-align: center; border-bottom: 3px double #000; padding-bottom: 18px; margin-bottom: 20px; }
            .school { font-size: 18px; font-weight: 700; }
            .title-main { font-family: 'Amiri', serif; font-size: 30px; font-weight: 700; margin: 8px 0; }
            .meta-row { display: flex; justify-content: center; flex-wrap: wrap; gap: 8px; margin-top: 14px; font-size: 13px; }
            .meta-item { background: #f5f5f5; padding: 5px 12px; border-radius: 4px; border: 1px solid #ddd; }
            .student-info { display: grid; grid-template-columns: 1fr 1fr; gap: 12px; margin: 18px 0; padding: 14px; background: #fafafa; border: 1px solid #ddd; border-radius: 6px; }
            .student-info > div { border-bottom: 1px dashed #999; padding-bottom: 4px; }
            .instructions { background: #fff8e1; border-right: 4px solid #f9a825; padding: 10px 16px; margin: 14px 0; font-size: 14px; border-radius: 4px; }
            .exam-section { margin: 24px 0; }
            .section-title { font-size: 17px; font-weight: 700; color: #fff; background: linear-gradient(135deg,#1976d2,#1565c0); padding: 10px 16px; border-radius: 6px; margin-bottom: 14px; }
            .question { margin: 16px 0; padding: 12px 14px; border-radius: 6px; border: 1px solid #eee; page-break-inside: avoid; }
            .q-header { font-size: 15.5px; font-weight: 600; margin-bottom: 8px; }
            .q-num { color: #1976d2; font-weight: 700; margin-left: 6px; }
            .qmark { color: #d32f2f; font-size: 12px; font-weight: 600; background: #ffebee; padding: 2px 8px; border-radius: 3px; margin-right: 8px; }
            .q-options { padding-right: 24px; margin-top: 6px; }
            .q-option { margin: 6px 0; font-size: 15px; }
            .q-circle { color: #999; margin-left: 4px; }
            .q-letter { color: #1976d2; font-weight: 700; }
            .q-answer-lines { margin-top: 8px; padding-right: 12px; }
            .line { border-bottom: 1px solid #999; height: 32px; margin: 8px 0; }
            .tf-box { padding-right: 24px; margin-top: 6px; font-size: 15px; }
            .tf-opt { display: inline-block; margin-left: 30px; font-weight: 600; }
            .match-table { width: 100%; border-collapse: collapse; margin-top: 8px; }
            .match-table th { background: #e3f2fd; padding: 8px; border: 1px solid #90caf9; font-size: 14px; }
            .match-table td { padding: 10px 12px; border: 1px solid #ddd; font-size: 14px; }
            .signature { display: flex; justify-content: space-between; margin-top: 30px; font-size: 14px; }
            .sign-block { text-align: center; min-width: 160px; }
            .sign-line { border-top: 1px solid #000; margin-top: 36px; padding-top: 4px; }
            .footer { margin-top: 30px; padding-top: 14px; border-top: 2px solid #000; text-align: center; font-size: 13px; color: #555; }
            .answers-page { page-break-before: always; padding: 20px; }
            .answers-title { text-align: center; font-size: 22px; font-weight: 700; background: #4caf50; color: #fff; padding: 12px; border-radius: 6px; margin-bottom: 18px; }
            .ans-item { padding: 7px 14px; background: #f1f8e9; margin: 5px 0; border-right: 4px solid #4caf50; border-radius: 4px; font-size: 14px; }
            .controls { position: fixed; top: 12px; left: 12px; background: #fff; border: 1px solid #ccc; border-radius: 8px; padding: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.15); z-index: 9999; }
            .controls button { background: #1976d2; color: #fff; border: none; padding: 8px 16px; margin: 0 4px; border-radius: 6px; cursor: pointer; font-family: 'Cairo'; font-weight: 600; }
            .controls button.green { background: #43a047; }
            @media print { .controls { display: none !important; } body { padding: 15px; } }
        """
        
        meta_items = '<div class="meta-item">📚 ' + grade + '</div>'
        if semester: meta_items += '<div class="meta-item">📅 ' + semester + '</div>'
        if unit: meta_items += '<div class="meta-item">📦 ' + unit + '</div>'
        meta_items += '<div class="meta-item">📅 ' + date_str + '</div>'
        if show_marks: meta_items += '<div class="meta-item">⭐ الكلية: ' + str(total_marks) + '</div>'
        
        school_block = ('<div class="school">' + school_name + '</div>') if school_name else ''
        class_text = class_name if class_name else "..................."
        teacher_sign = "المعلم " + (teacher_name if teacher_name else "")
        
        # ملخص أنواع الأسئلة
        type_summary = []
        if final["essay"]: type_summary.append(str(len(final["essay"])) + " مقالي")
        if final["tf"]: type_summary.append(str(len(final["tf"])) + " صواب/خطأ")
        if final["match"]: type_summary.append(str(len(final["match"])) + " وصّل")
        if final["mcq"]: type_summary.append(str(len(final["mcq"])) + " اختياري")
        instructions_text = "اقرأ الأسئلة بعناية. الاختبار يحتوي: " + "، ".join(type_summary) + "."
        
        answers_block = ""
        if include_answers:
            answers_block = '<div class="answers-page"><div class="answers-title">📋 ورقة الإجابات النموذجية</div>' + answers_html + '</div>'
        
        full_html = (
            '<!DOCTYPE html>\n<html lang="ar" dir="rtl">\n<head>\n<meta charset="UTF-8">\n'
            '<title>' + title + '</title>\n'
            '<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
            '<style>' + css_styles + '</style>\n</head>\n<body>\n'
            '<div class="controls"><button onclick="window.print()">🖨️ طباعة / PDF</button>'
            '<button class="green" onclick="window.close()">✕ إغلاق</button></div>\n'
            '<div class="header">' + school_block + '<div class="title-main">' + title + '</div>'
            '<div class="meta-row">' + meta_items + '</div></div>\n'
            '<div class="student-info">'
            '<div><b>اسم الطالب:</b> ...........................</div>'
            '<div><b>الفصل:</b> ' + class_text + '</div>'
            '<div><b>رقم الجلوس:</b> ..................</div>'
            '<div><b>الدرجة:</b> ......... من ' + str(total_marks) + '</div></div>\n'
            '<div class="instructions"><b>📌 تعليمات:</b> ' + instructions_text + '</div>\n'
            + sections_html +
            '<div class="signature"><div class="sign-block"><div class="sign-line">' + teacher_sign + '</div></div>'
            '<div class="sign-block"><div class="sign-line">المراجع</div></div></div>\n'
            '<div class="footer">🎓 إمبراطورية الرياضيات &nbsp;·&nbsp; مع تمنياتي بالتوفيق</div>\n'
            + answers_block + '\n</body>\n</html>'
        )
        
        return HTMLResponse(content=full_html)
    
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ════════════════════════════════════════════════════════════════════════════
# 🎛️ EXAM EDITOR — محرّر الاختبار التفاعلي (جمع → معاينة/تعديل → طباعة)
# ════════════════════════════════════════════════════════════════════════════
@app.post("/api/teacher/exam_generator/prepare_questions")
async def teacher_exam_prepare_questions(
    grade: str           = Form(...),
    semester: str        = Form(default=""),
    unit: str            = Form(default=""),
    lesson: str          = Form(default=""),
    lessons: str         = Form(default=""),
    n_essay: int         = Form(default=3),
    n_tf: int            = Form(default=5),
    n_match: int         = Form(default=1),
    n_mcq: int           = Form(default=4),
    use_ai: bool         = Form(default=True),
):
    """🎛️ يجمع الأسئلة ويُرجعها JSON للمحرّر التفاعلي (قبل الطباعة)"""
    try:
        variants = _grade_variants(grade) or [grade.strip()]
        lessons_list = []
        if lessons.strip():
            lessons_list = [l.strip() for l in lessons.split("|") if l.strip()]
        elif lesson.strip():
            lessons_list = [lesson.strip()]
        
        bank_questions = []
        seen_ids = set()
        for v in variants:
            if lessons_list:
                for L in lessons_list:
                    res = supabase.table("questions").select("*").eq("grade", v.strip()).ilike("lesson", L).execute()
                    for row in (res.data or []):
                        if row.get("id") not in seen_ids:
                            seen_ids.add(row.get("id"))
                            bank_questions.append(row)
            else:
                res = supabase.table("questions").select("*").eq("grade", v.strip()).execute()
                for row in (res.data or []):
                    if row.get("id") not in seen_ids:
                        seen_ids.add(row.get("id"))
                        bank_questions.append(row)
        
        def _norm(s):
            return str(s or "").strip().lower()
        if semester.strip():
            ns = _norm(semester)
            bank_questions = [q for q in bank_questions if ns in _norm(q.get("semester")) or _norm(q.get("semester")) in ns]
        if unit.strip():
            nu = _norm(unit)
            bank_questions = [q for q in bank_questions if nu in _norm(q.get("unit")) or _norm(q.get("unit")) in nu]
        
        import random
        random.shuffle(bank_questions)
        
        def detect_type(q):
            t = _norm(q.get("q_type") or q.get("type"))
            opts = q.get("options") or ""
            ans = _norm(q.get("answer") or q.get("correct_answer") or q.get("correct"))
            if "true_false" in t or "صواب" in t or ans in ("صواب", "خطأ", "true", "false"):
                return "tf"
            if "multiple" in t or "اختياري" in t or (opts and len(str(opts)) > 3):
                return "mcq"
            return "essay"
        
        pools = {"essay": [], "tf": [], "mcq": [], "match": []}
        for q in bank_questions:
            pools[detect_type(q)].append(q)
        
        def normalize_bank_q(q, typ):
            opts = q.get("options") or ""
            if isinstance(opts, str):
                opts = [o.strip() for o in opts.replace("،", ",").replace("|", ",").split(",") if o.strip()]
            return {
                "type": typ,
                "q": q.get("question") or "",
                "options": opts,
                "answer": q.get("answer") or q.get("correct_answer") or q.get("correct") or "",
                "pairs": q.get("pairs") or [],
                "source": "bank"
            }
        
        selected = {"essay": [], "tf": [], "mcq": [], "match": []}
        needs_ai = {}
        for typ, want in [("essay", n_essay), ("tf", n_tf), ("mcq", n_mcq), ("match", n_match)]:
            available = [normalize_bank_q(q, typ) for q in pools[typ][:want]]
            selected[typ] = available
            shortage = want - len(available)
            if shortage > 0:
                needs_ai[typ] = shortage
        
        # توليد الناقص بالـ AI
        if use_ai and needs_ai:
            api_key = os.getenv("GEMINI_API_KEY", "").strip()
            if api_key:
                ctx = f"الصف: {grade}"
                if unit: ctx += f" | الوحدة: {unit}"
                if lessons_list: ctx += f" | الدروس: {'، '.join(lessons_list)}"
                type_names = {"essay": "مقالي (يتطلب شرح/حل)", "tf": "صواب أو خطأ", "mcq": "اختيار من متعدد (4 خيارات)", "match": "وصّل من عمودين"}
                
                for typ, cnt in needs_ai.items():
                    prompt = f"""أنت معلم رياضيات خبير في سلطنة عُمان. أنشئ {cnt} سؤال من نوع "{type_names[typ]}".\n\nالسياق: {ctx}\n\nأخرج JSON فقط:"""
                    if typ == "match":
                        prompt += """\n{"questions":[{"type":"match","q":"وصّل كل عبارة بما يناسبها","pairs":[{"left":"عبارة","right":"مقابلها"}],"answer":"1-أ"}]}"""
                    elif typ == "tf":
                        prompt += """\n{"questions":[{"type":"tf","q":"العبارة","answer":"صواب"}]}"""
                    elif typ == "mcq":
                        prompt += """\n{"questions":[{"type":"mcq","q":"السؤال","options":["أ","ب","ج","د"],"answer":"أ"}]}"""
                    else:
                        prompt += """\n{"questions":[{"type":"essay","q":"السؤال المقالي","answer":"نموذج الإجابة"}]}"""
                    try:
                        import httpx
                        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
                        with httpx.Client(timeout=45.0) as client:
                            r = client.post(url, json={"contents":[{"parts":[{"text":prompt}]}],"generationConfig":{"temperature":0.8,"maxOutputTokens":3000,"responseMimeType":"application/json"}})
                            if r.status_code == 200:
                                raw = r.json().get("candidates",[{}])[0].get("content",{}).get("parts",[{}])[0].get("text","").strip()
                                if raw.startswith("```"):
                                    raw = raw.split("\n",1)[1] if "\n" in raw else raw
                                    raw = raw.rsplit("```",1)[0] if "```" in raw else raw
                                import json as jl
                                parsed = jl.loads(raw)
                                for q in parsed.get("questions", []):
                                    q["source"] = "ai"
                                    if not isinstance(q.get("options"), list):
                                        q["options"] = []
                                    if not isinstance(q.get("pairs"), list):
                                        q["pairs"] = []
                                    selected[typ].append(q)
                    except Exception as e:
                        print(f"[prepare AI] {typ}: {str(e)[:100]}")
        
        # دمج بترتيب: مقالي، صح/خطأ، وصّل، اختياري + درجة افتراضية
        default_marks = {"essay": 3, "tf": 1, "match": 4, "mcq": 1}
        ordered = []
        for typ in ["essay", "tf", "match", "mcq"]:
            for q in selected[typ]:
                q["marks"] = default_marks.get(typ, 1)
                ordered.append(q)
        
        return {
            "status": "ok",
            "questions": ordered,
            "counts": {t: len(selected[t]) for t in selected},
            "total": len(ordered)
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/teacher/exam_generator/build_from_questions")
async def teacher_exam_build_from_questions(
    title: str           = Form(default="اختبار رياضيات"),
    school_name: str     = Form(default=""),
    teacher_name: str    = Form(default=""),
    class_name: str      = Form(default=""),
    grade: str           = Form(default=""),
    semester: str        = Form(default=""),
    unit: str            = Form(default=""),
    questions_json: str  = Form(...),
    include_answers: bool = Form(default=True),
    show_marks: bool     = Form(default=True),
):
    """🖨️ يبني ورقة الاختبار من أسئلة جاهزة (بعد تعديل المعلم في المحرّر)"""
    try:
        import json as jl
        questions = jl.loads(questions_json)
        if not isinstance(questions, list) or not questions:
            raise HTTPException(status_code=400, detail="لا توجد أسئلة")
        
        from datetime import datetime
        date_str = datetime.now().strftime("%Y/%m/%d")
        arabic_letters = ["أ", "ب", "ج", "د", "هـ", "و"]
        
        def esc(s):
            return str(s or "").replace("<", "&lt;").replace(">", "&gt;")
        
        # المجموع الكلي = مجموع درجات الأسئلة
        total_marks = sum(float(q.get("marks", 1) or 0) for q in questions)
        
        # نجمّع حسب النوع للأقسام
        by_type = {"essay": [], "tf": [], "match": [], "mcq": []}
        for q in questions:
            t = q.get("type", "essay")
            if t not in by_type: t = "essay"
            by_type[t].append(q)
        
        section_defs = [
            ("essay", "✍️ السؤال الأول: أجب عن الأسئلة التالية"),
            ("tf", "✅ السؤال الثاني: ضع علامة (✓) أو (✗)"),
            ("match", "🔗 السؤال الثالث: صِل من العمود (أ) ما يناسبه من العمود (ب)"),
            ("mcq", "🔵 السؤال الرابع: اختر الإجابة الصحيحة"),
        ]
        
        sections_html = ""
        answers_html = ""
        n = 0
        for typ, sec_title in section_defs:
            qs = by_type[typ]
            if not qs:
                continue
            sections_html += '<div class="exam-section"><div class="section-title">' + sec_title + '</div>'
            for q in qs:
                n += 1
                qtext = esc(q.get("q"))
                marks = q.get("marks", 1)
                mark_label = ('<span class="qmark">[' + str(marks) + ' د]</span>') if show_marks else ""
                blk = '<div class="question"><div class="q-header"><span class="q-num">' + str(n) + '.</span> <span class="q-text">' + qtext + '</span>' + mark_label + '</div>'
                
                if typ == "mcq":
                    opts = q.get("options") or []
                    blk += '<div class="q-options">'
                    for idx, opt in enumerate(opts[:6]):
                        letter = arabic_letters[idx] if idx < len(arabic_letters) else str(idx+1)
                        blk += '<div class="q-option"><span class="q-circle">⃝</span> <span class="q-letter">(' + letter + ')</span> ' + esc(opt) + '</div>'
                    blk += '</div>'
                elif typ == "tf":
                    blk += '<div class="tf-box"><span class="tf-opt">صواب ( )</span> <span class="tf-opt">خطأ ( )</span></div>'
                elif typ == "match":
                    pairs = q.get("pairs") or []
                    if pairs:
                        import random as rnd
                        rights = [p.get("right","") for p in pairs]
                        rnd.shuffle(rights)
                        blk += '<table class="match-table"><tr><th>العمود (أ)</th><th>العمود (ب)</th></tr>'
                        for pi, p in enumerate(pairs):
                            r = rights[pi] if pi < len(rights) else ""
                            blk += '<tr><td>' + str(pi+1) + '. ' + esc(p.get("left","")) + ' (....)</td><td>' + (arabic_letters[pi] if pi < len(arabic_letters) else str(pi+1)) + '. ' + esc(r) + '</td></tr>'
                        blk += '</table>'
                else:
                    blk += '<div class="q-answer-lines"><div class="line"></div><div class="line"></div><div class="line"></div></div>'
                blk += '</div>'
                sections_html += blk
                answers_html += '<div class="ans-item"><strong>' + str(n) + '.</strong> ' + esc(q.get("answer")) + '</div>'
            sections_html += '</div>'
        
        css_styles = """
            * { box-sizing: border-box; margin: 0; padding: 0; }
            body { font-family: 'Cairo', sans-serif; background: #fff; color: #000; padding: 30px 25px; max-width: 800px; margin: 0 auto; line-height: 1.9; }
            .header { text-align: center; border-bottom: 3px double #000; padding-bottom: 18px; margin-bottom: 20px; }
            .school { font-size: 18px; font-weight: 700; }
            .title-main { font-family: 'Amiri', serif; font-size: 30px; font-weight: 700; margin: 8px 0; }
            .meta-row { display: flex; justify-content: center; flex-wrap: wrap; gap: 8px; margin-top: 14px; font-size: 13px; }
            .meta-item { background: #f5f5f5; padding: 5px 12px; border-radius: 4px; border: 1px solid #ddd; }
            .student-info { display: grid; grid-template-columns: 1fr 1fr; gap: 12px; margin: 18px 0; padding: 14px; background: #fafafa; border: 1px solid #ddd; border-radius: 6px; }
            .student-info > div { border-bottom: 1px dashed #999; padding-bottom: 4px; }
            .instructions { background: #fff8e1; border-right: 4px solid #f9a825; padding: 10px 16px; margin: 14px 0; font-size: 14px; border-radius: 4px; }
            .exam-section { margin: 24px 0; }
            .section-title { font-size: 17px; font-weight: 700; color: #fff; background: linear-gradient(135deg,#1976d2,#1565c0); padding: 10px 16px; border-radius: 6px; margin-bottom: 14px; }
            .question { margin: 16px 0; padding: 12px 14px; border-radius: 6px; border: 1px solid #eee; page-break-inside: avoid; }
            .q-header { font-size: 15.5px; font-weight: 600; margin-bottom: 8px; }
            .q-num { color: #1976d2; font-weight: 700; margin-left: 6px; }
            .qmark { color: #d32f2f; font-size: 12px; font-weight: 600; background: #ffebee; padding: 2px 8px; border-radius: 3px; margin-right: 8px; }
            .q-options { padding-right: 24px; margin-top: 6px; }
            .q-option { margin: 6px 0; font-size: 15px; }
            .q-circle { color: #999; margin-left: 4px; }
            .q-letter { color: #1976d2; font-weight: 700; }
            .q-answer-lines { margin-top: 8px; padding-right: 12px; }
            .line { border-bottom: 1px solid #999; height: 32px; margin: 8px 0; }
            .tf-box { padding-right: 24px; margin-top: 6px; font-size: 15px; }
            .tf-opt { display: inline-block; margin-left: 30px; font-weight: 600; }
            .match-table { width: 100%; border-collapse: collapse; margin-top: 8px; }
            .match-table th { background: #e3f2fd; padding: 8px; border: 1px solid #90caf9; font-size: 14px; }
            .match-table td { padding: 10px 12px; border: 1px solid #ddd; font-size: 14px; }
            .signature { display: flex; justify-content: space-between; margin-top: 30px; font-size: 14px; }
            .sign-block { text-align: center; min-width: 160px; }
            .sign-line { border-top: 1px solid #000; margin-top: 36px; padding-top: 4px; }
            .footer { margin-top: 30px; padding-top: 14px; border-top: 2px solid #000; text-align: center; font-size: 13px; color: #555; }
            .answers-page { page-break-before: always; padding: 20px; }
            .answers-title { text-align: center; font-size: 22px; font-weight: 700; background: #4caf50; color: #fff; padding: 12px; border-radius: 6px; margin-bottom: 18px; }
            .ans-item { padding: 7px 14px; background: #f1f8e9; margin: 5px 0; border-right: 4px solid #4caf50; border-radius: 4px; font-size: 14px; }
            .controls { position: fixed; top: 12px; left: 12px; background: #fff; border: 1px solid #ccc; border-radius: 8px; padding: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.15); z-index: 9999; }
            .controls button { background: #1976d2; color: #fff; border: none; padding: 8px 16px; margin: 0 4px; border-radius: 6px; cursor: pointer; font-family: 'Cairo'; font-weight: 600; }
            .controls button.green { background: #43a047; }
            @media print { .controls { display: none !important; } body { padding: 15px; } }
        """
        
        meta_items = ''
        if grade: meta_items += '<div class="meta-item">📚 ' + esc(grade) + '</div>'
        if semester: meta_items += '<div class="meta-item">📅 ' + esc(semester) + '</div>'
        if unit: meta_items += '<div class="meta-item">📦 ' + esc(unit) + '</div>'
        meta_items += '<div class="meta-item">📅 ' + date_str + '</div>'
        if show_marks: meta_items += '<div class="meta-item">⭐ الكلية: ' + str(int(total_marks)) + '</div>'
        
        school_block = ('<div class="school">' + esc(school_name) + '</div>') if school_name else ''
        class_text = esc(class_name) if class_name else "..................."
        teacher_sign = "المعلم " + esc(teacher_name)
        
        answers_block = ""
        if include_answers:
            answers_block = '<div class="answers-page"><div class="answers-title">📋 ورقة الإجابات النموذجية</div>' + answers_html + '</div>'
        
        full_html = (
            '<!DOCTYPE html>\n<html lang="ar" dir="rtl">\n<head>\n<meta charset="UTF-8">\n'
            '<title>' + esc(title) + '</title>\n'
            '<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
            '<style>' + css_styles + '</style>\n</head>\n<body>\n'
            '<div class="controls"><button onclick="window.print()">🖨️ طباعة / PDF</button>'
            '<button class="green" onclick="window.close()">✕ إغلاق</button></div>\n'
            '<div class="header">' + school_block + '<div class="title-main">' + esc(title) + '</div>'
            '<div class="meta-row">' + meta_items + '</div></div>\n'
            '<div class="student-info">'
            '<div><b>اسم الطالب:</b> ...........................</div>'
            '<div><b>الفصل:</b> ' + class_text + '</div>'
            '<div><b>رقم الجلوس:</b> ..................</div>'
            '<div><b>الدرجة:</b> ......... من ' + str(int(total_marks)) + '</div></div>\n'
            '<div class="instructions"><b>📌 تعليمات:</b> اقرأ الأسئلة بعناية وأجب عن الجميع.</div>\n'
            + sections_html +
            '<div class="signature"><div class="sign-block"><div class="sign-line">' + teacher_sign + '</div></div>'
            '<div class="sign-block"><div class="sign-line">المراجع</div></div></div>\n'
            '<div class="footer">🎓 إمبراطورية الرياضيات &nbsp;·&nbsp; مع تمنياتي بالتوفيق</div>\n'
            + answers_block + '\n</body>\n</html>'
        )
        return HTMLResponse(content=full_html)
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



async def teacher_certificate_build(
    template: str       = Form(default="gold"),     # gold | royal | modern | classic | floral
    student_name: str   = Form(...),
    achievement: str    = Form(default="التميّز في الرياضيات"),
    grade: str          = Form(default=""),
    school_name: str    = Form(default=""),
    teacher_name: str   = Form(default=""),
    date: str           = Form(default=""),
    extra_text: str     = Form(default=""),
):
    """🏆 بناء شهادة تقدير قابلة للطباعة"""
    try:
        from datetime import datetime
        if not date:
            date = datetime.now().strftime("%Y/%m/%d")
        
        # تصاميم متنوعة
        templates_data = {
            "gold": {
                "bg_color": "#fefae0",
                "border": "12px double #d4af37",
                "title_color": "#8b6914",
                "name_color": "#5d4e1f",
                "icon": "👑",
                "decoration": "⚜️"
            },
            "royal": {
                "bg_color": "#f0e8ff",
                "border": "10px double #5b3e8a",
                "title_color": "#3d2666",
                "name_color": "#2a1a47",
                "icon": "🎖️",
                "decoration": "♛"
            },
            "modern": {
                "bg_color": "#f0f9ff",
                "border": "8px solid #0ea5e9",
                "title_color": "#0369a1",
                "name_color": "#082f49",
                "icon": "🏆",
                "decoration": "★"
            },
            "classic": {
                "bg_color": "#fff5e6",
                "border": "10px ridge #c9742d",
                "title_color": "#7a4317",
                "name_color": "#4a2607",
                "icon": "📜",
                "decoration": "❦"
            },
            "floral": {
                "bg_color": "#fff0f5",
                "border": "8px double #d63384",
                "title_color": "#831843",
                "name_color": "#500724",
                "icon": "🌸",
                "decoration": "✿"
            }
        }
        
        t = templates_data.get(template, templates_data["gold"])
        
        css = """
            * { box-sizing: border-box; margin: 0; padding: 0; }
            body { font-family: 'Cairo', sans-serif; padding: 0; margin: 0; background: #fff; }
            .cert-page {
                width: 297mm; height: 210mm;
                padding: 30px;
                background: """ + t["bg_color"] + """;
                border: """ + t["border"] + """;
                margin: 20px auto;
                box-sizing: border-box;
                position: relative;
                page-break-after: always;
            }
            .cert-corner { position: absolute; font-size: 60px; opacity: 0.15; color: """ + t["title_color"] + """; }
            .cert-corner.tl { top: 30px; left: 30px; }
            .cert-corner.tr { top: 30px; right: 30px; }
            .cert-corner.bl { bottom: 30px; left: 30px; }
            .cert-corner.br { bottom: 30px; right: 30px; }
            .cert-content {
                text-align: center; height: 100%;
                display: flex; flex-direction: column;
                justify-content: center; align-items: center;
                position: relative; z-index: 2;
            }
            .cert-icon { font-size: 80px; margin-bottom: 14px; }
            .cert-title {
                font-family: 'Amiri', serif;
                font-size: 56px; font-weight: 700;
                color: """ + t["title_color"] + """;
                margin-bottom: 8px;
                text-shadow: 2px 2px 4px rgba(0,0,0,0.1);
            }
            .cert-subtitle {
                font-size: 22px;
                color: """ + t["title_color"] + """;
                margin-bottom: 30px;
                font-weight: 600;
            }
            .cert-text {
                font-size: 22px; line-height: 2;
                margin: 12px 0;
                color: #333;
            }
            .cert-name {
                font-family: 'Amiri', serif;
                font-size: 60px;
                font-weight: 700;
                color: """ + t["name_color"] + """;
                margin: 20px 0;
                padding: 14px 40px;
                border-bottom: 3px solid """ + t["title_color"] + """;
                display: inline-block;
            }
            .cert-achievement {
                font-size: 24px;
                color: """ + t["title_color"] + """;
                font-weight: 600;
                margin: 14px 0;
                padding: 10px 24px;
                background: rgba(255,255,255,0.5);
                border-radius: 8px;
                display: inline-block;
            }
            .cert-extra {
                font-size: 18px;
                color: #555;
                margin: 14px 30px;
                font-style: italic;
            }
            .cert-footer {
                display: flex;
                justify-content: space-around;
                width: 100%;
                margin-top: 40px;
                font-size: 16px;
            }
            .cert-sign-block {
                text-align: center;
                min-width: 200px;
            }
            .cert-sign-line {
                border-top: 2px solid """ + t["title_color"] + """;
                margin-top: 60px;
                padding-top: 6px;
                color: """ + t["name_color"] + """;
                font-weight: 600;
            }
            .controls { position: fixed; top: 12px; left: 12px; background: #fff; border: 1px solid #ccc; border-radius: 8px; padding: 10px; z-index: 9999; box-shadow: 0 2px 8px rgba(0,0,0,0.15); }
            .controls button { background: #1976d2; color: #fff; border: none; padding: 8px 16px; margin: 0 4px; border-radius: 6px; cursor: pointer; font-family: 'Cairo', sans-serif; font-weight: 600; }
            @media print {
                .controls { display: none !important; }
                @page { size: A4 landscape; margin: 0; }
                body { margin: 0; }
                .cert-page { margin: 0; box-shadow: none; }
            }
        """
        
        extra_block = ('<div class="cert-extra">' + extra_text + '</div>') if extra_text else ''
        teacher_block = teacher_name if teacher_name else "الأستاذ"
        school_block = school_name if school_name else "المدرسة"
        grade_text = ('في ' + grade) if grade else ''
        
        full_html = (
            '<!DOCTYPE html>\n'
            '<html lang="ar" dir="rtl">\n'
            '<head>\n'
            '<meta charset="UTF-8">\n'
            '<title>شهادة تقدير</title>\n'
            '<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
            '<style>' + css + '</style>\n'
            '</head>\n'
            '<body>\n'
            '<div class="controls">'
            '<button onclick="window.print()">🖨️ طباعة / حفظ PDF</button>'
            '<button onclick="window.close()" style="background:#43a047;">✕ إغلاق</button>'
            '</div>\n'
            '<div class="cert-page">'
            '<div class="cert-corner tl">' + t["decoration"] + '</div>'
            '<div class="cert-corner tr">' + t["decoration"] + '</div>'
            '<div class="cert-corner bl">' + t["decoration"] + '</div>'
            '<div class="cert-corner br">' + t["decoration"] + '</div>'
            '<div class="cert-content">'
            '<div class="cert-icon">' + t["icon"] + '</div>'
            '<div class="cert-title">شهادة تقدير</div>'
            '<div class="cert-subtitle">Certificate of Excellence</div>'
            '<div class="cert-text">تُمنح هذه الشهادة إلى الطالب/ـة المتميّز/ـة</div>'
            '<div class="cert-name">' + student_name + '</div>'
            '<div class="cert-text">' + grade_text + '</div>'
            '<div class="cert-text">تقديراً لـ</div>'
            '<div class="cert-achievement">' + achievement + '</div>'
            + extra_block +
            '<div class="cert-footer">'
            '<div class="cert-sign-block"><div class="cert-sign-line">' + teacher_block + '</div></div>'
            '<div class="cert-sign-block"><div class="cert-sign-line">📅 ' + date + '</div></div>'
            '<div class="cert-sign-block"><div class="cert-sign-line">' + school_block + '</div></div>'
            '</div>'
            '</div>'
            '</div>\n'
            '</body>\n</html>'
        )
        
        from fastapi.responses import HTMLResponse
        return HTMLResponse(content=full_html)
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ في بناء الشهادة: {str(e)[:200]}")


@app.post("/api/teacher/certificate/batch")
async def teacher_certificate_batch(
    template: str       = Form(default="gold"),
    student_names: str  = Form(...),    # أسماء مفصولة بأسطر جديدة
    achievement: str    = Form(default="التميّز في الرياضيات"),
    grade: str          = Form(default=""),
    school_name: str    = Form(default=""),
    teacher_name: str   = Form(default=""),
    date: str           = Form(default=""),
):
    """🏆 طباعة شهادات متعددة دفعة واحدة"""
    names = [n.strip() for n in student_names.replace("،", "\n").split("\n") if n.strip()]
    if not names:
        raise HTTPException(status_code=400, detail="أدخل أسماء الطلاب")
    if len(names) > 50:
        raise HTTPException(status_code=400, detail="الحد الأقصى 50 شهادة دفعة واحدة")
    
    # نبني نفس HTML الـ single لكن مع تكرار صفحات
    # نحاول استخدام نفس الدالة لكن مع تجميع
    from datetime import datetime
    if not date:
        date = datetime.now().strftime("%Y/%m/%d")
    
    templates_data = {
        "gold":    {"bg": "#fefae0", "border": "12px double #d4af37", "title_c": "#8b6914", "name_c": "#5d4e1f", "icon": "👑", "deco": "⚜️"},
        "royal":   {"bg": "#f0e8ff", "border": "10px double #5b3e8a", "title_c": "#3d2666", "name_c": "#2a1a47", "icon": "🎖️", "deco": "♛"},
        "modern":  {"bg": "#f0f9ff", "border": "8px solid #0ea5e9",   "title_c": "#0369a1", "name_c": "#082f49", "icon": "🏆", "deco": "★"},
        "classic": {"bg": "#fff5e6", "border": "10px ridge #c9742d",   "title_c": "#7a4317", "name_c": "#4a2607", "icon": "📜", "deco": "❦"},
        "floral":  {"bg": "#fff0f5", "border": "8px double #d63384",  "title_c": "#831843", "name_c": "#500724", "icon": "🌸", "deco": "✿"},
    }
    t = templates_data.get(template, templates_data["gold"])
    
    pages_html = ""
    for name in names:
        grade_text = ('في ' + grade) if grade else ''
        page = (
            '<div class="cert-page">'
            '<div class="cert-corner tl">' + t["deco"] + '</div>'
            '<div class="cert-corner tr">' + t["deco"] + '</div>'
            '<div class="cert-corner bl">' + t["deco"] + '</div>'
            '<div class="cert-corner br">' + t["deco"] + '</div>'
            '<div class="cert-content">'
            '<div class="cert-icon">' + t["icon"] + '</div>'
            '<div class="cert-title">شهادة تقدير</div>'
            '<div class="cert-subtitle">Certificate of Excellence</div>'
            '<div class="cert-text">تُمنح هذه الشهادة إلى الطالب/ـة المتميّز/ـة</div>'
            '<div class="cert-name">' + name + '</div>'
            '<div class="cert-text">' + grade_text + '</div>'
            '<div class="cert-text">تقديراً لـ</div>'
            '<div class="cert-achievement">' + achievement + '</div>'
            '<div class="cert-footer">'
            '<div class="cert-sign-block"><div class="cert-sign-line">' + (teacher_name or "الأستاذ") + '</div></div>'
            '<div class="cert-sign-block"><div class="cert-sign-line">📅 ' + date + '</div></div>'
            '<div class="cert-sign-block"><div class="cert-sign-line">' + (school_name or "المدرسة") + '</div></div>'
            '</div>'
            '</div>'
            '</div>'
        )
        pages_html += page
    
    css = """
        * { box-sizing: border-box; margin: 0; padding: 0; }
        body { font-family: 'Cairo', sans-serif; padding: 0; margin: 0; background: #fff; }
        .cert-page { width: 297mm; height: 210mm; padding: 30px; background: """ + t["bg"] + """; border: """ + t["border"] + """; margin: 20px auto; box-sizing: border-box; position: relative; page-break-after: always; }
        .cert-corner { position: absolute; font-size: 60px; opacity: 0.15; color: """ + t["title_c"] + """; }
        .cert-corner.tl { top: 30px; left: 30px; }
        .cert-corner.tr { top: 30px; right: 30px; }
        .cert-corner.bl { bottom: 30px; left: 30px; }
        .cert-corner.br { bottom: 30px; right: 30px; }
        .cert-content { text-align: center; height: 100%; display: flex; flex-direction: column; justify-content: center; align-items: center; position: relative; z-index: 2; }
        .cert-icon { font-size: 80px; margin-bottom: 14px; }
        .cert-title { font-family: 'Amiri', serif; font-size: 56px; font-weight: 700; color: """ + t["title_c"] + """; margin-bottom: 8px; }
        .cert-subtitle { font-size: 22px; color: """ + t["title_c"] + """; margin-bottom: 30px; font-weight: 600; }
        .cert-text { font-size: 22px; line-height: 2; margin: 12px 0; color: #333; }
        .cert-name { font-family: 'Amiri', serif; font-size: 60px; font-weight: 700; color: """ + t["name_c"] + """; margin: 20px 0; padding: 14px 40px; border-bottom: 3px solid """ + t["title_c"] + """; display: inline-block; }
        .cert-achievement { font-size: 24px; color: """ + t["title_c"] + """; font-weight: 600; margin: 14px 0; padding: 10px 24px; background: rgba(255,255,255,0.5); border-radius: 8px; display: inline-block; }
        .cert-footer { display: flex; justify-content: space-around; width: 100%; margin-top: 40px; font-size: 16px; }
        .cert-sign-block { text-align: center; min-width: 200px; }
        .cert-sign-line { border-top: 2px solid """ + t["title_c"] + """; margin-top: 60px; padding-top: 6px; color: """ + t["name_c"] + """; font-weight: 600; }
        .controls { position: fixed; top: 12px; left: 12px; background: #fff; border: 1px solid #ccc; border-radius: 8px; padding: 10px; z-index: 9999; box-shadow: 0 2px 8px rgba(0,0,0,0.15); }
        .controls button { background: #1976d2; color: #fff; border: none; padding: 8px 16px; margin: 0 4px; border-radius: 6px; cursor: pointer; font-family: 'Cairo', sans-serif; font-weight: 600; }
        .controls .info { display: inline-block; padding: 6px 14px; color: #1565c0; font-size: 13px; margin-right: 10px; }
        @media print {
            .controls { display: none !important; }
            @page { size: A4 landscape; margin: 0; }
            body { margin: 0; }
            .cert-page { margin: 0; }
        }
    """
    
    full_html = (
        '<!DOCTYPE html>\n'
        '<html lang="ar" dir="rtl">\n'
        '<head>\n'
        '<meta charset="UTF-8">\n'
        '<title>شهادات دفعة (' + str(len(names)) + ')</title>\n'
        '<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
        '<style>' + css + '</style>\n'
        '</head>\n<body>\n'
        '<div class="controls">'
        '<span class="info">📊 ' + str(len(names)) + ' شهادة جاهزة</span>'
        '<button onclick="window.print()">🖨️ طباعة الكل / حفظ PDF</button>'
        '<button onclick="window.close()" style="background:#43a047;">✕ إغلاق</button>'
        '</div>\n'
        + pages_html +
        '\n</body>\n</html>'
    )
    
    from fastapi.responses import HTMLResponse
    return HTMLResponse(content=full_html)



# ═══════════════════════════════════════════════════════════════
# 🤖 TEACHER AI ASSISTANT — مساعد المعلم التربوي
# ═══════════════════════════════════════════════════════════════

@app.post("/api/teacher/ai_assistant")
async def teacher_ai_assistant(
    request: Request,
    message: str         = Form(...),
    history: str         = Form(default="[]"),     # JSON array of past messages
    grade: str           = Form(default=""),
    subject: str         = Form(default="الرياضيات"),
):
    """
    🤖 مساعد تربوي ذكي — يجيب أسئلة المعلم في:
    - شرح المفاهيم بطرق مختلفة
    - اقتراح أنشطة تعليمية
    - تحسين تحضير
    - معالجة صعوبات الطلاب
    - إعطاء أمثلة حياتية
    """
    # rate limit
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=20, window_seconds=60):
        raise HTTPException(status_code=429, detail="طلبات كثيرة جداً، انتظر دقيقة")
    
    if not message.strip():
        raise HTTPException(status_code=400, detail="الرسالة فارغة")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ خدمة AI غير مُهيّأة (GEMINI_API_KEY مفقود)")
    
    # نُحدّد دور AI بدقة
    system_prompt = f"""أنت مساعد تربوي ذكي متخصص في تعليم {subject} للمراحل المدرسية.
دورك مساعدة المعلم على:
1. شرح المفاهيم الرياضية بطرق مبسّطة ومتنوّعة
2. اقتراح أنشطة تعليمية وألعاب صفية ممتعة
3. تقديم أمثلة حياتية واقعية
4. معالجة صعوبات تعلّم الطلاب
5. تحسين خطط الدروس والتحضيرات
6. اقتراح طرق تقويم متنوعة

ملاحظات مهمة:
- استخدم لغة عربية فصيحة وسهلة
- أعطِ إجابات عملية تطبيقية
- استخدم نقاط منظمة عند الحاجة
- اقترح أكثر من حل/طريقة عند الإمكان
- كن إيجابياً ومُلهماً
{f'- المعلم يدرّس {grade}' if grade else ''}
"""
    
    # نبني تاريخ المحادثة
    contents = []
    try:
        import json as _json
        past = _json.loads(history) if history else []
        for msg in past[-10:]:  # آخر 10 رسائل فقط
            role = "user" if msg.get("role") == "user" else "model"
            text = msg.get("content", "").strip()
            if text:
                contents.append({"role": role, "parts": [{"text": text}]})
    except Exception:
        pass
    
    contents.append({"role": "user", "parts": [{"text": message}]})
    
    try:
        import httpx
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
        
        payload = {
            "system_instruction": {
                "parts": [{"text": system_prompt}]
            },
            "contents": contents,
            "generationConfig": {
                "temperature": 0.7,
                "maxOutputTokens": 1500,
            }
        }
        
        async with httpx.AsyncClient(timeout=60.0) as client:
            res = await client.post(url, json=payload)
            
            if res.status_code != 200:
                # نحاول flash-lite كاحتياط
                fallback_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-lite:generateContent?key={api_key}"
                res = await client.post(fallback_url, json=payload)
            
            if res.status_code != 200:
                err_data = res.json() if res.content else {}
                err_msg = (err_data.get("error", {}).get("message", "") or "")[:200]
                raise HTTPException(status_code=503, detail=f"AI رفض الاستجابة: {err_msg}")
            
            data = res.json()
            
            try:
                text = data["candidates"][0]["content"]["parts"][0]["text"]
            except (KeyError, IndexError):
                raise HTTPException(status_code=500, detail="استجابة AI فارغة")
            
            return {
                "status": "ok",
                "response": text,
                "model": "gemini-2.5-flash"
            }
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"[teacher_ai] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ═══════════════════════════════════════════════════════════════
# 📓 LESSON PLAN GENERATOR — منشئ التحضيرات بالـ AI
# ═══════════════════════════════════════════════════════════════

@app.post("/api/teacher/generate_lesson_plan")
async def generate_lesson_plan(
    request: Request,
    grade: str           = Form(...),
    lesson_name: str     = Form(...),
    duration: int        = Form(default=45),    # بالدقائق
    objectives_focus: str = Form(default=""),    # محاور تركيز خاصة
    style: str           = Form(default="standard"),  # standard | interactive | discovery
):
    """
    📓 يولّد تحضير درس متكامل بالـ AI ثم يبني HTML احترافي للطباعة
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=10, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    if not lesson_name.strip():
        raise HTTPException(status_code=400, detail="اسم الدرس مطلوب")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY غير مُهيّأ")
    
    style_hints = {
        "standard":    "تحضير تقليدي منظم بأسلوب وزاري",
        "interactive": "تحضير يركّز على التعلّم النشط والأنشطة التفاعلية",
        "discovery":   "تحضير بأسلوب التعلّم بالاكتشاف والاستقصاء",
    }
    style_text = style_hints.get(style, style_hints["standard"])
    
    prompt = f"""أنشئ تحضير درس متكامل واحترافي بصيغة JSON صارمة بالمواصفات التالية:

📚 المعلومات الأساسية:
- الصف: {grade}
- الدرس: {lesson_name}
- المدة: {duration} دقيقة
- الأسلوب: {style_text}
{f'- محاور التركيز: {objectives_focus}' if objectives_focus else ''}

⚙️ المطلوب JSON بالحقول التالية فقط (لا تُضف أي نص خارج JSON):

{{
  "general_objectives": ["هدف عام 1", "هدف عام 2", "هدف عام 3"],
  "behavioral_objectives": ["أن يتمكن الطالب من...", "أن يستنتج...", "أن يحلّ..."],
  "tools": ["أداة 1", "أداة 2"],
  "introduction": {{
    "duration": 5,
    "activities": ["نشاط التمهيد 1", "نشاط 2"]
  }},
  "presentation": {{
    "duration": 25,
    "steps": [
      {{"title": "عنوان الخطوة", "content": "شرح مفصّل + أمثلة"}}
    ]
  }},
  "activities": [
    {{"title": "نشاط 1", "type": "فردي/جماعي", "description": "...", "duration": 5}}
  ],
  "evaluation": {{
    "duration": 10,
    "questions": ["سؤال تقويمي 1", "سؤال 2", "سؤال 3"]
  }},
  "homework": "وصف الواجب المنزلي",
  "common_difficulties": ["صعوبة 1 وحلّها", "صعوبة 2"],
  "real_life_examples": ["مثال حياتي 1", "مثال 2"]
}}

⚠️ أعطِ JSON فقط بدون أي نص قبله أو بعده، بدون ```json أو ```.
"""
    
    try:
        import httpx
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={api_key}"
        
        payload = {
            "contents": [{"role": "user", "parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.7,
                "maxOutputTokens": 3000,
                "responseMimeType": "application/json",
            }
        }
        
        async with httpx.AsyncClient(timeout=90.0) as client:
            res = await client.post(url, json=payload)
            
            if res.status_code != 200:
                fallback_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash-lite:generateContent?key={api_key}"
                res = await client.post(fallback_url, json=payload)
            
            if res.status_code != 200:
                raise HTTPException(status_code=503, detail="AI لم يستجب")
            
            data = res.json()
            
            try:
                raw_text = data["candidates"][0]["content"]["parts"][0]["text"].strip()
                # تنظيف JSON
                if raw_text.startswith("```"):
                    raw_text = raw_text.split("```")[1]
                    if raw_text.startswith("json"):
                        raw_text = raw_text[4:]
                raw_text = raw_text.strip()
                
                import json as _json
                plan = _json.loads(raw_text)
            except Exception as e:
                raise HTTPException(status_code=500, detail=f"AI أخرج JSON غير صالح: {str(e)[:100]}")
            
            return {
                "status": "ok",
                "plan": plan,
                "meta": {
                    "grade": grade,
                    "lesson_name": lesson_name,
                    "duration": duration,
                    "style": style,
                }
            }
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"[lesson_plan] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/teacher/lesson_plan_pdf")
async def lesson_plan_pdf(
    plan_json: str       = Form(...),
    grade: str           = Form(...),
    lesson_name: str     = Form(...),
    duration: int        = Form(default=45),
    teacher_name: str    = Form(default=""),
    school_name: str     = Form(default=""),
    date: str            = Form(default=""),
):
    """📄 يحوّل تحضير JSON إلى HTML احترافي للطباعة"""
    try:
        import json as _json
        plan = _json.loads(plan_json)
        
        from datetime import datetime
        if not date:
            date = datetime.now().strftime("%Y/%m/%d")
        
        # نبني الأقسام
        def render_list(items, css_class=""):
            if not items: return ""
            html = '<ul class="' + css_class + '">'
            for it in items:
                if isinstance(it, dict):
                    html += '<li>'
                    if it.get("title"):
                        html += '<b>' + str(it["title"]) + ':</b> '
                    if it.get("description"):
                        html += str(it["description"])
                    if it.get("content"):
                        html += str(it["content"])
                    html += '</li>'
                else:
                    html += '<li>' + str(it) + '</li>'
            html += '</ul>'
            return html
        
        # الأهداف
        gen_obj = render_list(plan.get("general_objectives", []))
        beh_obj = render_list(plan.get("behavioral_objectives", []))
        tools = render_list(plan.get("tools", []))
        
        # التمهيد
        intro = plan.get("introduction", {})
        intro_html = ""
        if intro:
            intro_html = '<div class="time-badge">⏱️ ' + str(intro.get("duration", 5)) + ' دقائق</div>'
            intro_html += render_list(intro.get("activities", []))
        
        # العرض
        present = plan.get("presentation", {})
        present_html = ""
        if present:
            present_html = '<div class="time-badge">⏱️ ' + str(present.get("duration", 25)) + ' دقيقة</div>'
            for step in present.get("steps", []):
                if isinstance(step, dict):
                    present_html += '<div class="step-block">'
                    present_html += '<h4>📌 ' + str(step.get("title", "")) + '</h4>'
                    present_html += '<p>' + str(step.get("content", "")) + '</p>'
                    present_html += '</div>'
                else:
                    present_html += '<p>' + str(step) + '</p>'
        
        # الأنشطة
        activities = plan.get("activities", [])
        activities_html = ""
        for act in activities:
            if isinstance(act, dict):
                activities_html += '<div class="activity-block">'
                activities_html += '<div class="act-header">'
                activities_html += '<span class="act-title">🎯 ' + str(act.get("title", "")) + '</span>'
                if act.get("type"):
                    activities_html += '<span class="act-type">' + str(act["type"]) + '</span>'
                if act.get("duration"):
                    activities_html += '<span class="act-time">⏱️ ' + str(act["duration"]) + ' د</span>'
                activities_html += '</div>'
                activities_html += '<p>' + str(act.get("description", "")) + '</p>'
                activities_html += '</div>'
        
        # التقويم
        evaluation = plan.get("evaluation", {})
        eval_html = ""
        if evaluation:
            eval_html = '<div class="time-badge">⏱️ ' + str(evaluation.get("duration", 10)) + ' دقائق</div>'
            eval_html += render_list(evaluation.get("questions", []))
        
        # الواجب
        homework = plan.get("homework", "")
        difficulties = render_list(plan.get("common_difficulties", []))
        examples = render_list(plan.get("real_life_examples", []))
        
        css = """
            * { box-sizing: border-box; margin: 0; padding: 0; }
            body { font-family: 'Cairo', sans-serif; padding: 25px; max-width: 850px; margin: 0 auto; line-height: 1.8; color: #000; background: #fff; }
            .header { text-align: center; padding-bottom: 18px; margin-bottom: 24px; border-bottom: 3px double #1565c0; }
            .school { font-size: 17px; font-weight: 700; color: #555; margin-bottom: 4px; }
            h1 { font-family: 'Amiri', serif; font-size: 34px; color: #1565c0; margin: 6px 0; }
            .meta-grid { display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap: 12px; margin-top: 14px; }
            .meta-item { background: #e3f2fd; padding: 8px 14px; border-radius: 6px; font-size: 14px; border-right: 4px solid #1565c0; }
            section { margin: 22px 0; padding: 18px; background: #fafafa; border: 1px solid #eee; border-radius: 8px; page-break-inside: avoid; }
            section h2 { font-family: 'Amiri', serif; color: #1565c0; font-size: 22px; margin-bottom: 12px; padding-bottom: 6px; border-bottom: 2px solid #1565c0; }
            section ul { padding-right: 22px; }
            section li { margin: 6px 0; line-height: 1.9; }
            .time-badge { display: inline-block; background: #fff3e0; color: #e65100; padding: 4px 12px; border-radius: 12px; font-weight: 700; font-size: 13px; margin-bottom: 10px; }
            .step-block { margin: 12px 0; padding: 12px 14px; background: #fff; border-right: 4px solid #1976d2; border-radius: 4px; }
            .step-block h4 { color: #1976d2; margin-bottom: 6px; font-size: 16px; }
            .step-block p { font-size: 14px; line-height: 1.8; }
            .activity-block { margin: 10px 0; padding: 12px 14px; background: #fff; border-right: 4px solid #f57c00; border-radius: 4px; }
            .act-header { display: flex; gap: 10px; align-items: center; margin-bottom: 6px; flex-wrap: wrap; }
            .act-title { font-weight: 700; color: #e65100; font-size: 15px; }
            .act-type { background: #fff3e0; color: #e65100; padding: 2px 10px; border-radius: 10px; font-size: 12px; }
            .act-time { background: #e3f2fd; color: #1565c0; padding: 2px 10px; border-radius: 10px; font-size: 12px; }
            .homework-box { background: #fff8e1; border: 2px solid #ffa726; padding: 14px; border-radius: 8px; }
            .signature { display: flex; justify-content: space-between; margin-top: 40px; font-size: 14px; }
            .sign-block { text-align: center; min-width: 200px; }
            .sign-line { border-top: 2px solid #000; margin-top: 50px; padding-top: 6px; }
            .controls { position: fixed; top: 12px; left: 12px; background: #fff; border: 1px solid #ccc; border-radius: 8px; padding: 10px; box-shadow: 0 2px 8px rgba(0,0,0,0.15); z-index: 9999; }
            .controls button { background: #1976d2; color: #fff; border: none; padding: 8px 16px; margin: 0 4px; border-radius: 6px; cursor: pointer; font-family: 'Cairo', sans-serif; font-weight: 600; }
            .controls button.green { background: #43a047; }
            @media print { .controls { display: none !important; } body { padding: 15px; } @page { size: A4; margin: 1cm; } }
        """
        
        # نُغلّف الأقسام
        def section_block(title, content):
            if not content or content == "":
                return ""
            return '<section><h2>' + title + '</h2>' + content + '</section>'
        
        sections_html = ""
        sections_html += section_block("🎯 الأهداف العامة", gen_obj)
        sections_html += section_block("📋 الأهداف السلوكية", beh_obj)
        sections_html += section_block("🛠️ الوسائل والأدوات", tools)
        sections_html += section_block("🚀 التمهيد", intro_html)
        sections_html += section_block("📚 العرض والشرح", present_html)
        sections_html += section_block("📝 الأنشطة والتطبيقات", activities_html)
        sections_html += section_block("✅ التقويم", eval_html)
        if homework:
            sections_html += '<section><h2>📌 الواجب المنزلي</h2><div class="homework-box">' + homework + '</div></section>'
        sections_html += section_block("⚠️ صعوبات شائعة وحلولها", difficulties)
        sections_html += section_block("🌟 أمثلة من الحياة", examples)
        
        full_html = (
            '<!DOCTYPE html>\n<html lang="ar" dir="rtl">\n<head>\n'
            '<meta charset="UTF-8">\n<title>' + lesson_name + '</title>\n'
            '<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
            '<style>' + css + '</style>\n</head>\n<body>\n'
            '<div class="controls">'
            '<button onclick="window.print()">🖨️ طباعة / حفظ PDF</button>'
            '<button class="green" onclick="window.close()">✕ إغلاق</button>'
            '</div>\n'
            '<div class="header">'
            + (('<div class="school">' + school_name + '</div>') if school_name else '') +
            '<h1>📓 ' + lesson_name + '</h1>'
            '<div class="meta-grid">'
            '<div class="meta-item">📚 ' + grade + '</div>'
            '<div class="meta-item">⏱️ ' + str(duration) + ' دقيقة</div>'
            '<div class="meta-item">📅 ' + date + '</div>'
            + (('<div class="meta-item">👨‍🏫 ' + teacher_name + '</div>') if teacher_name else '') +
            '</div>'
            '</div>\n'
            + sections_html +
            '<div class="signature">'
            '<div class="sign-block"><div class="sign-line">المعلم ' + (teacher_name or '') + '</div></div>'
            '<div class="sign-block"><div class="sign-line">المراجع</div></div>'
            '</div>\n'
            '</body>\n</html>'
        )
        
        from fastapi.responses import HTMLResponse
        return HTMLResponse(content=full_html)
    
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ في بناء PDF: {str(e)[:200]}")


# ═══════════════════════════════════════════════════════════════
# 🎬 SHORT VIDEOS LIBRARY — مكتبة الفيديوهات للمعلمين
# ═══════════════════════════════════════════════════════════════

@app.get("/api/teacher/videos/list")
async def teacher_videos_list(
    grade: str = "",
    unit: str = "",
    lesson: str = "",
    search: str = "",
):
    """🎬 قائمة الفيديوهات التعليمية المتاحة للمعلمين"""
    try:
        # محاولة جلب من جدول short_videos
        try:
            res = supabase.table("short_videos").select("*").eq("is_active", True).order("created_at", desc=True).execute()
            all_videos = res.data or []
        except Exception:
            # لو الجدول غير موجود
            return {"status": "ok", "videos": [], "count": 0, "warning": "جدول الفيديوهات غير موجود بعد"}
        
        # فلترة
        def _norm(s): return str(s or "").strip().lower()
        
        if grade.strip():
            ng = _norm(grade)
            all_videos = [v for v in all_videos if 
                _norm(v.get("grade")) == ng or ng in _norm(v.get("grade")) or _norm(v.get("grade")) in ng]
        
        if unit.strip():
            nu = _norm(unit)
            all_videos = [v for v in all_videos if _norm(v.get("unit")) == nu or nu in _norm(v.get("unit"))]
        
        if lesson.strip():
            nl = _norm(lesson)
            all_videos = [v for v in all_videos if _norm(v.get("lesson")) == nl or nl in _norm(v.get("lesson"))]
        
        if search.strip():
            ns = _norm(search)
            all_videos = [v for v in all_videos if 
                ns in _norm(v.get("title")) or ns in _norm(v.get("description"))]
        
        return {
            "status": "ok",
            "videos": all_videos,
            "count": len(all_videos),
        }
    
    except Exception as e:
        print(f"[videos] error: {e}")
        return {"status": "error", "videos": [], "count": 0, "error": str(e)[:200]}


@app.post("/api/teacher/videos/{video_id}/track_view")
async def teacher_track_video_view(video_id: int):
    """📊 تتبّع مشاهدة فيديو من المعلم"""
    try:
        # ننقّب الـ view_count
        cur = supabase.table("short_videos").select("view_count").eq("id", video_id).maybe_single().execute()
        if cur and cur.data:
            new_count = (cur.data.get("view_count") or 0) + 1
            supabase.table("short_videos").update({"view_count": new_count}).eq("id", video_id).execute()
            return {"status": "ok", "view_count": new_count}
    except Exception as e:
        print(f"[track_view] error: {e}")
    return {"status": "ok"}



# ═══════════════════════════════════════════════════════════════
# 🏆 PROFESSIONAL CERTIFICATES — قوالب احترافية حديثة
# ═══════════════════════════════════════════════════════════════

# 🎨 8 قوالب احترافية متنوعة (4 جديدة + 4 كلاسيكية محسّنة)
CERT_TEMPLATES = {
    # ───── الكلاسيكية المحسّنة ─────
    "gold": {
        "name": "الذهبي الملكي",
        "primary": "#d4af37",
        "secondary": "#8b6914",
        "bg": "linear-gradient(135deg, #fefae0 0%, #fff3c4 50%, #fefae0 100%)",
        "title_color": "#5d4e1f",
        "name_color": "#3d2e0f",
        "accent": "#b8860b",
        "border": "double 14px",
        "icon": "👑",
        "decoration": "⚜️",
        "pattern": "royal",
    },
    "modern_blue": {
        "name": "الأزرق العصري",
        "primary": "#0ea5e9",
        "secondary": "#0c4a6e",
        "bg": "linear-gradient(135deg, #f0f9ff 0%, #e0f2fe 100%)",
        "title_color": "#0369a1",
        "name_color": "#082f49",
        "accent": "#0284c7",
        "border": "solid 8px",
        "icon": "🏆",
        "decoration": "★",
        "pattern": "modern",
    },
    "elegant_dark": {
        "name": "الأنيق الداكن",
        "primary": "#1e293b",
        "secondary": "#334155",
        "bg": "linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%)",
        "title_color": "#0f172a",
        "name_color": "#020617",
        "accent": "#475569",
        "border": "solid 3px",
        "icon": "🎖️",
        "decoration": "◆",
        "pattern": "minimal",
    },
    "vibrant_purple": {
        "name": "البنفسجي الحيوي",
        "primary": "#9333ea",
        "secondary": "#581c87",
        "bg": "linear-gradient(135deg, #faf5ff 0%, #f3e8ff 100%)",
        "title_color": "#6b21a8",
        "name_color": "#3b0764",
        "accent": "#a855f7",
        "border": "double 10px",
        "icon": "🌟",
        "decoration": "✦",
        "pattern": "modern",
    },
    # ───── الجديدة الإبداعية ─────
    "geometric_emerald": {
        "name": "الزمردي الهندسي",
        "primary": "#059669",
        "secondary": "#064e3b",
        "bg": "linear-gradient(135deg, #ecfdf5 0%, #d1fae5 100%)",
        "title_color": "#065f46",
        "name_color": "#022c22",
        "accent": "#10b981",
        "border": "solid 5px",
        "icon": "🌿",
        "decoration": "◊",
        "pattern": "geometric",
    },
    "sunset_orange": {
        "name": "غروب البرتقالي",
        "primary": "#ea580c",
        "secondary": "#7c2d12",
        "bg": "linear-gradient(135deg, #fff7ed 0%, #ffedd5 50%, #fed7aa 100%)",
        "title_color": "#9a3412",
        "name_color": "#431407",
        "accent": "#f97316",
        "border": "ridge 12px",
        "icon": "🔆",
        "decoration": "❋",
        "pattern": "sunset",
    },
    "minimal_white": {
        "name": "الأبيض البسيط",
        "primary": "#000000",
        "secondary": "#404040",
        "bg": "#ffffff",
        "title_color": "#000000",
        "name_color": "#000000",
        "accent": "#737373",
        "border": "solid 2px",
        "icon": "📜",
        "decoration": "—",
        "pattern": "minimal",
    },
    "rose_gold": {
        "name": "الذهبي الوردي",
        "primary": "#e11d48",
        "secondary": "#881337",
        "bg": "linear-gradient(135deg, #fff1f2 0%, #ffe4e6 50%, #fecdd3 100%)",
        "title_color": "#9f1239",
        "name_color": "#4c0519",
        "accent": "#f43f5e",
        "border": "double 8px",
        "icon": "🌹",
        "decoration": "❀",
        "pattern": "floral",
    },
}


def _build_cert_pattern_svg(pattern_type: str, primary: str, secondary: str) -> str:
    """🎨 يبني SVG decoration للنمط"""
    if pattern_type == "royal":
        return f'''
        <div class="cert-corner cert-corner-tl"><svg viewBox="0 0 80 80" width="80" height="80">
            <path d="M0,0 L80,0 L80,15 Q70,15 60,5 Q50,0 40,5 Q30,15 20,5 Q10,0 0,15 Z" fill="{primary}" opacity="0.3"/>
            <circle cx="15" cy="15" r="6" fill="{primary}" opacity="0.5"/>
        </svg></div>
        <div class="cert-corner cert-corner-tr"><svg viewBox="0 0 80 80" width="80" height="80" style="transform:scaleX(-1);">
            <path d="M0,0 L80,0 L80,15 Q70,15 60,5 Q50,0 40,5 Q30,15 20,5 Q10,0 0,15 Z" fill="{primary}" opacity="0.3"/>
            <circle cx="15" cy="15" r="6" fill="{primary}" opacity="0.5"/>
        </svg></div>
        <div class="cert-corner cert-corner-bl"><svg viewBox="0 0 80 80" width="80" height="80" style="transform:scaleY(-1);">
            <path d="M0,0 L80,0 L80,15 Q70,15 60,5 Q50,0 40,5 Q30,15 20,5 Q10,0 0,15 Z" fill="{primary}" opacity="0.3"/>
        </svg></div>
        <div class="cert-corner cert-corner-br"><svg viewBox="0 0 80 80" width="80" height="80" style="transform:scale(-1,-1);">
            <path d="M0,0 L80,0 L80,15 Q70,15 60,5 Q50,0 40,5 Q30,15 20,5 Q10,0 0,15 Z" fill="{primary}" opacity="0.3"/>
        </svg></div>
        '''
    elif pattern_type == "modern":
        return f'''
        <div class="cert-shape cert-shape-1" style="background:{primary};opacity:0.08;"></div>
        <div class="cert-shape cert-shape-2" style="background:{secondary};opacity:0.06;"></div>
        <div class="cert-stripe" style="background:linear-gradient(90deg,transparent,{primary},transparent);opacity:0.5;"></div>
        '''
    elif pattern_type == "geometric":
        return f'''
        <div class="cert-geo-pattern"></div>
        <svg style="position:absolute;top:30px;left:30px;" width="60" height="60" viewBox="0 0 60 60">
            <polygon points="30,5 55,20 55,40 30,55 5,40 5,20" fill="none" stroke="{primary}" stroke-width="2" opacity="0.4"/>
            <polygon points="30,15 45,25 45,35 30,45 15,35 15,25" fill="{primary}" opacity="0.2"/>
        </svg>
        <svg style="position:absolute;top:30px;right:30px;" width="60" height="60" viewBox="0 0 60 60">
            <polygon points="30,5 55,20 55,40 30,55 5,40 5,20" fill="none" stroke="{primary}" stroke-width="2" opacity="0.4"/>
            <polygon points="30,15 45,25 45,35 30,45 15,35 15,25" fill="{primary}" opacity="0.2"/>
        </svg>
        '''
    elif pattern_type == "sunset":
        return f'''
        <div class="cert-sun" style="background:radial-gradient(circle,{primary} 0%,transparent 70%);"></div>
        <div class="cert-rays"></div>
        '''
    elif pattern_type == "floral":
        return f'''
        <div class="cert-floral cert-floral-tl">❀</div>
        <div class="cert-floral cert-floral-tr">❀</div>
        <div class="cert-floral cert-floral-bl">❀</div>
        <div class="cert-floral cert-floral-br">❀</div>
        <div class="cert-floral-line"></div>
        '''
    else:  # minimal
        return f'''
        <div class="cert-minimal-line cert-line-top" style="background:{primary};"></div>
        <div class="cert-minimal-line cert-line-bottom" style="background:{primary};"></div>
        '''


def _build_cert_html(template_id: str, student_name: str, achievement: str,
                     grade: str = "", school_name: str = "", teacher_name: str = "",
                     date: str = "", extra_text: str = "",
                     custom_template_html: str = "") -> str:
    """🏗️ يبني HTML شهادة احترافية"""
    
    # 🎨 لو المعلم رفع قالباً مخصصاً، نستخدمه
    if custom_template_html:
        return _apply_custom_template(custom_template_html, {
            "{{name}}": student_name,
            "{{achievement}}": achievement,
            "{{grade}}": grade,
            "{{school}}": school_name,
            "{{teacher}}": teacher_name,
            "{{date}}": date,
            "{{extra}}": extra_text,
        })
    
    t = CERT_TEMPLATES.get(template_id, CERT_TEMPLATES["gold"])
    pattern_html = _build_cert_pattern_svg(t["pattern"], t["primary"], t["secondary"])
    
    grade_text = f'في {grade}' if grade else ''
    extra_block = f'<div class="cert-extra">{extra_text}</div>' if extra_text else ''
    
    css = f"""
        * {{ box-sizing: border-box; margin: 0; padding: 0; }}
        @page {{ size: A4 landscape; margin: 0; }}
        body {{ font-family: 'Cairo', sans-serif; padding: 0; margin: 0; background: #f0f0f0; }}
        
        .cert-page {{
            width: 297mm;
            height: 210mm;
            background: {t['bg']};
            margin: 20px auto;
            padding: 30px 50px 25px;
            box-sizing: border-box;
            position: relative;
            border: {t['border']} {t['primary']};
            overflow: hidden;
            page-break-after: always;
            box-shadow: 0 12px 40px rgba(0,0,0,0.15);
            display: flex;
            flex-direction: column;
            justify-content: space-between;
        }}
        
        /* الزخارف الزاوية */
        .cert-corner {{
            position: absolute;
            width: 80px;
            height: 80px;
        }}
        .cert-corner-tl {{ top: 20px; right: 20px; }}
        .cert-corner-tr {{ top: 20px; left: 20px; }}
        .cert-corner-bl {{ bottom: 20px; right: 20px; }}
        .cert-corner-br {{ bottom: 20px; left: 20px; }}
        
        /* الأشكال (modern pattern) */
        .cert-shape {{
            position: absolute;
            border-radius: 50%;
        }}
        .cert-shape-1 {{
            width: 350px; height: 350px;
            top: -120px; left: -120px;
        }}
        .cert-shape-2 {{
            width: 250px; height: 250px;
            bottom: -100px; right: -100px;
        }}
        .cert-stripe {{
            position: absolute;
            top: 50%; left: 0; right: 0;
            height: 2px;
            transform: translateY(-50%);
        }}
        
        /* النمط الهندسي */
        .cert-geo-pattern {{
            position: absolute;
            inset: 0;
            background-image: 
                linear-gradient(45deg, {t['primary']}10 25%, transparent 25%),
                linear-gradient(-45deg, {t['primary']}10 25%, transparent 25%);
            background-size: 30px 30px;
            opacity: 0.15;
        }}
        
        /* غروب */
        .cert-sun {{
            position: absolute;
            top: -120px; left: 50%;
            transform: translateX(-50%);
            width: 400px; height: 400px;
            border-radius: 50%;
            opacity: 0.4;
        }}
        .cert-rays {{
            position: absolute;
            inset: 0;
            background: repeating-linear-gradient(
                from 90deg at 50% 0%,
                transparent 0deg, transparent 20deg,
                {t['primary']}15 20deg, {t['primary']}15 25deg
            );
        }}
        
        /* الزهور */
        .cert-floral {{
            position: absolute;
            font-size: 60px;
            color: {t['primary']};
            opacity: 0.3;
        }}
        .cert-floral-tl {{ top: 20px; right: 30px; }}
        .cert-floral-tr {{ top: 20px; left: 30px; }}
        .cert-floral-bl {{ bottom: 20px; right: 30px; }}
        .cert-floral-br {{ bottom: 20px; left: 30px; }}
        .cert-floral-line {{
            position: absolute;
            top: 50%;
            left: 100px; right: 100px;
            height: 2px;
            background: linear-gradient(90deg, transparent, {t['primary']}, transparent);
            opacity: 0.4;
        }}
        
        /* الخطوط البسيطة */
        .cert-minimal-line {{
            position: absolute;
            left: 80px;
            right: 80px;
            height: 3px;
        }}
        .cert-line-top {{ top: 60px; }}
        .cert-line-bottom {{ bottom: 60px; }}
        
        /* المحتوى */
        .cert-content {{
            position: relative;
            z-index: 5;
            text-align: center;
            flex: 1;
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            min-height: 0;
            padding-bottom: 10px;
        }}
        
        .cert-icon {{
            font-size: 70px;
            margin-bottom: 12px;
            line-height: 1;
            filter: drop-shadow(0 4px 8px {t['primary']}40);
        }}
        
        .cert-title {{
            font-family: 'Amiri', serif;
            font-size: 52px;
            font-weight: 700;
            color: {t['title_color']};
            margin-bottom: 4px;
            letter-spacing: 2px;
            text-shadow: 1px 1px 2px rgba(0,0,0,0.05);
        }}
        
        .cert-subtitle {{
            font-size: 18px;
            color: {t['accent']};
            margin-bottom: 22px;
            font-weight: 600;
            letter-spacing: 4px;
            text-transform: uppercase;
        }}
        
        .cert-divider {{
            width: 200px;
            height: 3px;
            background: linear-gradient(90deg, transparent, {t['primary']}, transparent);
            margin: 12px auto 18px;
        }}
        
        .cert-text {{
            font-size: 19px;
            line-height: 2;
            margin: 8px 0;
            color: {t['secondary']};
            font-weight: 500;
        }}
        
        .cert-name {{
            font-family: 'Amiri', serif;
            font-size: 64px;
            font-weight: 700;
            color: {t['name_color']};
            margin: 14px 0;
            padding: 10px 50px;
            display: inline-block;
            position: relative;
        }}
        .cert-name::before, .cert-name::after {{
            content: '';
            position: absolute;
            top: 50%;
            width: 30px;
            height: 2px;
            background: {t['primary']};
        }}
        .cert-name::before {{ right: 10px; }}
        .cert-name::after {{ left: 10px; }}
        
        .cert-name-line {{
            width: 60%;
            height: 2px;
            background: {t['primary']};
            margin: 0 auto 14px;
            opacity: 0.6;
        }}
        
        .cert-grade {{
            font-size: 20px;
            color: {t['accent']};
            font-weight: 600;
            margin: 6px 0;
        }}
        
        .cert-achievement {{
            font-size: 22px;
            color: {t['title_color']};
            font-weight: 700;
            margin: 14px 0;
            padding: 12px 32px;
            background: rgba(255,255,255,0.6);
            border-radius: 12px;
            display: inline-block;
            border: 2px solid {t['primary']}40;
            backdrop-filter: blur(4px);
        }}
        
        .cert-extra {{
            font-size: 16px;
            color: {t['secondary']};
            margin: 12px 40px;
            font-style: italic;
            opacity: 0.85;
        }}
        
        /* التذييل — داخل الإطار، مضمون أن لا يخرج */
        .cert-footer {{
            display: flex;
            justify-content: space-around;
            align-items: flex-end;
            width: 100%;
            max-width: 100%;
            margin: 0 auto;
            padding: 14px 20px 6px;
            font-size: 13px;
            gap: 16px;
            box-sizing: border-box;
            position: relative;
            z-index: 6;
            flex-shrink: 0;
        }}
        .cert-sign-block {{
            text-align: center;
            flex: 1 1 0;
            min-width: 0;
            max-width: 32%;
            overflow: hidden;
        }}
        .cert-sign-line {{
            border-top: 2px solid {t['primary']};
            margin-top: 20px;
            padding-top: 5px;
            color: {t['name_color']};
            font-weight: 600;
            font-size: 13px;
            white-space: nowrap;
            overflow: hidden;
            text-overflow: ellipsis;
            line-height: 1.4;
        }}
        .cert-sign-label {{
            font-size: 10px;
            color: {t['secondary']};
            opacity: 0.7;
            margin-top: 2px;
            display: block;
        }}
        
        /* الختم */
        .cert-seal {{
            position: absolute;
            bottom: 30px;
            right: 50px;
            width: 80px;
            height: 80px;
            border: 4px double {t['primary']};
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 32px;
            color: {t['primary']};
            background: rgba(255,255,255,0.7);
            transform: rotate(-15deg);
            opacity: 0.7;
            z-index: 4;
        }}
        
        /* أزرار التحكم */
        .controls {{
            position: fixed;
            top: 12px;
            left: 12px;
            background: #fff;
            border: 1px solid #ccc;
            border-radius: 10px;
            padding: 12px;
            z-index: 9999;
            box-shadow: 0 4px 16px rgba(0,0,0,0.2);
        }}
        .controls button {{
            background: linear-gradient(135deg, #1976d2, #1565c0);
            color: #fff;
            border: none;
            padding: 10px 20px;
            margin: 0 4px;
            border-radius: 8px;
            cursor: pointer;
            font-family: 'Cairo', sans-serif;
            font-weight: 700;
            font-size: 14px;
            box-shadow: 0 2px 6px rgba(0,0,0,0.15);
        }}
        .controls button.green {{
            background: linear-gradient(135deg, #43a047, #2e7d32);
        }}
        .controls button:hover {{ transform: translateY(-2px); }}
        
        @media print {{
            .controls {{ display: none !important; }}
            body {{ background: #fff; margin: 0; padding: 0; }}
            .cert-page {{ margin: 0 !important; box-shadow: none !important; }}
        }}
    """
    
    full_html = (
        f'<!DOCTYPE html>\n<html lang="ar" dir="rtl">\n<head>\n'
        f'<meta charset="UTF-8">\n<title>شهادة - {student_name}</title>\n'
        f'<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">\n'
        f'<style>{css}</style>\n</head>\n<body>\n'
        f'<div class="controls">'
        f'<button onclick="window.print()">🖨️ طباعة / حفظ PDF</button>'
        f'<button class="green" onclick="window.close()">✕ إغلاق</button>'
        f'</div>\n'
        f'<div class="cert-page">'
        f'{pattern_html}'
        f'<div class="cert-content">'
        f'<div class="cert-icon">{t["icon"]}</div>'
        f'<div class="cert-title">شهادة تقدير</div>'
        f'<div class="cert-subtitle">Certificate of Excellence</div>'
        f'<div class="cert-divider"></div>'
        f'<div class="cert-text">تُمنح هذه الشهادة بكل فخر إلى الطالب/ـة المتميّز/ـة</div>'
        f'<div class="cert-name">{student_name}</div>'
        f'<div class="cert-name-line"></div>'
        f'{f"<div class=\"cert-grade\">{grade_text}</div>" if grade else ""}'
        f'<div class="cert-text">تقديراً لـ</div>'
        f'<div class="cert-achievement">{achievement}</div>'
        f'{extra_block}'
        f'<div class="cert-footer">'
        f'<div class="cert-sign-block"><div class="cert-sign-line">{teacher_name or "المعلم"}<span class="cert-sign-label">توقيع المعلم</span></div></div>'
        f'<div class="cert-sign-block"><div class="cert-sign-line">📅 {date}<span class="cert-sign-label">التاريخ</span></div></div>'
        f'<div class="cert-sign-block"><div class="cert-sign-line">{school_name or "المدرسة"}<span class="cert-sign-label">المدرسة</span></div></div>'
        f'</div>'
        f'</div>'
        f'<div class="cert-seal">{t["icon"]}</div>'
        f'</div>\n'
        f'</body>\n</html>'
    )
    return full_html


def _apply_custom_template(template_html: str, replacements: dict) -> str:
    """🎨 يطبّق متغيرات على قالب HTML مخصص"""
    result = template_html
    for placeholder, value in replacements.items():
        result = result.replace(placeholder, str(value or ""))
    
    # نُضيف أزرار التحكم لو لم تكن موجودة
    if 'window.print()' not in result:
        controls = '''
<div style="position:fixed;top:12px;left:12px;background:#fff;border:1px solid #ccc;border-radius:10px;padding:12px;z-index:9999;box-shadow:0 4px 16px rgba(0,0,0,0.2);">
    <button onclick="window.print()" style="background:linear-gradient(135deg,#1976d2,#1565c0);color:#fff;border:none;padding:10px 20px;margin:0 4px;border-radius:8px;cursor:pointer;font-family:Cairo,sans-serif;font-weight:700;">🖨️ طباعة / حفظ PDF</button>
    <button onclick="window.close()" style="background:linear-gradient(135deg,#43a047,#2e7d32);color:#fff;border:none;padding:10px 20px;margin:0 4px;border-radius:8px;cursor:pointer;font-family:Cairo,sans-serif;font-weight:700;">✕ إغلاق</button>
</div>
<style>@media print { div[style*="position:fixed"] { display:none !important; } }</style>
'''
        if '</body>' in result:
            result = result.replace('</body>', controls + '</body>')
        else:
            result += controls
    
    return result


# ═══════════════════════════════════════════════════════════════
# 🏆 الـ Endpoints الجديدة (تستبدل القديمة)
# ═══════════════════════════════════════════════════════════════

@app.get("/api/teacher/certificate/templates")
async def teacher_cert_templates():
    """📋 قائمة القوالب المتاحة"""
    return {
        "templates": [
            {"id": k, "name": v["name"], "icon": v["icon"], 
             "primary": v["primary"], "bg": v["bg"]}
            for k, v in CERT_TEMPLATES.items()
        ]
    }


@app.post("/api/teacher/certificate/build_v2")
async def teacher_certificate_build_v2(
    template: str       = Form(default="gold"),
    student_name: str   = Form(...),
    achievement: str    = Form(default="التميّز في الرياضيات"),
    grade: str          = Form(default=""),
    school_name: str    = Form(default=""),
    teacher_name: str   = Form(default=""),
    date: str           = Form(default=""),
    extra_text: str     = Form(default=""),
    custom_template: str = Form(default=""),  # HTML مخصص اختياري
):
    """🏆 بناء شهادة بالقالب المختار أو قالب مخصص"""
    try:
        from datetime import datetime
        if not date:
            date = datetime.now().strftime("%Y/%m/%d")
        
        if not student_name.strip():
            raise HTTPException(status_code=400, detail="اسم الطالب مطلوب")
        
        full_html = _build_cert_html(
            template_id=template,
            student_name=student_name.strip(),
            achievement=achievement.strip(),
            grade=grade.strip(),
            school_name=school_name.strip(),
            teacher_name=teacher_name.strip(),
            date=date.strip(),
            extra_text=extra_text.strip(),
            custom_template_html=custom_template.strip() if custom_template else "",
        )
        
        from fastapi.responses import HTMLResponse
        return HTMLResponse(content=full_html)
    
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/teacher/certificate/batch_v2")
async def teacher_certificate_batch_v2(
    template: str       = Form(default="gold"),
    student_names: str  = Form(...),
    achievement: str    = Form(default="التميّز في الرياضيات"),
    grade: str          = Form(default=""),
    school_name: str    = Form(default=""),
    teacher_name: str   = Form(default=""),
    date: str           = Form(default=""),
    custom_template: str = Form(default=""),
):
    """🏆 بناء شهادات متعددة"""
    names = [n.strip() for n in student_names.replace("،", "\n").split("\n") if n.strip()]
    if not names:
        raise HTTPException(status_code=400, detail="أدخل أسماء الطلاب")
    if len(names) > 100:
        raise HTTPException(status_code=400, detail="الحد الأقصى 100 شهادة")
    
    from datetime import datetime
    if not date:
        date = datetime.now().strftime("%Y/%m/%d")
    
    # نبني كل صفحة
    pages = []
    for name in names:
        if custom_template:
            page_html = _apply_custom_template(custom_template, {
                "{{name}}": name, "{{achievement}}": achievement,
                "{{grade}}": grade, "{{school}}": school_name,
                "{{teacher}}": teacher_name, "{{date}}": date,
            })
            # نستخرج body فقط
            if '<body' in page_html:
                start = page_html.find('<body')
                start = page_html.find('>', start) + 1
                end = page_html.find('</body>')
                pages.append(page_html[start:end])
            else:
                pages.append(page_html)
        else:
            single = _build_cert_html(
                template_id=template, student_name=name, achievement=achievement,
                grade=grade, school_name=school_name, teacher_name=teacher_name, date=date,
            )
            # استخرج body
            start = single.find('<body')
            start = single.find('>', start) + 1
            end = single.find('</body>')
            body_content = single[start:end]
            # احذف controls
            controls_start = body_content.find('<div class="controls">')
            if controls_start >= 0:
                controls_end = body_content.find('</div>', controls_start) + 6
                body_content = body_content[:controls_start] + body_content[controls_end:]
            pages.append(body_content)
    
    # نستخرج <head> من أول صفحة (للـ CSS)
    if not custom_template:
        first_full = _build_cert_html(
            template_id=template, student_name="x", achievement=achievement,
            grade=grade, school_name=school_name, teacher_name=teacher_name, date=date,
        )
        head_start = first_full.find('<head>')
        head_end = first_full.find('</head>') + 7
        head_content = first_full[head_start:head_end]
    else:
        head_content = '''<head>
<meta charset="UTF-8">
<title>شهادات</title>
<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&display=swap" rel="stylesheet">
</head>'''
    
    final_html = (
        '<!DOCTYPE html><html lang="ar" dir="rtl">'
        + head_content +
        '<body style="margin:0;padding:0;background:#f0f0f0;">'
        '<div style="position:fixed;top:12px;left:12px;background:#fff;border:1px solid #ccc;border-radius:10px;padding:12px;z-index:9999;box-shadow:0 4px 16px rgba(0,0,0,0.2);">'
        f'<span style="display:inline-block;padding:6px 14px;color:#1565c0;font-size:13px;font-weight:600;">📊 {len(names)} شهادة</span>'
        '<button onclick="window.print()" style="background:linear-gradient(135deg,#1976d2,#1565c0);color:#fff;border:none;padding:10px 20px;border-radius:8px;cursor:pointer;font-family:Cairo;font-weight:700;">🖨️ طباعة الكل / حفظ PDF</button>'
        '<button onclick="window.close()" style="background:linear-gradient(135deg,#43a047,#2e7d32);color:#fff;border:none;padding:10px 20px;margin-right:6px;border-radius:8px;cursor:pointer;font-family:Cairo;font-weight:700;">✕ إغلاق</button>'
        '</div>'
        '<style>@media print { div[style*="position:fixed"] { display:none !important; } }</style>'
        + ''.join(pages) +
        '</body></html>'
    )
    
    from fastapi.responses import HTMLResponse
    return HTMLResponse(content=final_html)



# ════════════════════════════════════════════════════════════════
# 🏛️ TEACHER CLASSROOMS — نظام الفصول الخاصة بالمعلم
# ════════════════════════════════════════════════════════════════

import secrets as _crypto_secrets
import string as _string

def _generate_classroom_code():
    """🎫 يُولّد كود فصل فريد: XXX-1234"""
    letters = ''.join(_crypto_secrets.choice(_string.ascii_uppercase) for _ in range(3))
    digits = ''.join(_crypto_secrets.choice(_string.digits) for _ in range(4))
    return f"{letters}-{digits}"


# ═══════════ Endpoints للمعلم ═══════════


# ═══════════════════════════════════════════════════════════════
# 🖼️ شهادة من صورة فارغة (يرفعها المعلم)
# ═══════════════════════════════════════════════════════════════

@app.post("/api/teacher/certificate/build_from_image")
async def teacher_certificate_from_image(
    image_data: str       = Form(...),       # base64 image
    student_name: str     = Form(...),
    achievement: str      = Form(default="التميّز في الرياضيات"),
    grade: str            = Form(default=""),
    school_name: str      = Form(default=""),
    teacher_name: str     = Form(default=""),
    date: str             = Form(default=""),
    
    # مواقع النصوص (نسبة مئوية من الصورة)
    name_x: float         = Form(default=50),
    name_y: float         = Form(default=45),
    name_size: int        = Form(default=64),
    name_color: str       = Form(default="#000000"),
    name_font: str        = Form(default="Amiri"),
    
    achievement_x: float  = Form(default=50),
    achievement_y: float  = Form(default=58),
    achievement_size: int = Form(default=24),
    achievement_color: str = Form(default="#444"),
    
    teacher_x: float      = Form(default=20),
    teacher_y: float      = Form(default=85),
    teacher_size: int     = Form(default=18),
    
    school_x: float       = Form(default=80),
    school_y: float       = Form(default=85),
    school_size: int      = Form(default=18),
    
    date_x: float         = Form(default=50),
    date_y: float         = Form(default=92),
    date_size: int        = Form(default=16),
):
    """🖼️ ينشئ شهادة من صورة فارغة + كتابة عليها"""
    try:
        if not student_name.strip():
            raise HTTPException(status_code=400, detail="اسم الطالب مطلوب")
        
        from datetime import datetime
        if not date:
            date = datetime.now().strftime("%Y/%m/%d")
        
        # بناء النصوص (overlays)
        overlays_html = ""
        
        # اسم الطالب
        overlays_html += (
            f'<div class="text-overlay" style="right:{100-name_x}%; top:{name_y}%; '
            f'font-family:\'{name_font}\',serif; font-size:{name_size}px; '
            f'color:{name_color}; font-weight:bold;">'
            f'{student_name}</div>\n'
        )
        
        # سبب التقدير
        if achievement:
            overlays_html += (
                f'<div class="text-overlay" style="right:{100-achievement_x}%; top:{achievement_y}%; '
                f'font-size:{achievement_size}px; color:{achievement_color}; font-weight:600;">'
                f'{achievement}</div>\n'
            )
        
        # الصف (تحت الاسم)
        if grade:
            overlays_html += (
                f'<div class="text-overlay" style="right:{100-name_x}%; top:{name_y + 7}%; '
                f'font-size:{achievement_size - 4}px; color:#666;">'
                f'{grade}</div>\n'
            )
        
        # المعلم
        if teacher_name:
            overlays_html += (
                f'<div class="text-overlay" style="right:{100-teacher_x}%; top:{teacher_y}%; '
                f'font-size:{teacher_size}px; color:#444; font-weight:600;">'
                f'{teacher_name}</div>\n'
            )
        
        # المدرسة
        if school_name:
            overlays_html += (
                f'<div class="text-overlay" style="right:{100-school_x}%; top:{school_y}%; '
                f'font-size:{school_size}px; color:#444; font-weight:600;">'
                f'{school_name}</div>\n'
            )
        
        # التاريخ
        overlays_html += (
            f'<div class="text-overlay" style="right:{100-date_x}%; top:{date_y}%; '
            f'font-size:{date_size}px; color:#666;">'
            f'📅 {date}</div>\n'
        )
        
        # بناء HTML الكامل (multi-line string واحد نظيف)
        full_html = f"""<!DOCTYPE html>
<html lang="ar" dir="rtl">
<head>
<meta charset="UTF-8">
<title>شهادة - {student_name}</title>
<link href="https://fonts.googleapis.com/css2?family=Cairo:wght@400;600;700;900&family=Amiri:wght@700&family=Tajawal:wght@500;700&display=swap" rel="stylesheet">
<style>
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
@page {{ size: A4 landscape; margin: 0; }}
body {{ font-family: "Cairo", sans-serif; padding: 0; margin: 0; background: #f0f0f0; }}

.cert-page {{
    width: 297mm;
    height: 210mm;
    margin: 20px auto;
    position: relative;
    background-image: url("{image_data}");
    background-size: 100% 100%;
    background-repeat: no-repeat;
    background-position: center;
    box-shadow: 0 12px 40px rgba(0,0,0,0.2);
    page-break-after: always;
}}

.text-overlay {{
    position: absolute;
    transform: translate(50%, -50%);
    text-align: center;
    white-space: nowrap;
    text-shadow: 0 1px 3px rgba(255,255,255,0.6);
}}

.controls {{
    position: fixed;
    top: 12px;
    left: 12px;
    background: #fff;
    border: 1px solid #ccc;
    border-radius: 10px;
    padding: 12px;
    z-index: 9999;
    box-shadow: 0 4px 16px rgba(0,0,0,0.2);
}}

.controls button {{
    background: linear-gradient(135deg, #1976d2, #1565c0);
    color: #fff;
    border: none;
    padding: 10px 20px;
    margin: 0 4px;
    border-radius: 8px;
    cursor: pointer;
    font-family: "Cairo";
    font-weight: 700;
    font-size: 14px;
}}

.controls button.green {{
    background: linear-gradient(135deg, #43a047, #2e7d32);
}}

@media print {{
    .controls {{ display: none !important; }}
    body {{ background: #fff; margin: 0; padding: 0; }}
    .cert-page {{ margin: 0 !important; box-shadow: none !important; }}
}}
</style>
</head>
<body>

<div class="controls">
    <button onclick="window.print()">🖨️ طباعة / حفظ PDF</button>
    <button class="green" onclick="window.close()">✕ إغلاق</button>
</div>

<div class="cert-page">
{overlays_html}
</div>

</body>
</html>"""
        
        from fastapi.responses import HTMLResponse
        return HTMLResponse(content=full_html)
    
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



@app.post("/api/teacher/classrooms/create")
async def teacher_create_classroom(
    teacher_id: int     = Form(...),
    name: str           = Form(...),
    description: str    = Form(default=""),
    grade: str          = Form(default=""),
    cover_color: str    = Form(default="#1976d2"),
    cover_emoji: str    = Form(default="📚"),
    max_students: int   = Form(default=50),
    requires_approval: bool = Form(default=False),
):
    """🏛️ إنشاء فصل دراسي جديد"""
    if not name.strip():
        raise HTTPException(status_code=400, detail="اسم الفصل مطلوب")
    if len(name) > 120:
        raise HTTPException(status_code=400, detail="اسم الفصل طويل جداً")
    
    # نتأكد من وجود المعلم
    try:
        t = supabase.table("teachers").select("id").eq("id", teacher_id).maybe_single().execute()
        if not t or not t.data:
            raise HTTPException(status_code=404, detail="المعلم غير موجود")
    except HTTPException:
        raise
    except Exception:
        pass  # نسمح بالاستمرار حتى لو لم نتمكن من التحقق
    
    # نولّد كوداً فريداً
    code = None
    for _ in range(10):
        candidate = _generate_classroom_code()
        try:
            existing = supabase.table("teacher_classrooms").select("id").eq("code", candidate).maybe_single().execute()
            if not existing or not existing.data:
                code = candidate
                break
        except Exception:
            code = candidate
            break
    
    if not code:
        raise HTTPException(status_code=500, detail="فشل توليد كود فريد")
    
    try:
        res = supabase.table("teacher_classrooms").insert({
            "owner_teacher_id": teacher_id,
            "name": name.strip(),
            "code": code,
            "description": description.strip()[:1000],
            "grade": grade.strip()[:80],
            "cover_color": cover_color[:20],
            "cover_emoji": cover_emoji[:10],
            "max_students": max(5, min(int(max_students), 200)),
            "requires_approval": bool(requires_approval),
            "is_active": True,
        }).execute()
        
        if not res.data:
            raise HTTPException(status_code=500, detail="فشل إنشاء الفصل (لا بيانات راجعة)")
        
        return {
            "status": "ok",
            "classroom": res.data[0],
            "message": f"✅ تم إنشاء الفصل! الكود: {code}"
        }
    except HTTPException:
        raise
    except Exception as e:
        err_str = str(e).lower()
        # رسائل واضحة عن المشاكل الشائعة
        if "relation" in err_str or "does not exist" in err_str or "table" in err_str:
            raise HTTPException(status_code=500, detail="❌ جدول الفصول غير موجود! يجب تشغيل classrooms_migration.sql في Supabase أولاً")
        if "violates" in err_str or "constraint" in err_str:
            raise HTTPException(status_code=400, detail=f"❌ خطأ في البيانات: {str(e)[:150]}")
        print(f"[create_classroom] error: {e}")
        raise HTTPException(status_code=500, detail=f"❌ خطأ: {str(e)[:200]}")


@app.get("/api/teacher/classrooms/my")
async def teacher_my_classrooms(teacher_id: int):
    """📋 قائمة فصول المعلم"""
    try:
        res = supabase.table("teacher_classrooms").select("*").eq(
            "owner_teacher_id", teacher_id
        ).order("created_at", desc=True).execute()
        
        classrooms = res.data or []
        
        # نُضيف عدد الطلاب لكل فصل
        for c in classrooms:
            try:
                count_res = supabase.table("classroom_members").select(
                    "id", count="exact"
                ).eq("classroom_id", c["id"]).eq("status", "active").execute()
                c["members_count"] = count_res.count or 0
            except Exception:
                c["members_count"] = 0
        
        return {"status": "ok", "classrooms": classrooms, "count": len(classrooms)}
    except Exception as e:
        print(f"[my_classrooms] error: {e}")
        return {"status": "error", "classrooms": [], "count": 0}


@app.get("/api/teacher/classrooms/{classroom_id}")
async def teacher_classroom_details(classroom_id: int, teacher_id: int):
    """📊 تفاصيل فصل (للمعلم المالك فقط)"""
    try:
        # نتأكد من الملكية
        res = supabase.table("teacher_classrooms").select("*").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود أو ليس لك")
        
        classroom = res.data
        
        # الأعضاء
        members_res = supabase.table("classroom_members").select(
            "*"
        ).eq("classroom_id", classroom_id).order("joined_at", desc=False).execute()
        members = members_res.data or []
        
        # نضيف بيانات الطلاب
        if members:
            student_ids = [m["student_id"] for m in members]
            try:
                students_res = supabase.table("students").select(
                    "id, full_name, username, grade, total_points"
                ).in_("id", student_ids).execute()
                students_map = {s["id"]: s for s in (students_res.data or [])}
                for m in members:
                    s = students_map.get(m["student_id"], {})
                    m["full_name"] = s.get("full_name", f"طالب #{m['student_id']}")
                    m["username"] = s.get("username", "")
                    m["grade"] = s.get("grade", "")
                    m["total_points"] = s.get("total_points", 0)
            except Exception:
                pass
        
        return {
            "status": "ok",
            "classroom": classroom,
            "members": members,
            "members_count": len(members),
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"[classroom_details] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.put("/api/teacher/classrooms/{classroom_id}")
async def teacher_update_classroom(
    classroom_id: int,
    teacher_id: int     = Form(...),
    name: str           = Form(default=""),
    description: str    = Form(default=""),
    grade: str          = Form(default=""),
    cover_color: str    = Form(default=""),
    cover_emoji: str    = Form(default=""),
    is_active: bool     = Form(default=True),
):
    """✏️ تعديل بيانات الفصل"""
    try:
        # تحقق من الملكية
        res = supabase.table("teacher_classrooms").select("id").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود أو ليس لك")
        
        update_data = {"is_active": is_active}
        if name.strip(): update_data["name"] = name.strip()[:120]
        if description: update_data["description"] = description.strip()[:1000]
        if grade: update_data["grade"] = grade.strip()[:80]
        if cover_color: update_data["cover_color"] = cover_color[:20]
        if cover_emoji: update_data["cover_emoji"] = cover_emoji[:10]
        
        supabase.table("teacher_classrooms").update(update_data).eq("id", classroom_id).execute()
        return {"status": "ok", "message": "✅ تم التحديث"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/teacher/classrooms/{classroom_id}")
async def teacher_delete_classroom(classroom_id: int, teacher_id: int):
    """🗑️ حذف الفصل (مع كل بياناته)"""
    try:
        res = supabase.table("teacher_classrooms").select("id").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود")
        
        supabase.table("teacher_classrooms").delete().eq("id", classroom_id).execute()
        return {"status": "ok", "message": "🗑️ تم الحذف"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ═══════════ إدارة الأعضاء ═══════════

@app.delete("/api/teacher/classrooms/{classroom_id}/members/{student_id}")
async def teacher_remove_member(classroom_id: int, student_id: int, teacher_id: int):
    """🚫 حذف عضو من الفصل"""
    try:
        # تحقق من الملكية
        res = supabase.table("teacher_classrooms").select("id").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود")
        
        supabase.table("classroom_members").delete().eq(
            "classroom_id", classroom_id
        ).eq("student_id", student_id).execute()
        return {"status": "ok", "message": "🚫 تم الإزالة"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ═══════════ الإعلانات ═══════════

@app.post("/api/teacher/classrooms/{classroom_id}/announcements")
async def teacher_post_announcement(
    classroom_id: int,
    teacher_id: int     = Form(...),
    title: str          = Form(...),
    body: str           = Form(default=""),
    type: str           = Form(default="info"),
    pinned: bool        = Form(default=False),
):
    """📢 نشر إعلان للفصل"""
    try:
        # تحقق من الملكية
        res = supabase.table("teacher_classrooms").select("id").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود")
        
        if type not in ("info", "warning", "exam", "celebration"):
            type = "info"
        
        ins = supabase.table("classroom_announcements").insert({
            "classroom_id": classroom_id,
            "title": title.strip()[:200],
            "body": body.strip()[:5000],
            "type": type,
            "pinned": bool(pinned),
        }).execute()
        
        return {"status": "ok", "announcement": ins.data[0] if ins.data else None}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teacher/classrooms/{classroom_id}/announcements")
async def teacher_list_announcements(classroom_id: int):
    """📋 قائمة إعلانات الفصل"""
    try:
        res = supabase.table("classroom_announcements").select("*").eq(
            "classroom_id", classroom_id
        ).order("pinned", desc=True).order("created_at", desc=True).limit(50).execute()
        return {"status": "ok", "announcements": res.data or []}
    except Exception as e:
        return {"status": "error", "announcements": [], "error": str(e)[:200]}


@app.delete("/api/teacher/classrooms/announcements/{announcement_id}")
async def teacher_delete_announcement(announcement_id: int, teacher_id: int):
    """🗑️ حذف إعلان"""
    try:
        # نتحقق أن الإعلان من فصل يملكه المعلم
        ann = supabase.table("classroom_announcements").select("classroom_id").eq(
            "id", announcement_id
        ).maybe_single().execute()
        if not ann or not ann.data:
            raise HTTPException(status_code=404, detail="الإعلان غير موجود")
        
        cls = supabase.table("teacher_classrooms").select("id").eq(
            "id", ann.data["classroom_id"]
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not cls or not cls.data:
            raise HTTPException(status_code=403, detail="ليست لك صلاحية")
        
        supabase.table("classroom_announcements").delete().eq("id", announcement_id).execute()
        return {"status": "ok"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ═══════════ اختبارات الفصل ═══════════

@app.post("/api/teacher/classrooms/{classroom_id}/exams")
async def teacher_create_exam(
    classroom_id: int,
    teacher_id: int     = Form(...),
    title: str          = Form(...),
    description: str    = Form(default=""),
    questions_json: str = Form(...),       # JSON array
    total_marks: int    = Form(default=20),
    duration_min: int   = Form(default=30),
    scheduled_at: str   = Form(default=""),
    expires_at: str     = Form(default=""),
    is_published: bool  = Form(default=False),
):
    """📝 إنشاء اختبار للفصل"""
    try:
        # تحقق من الملكية
        res = supabase.table("teacher_classrooms").select("id").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود")
        
        # نتحقق من JSON
        import json as _json
        try:
            questions = _json.loads(questions_json)
            if not isinstance(questions, list) or not questions:
                raise ValueError("لا أسئلة")
        except Exception:
            raise HTTPException(status_code=400, detail="JSON الأسئلة غير صالح")
        
        data = {
            "classroom_id": classroom_id,
            "title": title.strip()[:200],
            "description": description.strip()[:1000],
            "questions_json": questions,
            "total_marks": max(1, min(int(total_marks), 200)),
            "duration_min": max(5, min(int(duration_min), 240)),
            "is_published": bool(is_published),
        }
        if scheduled_at.strip(): data["scheduled_at"] = scheduled_at.strip()
        if expires_at.strip(): data["expires_at"] = expires_at.strip()
        
        ins = supabase.table("classroom_exams").insert(data).execute()
        return {"status": "ok", "exam": ins.data[0] if ins.data else None}
    except HTTPException:
        raise
    except Exception as e:
        print(f"[create_exam] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teacher/classrooms/{classroom_id}/exams")
async def teacher_list_exams(classroom_id: int):
    """📋 قائمة اختبارات الفصل"""
    try:
        res = supabase.table("classroom_exams").select("*").eq(
            "classroom_id", classroom_id
        ).order("created_at", desc=True).execute()
        
        exams = res.data or []
        # نضيف عدد المُجيبين
        for e in exams:
            try:
                count_res = supabase.table("classroom_exam_results").select(
                    "id", count="exact"
                ).eq("exam_id", e["id"]).execute()
                e["submissions_count"] = count_res.count or 0
            except Exception:
                e["submissions_count"] = 0
        
        return {"status": "ok", "exams": exams}
    except Exception as e:
        return {"status": "error", "exams": [], "error": str(e)[:200]}


@app.get("/api/teacher/classrooms/exams/{exam_id}/results")
async def teacher_exam_results(exam_id: int, teacher_id: int):
    """📊 نتائج اختبار"""
    try:
        # تحقق
        e = supabase.table("classroom_exams").select("classroom_id, title, total_marks").eq(
            "id", exam_id
        ).maybe_single().execute()
        if not e or not e.data:
            raise HTTPException(status_code=404, detail="الاختبار غير موجود")
        
        cls = supabase.table("teacher_classrooms").select("id").eq(
            "id", e.data["classroom_id"]
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not cls or not cls.data:
            raise HTTPException(status_code=403, detail="ليست لك صلاحية")
        
        # نتائج
        results_res = supabase.table("classroom_exam_results").select("*").eq(
            "exam_id", exam_id
        ).order("score", desc=True).execute()
        results = results_res.data or []
        
        # بيانات الطلاب
        if results:
            student_ids = [r["student_id"] for r in results]
            try:
                students_res = supabase.table("students").select(
                    "id, full_name, grade"
                ).in_("id", student_ids).execute()
                students_map = {s["id"]: s for s in (students_res.data or [])}
                for r in results:
                    s = students_map.get(r["student_id"], {})
                    r["full_name"] = s.get("full_name", f"طالب #{r['student_id']}")
                    r["grade"] = s.get("grade", "")
            except Exception:
                pass
        
        # إحصاءات
        scores = [float(r["score"] or 0) for r in results]
        stats = {
            "submissions": len(results),
            "average": round(sum(scores) / len(scores), 2) if scores else 0,
            "highest": max(scores) if scores else 0,
            "lowest": min(scores) if scores else 0,
        }
        
        return {
            "status": "ok",
            "exam": e.data,
            "results": results,
            "stats": stats,
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/teacher/classrooms/exams/{exam_id}")
async def teacher_delete_exam(exam_id: int, teacher_id: int):
    """🗑️ حذف اختبار"""
    try:
        e = supabase.table("classroom_exams").select("classroom_id").eq("id", exam_id).maybe_single().execute()
        if not e or not e.data:
            raise HTTPException(status_code=404, detail="الاختبار غير موجود")
        cls = supabase.table("teacher_classrooms").select("id").eq(
            "id", e.data["classroom_id"]
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not cls or not cls.data:
            raise HTTPException(status_code=403, detail="ليست لك صلاحية")
        
        supabase.table("classroom_exams").delete().eq("id", exam_id).execute()
        return {"status": "ok"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ═══════════ الحصص الافتراضية ═══════════

@app.post("/api/teacher/classrooms/{classroom_id}/sessions")
async def teacher_create_session(
    classroom_id: int,
    teacher_id: int     = Form(...),
    title: str          = Form(...),
    description: str    = Form(default=""),
    link: str           = Form(...),
    scheduled_at: str   = Form(...),
    duration_min: int   = Form(default=60),
):
    """🎥 إنشاء حصة افتراضية"""
    try:
        res = supabase.table("teacher_classrooms").select("id").eq(
            "id", classroom_id
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود")
        
        if not link.strip().startswith(("http://", "https://")):
            raise HTTPException(status_code=400, detail="رابط الحصة غير صالح")
        
        ins = supabase.table("classroom_live_sessions").insert({
            "classroom_id": classroom_id,
            "title": title.strip()[:200],
            "description": description.strip()[:1000],
            "link": link.strip(),
            "scheduled_at": scheduled_at,
            "duration_min": max(15, min(int(duration_min), 240)),
            "status": "scheduled",
        }).execute()
        return {"status": "ok", "session": ins.data[0] if ins.data else None}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teacher/classrooms/{classroom_id}/sessions")
async def teacher_list_sessions(classroom_id: int):
    """📋 قائمة الحصص الافتراضية"""
    try:
        res = supabase.table("classroom_live_sessions").select("*").eq(
            "classroom_id", classroom_id
        ).order("scheduled_at", desc=False).execute()
        return {"status": "ok", "sessions": res.data or []}
    except Exception as e:
        return {"status": "error", "sessions": [], "error": str(e)[:200]}


@app.delete("/api/teacher/classrooms/sessions/{session_id}")
async def teacher_delete_session(session_id: int, teacher_id: int):
    """🗑️ حذف حصة"""
    try:
        s = supabase.table("classroom_live_sessions").select("classroom_id").eq("id", session_id).maybe_single().execute()
        if not s or not s.data:
            raise HTTPException(status_code=404, detail="الحصة غير موجودة")
        cls = supabase.table("teacher_classrooms").select("id").eq(
            "id", s.data["classroom_id"]
        ).eq("owner_teacher_id", teacher_id).maybe_single().execute()
        if not cls or not cls.data:
            raise HTTPException(status_code=403, detail="ليست لك صلاحية")
        
        supabase.table("classroom_live_sessions").delete().eq("id", session_id).execute()
        return {"status": "ok"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ════════════════════════════════════════════════════════════════
# 🎓 Endpoints للطالب
# ════════════════════════════════════════════════════════════════

@app.post("/api/student/classrooms/join")
async def student_join_classroom(
    student_id: int     = Form(...),
    teacher_code: str   = Form(...),       # كود الفصل
):
    """
    🎓 انضمام الطالب لفصل
    شروط صارمة:
    1. الطالب موجود في المنصة
    2. كود الفصل صالح
    3. الفصل نشط
    4. لم يتجاوز الحد الأقصى
    5. لم ينضم سابقاً
    """
    if not teacher_code.strip():
        raise HTTPException(status_code=400, detail="كود الفصل مطلوب")
    
    code = teacher_code.strip().upper()
    
    try:
        # 1. الطالب موجود
        student_res = supabase.table("students").select(
            "id, full_name, grade"
        ).eq("id", student_id).maybe_single().execute()
        if not student_res or not student_res.data:
            raise HTTPException(status_code=403, detail="❌ يجب أن تكون مسجلاً في المنصة")
        student = student_res.data
        
        # 2. الفصل موجود ونشط
        cls_res = supabase.table("teacher_classrooms").select("*").eq(
            "code", code
        ).eq("is_active", True).maybe_single().execute()
        if not cls_res or not cls_res.data:
            raise HTTPException(status_code=404, detail="❌ كود الفصل غير صحيح أو الفصل غير نشط")
        classroom = cls_res.data
        
        # 3. لم ينضم سابقاً
        existing = supabase.table("classroom_members").select("id, status").eq(
            "classroom_id", classroom["id"]
        ).eq("student_id", student_id).maybe_single().execute()
        if existing and existing.data:
            if existing.data.get("status") == "banned":
                raise HTTPException(status_code=403, detail="❌ تم حظرك من هذا الفصل")
            raise HTTPException(status_code=400, detail="❌ أنت منضم بالفعل لهذا الفصل")
        
        # 4. الحد الأقصى
        count_res = supabase.table("classroom_members").select(
            "id", count="exact"
        ).eq("classroom_id", classroom["id"]).eq("status", "active").execute()
        if (count_res.count or 0) >= classroom.get("max_students", 50):
            raise HTTPException(status_code=400, detail="❌ الفصل ممتلئ")
        
        # 5. انضم
        status = "pending" if classroom.get("requires_approval") else "active"
        ins = supabase.table("classroom_members").insert({
            "classroom_id": classroom["id"],
            "student_id": student_id,
            "status": status,
        }).execute()
        
        msg = "✅ انضممت بنجاح!" if status == "active" else "⏳ طلبك قيد الموافقة"
        return {
            "status": "ok",
            "membership": ins.data[0] if ins.data else None,
            "classroom": classroom,
            "message": msg,
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"[join_classroom] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/student/classrooms/my")
async def student_my_classrooms(student_id: int):
    """📋 قائمة فصول الطالب"""
    try:
        # الفصول التي ينتمي لها
        members_res = supabase.table("classroom_members").select(
            "classroom_id, status, joined_at"
        ).eq("student_id", student_id).eq("status", "active").execute()
        members = members_res.data or []
        
        if not members:
            return {"status": "ok", "classrooms": [], "count": 0}
        
        classroom_ids = [m["classroom_id"] for m in members]
        cls_res = supabase.table("teacher_classrooms").select("*").in_(
            "id", classroom_ids
        ).eq("is_active", True).execute()
        classrooms = cls_res.data or []
        
        # نضيف اسم المعلم لكل فصل
        teacher_ids = list(set([c["owner_teacher_id"] for c in classrooms]))
        if teacher_ids:
            try:
                teachers_res = supabase.table("teachers").select("id, full_name").in_(
                    "id", teacher_ids
                ).execute()
                teachers_map = {t["id"]: t.get("full_name", "") for t in (teachers_res.data or [])}
                for c in classrooms:
                    c["teacher_name"] = teachers_map.get(c["owner_teacher_id"], "")
            except Exception:
                pass
        
        return {"status": "ok", "classrooms": classrooms, "count": len(classrooms)}
    except Exception as e:
        return {"status": "error", "classrooms": [], "count": 0, "error": str(e)[:200]}


@app.get("/api/student/classrooms/{classroom_id}/feed")
async def student_classroom_feed(classroom_id: int, student_id: int):
    """
    📰 ملخّص الفصل للطالب:
    - معلومات الفصل
    - الإعلانات
    - الاختبارات النشطة
    - الحصص القادمة
    """
    try:
        # تحقق من العضوية
        member = supabase.table("classroom_members").select("id, status").eq(
            "classroom_id", classroom_id
        ).eq("student_id", student_id).eq("status", "active").maybe_single().execute()
        if not member or not member.data:
            raise HTTPException(status_code=403, detail="❌ لست عضواً في هذا الفصل")
        
        # الفصل
        cls_res = supabase.table("teacher_classrooms").select("*").eq(
            "id", classroom_id
        ).maybe_single().execute()
        if not cls_res or not cls_res.data:
            raise HTTPException(status_code=404, detail="الفصل غير موجود")
        classroom = cls_res.data
        
        # اسم المعلم
        try:
            t_res = supabase.table("teachers").select("full_name").eq(
                "id", classroom["owner_teacher_id"]
            ).maybe_single().execute()
            classroom["teacher_name"] = (t_res.data or {}).get("full_name", "") if t_res else ""
        except Exception:
            classroom["teacher_name"] = ""
        
        # الإعلانات
        try:
            ann_res = supabase.table("classroom_announcements").select("*").eq(
                "classroom_id", classroom_id
            ).order("pinned", desc=True).order("created_at", desc=True).limit(20).execute()
            announcements = ann_res.data or []
        except Exception:
            announcements = []
        
        # الاختبارات (المنشورة فقط)
        try:
            exams_res = supabase.table("classroom_exams").select(
                "id, title, description, total_marks, duration_min, scheduled_at, expires_at, is_published"
            ).eq("classroom_id", classroom_id).eq("is_published", True).order("created_at", desc=True).execute()
            exams = exams_res.data or []
            
            # هل الطالب أجاب على كل منها؟
            if exams:
                exam_ids = [e["id"] for e in exams]
                results_res = supabase.table("classroom_exam_results").select(
                    "exam_id, score, max_score, submitted_at"
                ).in_("exam_id", exam_ids).eq("student_id", student_id).execute()
                results_map = {r["exam_id"]: r for r in (results_res.data or [])}
                for e in exams:
                    if e["id"] in results_map:
                        e["my_result"] = results_map[e["id"]]
        except Exception:
            exams = []
        
        # الحصص الافتراضية القادمة
        try:
            from datetime import datetime, timedelta
            now = datetime.utcnow().isoformat()
            sessions_res = supabase.table("classroom_live_sessions").select("*").eq(
                "classroom_id", classroom_id
            ).neq("status", "cancelled").order("scheduled_at", desc=False).limit(20).execute()
            sessions = sessions_res.data or []
        except Exception:
            sessions = []
        
        return {
            "status": "ok",
            "classroom": classroom,
            "announcements": announcements,
            "exams": exams,
            "sessions": sessions,
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"[student_feed] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/student/classrooms/exams/{exam_id}/submit")
async def student_submit_exam(
    exam_id: int,
    student_id: int     = Form(...),
    answers_json: str   = Form(...),
    time_taken_sec: int = Form(default=0),
):
    """📝 تسليم اختبار من الطالب — تصحيح تلقائي"""
    try:
        # 1. الاختبار موجود ومنشور
        exam_res = supabase.table("classroom_exams").select("*").eq(
            "id", exam_id
        ).eq("is_published", True).maybe_single().execute()
        if not exam_res or not exam_res.data:
            raise HTTPException(status_code=404, detail="❌ الاختبار غير متاح")
        exam = exam_res.data
        
        # 2. الطالب عضو في الفصل
        member = supabase.table("classroom_members").select("id").eq(
            "classroom_id", exam["classroom_id"]
        ).eq("student_id", student_id).eq("status", "active").maybe_single().execute()
        if not member or not member.data:
            raise HTTPException(status_code=403, detail="❌ لست عضواً في هذا الفصل")
        
        # 3. لم يحاول قبل
        existing = supabase.table("classroom_exam_results").select("id").eq(
            "exam_id", exam_id
        ).eq("student_id", student_id).maybe_single().execute()
        if existing and existing.data:
            raise HTTPException(status_code=400, detail="❌ سبق أن أجبت على هذا الاختبار")
        
        # 4. حساب الدرجة
        import json as _json
        try:
            answers = _json.loads(answers_json)
        except Exception:
            raise HTTPException(status_code=400, detail="JSON الإجابات غير صالح")
        
        questions = exam.get("questions_json", [])
        if isinstance(questions, str):
            questions = _json.loads(questions)
        
        total_q = len(questions)
        correct = 0
        for i, q in enumerate(questions):
            student_ans = str(answers.get(str(i), "")).strip().lower()
            correct_ans = str(q.get("answer", "")).strip().lower()
            if student_ans and student_ans == correct_ans:
                correct += 1
        
        max_score = exam.get("total_marks", 20)
        score = round((correct / total_q) * max_score, 2) if total_q else 0
        
        # 5. حفظ
        ins = supabase.table("classroom_exam_results").insert({
            "exam_id": exam_id,
            "student_id": student_id,
            "score": score,
            "max_score": max_score,
            "answers_json": answers,
            "time_taken_sec": int(time_taken_sec),
        }).execute()
        
        return {
            "status": "ok",
            "score": score,
            "max_score": max_score,
            "correct": correct,
            "total_questions": total_q,
            "percentage": round((correct / total_q) * 100, 1) if total_q else 0,
        }
    except HTTPException:
        raise
    except Exception as e:
        print(f"[submit_exam] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/student/classrooms/{classroom_id}/leave")
async def student_leave_classroom(classroom_id: int, student_id: int):
    """🚪 مغادرة الفصل"""
    try:
        supabase.table("classroom_members").delete().eq(
            "classroom_id", classroom_id
        ).eq("student_id", student_id).execute()
        return {"status": "ok", "message": "🚪 غادرت الفصل"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")



# ════════════════════════════════════════════════════════════════
# 📰 NEWS TICKER + TEACHER NOTIFICATIONS — نظام كامل
# ════════════════════════════════════════════════════════════════

# ─── الأشرطة الإخبارية ───

@app.get("/api/tickers/active")
async def get_active_tickers(audience: str = "student"):
    """📰 الأشرطة النشطة لجمهور محدد (student أو teacher)"""
    try:
        if audience not in ("student", "teacher"):
            audience = "student"
        
        from datetime import datetime
        now = datetime.utcnow().isoformat()
        
        res = supabase.table("news_tickers").select("*").eq(
            "audience", audience
        ).eq("is_active", True).order("priority", desc=True).order("created_at", desc=True).execute()
        
        tickers = res.data or []
        
        # نفلتر منتهية الصلاحية
        active = []
        for t in tickers:
            exp = t.get("expires_at")
            if exp and exp < now:
                continue
            active.append(t)
        
        return {"status": "ok", "tickers": active, "count": len(active)}
    except Exception as e:
        # لو الجدول غير موجود، نُرجع قائمة فارغة
        if "relation" in str(e).lower() or "does not exist" in str(e).lower():
            return {"status": "ok", "tickers": [], "count": 0, "warning": "جدول الأشرطة غير موجود"}
        return {"status": "error", "tickers": [], "count": 0, "error": str(e)[:200]}


@app.post("/api/admin/tickers/create")
async def admin_create_ticker(
    request: Request,
    audience: str        = Form(...),       # student | teacher
    text: str            = Form(...),
    icon: str            = Form(default="📢"),
    bg_color: str        = Form(default="#1565c0"),
    text_color: str      = Form(default="#ffffff"),
    speed: int           = Form(default=50),
    priority: int        = Form(default=0),
    expires_at: str      = Form(default=""),
):
    """📰 إنشاء شريط إخباري جديد (Admin only)"""
    if audience not in ("student", "teacher"):
        raise HTTPException(status_code=400, detail="audience يجب أن يكون student أو teacher")
    if not text.strip():
        raise HTTPException(status_code=400, detail="نص الشريط مطلوب")
    
    try:
        data = {
            "audience": audience,
            "text": text.strip()[:500],
            "icon": icon[:20],
            "bg_color": bg_color[:20],
            "text_color": text_color[:20],
            "speed": max(20, min(int(speed), 200)),
            "priority": int(priority),
            "is_active": True,
        }
        if expires_at.strip():
            data["expires_at"] = expires_at.strip()
        
        res = supabase.table("news_tickers").insert(data).execute()
        return {
            "status": "ok",
            "ticker": res.data[0] if res.data else None,
            "message": "✅ تم إنشاء الشريط"
        }
    except Exception as e:
        err = str(e).lower()
        if "relation" in err or "does not exist" in err:
            raise HTTPException(status_code=500, detail="❌ جدول news_tickers غير موجود — شغّل news_system_migration.sql")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/admin/tickers/list")
async def admin_list_tickers(audience: str = ""):
    """📋 قائمة الأشرطة (للأدمن)"""
    try:
        q = supabase.table("news_tickers").select("*")
        if audience in ("student", "teacher"):
            q = q.eq("audience", audience)
        res = q.order("priority", desc=True).order("created_at", desc=True).execute()
        return {"status": "ok", "tickers": res.data or []}
    except Exception as e:
        return {"status": "error", "tickers": [], "error": str(e)[:200]}


@app.put("/api/admin/tickers/{ticker_id}")
async def admin_update_ticker(
    ticker_id: int,
    text: str            = Form(default=""),
    icon: str            = Form(default=""),
    bg_color: str        = Form(default=""),
    text_color: str      = Form(default=""),
    speed: int           = Form(default=0),
    priority: int        = Form(default=-999),
    is_active: bool      = Form(default=True),
):
    """✏️ تعديل شريط"""
    try:
        update_data = {"is_active": is_active}
        if text.strip(): update_data["text"] = text.strip()[:500]
        if icon: update_data["icon"] = icon[:20]
        if bg_color: update_data["bg_color"] = bg_color[:20]
        if text_color: update_data["text_color"] = text_color[:20]
        if speed > 0: update_data["speed"] = max(20, min(int(speed), 200))
        if priority != -999: update_data["priority"] = int(priority)
        
        supabase.table("news_tickers").update(update_data).eq("id", ticker_id).execute()
        return {"status": "ok", "message": "✅ تم التحديث"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/admin/tickers/{ticker_id}")
async def admin_delete_ticker(ticker_id: int):
    """🗑️ حذف شريط"""
    try:
        supabase.table("news_tickers").delete().eq("id", ticker_id).execute()
        return {"status": "ok"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ─── إشعارات المعلمين ───

@app.get("/api/teacher/notifications")
async def get_teacher_notifications(teacher_id: int):
    """📨 إشعارات المعلم (الخاصة + العامة)"""
    try:
        from datetime import datetime
        now = datetime.utcnow().isoformat()
        
        # الخاصة بالمعلم
        own = supabase.table("teacher_notifications").select("*").eq(
            "teacher_id", teacher_id
        ).order("pinned", desc=True).order("created_at", desc=True).limit(50).execute()
        
        # العامة (teacher_id = NULL)
        public = supabase.table("teacher_notifications").select("*").is_(
            "teacher_id", "null"
        ).order("pinned", desc=True).order("created_at", desc=True).limit(20).execute()
        
        all_notifs = list(own.data or []) + list(public.data or [])
        
        # نفلتر منتهية الصلاحية
        active = []
        for n in all_notifs:
            exp = n.get("expires_at")
            if exp and exp < now:
                continue
            active.append(n)
        
        # نرتب بالأولوية
        active.sort(key=lambda x: (not x.get("pinned"), -(int(x.get("id") or 0))))
        
        unread_count = sum(1 for n in active if not n.get("is_read"))
        
        return {
            "status": "ok",
            "notifications": active[:50],
            "unread_count": unread_count,
        }
    except Exception as e:
        if "relation" in str(e).lower():
            return {"status": "ok", "notifications": [], "unread_count": 0, "warning": "جدول الإشعارات غير موجود"}
        return {"status": "error", "notifications": [], "unread_count": 0, "error": str(e)[:200]}


@app.post("/api/teacher/notifications/{notif_id}/read")
async def mark_notification_read(notif_id: int):
    """✅ تعليم إشعار كمقروء"""
    try:
        supabase.table("teacher_notifications").update({"is_read": True}).eq("id", notif_id).execute()
        return {"status": "ok"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/admin/notifications/send")
async def admin_send_notification(
    request: Request,
    teacher_id: int      = Form(default=0),       # 0 = للجميع
    title: str           = Form(...),
    body: str            = Form(default=""),
    type: str            = Form(default="info"),
    icon: str            = Form(default="📨"),
    pinned: bool         = Form(default=False),
    expires_at: str      = Form(default=""),
):
    """📨 إرسال إشعار للمعلمين (Admin)"""
    if not title.strip():
        raise HTTPException(status_code=400, detail="العنوان مطلوب")
    if type not in ("info", "warning", "success", "urgent"):
        type = "info"
    
    try:
        data = {
            "title": title.strip()[:200],
            "body": body.strip()[:5000],
            "type": type,
            "icon": icon[:20],
            "pinned": bool(pinned),
            "is_read": False,
        }
        # 0 = للجميع → NULL في DB
        if teacher_id and teacher_id > 0:
            data["teacher_id"] = teacher_id
        if expires_at.strip():
            data["expires_at"] = expires_at.strip()
        
        res = supabase.table("teacher_notifications").insert(data).execute()
        target = "للجميع" if not (teacher_id and teacher_id > 0) else f"للمعلم #{teacher_id}"
        return {
            "status": "ok",
            "notification": res.data[0] if res.data else None,
            "message": f"✅ أُرسل الإشعار {target}"
        }
    except Exception as e:
        err = str(e).lower()
        if "relation" in err:
            raise HTTPException(status_code=500, detail="❌ جدول الإشعارات غير موجود — شغّل news_system_migration.sql")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/admin/notifications/list")
async def admin_list_notifications(limit: int = 100):
    """📋 كل الإشعارات (للأدمن)"""
    try:
        res = supabase.table("teacher_notifications").select("*").order(
            "created_at", desc=True
        ).limit(min(limit, 500)).execute()
        return {"status": "ok", "notifications": res.data or []}
    except Exception as e:
        return {"status": "error", "notifications": [], "error": str(e)[:200]}


@app.delete("/api/admin/notifications/{notif_id}")
async def admin_delete_notification(notif_id: int):
    """🗑️ حذف إشعار"""
    try:
        supabase.table("teacher_notifications").delete().eq("id", notif_id).execute()
        return {"status": "ok"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")




# ═══════════════════════════════════════════════════════════════
# 📧 TEACHER PROFILE — تحديث الملف الشخصي + جمع الإيميلات
# ═══════════════════════════════════════════════════════════════

@app.post("/api/teacher/profile/update")
async def teacher_profile_update(
    teacher_id: int     = Form(...),
    full_name: str      = Form(default=""),
    email: str          = Form(default=""),
    phone: str          = Form(default=""),
    school: str         = Form(default=""),
    subject: str        = Form(default=""),
    city: str           = Form(default=""),
):
    """📝 تحديث ملف المعلم"""
    try:
        update_data = {}
        if full_name.strip(): update_data["full_name"] = full_name.strip()[:120]
        if email.strip(): update_data["email"] = email.strip().lower()[:200]
        if phone.strip(): update_data["phone"] = phone.strip()[:30]
        if school.strip(): update_data["school"] = school.strip()[:200]
        if subject.strip(): update_data["subject"] = subject.strip()[:100]
        if city.strip(): update_data["city"] = city.strip()[:100]
        
        if not update_data:
            raise HTTPException(status_code=400, detail="لا بيانات للتحديث")
        
        supabase.table("teachers").update(update_data).eq("id", teacher_id).execute()
        return {"status": "ok", "message": "✅ تم التحديث", "updated_fields": list(update_data.keys())}
    except HTTPException:
        raise
    except Exception as e:
        err = str(e).lower()
        if "column" in err and ("email" in err or "phone" in err or "school" in err or "subject" in err or "city" in err):
            raise HTTPException(status_code=500, detail="❌ الأعمدة الجديدة غير موجودة في قاعدة البيانات — شغّل teachers_email_migration.sql")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teacher/profile")
async def teacher_profile_get(teacher_id: int):
    """📋 جلب بيانات المعلم"""
    try:
        res = supabase.table("teachers").select(
            "id, full_name, username, email, phone, school, subject, city"
        ).eq("id", teacher_id).maybe_single().execute()
        
        if not res or not res.data:
            raise HTTPException(status_code=404, detail="المعلم غير موجود")
        
        return {"status": "ok", "profile": res.data}
    except HTTPException:
        raise
    except Exception as e:
        # لو الحقول غير موجودة، نُرجع ما نقدر
        try:
            res = supabase.table("teachers").select("id, full_name, username").eq("id", teacher_id).maybe_single().execute()
            if res and res.data:
                return {"status": "ok", "profile": res.data, "warning": "بعض الحقول غير متاحة"}
        except Exception:
            pass
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/admin/teachers/emails")
async def admin_get_teacher_emails():
    """📧 قائمة إيميلات المعلمين للأدمن"""
    try:
        res = supabase.table("teachers").select(
            "id, full_name, username, email, phone, school, subject, city, created_at"
        ).order("created_at", desc=True).execute()
        
        teachers = res.data or []
        
        # نُحسب الإحصاءات
        with_email = sum(1 for t in teachers if t.get("email"))
        with_phone = sum(1 for t in teachers if t.get("phone"))
        with_school = sum(1 for t in teachers if t.get("school"))
        
        return {
            "status": "ok",
            "teachers": teachers,
            "stats": {
                "total": len(teachers),
                "with_email": with_email,
                "with_phone": with_phone,
                "with_school": with_school,
                "completion": round((with_email / max(len(teachers), 1)) * 100, 1),
            }
        }
    except Exception as e:
        # لو الحقول غير موجودة، نُرجع الأساسيات
        try:
            res = supabase.table("teachers").select("id, full_name, username, created_at").execute()
            return {
                "status": "ok",
                "teachers": res.data or [],
                "stats": {"total": len(res.data or []), "with_email": 0, "with_phone": 0, "with_school": 0, "completion": 0},
                "warning": "❌ حقول email/phone غير موجودة — شغّل teachers_email_migration.sql"
            }
        except Exception:
            return {"status": "error", "teachers": [], "stats": {}, "error": str(e)[:200]}


@app.get("/api/admin/teachers/emails/export")
async def admin_export_teacher_emails():
    """📤 تصدير الإيميلات كـ CSV"""
    try:
        res = supabase.table("teachers").select(
            "id, full_name, username, email, phone, school, subject, city"
        ).execute()
        
        teachers = res.data or []
        
        # نبني CSV
        import csv
        from io import StringIO
        
        buf = StringIO()
        writer = csv.writer(buf)
        writer.writerow(["ID", "الاسم", "اسم المستخدم", "الإيميل", "الجوال", "المدرسة", "المادة", "المدينة"])
        for t in teachers:
            writer.writerow([
                t.get("id", ""),
                t.get("full_name", ""),
                t.get("username", ""),
                t.get("email", ""),
                t.get("phone", ""),
                t.get("school", ""),
                t.get("subject", ""),
                t.get("city", ""),
            ])
        
        csv_content = buf.getvalue()
        # BOM للعربية
        csv_content = '\ufeff' + csv_content
        
        from fastapi.responses import Response
        return Response(
            content=csv_content.encode('utf-8'),
            media_type="text/csv; charset=utf-8",
            headers={
                "Content-Disposition": "attachment; filename=teachers_emails.csv"
            }
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")




# ════════════════════════════════════════════════════════════════
# 🚀 COMPLETE AI LESSON PACK — حزمة الدرس الكاملة بالذكاء الاصطناعي
# يولّد في طلب واحد: تحضير + أسئلة + ورقة عمل + واجب + بطاقات + تحدي
# ════════════════════════════════════════════════════════════════

LESSON_PACK_PROMPT = """أنت معلم رياضيات خبير في المنهج العماني. مهمتك إنتاج حزمة تعليمية متكاملة لدرس واحد.

أعد الإجابة بصيغة JSON صالحة فقط، بدون أي نص خارجي قبل أو بعد. الصيغة المطلوبة:

{
  "title": "عنوان الدرس",
  "subtitle": "وصف موجز في سطر واحد",
  
  "lesson_plan": {
    "objectives": ["هدف 1", "هدف 2", "هدف 3", "هدف 4"],
    "vocabulary": ["مصطلح: تعريفه", "مصطلح: تعريفه", "مصطلح: تعريفه"],
    "materials": ["وسيلة 1", "وسيلة 2", "وسيلة 3"],
    "introduction": "فقرة التمهيد (3-4 أسطر) - كيف يبدأ المعلم الدرس",
    "main_content": "شرح المحتوى الأساسي (8-10 أسطر) - الأفكار الرئيسية بالتفصيل",
    "examples": [
      {"question": "مثال 1", "solution": "الحل خطوة بخطوة"},
      {"question": "مثال 2", "solution": "الحل خطوة بخطوة"},
      {"question": "مثال 3", "solution": "الحل خطوة بخطوة"}
    ],
    "common_mistakes": ["خطأ شائع 1", "خطأ شائع 2", "خطأ شائع 3"],
    "closure": "ختام الدرس (2-3 أسطر)"
  },
  
  "exam_questions": [
    {"q": "السؤال 1 (سهل)", "answer": "الإجابة", "difficulty": "easy"},
    {"q": "السؤال 2 (سهل)", "answer": "الإجابة", "difficulty": "easy"},
    {"q": "السؤال 3 (سهل)", "answer": "الإجابة", "difficulty": "easy"},
    {"q": "السؤال 4 (سهل)", "answer": "الإجابة", "difficulty": "easy"},
    {"q": "السؤال 5 (متوسط)", "answer": "الإجابة", "difficulty": "medium"},
    {"q": "السؤال 6 (متوسط)", "answer": "الإجابة", "difficulty": "medium"},
    {"q": "السؤال 7 (متوسط)", "answer": "الإجابة", "difficulty": "medium"},
    {"q": "السؤال 8 (متوسط)", "answer": "الإجابة", "difficulty": "medium"},
    {"q": "السؤال 9 (متوسط)", "answer": "الإجابة", "difficulty": "medium"},
    {"q": "السؤال 10 (صعب)", "answer": "الإجابة", "difficulty": "hard"}
  ],
  
  "worksheet": {
    "title": "ورقة عمل صفية",
    "instructions": "تعليمات للطلاب",
    "questions": [
      "سؤال 1 — ورقة العمل",
      "سؤال 2 — ورقة العمل",
      "سؤال 3 — ورقة العمل",
      "سؤال 4 — ورقة العمل",
      "سؤال 5 — ورقة العمل"
    ]
  },
  
  "homework": {
    "title": "الواجب المنزلي",
    "instructions": "تعليمات الواجب",
    "questions": [
      {"q": "سؤال الواجب 1", "answer": "الإجابة"},
      {"q": "سؤال الواجب 2", "answer": "الإجابة"},
      {"q": "سؤال الواجب 3", "answer": "الإجابة"},
      {"q": "سؤال الواجب 4", "answer": "الإجابة"},
      {"q": "سؤال الواجب 5", "answer": "الإجابة"}
    ]
  },
  
  "flashcards": [
    {"front": "السؤال 1", "back": "الإجابة 1"},
    {"front": "السؤال 2", "back": "الإجابة 2"},
    {"front": "السؤال 3", "back": "الإجابة 3"},
    {"front": "السؤال 4", "back": "الإجابة 4"},
    {"front": "السؤال 5", "back": "الإجابة 5"},
    {"front": "السؤال 6", "back": "الإجابة 6"},
    {"front": "السؤال 7", "back": "الإجابة 7"},
    {"front": "السؤال 8", "back": "الإجابة 8"},
    {"front": "السؤال 9", "back": "الإجابة 9"},
    {"front": "السؤال 10", "back": "الإجابة 10"}
  ],
  
  "challenge": {
    "title": "تحدي اليوم",
    "question": "سؤال تحدّي صعب ومُمتع",
    "hint": "تلميح للحل",
    "answer": "الإجابة الكاملة مع الشرح"
  },
  
  "video_suggestions": [
    {"title": "عنوان فيديو مقترح", "search": "كلمات بحث في يوتيوب"},
    {"title": "عنوان فيديو مقترح", "search": "كلمات بحث في يوتيوب"},
    {"title": "عنوان فيديو مقترح", "search": "كلمات بحث في يوتيوب"}
  ],
  
  "mindmap": {
    "central": "الموضوع المركزي",
    "branches": [
      {"name": "فرع 1", "items": ["نقطة", "نقطة", "نقطة"]},
      {"name": "فرع 2", "items": ["نقطة", "نقطة", "نقطة"]},
      {"name": "فرع 3", "items": ["نقطة", "نقطة", "نقطة"]},
      {"name": "فرع 4", "items": ["نقطة", "نقطة", "نقطة"]}
    ]
  }
}

تأكد من:
- الإجابات صحيحة رياضياً
- التدرّج من السهل للصعب
- المحتوى مناسب لعمر الطلاب
- اللغة العربية فصحى وواضحة
- JSON صالح 100% بدون أي markdown أو commentaries
"""


@app.post("/api/teacher/lesson_pack/generate")
async def teacher_lesson_pack_generate(
    grade: str            = Form(...),       # مثال: الصف السادس
    lesson_name: str      = Form(...),       # مثال: الكسور العشرية
    semester: str         = Form(default=""),
    unit: str             = Form(default=""),
    extra_notes: str      = Form(default=""),
):
    """🚀 يولّد حزمة درس كاملة (تحضير + أسئلة + ورقة عمل + واجب + بطاقات + تحدي) في طلب واحد"""
    if not lesson_name.strip():
        raise HTTPException(status_code=400, detail="اسم الدرس مطلوب")
    if not grade.strip():
        raise HTTPException(status_code=400, detail="الصف مطلوب")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(
            status_code=503,
            detail="❌ خدمة AI غير مُهيّأة — يجب إضافة GEMINI_API_KEY في إعدادات الخادم"
        )
    
    # بناء الـ prompt
    user_request = f"""
أنشئ حزمة درس كاملة للمعلومات التالية:

📚 الصف: {grade.strip()}
📖 اسم الدرس: {lesson_name.strip()}
"""
    if semester.strip():
        user_request += f"📅 الفصل: {semester.strip()}\n"
    if unit.strip():
        user_request += f"📦 الوحدة: {unit.strip()}\n"
    if extra_notes.strip():
        user_request += f"\n📝 ملاحظات إضافية: {extra_notes.strip()[:500]}"
    
    full_prompt = LESSON_PACK_PROMPT + "\n\n" + user_request
    
    # استدعاء Gemini (نستخدم Pro للجودة العالية - يحتاج المزيد من tokens)
    import httpx
    import json as json_lib
    
    # نُحاول flash أولاً ثم pro لو فشل
    models_to_try = ["gemini-2.5-flash", "gemini-2.5-pro", "gemini-2.0-flash"]
    
    last_error = None
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": full_prompt}]}],
                "generationConfig": {
                    "temperature": 0.7,
                    "maxOutputTokens": 8000,  # كبير لأن المحتوى كبير
                    "topP": 0.9,
                    "responseMimeType": "application/json",  # نطلب JSON صريحة
                },
                "safetySettings": [
                    {"category": "HARM_CATEGORY_HARASSMENT",        "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                    {"category": "HARM_CATEGORY_HATE_SPEECH",       "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                    {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_LOW_AND_ABOVE"},
                    {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_MEDIUM_AND_ABOVE"},
                ]
            }
            
            async with httpx.AsyncClient(timeout=120) as client:  # 2 دقائق timeout
                res = await client.post(url, json=payload)
                
                if res.status_code == 429:  # Rate limited — نحاول الموديل التالي
                    last_error = f"{model_name}: تجاوز الحصة"
                    continue
                
                if res.status_code != 200:
                    last_error = f"{model_name}: HTTP {res.status_code}"
                    continue
                
                data = res.json()
                
                if "candidates" not in data or not data["candidates"]:
                    last_error = f"{model_name}: لا استجابة"
                    continue
                
                cand = data["candidates"][0]
                if "content" not in cand or "parts" not in cand["content"]:
                    last_error = f"{model_name}: استجابة فارغة"
                    continue
                
                text_response = cand["content"]["parts"][0].get("text", "").strip()
                if not text_response:
                    last_error = f"{model_name}: نص فارغ"
                    continue
                
                # نُحاول parse JSON
                # نُزيل أي markdown wrappers لو كانت
                if text_response.startswith("```"):
                    text_response = text_response.split("```")[1]
                    if text_response.startswith("json"):
                        text_response = text_response[4:]
                    text_response = text_response.strip()
                if text_response.endswith("```"):
                    text_response = text_response[:-3].strip()
                
                try:
                    pack = json_lib.loads(text_response)
                except json_lib.JSONDecodeError as je:
                    # نحاول استخراج JSON من النص
                    start_idx = text_response.find("{")
                    end_idx = text_response.rfind("}")
                    if start_idx >= 0 and end_idx > start_idx:
                        try:
                            pack = json_lib.loads(text_response[start_idx:end_idx+1])
                        except Exception:
                            last_error = f"{model_name}: JSON تالف"
                            continue
                    else:
                        last_error = f"{model_name}: لا JSON"
                        continue
                
                # نجاح!
                return {
                    "status": "ok",
                    "model_used": model_name,
                    "grade": grade,
                    "lesson_name": lesson_name,
                    "pack": pack,
                    "message": f"✅ تم توليد الحزمة الكاملة (موديل: {model_name})"
                }
        
        except httpx.TimeoutException:
            last_error = f"{model_name}: انتهت المهلة"
            continue
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:100]}"
            continue
    
    # كل الموديلات فشلت
    raise HTTPException(
        status_code=502,
        detail=f"❌ فشل توليد الحزمة من كل الموديلات. آخر خطأ: {last_error}"
    )


@app.post("/api/teacher/lesson_pack/save_questions")
async def teacher_lesson_pack_save_questions(
    grade: str            = Form(...),
    lesson_name: str      = Form(...),
    questions_json: str   = Form(...),       # JSON من exam_questions
    semester: str         = Form(default=""),
    unit: str             = Form(default=""),
):
    """💾 حفظ الأسئلة المُولّدة في بنك الأسئلة"""
    if not questions_json.strip():
        raise HTTPException(status_code=400, detail="لا أسئلة")
    
    try:
        import json as json_lib
        questions = json_lib.loads(questions_json)
        
        if not isinstance(questions, list) or not questions:
            raise HTTPException(status_code=400, detail="صيغة الأسئلة غير صحيحة")
        
        # نُحضّر بيانات الإدراج
        to_insert = []
        for q in questions[:30]:  # حد أقصى 30
            if not isinstance(q, dict):
                continue
            qtext = (q.get("q") or q.get("question") or "").strip()
            ans = (q.get("answer") or q.get("a") or "").strip()
            if not qtext or not ans:
                continue
            
            to_insert.append({
                "grade": grade[:80],
                "semester": semester[:80] if semester else "غير محدد",
                "unit": unit[:120] if unit else "غير محدد",
                "lesson": lesson_name[:120],
                "question": qtext[:2000],
                "answer": ans[:1000],
            })
        
        if not to_insert:
            raise HTTPException(status_code=400, detail="لا أسئلة صالحة للحفظ")
        
        res = supabase.table("questions").insert(to_insert).execute()
        saved = len(res.data or [])
        
        return {
            "status": "ok",
            "saved": saved,
            "message": f"✅ حُفظ {saved} سؤال في بنك الأسئلة"
        }
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ════════════════════════════════════════════════════════════
# 🎮 GAMES WORKSHOP — Backend endpoints لتغذية الألعاب
# ════════════════════════════════════════════════════════════

@app.get("/api/teacher/games/filters")
async def games_get_filters(grade: str = "", semester: str = "", unit: str = ""):
    """
    📋 جلب القيم المتاحة للفلاتر (cascade dropdown)
    
    منطق متدرّج:
    - بدون grade   → يرجع كل الصفوف المتاحة
    - مع grade فقط → يرجع الفصول والوحدات والدروس لذلك الصف
    - مع grade + semester → يرجع الوحدات والدروس لهذا الفصل
    - مع grade + semester + unit → يرجع الدروس لهذه الوحدة
    """
    try:
        query = supabase.table("questions").select("grade,semester,unit,lesson,subject")
        if grade.strip():
            query = query.eq("grade", grade.strip())
        if semester.strip():
            query = query.eq("semester", semester.strip())
        if unit.strip():
            query = query.eq("unit", unit.strip())
        
        # نجلب حتى 5000 سجل لاستخراج القيم الفريدة
        res = query.limit(5000).execute()
        rows = res.data or []
        
        # نستخرج القيم الفريدة المرتّبة
        grades_set, semesters_set, units_set, lessons_set, subjects_set = set(), set(), set(), set(), set()
        for r in rows:
            if r.get("grade"):
                grades_set.add(r["grade"])
            if r.get("semester"):
                semesters_set.add(r["semester"])
            if r.get("unit"):
                units_set.add(r["unit"])
            if r.get("lesson"):
                lessons_set.add(r["lesson"])
            if r.get("subject"):
                subjects_set.add(r["subject"])
        
        return {
            "status": "ok",
            "grades":    sorted(grades_set),
            "semesters": sorted(semesters_set),
            "units":     sorted(units_set),
            "lessons":   sorted(lessons_set),
            "subjects":  sorted(subjects_set),
            "total_rows": len(rows)
        }
    except Exception as e:
        print(f"[games/filters] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teacher/games/bank")
async def games_get_from_bank(
    grade: str = "",
    semester: str = "",
    unit: str = "",
    lesson: str = "",
    q_type: str = "",
    limit: int = 50,
    random_order: bool = True
):
    """
    📚 جلب الأسئلة من بنك الأسئلة بفلاتر متعددة - للاستخدام في الألعاب
    
    Parameters:
    - grade, semester, unit, lesson: فلاتر التصنيف
    - q_type: نوع السؤال (multiple_choice / true_false / written / ...)
    - limit: العدد المطلوب (افتراضي 50، أقصى 200)
    - random_order: ترتيب عشوائي (افتراضي True)
    
    Returns:
    - قائمة الأسئلة بصيغة موحدة للألعاب: { question, answer, options[], q_type, lesson, grade }
    """
    limit = max(1, min(200, limit))
    
    try:
        query = supabase.table("questions").select("id,question,answer,options,q_type,grade,semester,unit,lesson,subject")
        if grade.strip():
            query = query.eq("grade", grade.strip())
        if semester.strip():
            query = query.eq("semester", semester.strip())
        if unit.strip():
            query = query.eq("unit", unit.strip())
        if lesson.strip():
            query = query.eq("lesson", lesson.strip())
        if q_type.strip():
            query = query.eq("q_type", q_type.strip())
        
        # نجلب أكثر من المطلوب للسماح بالخلط
        fetch_limit = min(500, limit * 3) if random_order else limit
        res = query.limit(fetch_limit).execute()
        rows = res.data or []
        
        # خلط عشوائي
        if random_order:
            import random as _random
            _random.shuffle(rows)
        
        # نأخذ أول N
        rows = rows[:limit]
        
        # تحويل options من string إلى list (لو كانت محفوظة كنص)
        for r in rows:
            opts = r.get("options", "")
            if isinstance(opts, str) and opts.strip():
                # نقسّمها بـ | أو ; أو \n
                for sep in ['|', ';', '\n', ',']:
                    if sep in opts:
                        r["options"] = [o.strip() for o in opts.split(sep) if o.strip()]
                        break
                else:
                    r["options"] = [opts.strip()] if opts.strip() else []
            elif not isinstance(opts, list):
                r["options"] = []
        
        return {
            "status": "ok",
            "count": len(rows),
            "questions": rows
        }
    except Exception as e:
        print(f"[games/bank] error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/teacher/games/ai_generate")
async def games_ai_generate(
    request: Request,
    grade: str        = Form(default=""),
    lesson: str       = Form(default=""),
    topic: str        = Form(...),                  # الموضوع (إجباري)
    count: int        = Form(default=10),           # عدد الأسئلة
    difficulty: str   = Form(default="medium"),     # easy / medium / hard
    q_format: str     = Form(default="mcq"),        # mcq (4 خيارات) / short (إجابة قصيرة) / pairs (سؤال↔جواب)
    language: str     = Form(default="ar")          # ar / en
):
    """
    🤖 توليد أسئلة بالذكاء الاصطناعي للألعاب التعليمية
    
    q_format:
    - mcq:   4 خيارات (مناسب لـ Quiz Battle)
    - short: سؤال + إجابة قصيرة (مناسب لـ Lucky Wheel, Beat the Clock, Team Battle)
    - pairs: سؤال + جواب (مناسب لـ Memory Match)
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=10, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    if not topic.strip():
        raise HTTPException(status_code=400, detail="الموضوع مطلوب")
    
    count = max(3, min(40, count))  # بين 3 و 40 سؤال
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ خدمة AI غير مُهيّأة (GEMINI_API_KEY مفقود)")
    
    # بناء الـ prompt حسب الـ format
    diff_ar = {"easy": "سهل", "medium": "متوسط", "hard": "صعب"}.get(difficulty, "متوسط")
    
    context_parts = []
    if grade.strip():
        context_parts.append(f"الصف: {grade.strip()}")
    if lesson.strip():
        context_parts.append(f"الدرس: {lesson.strip()}")
    context_parts.append(f"الموضوع: {topic.strip()}")
    context_parts.append(f"المستوى: {diff_ar}")
    context = "\n".join(context_parts)
    
    if q_format == "mcq":
        format_instructions = f"""
أنشئ {count} سؤال اختيار من متعدد (٤ خيارات) عن الموضوع التالي:

{context}

أعد النتيجة كـ JSON صحيح فقط (بدون أي شرح) بالشكل:
{{
  "questions": [
    {{
      "q": "نص السؤال",
      "a": ["الخيار الأول", "الخيار الثاني", "الخيار الثالث", "الخيار الرابع"],
      "correct": 0,
      "points": 10
    }}
  ]
}}

⚠️ قواعد مهمة:
- correct رقم من ٠ إلى ٣ (موقع الإجابة الصحيحة)
- اجعل الخيارات الخاطئة معقولة (errors شائعة) ليست عشوائية
- استخدم اللغة العربية الفصحى
- اجعل الأسئلة متنوعة وغير مكررة
- استخدم الأرقام العربية (٠١٢٣٤٥٦٧٨٩) في النص
- points = 10 لكل سؤال (المعلم يعدّلها لاحقاً)
"""
    elif q_format == "short":
        format_instructions = f"""
أنشئ {count} سؤال بإجابة قصيرة عن الموضوع التالي:

{context}

أعد النتيجة كـ JSON صحيح فقط (بدون أي شرح) بالشكل:
{{
  "questions": [
    {{ "q": "نص السؤال", "a": "الإجابة القصيرة" }}
  ]
}}

⚠️ قواعد مهمة:
- الإجابة كلمة أو عدد أو جملة قصيرة جداً (للحساب الذهني والمراجعة السريعة)
- الأسئلة قصيرة وواضحة
- استخدم الأرقام العربية (٠١٢٣٤٥٦٧٨٩)
- تنوّع في الأسئلة (جمع، طرح، ضرب، قسمة، كسور، نسب...)
"""
    else:  # pairs
        format_instructions = f"""
أنشئ {count} زوج (سؤال ↔ جواب) للمطابقة عن الموضوع التالي:

{context}

أعد النتيجة كـ JSON صحيح فقط (بدون أي شرح) بالشكل:
{{
  "questions": [
    {{ "a": "العنصر الأول (سؤال أو عملية)", "b": "العنصر الثاني (الجواب أو الناتج)" }}
  ]
}}

⚠️ قواعد مهمة:
- كل زوج يكون قصيراً جداً (يكتب على بطاقة صغيرة)
- مناسب للعبة مطابقة الذاكرة
- استخدم الأرقام العربية
"""
    
    # استدعاء Gemini
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": format_instructions}]}],
                "generationConfig": {
                    "temperature": 0.85,
                    "maxOutputTokens": 4000,
                    "topP": 0.9,
                    "responseMimeType": "application/json"
                }
            }
            with httpx.Client(timeout=45.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    text = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    # تنظيف الـ markdown إن وجد
                    if text.startswith("```"):
                        text = text.split("\n", 1)[1] if "\n" in text else text
                        text = text.rsplit("```", 1)[0] if "```" in text else text
                    
                    parsed = json_lib.loads(text)
                    questions = parsed.get("questions", [])
                    if not isinstance(questions, list) or len(questions) == 0:
                        raise ValueError("استجابة AI فارغة")
                    
                    return {
                        "status": "ok",
                        "count": len(questions),
                        "questions": questions,
                        "model": model_name,
                        "format": q_format
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:200]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON parse error: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:200]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل توليد الأسئلة بعد محاولة كل النماذج. آخر خطأ: {last_error}")


# ════════════════════════════════════════════════════════════════════════════
# 🏭 PREP EXTRACTOR — مصنع الأسئلة الشخصي (Personal Question Factory)
# ════════════════════════════════════════════════════════════════════════════
# استخدام شخصي للمدير: استخراج PDF + Gemini + ربط بالمنهج + حفظ في البنك
# ════════════════════════════════════════════════════════════════════════════

@app.post("/api/prep/pdf_info")
async def prep_pdf_info(
    pdf_file: UploadFile = File(...),
    admin = Depends(get_current_admin)
):
    """
    📄 معلومات الـ PDF: عدد الصفحات + معاينة سريعة لأول صفحة
    """
    try:
        content = await pdf_file.read()
        if len(content) > 50 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="حجم الملف أكبر من 50 ميجابايت")
        
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader  # type: ignore
            except ImportError:
                raise HTTPException(status_code=503, detail="❌ مكتبة PDF غير مثبتة. ثبّتها: pip install pypdf")
        
        import io
        reader = PdfReader(io.BytesIO(content))
        page_count = len(reader.pages)
        
        # معاينة أول صفحة
        preview = ""
        try:
            if page_count > 0:
                preview = (reader.pages[0].extract_text() or "")[:500]
        except Exception:
            pass
        
        return {
            "status": "ok",
            "page_count": page_count,
            "file_size": len(content),
            "filename": pdf_file.filename,
            "preview": preview
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ في قراءة PDF: {str(e)[:200]}")


@app.post("/api/prep/extract_text")
async def prep_extract_text(
    pdf_file: UploadFile = File(...),
    page_start: int = Form(default=1),
    page_end: int = Form(default=15),
    admin = Depends(get_current_admin)
):
    """
    📖 استخراج النص من نطاق صفحات محدد في PDF
    معالجة قوية مع بدائل لـ PDF العربية والمعقدة
    """
    try:
        content = await pdf_file.read()
        if not content or len(content) == 0:
            raise HTTPException(status_code=400, detail="الملف فارغ أو لم يُرفع")
        
        if len(content) > 50 * 1024 * 1024:
            raise HTTPException(status_code=413, detail="حجم الملف أكبر من 50 ميجابايت")
        
        # محاولة استيراد المكتبات
        try:
            from pypdf import PdfReader
            pdf_lib = "pypdf"
        except ImportError:
            try:
                from PyPDF2 import PdfReader  # type: ignore
                pdf_lib = "PyPDF2"
            except ImportError:
                raise HTTPException(
                    status_code=503,
                    detail="❌ مكتبة PDF غير مثبتة. شغّل: pip install pypdf"
                )
        
        import io
        
        # قراءة الـ PDF مع معالجة الـ encryption
        try:
            reader = PdfReader(io.BytesIO(content))
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"❌ فشل قراءة PDF: قد يكون الملف تالفاً. ({str(e)[:100]})"
            )
        
        # فحص التشفير
        if hasattr(reader, 'is_encrypted') and reader.is_encrypted:
            try:
                reader.decrypt("")
            except Exception:
                raise HTTPException(
                    status_code=400,
                    detail="❌ هذا الـ PDF مُشفّر بكلمة مرور. يرجى فك التشفير أولاً."
                )
        
        try:
            total_pages = len(reader.pages)
        except Exception as e:
            raise HTTPException(
                status_code=400,
                detail=f"❌ تعذّر قراءة عدد الصفحات: {str(e)[:100]}"
            )
        
        if total_pages == 0:
            raise HTTPException(status_code=400, detail="❌ الـ PDF لا يحتوي على صفحات")
        
        # تنظيف الحدود
        p_start = max(1, min(page_start, total_pages))
        p_end = max(p_start, min(page_end, total_pages))
        
        parts = []
        failed_pages = []
        empty_pages = []
        
        for i in range(p_start - 1, p_end):
            try:
                page = reader.pages[i]
                
                # محاولة استخراج النص بطرق مختلفة
                text = None
                
                # الطريقة 1: extract_text العادية
                try:
                    text = page.extract_text()
                except Exception:
                    pass
                
                # الطريقة 2: extract_text بـ extraction_mode='layout' (للـ pypdf الحديث)
                if not text or len(text.strip()) < 10:
                    try:
                        text = page.extract_text(extraction_mode="layout")
                    except Exception:
                        pass
                
                # الطريقة 3: تجاهل أخطاء التشفير
                if not text or len(text.strip()) < 10:
                    try:
                        text = page.extract_text(0)  # mode 0 = افتراضي بدون layout
                    except Exception:
                        pass
                
                # تنظيف النص
                if text:
                    text = text.strip()
                    # إزالة المسافات الزائدة المتتالية
                    import re
                    text = re.sub(r'\n{3,}', '\n\n', text)
                    text = re.sub(r' {2,}', ' ', text)
                
                if text and len(text) >= 5:
                    parts.append(f"[صفحة {i + 1}]\n{text}")
                else:
                    empty_pages.append(i + 1)
            except Exception as e:
                failed_pages.append(i + 1)
                print(f"[prep_extract] فشل في صفحة {i+1}: {str(e)[:100]}")
                continue
        
        full_text = "\n\n".join(parts)
        
        # رسائل تشخيصية مفيدة
        if not full_text or len(full_text) < 30:
            # محاولة فهم السبب
            diagnostics = []
            if empty_pages:
                diagnostics.append(f"الصفحات {empty_pages[:5]} لا تحتوي نصاً")
            if failed_pages:
                diagnostics.append(f"فشلت قراءة الصفحات {failed_pages[:5]}")
            
            diag_msg = " | ".join(diagnostics) if diagnostics else "النص الناتج فارغ"
            
            raise HTTPException(
                status_code=422,
                detail=(
                    f"❌ لم نتمكن من استخراج نص من النطاق المحدد ({diag_msg}). "
                    f"الأسباب المحتملة:\n"
                    f"1. الـ PDF عبارة عن صور ممسوحة ضوئياً (يحتاج OCR)\n"
                    f"2. الـ PDF يستخدم خطوطاً غير قياسية\n"
                    f"3. المحتوى كله صور بدون نص\n\n"
                    f"💡 جرّب: صفحات أخرى أو ملف PDF آخر"
                )
            )
        
        return {
            "status": "ok",
            "text": full_text,
            "char_count": len(full_text),
            "pages_extracted": len(parts),
            "pages_requested": p_end - p_start + 1,
            "empty_pages": empty_pages,
            "failed_pages": failed_pages,
            "page_start": p_start,
            "page_end": p_end,
            "total_pages": total_pages,
            "pdf_library": pdf_lib
        }
    except HTTPException:
        raise
    except Exception as e:
        # خطأ غير متوقع - نطبع تفاصيل أكثر
        import traceback
        print(f"[prep_extract] خطأ غير متوقع:\n{traceback.format_exc()}")
        raise HTTPException(
            status_code=500,
            detail=f"❌ خطأ في استخراج النص: {str(e)[:200]}"
        )


@app.post("/api/prep/ai_generate")
async def prep_ai_generate(
    request: Request,
    text: str = Form(...),
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    n_mcq: int = Form(default=5),
    n_tf: int = Form(default=3),
    n_short: int = Form(default=2),
    difficulty: str = Form(default="medium"),
    admin = Depends(get_current_admin)
):
    """
    🤖 توليد أسئلة بـ Gemini من نص PDF مستخرج
    أنواع: اختياري (mcq) + صواب/خطأ (tf) + إجابة قصيرة (short)
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=15, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    text = text.strip()
    if not text or len(text) < 50:
        raise HTTPException(status_code=400, detail="النص قصير جداً (أقل من 50 حرف)")
    
    total = n_mcq + n_tf + n_short
    if total < 1:
        raise HTTPException(status_code=400, detail="حدد عدداً للأسئلة")
    if total > 50:
        raise HTTPException(status_code=400, detail="الحد الأقصى 50 سؤال في المرة الواحدة")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    diff_ar = {"easy": "سهلة ومباشرة", "medium": "متوسطة الصعوبة", "hard": "تحليلية وتطبيقية"}.get(difficulty, "متوسطة")
    
    context_parts = []
    if grade: context_parts.append(f"الصف: {grade}")
    if semester: context_parts.append(f"الفصل: {semester}")
    if unit: context_parts.append(f"الوحدة: {unit}")
    if lesson: context_parts.append(f"الدرس: {lesson}")
    context = " | ".join(context_parts) if context_parts else "غير محدد"
    
    prompt = f"""أنت خبير في تعليم الرياضيات للمرحلة الأساسية في سلطنة عُمان.

المهمة: أنشئ {total} سؤالاً من محتوى الدرس التالي.

السياق: {context}
الصعوبة: {diff_ar}
التوزيع المطلوب:
- {n_mcq} سؤال اختياري (4 خيارات)
- {n_tf} سؤال صواب/خطأ
- {n_short} سؤال إجابة قصيرة

محتوى الدرس:
---
{text[:5000]}
---

قواعد صارمة:
- العربية الفصحى، أرقام عربية (٠١٢٣٤٥٦٧٨٩)
- للاختياري: 4 خيارات، الإجابة تطابق أحد الخيارات نصاً
- للصواب/خطأ: الإجابة "صواب" أو "خطأ" فقط
- للإجابة القصيرة: إجابة مختصرة (كلمة أو رقم أو جملة قصيرة)
- الأسئلة متنوعة وغير مكررة
- ربط الأسئلة بمحتوى الدرس الفعلي

أخرج JSON فقط بهذا الشكل بدون أي نص آخر:
{{
  "questions": [
    {{"type": "mcq", "q": "نص السؤال", "options": ["خيار1", "خيار2", "خيار3", "خيار4"], "answer": "خيار1"}},
    {{"type": "tf", "q": "نص السؤال", "options": [], "answer": "صواب"}},
    {{"type": "short", "q": "نص السؤال", "options": [], "answer": "الإجابة"}}
  ]
}}"""
    
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": 0.8,
                    "maxOutputTokens": 6000,
                    "topP": 0.9,
                    "responseMimeType": "application/json"
                }
            }
            with httpx.Client(timeout=60.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw_text = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    
                    if raw_text.startswith("```"):
                        raw_text = raw_text.split("\n", 1)[1] if "\n" in raw_text else raw_text
                        raw_text = raw_text.rsplit("```", 1)[0] if "```" in raw_text else raw_text
                    
                    parsed = json_lib.loads(raw_text)
                    questions = parsed.get("questions", [])
                    if not isinstance(questions, list) or len(questions) == 0:
                        raise ValueError("استجابة AI فارغة")
                    
                    # تعديل صيغة الأنواع للتوافق العربي
                    type_map = {"mcq": "اختياري", "tf": "صواب/خطأ", "short": "إجابة قصيرة"}
                    for q in questions:
                        q["type"] = type_map.get(q.get("type", "mcq"), q.get("type", "اختياري"))
                        # تأكيد أن options قائمة
                        if not isinstance(q.get("options"), list):
                            q["options"] = []
                    
                    return {
                        "status": "ok",
                        "count": len(questions),
                        "questions": questions,
                        "model": model_name,
                        "context": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson}
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:200]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON parse: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:200]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل التوليد. آخر خطأ: {last_error}")


@app.post("/api/prep/save_to_bank")
async def prep_save_to_bank(
    questions_json: str = Form(...),
    grade: str = Form(...),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    subject: str = Form(default="الرياضيات"),
    admin = Depends(get_current_admin)
):
    """
    💾 حفظ مباشر للأسئلة المُولَّدة في بنك الأسئلة الرسمي للمنصة
    """
    try:
        import json as json_lib
        questions = json_lib.loads(questions_json)
        if not isinstance(questions, list) or len(questions) == 0:
            raise HTTPException(status_code=400, detail="لا توجد أسئلة لحفظها")
        
        if not grade.strip():
            raise HTTPException(status_code=400, detail="الصف مطلوب")
        
        # تحويل الصيغ
        type_reverse = {"اختياري": "multiple_choice", "صواب/خطأ": "true_false", "إجابة قصيرة": "short_answer"}
        
        rows = []
        for q in questions:
            q_text = (q.get("q") or q.get("question") or "").strip()
            if not q_text:
                continue
            
            q_type_ar = q.get("type", "اختياري")
            q_type_en = type_reverse.get(q_type_ar, "multiple_choice")
            
            options = q.get("options", [])
            if isinstance(options, list):
                options_str = ", ".join(str(o) for o in options if o)
            else:
                options_str = str(options)
            
            row = {
                "question": q_text,
                "answer": str(q.get("answer") or q.get("correct") or "").strip(),
                "options": options_str,
                "q_type": q_type_en,
                "grade": grade.strip(),
                "subject": subject.strip(),
            }
            if semester.strip(): row["semester"] = semester.strip()
            if unit.strip(): row["unit"] = unit.strip()
            if lesson.strip(): row["lesson"] = lesson.strip()
            
            rows.append(row)
        
        if not rows:
            raise HTTPException(status_code=400, detail="لا توجد أسئلة صالحة بعد التحقق")
        
        # حفظ في Supabase
        res = supabase.table("questions").insert(rows).execute()
        saved_count = len(res.data or [])
        
        return {
            "status": "ok",
            "saved": saved_count,
            "total_attempted": len(questions),
            "message": f"✅ تم حفظ {saved_count} سؤال في بنك الأسئلة"
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ في الحفظ: {str(e)[:200]}")


@app.get("/api/prep/sessions_history")
async def prep_sessions_history(admin = Depends(get_current_admin)):
    """
    📊 إحصاءات استخدام المصنع (آخر الأسئلة المضافة من المدير)
    """
    try:
        # نجلب آخر 50 سؤال أضافه المدير
        res = supabase.table("questions").select("id,question,grade,subject,unit,lesson,q_type,created_at").order("created_at", desc=True).limit(50).execute()
        rows = res.data or []
        
        # إحصاءات
        total_count_res = supabase.table("questions").select("id", count="exact").execute()
        total_count = total_count_res.count or 0
        
        return {
            "status": "ok",
            "total_in_bank": total_count,
            "recent": rows[:20]
        }
    except Exception as e:
        return {"status": "ok", "total_in_bank": 0, "recent": [], "warning": str(e)[:100]}


# ════════════════════════════════════════════════════════════════════════════
# 💾 SESSIONS — سجل جلسات المصنع + الاستئناف
# ════════════════════════════════════════════════════════════════════════════
# نستخدم جدول prep_sessions (يُنشأ تلقائياً إن لم يكن موجوداً)
# الـ schema المتوقع:
#   id (uuid, primary key)
#   admin_id (text, default 'admin')
#   title (text)
#   pdf_filename (text)
#   grade, semester, unit, lesson (text)
#   extracted_text (text, nullable)
#   questions (jsonb) - مصفوفة الأسئلة
#   meta (jsonb) - إعدادات (n_mcq, n_tf, n_short, difficulty)
#   created_at (timestamp default now())
#   updated_at (timestamp)
#   status (text: 'draft' | 'completed' | 'saved')
# ════════════════════════════════════════════════════════════════════════════

@app.post("/api/prep/session/save")
async def prep_session_save(
    title: str = Form(...),
    pdf_filename: str = Form(default=""),
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    extracted_text: str = Form(default=""),
    questions_json: str = Form(default="[]"),
    meta_json: str = Form(default="{}"),
    status: str = Form(default="draft"),
    session_id: str = Form(default=""),
    admin = Depends(get_current_admin)
):
    """
    💾 حفظ/تحديث جلسة مصنع. إذا session_id موجود → تحديث، وإلا → إنشاء جديد
    """
    try:
        import json as json_lib
        from datetime import datetime
        
        try:
            questions = json_lib.loads(questions_json) if questions_json else []
            meta = json_lib.loads(meta_json) if meta_json else {}
        except json_lib.JSONDecodeError:
            questions = []
            meta = {}
        
        # نختصر النص للحفظ (أول 50KB)
        extracted_short = extracted_text[:50000] if extracted_text else ""
        
        row = {
            "title": title.strip()[:200] or f"جلسة {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            "pdf_filename": pdf_filename[:200],
            "grade": grade,
            "semester": semester,
            "unit": unit,
            "lesson": lesson,
            "extracted_text": extracted_short,
            "questions": questions,
            "meta": meta,
            "status": status,
            "updated_at": datetime.utcnow().isoformat()
        }
        
        if session_id and session_id.strip():
            # تحديث
            res = supabase.table("prep_sessions").update(row).eq("id", session_id.strip()).execute()
            if not res.data:
                # إذا فشل التحديث، نُنشئ جديد
                res = supabase.table("prep_sessions").insert(row).execute()
        else:
            # إنشاء جديد
            res = supabase.table("prep_sessions").insert(row).execute()
        
        new_id = (res.data[0]["id"] if res.data else None)
        return {
            "status": "ok",
            "session_id": new_id,
            "message": "✅ حُفظت الجلسة"
        }
    except Exception as e:
        # ربما الجدول غير موجود - نعطي رسالة واضحة
        err_msg = str(e)[:300]
        if "does not exist" in err_msg or "relation" in err_msg.lower():
            return {
                "status": "error",
                "needs_table": True,
                "message": "⚠️ جدول prep_sessions غير موجود. يجب إنشاؤه أولاً.",
                "sql_to_run": """CREATE TABLE IF NOT EXISTS prep_sessions (
  id uuid DEFAULT gen_random_uuid() PRIMARY KEY,
  title text,
  pdf_filename text,
  grade text,
  semester text,
  unit text,
  lesson text,
  extracted_text text,
  questions jsonb DEFAULT '[]'::jsonb,
  meta jsonb DEFAULT '{}'::jsonb,
  status text DEFAULT 'draft',
  created_at timestamptz DEFAULT now(),
  updated_at timestamptz DEFAULT now()
);"""
            }
        raise HTTPException(status_code=500, detail=f"خطأ في الحفظ: {err_msg}")


@app.get("/api/prep/session/list")
async def prep_session_list(limit: int = 20, admin = Depends(get_current_admin)):
    """
    📋 قائمة الجلسات السابقة (أحدثها أولاً)
    """
    try:
        limit = max(1, min(50, limit))
        res = supabase.table("prep_sessions").select("id,title,pdf_filename,grade,semester,unit,lesson,status,created_at,updated_at,questions").order("updated_at", desc=True).limit(limit).execute()
        rows = res.data or []
        
        # نُلخّص: عدد الأسئلة فقط (لا نُرسل المحتوى كاملاً)
        for r in rows:
            qs = r.get("questions") or []
            if isinstance(qs, list):
                r["questions_count"] = len(qs)
            else:
                r["questions_count"] = 0
            # نُزيل المحتوى الثقيل من القائمة
            r.pop("questions", None)
        
        return {"status": "ok", "sessions": rows, "count": len(rows)}
    except Exception as e:
        err = str(e)[:200]
        if "does not exist" in err or "relation" in err.lower():
            return {"status": "ok", "sessions": [], "count": 0, "warning": "جدول الجلسات غير موجود"}
        return {"status": "ok", "sessions": [], "count": 0, "warning": err}


@app.get("/api/prep/session/{session_id}")
async def prep_session_get(session_id: str, admin = Depends(get_current_admin)):
    """
    📖 جلب جلسة كاملة بالـ id (للاستئناف)
    """
    try:
        res = supabase.table("prep_sessions").select("*").eq("id", session_id).single().execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="الجلسة غير موجودة")
        return {"status": "ok", "session": res.data}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/prep/session/{session_id}")
async def prep_session_delete(session_id: str, admin = Depends(get_current_admin)):
    """
    🗑️ حذف جلسة
    """
    try:
        supabase.table("prep_sessions").delete().eq("id", session_id).execute()
        return {"status": "ok", "message": "✅ حُذفت الجلسة"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ════════════════════════════════════════════════════════════════════════════
# 📊 BATCH PROCESSING — معالجة دفعات (100+ سؤال تلقائياً)
# ════════════════════════════════════════════════════════════════════════════
@app.post("/api/prep/ai_generate_batch")
async def prep_ai_generate_batch(
    request: Request,
    text: str = Form(...),
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    target_total: int = Form(default=100),
    n_mcq_per_batch: int = Form(default=10),
    n_tf_per_batch: int = Form(default=5),
    n_short_per_batch: int = Form(default=5),
    difficulty: str = Form(default="medium"),
    batch_num: int = Form(default=1),
    already_generated: int = Form(default=0),
    admin = Depends(get_current_admin)
):
    """
    🔄 توليد دفعة واحدة من ضمن مهمة كبيرة (يُستدعى عدة مرات من الـ frontend)
    
    الـ frontend يُكرّر النداء حتى يصل لـ target_total، مع زيادة batch_num وdifferent random_seed.
    
    نستخدم variation_seed في الـ prompt لضمان تنوع الدفعات.
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=30, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    text = text.strip()
    if not text or len(text) < 50:
        raise HTTPException(status_code=400, detail="النص قصير جداً")
    
    batch_size = n_mcq_per_batch + n_tf_per_batch + n_short_per_batch
    if batch_size < 1 or batch_size > 30:
        raise HTTPException(status_code=400, detail="حجم الدفعة بين 1 و 30")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    diff_ar = {"easy": "سهلة ومباشرة", "medium": "متوسطة الصعوبة", "hard": "تحليلية وتطبيقية"}.get(difficulty, "متوسطة")
    
    context_parts = []
    if grade: context_parts.append(f"الصف: {grade}")
    if semester: context_parts.append(f"الفصل: {semester}")
    if unit: context_parts.append(f"الوحدة: {unit}")
    if lesson: context_parts.append(f"الدرس: {lesson}")
    context = " | ".join(context_parts) if context_parts else "غير محدد"
    
    # variation hints لضمان تنوع الدفعات
    variation_hints = [
        "ركّز على الأسئلة المباشرة والتعريفات",
        "ركّز على الأمثلة التطبيقية والمسائل العملية",
        "ركّز على التحليل والمقارنة والاستنتاج",
        "ركّز على المهارات الحسابية والعمليات",
        "ركّز على الفهم والإدراك المفاهيمي",
        "ركّز على حل المشكلات المركّبة",
        "ركّز على التطبيقات الحياتية اليومية",
        "ركّز على الأسئلة العكسية (من الإجابة للسؤال)",
        "ركّز على الأسئلة التي تختبر سوء الفهم الشائع",
        "ركّز على المفاهيم المتقاطعة بين الوحدات"
    ]
    variation = variation_hints[(batch_num - 1) % len(variation_hints)]
    
    prompt = f"""أنت خبير في تعليم الرياضيات للمرحلة الأساسية في سلطنة عُمان.

المهمة: أنشئ دفعة {batch_num} من الأسئلة ({batch_size} سؤال) من محتوى الدرس.

⚠️ مهم: سبق توليد {already_generated} سؤال من نفس المحتوى. تجنّب التكرار تماماً.
🎯 توجيه الدفعة الحالية: {variation}

السياق: {context}
الصعوبة: {diff_ar}
التوزيع:
- {n_mcq_per_batch} اختياري (4 خيارات)
- {n_tf_per_batch} صواب/خطأ
- {n_short_per_batch} إجابة قصيرة

محتوى الدرس:
---
{text[:5000]}
---

قواعد:
- العربية الفصحى، أرقام عربية (٠-٩)
- الإجابة للاختياري تطابق أحد الخيارات نصاً
- الإجابة لصواب/خطأ: "صواب" أو "خطأ" فقط
- أسئلة جديدة كلياً، غير مكررة من الدفعات السابقة

أخرج JSON فقط:
{{"questions":[{{"type":"mcq","q":"...","options":["...","...","...","..."],"answer":"..."}},...]}}"""
    
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": 0.85 + (batch_num % 3) * 0.03,  # تغيير الإبداع بين الدفعات
                    "maxOutputTokens": 4000,
                    "topP": 0.9,
                    "responseMimeType": "application/json"
                }
            }
            with httpx.Client(timeout=45.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw_text = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    
                    if raw_text.startswith("```"):
                        raw_text = raw_text.split("\n", 1)[1] if "\n" in raw_text else raw_text
                        raw_text = raw_text.rsplit("```", 1)[0] if "```" in raw_text else raw_text
                    
                    parsed = json_lib.loads(raw_text)
                    questions = parsed.get("questions", [])
                    if not isinstance(questions, list) or len(questions) == 0:
                        raise ValueError("استجابة فارغة")
                    
                    type_map = {"mcq": "اختياري", "tf": "صواب/خطأ", "short": "إجابة قصيرة"}
                    for q in questions:
                        q["type"] = type_map.get(q.get("type", "mcq"), q.get("type", "اختياري"))
                        if not isinstance(q.get("options"), list):
                            q["options"] = []
                        q["_batch"] = batch_num  # نُعلّم الدفعة
                    
                    return {
                        "status": "ok",
                        "count": len(questions),
                        "questions": questions,
                        "model": model_name,
                        "batch_num": batch_num,
                        "variation": variation
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:150]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:150]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشلت الدفعة {batch_num}: {last_error}")


# ════════════════════════════════════════════════════════════════════════════
# 💬 AI CHAT — شات تفاعلي لتعديل/تحسين الأسئلة
# ════════════════════════════════════════════════════════════════════════════
@app.post("/api/prep/ai_chat_modify")
async def prep_ai_chat_modify(
    request: Request,
    questions_json: str = Form(...),
    instruction: str = Form(...),
    selected_indices: str = Form(default=""),
    admin = Depends(get_current_admin)
):
    """
    💬 المعلم يعطي تعليمات لتعديل الأسئلة:
    - "حسّن السؤال 3"
    - "اجعل كل الأسئلة أصعب"
    - "غيّر صياغة الأسئلة 1، 5، 7"
    - "أضف تنوّعاً في الأسئلة"
    - "صحّح الإملاء"
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=20, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة")
    
    if not instruction.strip():
        raise HTTPException(status_code=400, detail="لا توجد تعليمات")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="GEMINI_API_KEY مفقود")
    
    import json as json_lib
    
    try:
        all_questions = json_lib.loads(questions_json)
    except json_lib.JSONDecodeError:
        raise HTTPException(status_code=400, detail="JSON غير صحيح")
    
    if not isinstance(all_questions, list) or len(all_questions) == 0:
        raise HTTPException(status_code=400, detail="لا توجد أسئلة")
    
    # نحدد الأسئلة المستهدفة
    target_indices = []
    if selected_indices.strip():
        try:
            target_indices = [int(i) for i in selected_indices.split(",") if i.strip().isdigit()]
            target_indices = [i for i in target_indices if 0 <= i < len(all_questions)]
        except Exception:
            target_indices = []
    
    # إذا لم تُحدد أسئلة، نطبّق على الكل (حتى 30 سؤال كحد أقصى)
    if not target_indices:
        target_indices = list(range(min(len(all_questions), 30)))
    
    target_questions = [all_questions[i] for i in target_indices]
    
    prompt = f"""أنت محرّر خبير لأسئلة الرياضيات. مهمتك تعديل الأسئلة التالية حسب تعليمات المعلم.

📋 تعليمات المعلم:
"{instruction.strip()}"

🎯 الأسئلة المستهدفة ({len(target_questions)} سؤال):
{json_lib.dumps(target_questions, ensure_ascii=False, indent=2)}

⚠️ قواعد صارمة:
- أعد نفس عدد الأسئلة بنفس الترتيب
- حافظ على نفس البنية (type, q, options, answer)
- طبّق التعليمات بدقة
- العربية الفصحى، أرقام عربية (٠-٩)
- الإجابة تطابق أحد الخيارات (للاختياري)

أخرج JSON فقط بهذا الشكل:
{{"questions": [{{...}}, {{...}}, ...]}}"""
    
    import httpx
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": 0.7,
                    "maxOutputTokens": 6000,
                    "responseMimeType": "application/json"
                }
            }
            with httpx.Client(timeout=60.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```", 1)[0] if "```" in raw else raw
                    parsed = json_lib.loads(raw)
                    modified = parsed.get("questions", [])
                    if len(modified) != len(target_questions):
                        # اختلاف العدد - نعطي تحذير لكن نقبل
                        pass
                    
                    return {
                        "status": "ok",
                        "modified_count": len(modified),
                        "indices": target_indices,
                        "questions": modified,
                        "model": model_name
                    }
                else:
                    last_error = f"HTTP {resp.status_code}"
        except json_lib.JSONDecodeError:
            last_error = "JSON parse error"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:150]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل التعديل: {last_error}")


# ════════════════════════════════════════════════════════════════════════════
# 📝 WORKSHEET GENERATOR — مولّد أوراق العمل بـ AI
# ════════════════════════════════════════════════════════════════════════════
@app.post("/api/prep/worksheet_generate")
async def prep_worksheet_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    worksheet_type: str = Form(default="practice"),
    difficulty: str = Form(default="medium"),
    n_questions: int = Form(default=10),
    extra_text: str = Form(default=""),
    include_instructions: bool = Form(default=True),
    admin = Depends(get_current_admin)
):
    """
    📝 توليد ورقة عمل كاملة بـ AI:
    - عنوان + تعليمات + أسئلة متنوعة
    - أنواع: practice (تدريب) / review (مراجعة) / homework (واجب) / quiz (اختبار قصير)
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=15, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    if n_questions < 1 or n_questions > 30:
        raise HTTPException(status_code=400, detail="عدد الأسئلة بين 1 و 30")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    type_ar = {
        "practice": "ورقة تدريب",
        "review": "ورقة مراجعة",
        "homework": "واجب منزلي",
        "quiz": "اختبار قصير"
    }.get(worksheet_type, "ورقة تدريب")
    
    diff_ar = {"easy": "سهلة", "medium": "متوسطة", "hard": "متقدمة"}.get(difficulty, "متوسطة")
    
    context_parts = []
    if grade: context_parts.append(f"الصف: {grade}")
    if semester: context_parts.append(f"الفصل: {semester}")
    if unit: context_parts.append(f"الوحدة: {unit}")
    if lesson: context_parts.append(f"الدرس: {lesson}")
    context = " | ".join(context_parts) if context_parts else "رياضيات عامة"
    
    extra_context = f"\n\nمحتوى إضافي للاستناد إليه:\n{extra_text[:3000]}" if extra_text.strip() else ""
    
    prompt = f"""أنت خبير في إعداد أوراق العمل التعليمية للرياضيات في سلطنة عُمان.

المهمة: أنشئ {type_ar} متكاملة.

السياق: {context}
نوع الورقة: {type_ar}
الصعوبة: {diff_ar}
عدد الأسئلة: {n_questions}{extra_context}

أخرج JSON فقط بهذا الشكل:
{{
  "title": "عنوان ورقة العمل",
  "instructions": "تعليمات واضحة للطالب",
  "learning_objective": "الهدف التعليمي المختصر",
  "estimated_time": "الزمن المقدّر بالدقائق",
  "questions": [
    {{"num": 1, "type": "mcq", "q": "نص السؤال", "options": ["أ","ب","ج","د"], "answer": "أ", "marks": 2}},
    {{"num": 2, "type": "fill", "q": "أكمل: ... = ____", "answer": "الإجابة", "marks": 1}},
    {{"num": 3, "type": "solve", "q": "مسألة تتطلب حلاً", "answer": "خطوات الحل", "marks": 3}}
  ],
  "total_marks": 20
}}

قواعد:
- العربية الفصحى، أرقام عربية (٠-٩)
- نوّع الأسئلة (mcq اختياري، fill إكمال، solve حل مسائل، tf صواب/خطأ)
- لكل سؤال درجة (marks) مناسبة لصعوبته
- أسئلة متدرّجة من السهل للصعب
- مناسبة لنوع الورقة ({type_ar})"""
    
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {
                    "temperature": 0.8,
                    "maxOutputTokens": 6000,
                    "responseMimeType": "application/json"
                }
            }
            with httpx.Client(timeout=60.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```", 1)[0] if "```" in raw else raw
                    parsed = json_lib.loads(raw)
                    
                    if not parsed.get("questions"):
                        raise ValueError("لا توجد أسئلة")
                    
                    return {
                        "status": "ok",
                        "worksheet": parsed,
                        "model": model_name,
                        "type": worksheet_type,
                        "context": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson}
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:150]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:150]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل توليد ورقة العمل: {last_error}")


@app.post("/api/teacher/worksheet_generate")
async def teacher_worksheet_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    worksheet_type: str = Form(default="practice"),
    difficulty: str = Form(default="medium"),
    n_questions: int = Form(default=10),
    extra_text: str = Form(default=""),
    include_instructions: bool = Form(default=True),
):
    """📝 [معلم] توليد ورقة عمل — من لوحة المعلم (بلا حماية admin)"""
    # يستدعي نفس منطق دالة الأدمن مباشرةً (admin غير مطلوب)
    return await prep_worksheet_generate(
        request=request, grade=grade, semester=semester, unit=unit, lesson=lesson,
        worksheet_type=worksheet_type, difficulty=difficulty, n_questions=n_questions,
        extra_text=extra_text, include_instructions=include_instructions, admin="teacher"
    )


# ════════════════════════════════════════════════════════════════════════════
# 🎬 SLIDES GENERATOR — مولّد العروض التقديمية بـ AI
# ════════════════════════════════════════════════════════════════════════════
async def _generate_slides_core(grade, semester, unit, lesson, n_slides, n_questions, extra_text, ip):
    """🎬 المنطق الأساسي لتوليد العرض — مشترك بين الأدمن والمعلم"""
    if _is_rate_limited(ip, max_calls=15, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    if n_slides < 4 or n_slides > 20:
        raise HTTPException(status_code=400, detail="عدد الشرائح بين 4 و 20")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    context_parts = []
    if grade: context_parts.append(f"الصف: {grade}")
    if semester: context_parts.append(f"الفصل: {semester}")
    if unit: context_parts.append(f"الوحدة: {unit}")
    if lesson: context_parts.append(f"الدرس: {lesson}")
    context = " | ".join(context_parts) if context_parts else "رياضيات عامة"
    
    extra_context = f"\n\nمحتوى الدرس للاستناد إليه:\n{extra_text[:3000]}" if extra_text.strip() else ""
    
    prompt = f"""أنت خبير في إعداد العروض التقديمية التعليمية للرياضيات في سلطنة عُمان.

المهمة: أنشئ عرضاً تقديمياً من {n_slides} شريحة لدرس رياضيات، يتضمن {n_questions} سؤال تفاعلي.

السياق: {context}{extra_context}

أنواع الشرائح المطلوبة بالترتيب:
1. شريحة عنوان (title): عنوان الدرس
2. شريحة أهداف (objectives): أهداف الدرس كنقاط
3. شريحة خريطة ذهنية (mindmap): المفهوم المركزي وتتفرّع منه ٣-٥ أفكار رئيسية
3-؟. شرائح شرح (content): المفاهيم كنقاط واضحة
شرائح أمثلة (example): مثال محلول خطوة بخطوة
شريحة رسم بياني (chart): إذا كان الدرس يتضمّن أرقاماً أو نسباً أو مقارنات، أضف شريحة بيانات
{n_questions} شريحة سؤال تفاعلي (question): سؤال + إجابة مخفية
شريحة ختام (closing): تلخيص

أخرج JSON فقط بهذا الشكل:
{{
  "title": "عنوان الدرس",
  "slides": [
    {{"type": "title", "title": "عنوان الدرس", "subtitle": "الصف · الوحدة"}},
    {{"type": "objectives", "title": "أهداف الدرس", "points": ["هدف 1", "هدف 2", "هدف 3"]}},
    {{"type": "mindmap", "title": "خريطة المفهوم", "center": "المفهوم المركزي", "branches": ["فكرة 1", "فكرة 2", "فكرة 3", "فكرة 4"]}},
    {{"type": "content", "title": "عنوان الفكرة", "points": ["نقطة 1", "نقطة 2"]}},
    {{"type": "example", "title": "مثال محلول", "problem": "نص المسألة", "steps": ["الخطوة 1", "الخطوة 2"], "result": "النتيجة النهائية"}},
    {{"type": "chart", "title": "عنوان الرسم", "chart_label": "وصف ما تمثّله الأرقام", "data": [{{"label": "أ", "value": 30}}, {{"label": "ب", "value": 50}}, {{"label": "ج", "value": 20}}]}},
    {{"type": "question", "title": "سؤال تفاعلي", "question": "نص السؤال", "answer": "الإجابة الصحيحة"}},
    {{"type": "closing", "title": "الخلاصة", "points": ["ملخص 1", "ملخص 2"]}}
  ]
}}

قواعد:
- العربية الفصحى، أرقام عربية (٠-٩)
- كل شريحة محتواها مركّز ومناسب للعرض (ليس نصاً طويلاً)
- الأمثلة المحلولة واضحة خطوة بخطوة
- الأسئلة التفاعلية مناسبة لمستوى الطلاب
- **مهم:** أضف شريحة mindmap واحدة للمفهوم الرئيسي (٣-٥ تفرّعات)
- شريحة chart فقط إذا كان الدرس يحتوي أرقاماً/نسباً/مقارنات حقيقية (قيم data أرقام صحيحة موجبة)؛ إن لم يكن الدرس رقمياً فلا تضف chart
- المجموع تقريباً {n_slides} شريحة"""
    
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {"temperature": 0.8, "maxOutputTokens": 8000, "responseMimeType": "application/json"}
            }
            with httpx.Client(timeout=70.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```", 1)[0] if "```" in raw else raw
                    parsed = json_lib.loads(raw)
                    if not parsed.get("slides"):
                        raise ValueError("لا توجد شرائح")
                    return {
                        "status": "ok",
                        "presentation": parsed,
                        "model": model_name,
                        "context": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson}
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:150]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:150]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل توليد العرض: {last_error}")


@app.post("/api/prep/slides_generate")
async def prep_slides_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    n_slides: int = Form(default=8),
    n_questions: int = Form(default=3),
    extra_text: str = Form(default=""),
    admin = Depends(get_current_admin)
):
    """🎬 [أدمن] توليد عرض تقديمي — من داخل المصنع"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_slides_core(grade, semester, unit, lesson, n_slides, n_questions, extra_text, ip)


@app.post("/api/teacher/slides_generate")
async def teacher_slides_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    n_slides: int = Form(default=8),
    n_questions: int = Form(default=3),
    extra_text: str = Form(default=""),
):
    """🎬 [معلم] توليد عرض تقديمي — من لوحة المعلم (بلا حماية admin، مثل بقية أدوات المعلم)"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_slides_core(grade, semester, unit, lesson, n_slides, n_questions, extra_text, ip)


# ════════════════════════════════════════════════════════════════════════════
# 🎥 CINEMATIC LESSON — مولّد الدروس السينمائية المتحركة (GSAP)
# ════════════════════════════════════════════════════════════════════════════
async def _generate_cinematic_core(grade, semester, unit, lesson, teacher_name, n_examples, n_challenges, ip):
    """🎥 المنطق الأساسي لتوليد فيلم الدرس — نظام القوالب (AI يملأ المحتوى فقط)"""
    if _is_rate_limited(ip, max_calls=12, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    n_examples = max(1, min(n_examples, 8))
    n_challenges = max(0, min(n_challenges, 3))
    
    context_parts = []
    if grade: context_parts.append(f"الصف: {grade}")
    if semester: context_parts.append(f"الفصل: {semester}")
    if unit: context_parts.append(f"الوحدة: {unit}")
    if lesson: context_parts.append(f"الدرس: {lesson}")
    context = " | ".join(context_parts) if context_parts else "رياضيات عامة"
    
    prompt = f"""أنت خبير في إعداد الدروس التعليمية المتحركة للرياضيات في سلطنة عُمان.

المهمة: أنشئ محتوى "فيلم درس" متحرك من مشاهد متتابعة.

السياق: {context}

المشاهد المطلوبة بالترتيب:
1. مشهد مفهوم واحد (concept): يشرح الفكرة الأساسية + 3 مصطلحات للصناديق الملونة
2. عدد {n_examples} مشهد مثال محلول (example): مسألة + خطوات حل + نتيجة نهائية
{f'3. عدد {n_challenges} مشهد تحدّي (challenge): لغز أو سؤال محفّز للتفكير' if n_challenges > 0 else ''}
4. مشهد توصيات ختامي واحد (summary): 3 نصائح + جملة ختامية ملهمة

أخرج JSON فقط بهذا الشكل بالضبط:
{{
  "lesson_title": "عنوان الدرس المختصر",
  "lesson_code": "رمز الدرس مثل ٢٣-١",
  "scenes": [
    {{
      "type": "concept",
      "title": "ما هو الموضوع؟",
      "lines": ["جملة شرح 1", "جملة شرح 2", "جملة شرح 3"],
      "boxes": [
        {{"label": "مصطلح 1", "color": "green"}},
        {{"label": "مصطلح 2", "color": "pink"}},
        {{"label": "مصطلح 3", "color": "yellow"}}
      ]
    }},
    {{
      "type": "example",
      "title": "مثال ١: العنوان",
      "problem": "نص المسألة",
      "steps": ["الخطوة 1", "الخطوة 2"],
      "result": "النتيجة النهائية"
    }},
    {{
      "type": "challenge",
      "title": "لغز التحدي 👑",
      "lines": ["سطر اللغز 1", "سطر اللغز 2"],
      "hint": "السؤال المطلوب"
    }},
    {{
      "type": "summary",
      "title": "توصيات وتغذية راجعة 🌟",
      "tips": ["نصيحة 1", "نصيحة 2", "نصيحة 3"],
      "closing": "جملة ختامية ملهمة"
    }}
  ]
}}

قواعد صارمة:
- العربية الفصحى، أرقام عربية (٠-٩) في كل النصوص
- كل مثال: مسألة واضحة + خطوتان للحل على الأقل + نتيجة
- النصوص مركّزة ومناسبة للعرض المتحرك (ليست فقرات طويلة)
- ألوان الصناديق من: green, pink, yellow, cyan فقط
- التزم بالترتيب: concept ثم examples ثم challenges ثم summary"""
    
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {"temperature": 0.85, "maxOutputTokens": 8000, "responseMimeType": "application/json"}
            }
            with httpx.Client(timeout=75.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```", 1)[0] if "```" in raw else raw
                    parsed = json_lib.loads(raw)
                    if not parsed.get("scenes"):
                        raise ValueError("لا توجد مشاهد")
                    return {
                        "status": "ok",
                        "movie": parsed,
                        "model": model_name,
                        "context": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson, "teacher_name": teacher_name}
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:150]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:150]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل توليد الفيلم: {last_error}")


@app.post("/api/prep/cinematic_generate")
async def prep_cinematic_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    teacher_name: str = Form(default="أ. رشدي سيد"),
    n_examples: int = Form(default=4),
    n_challenges: int = Form(default=1),
    admin = Depends(get_current_admin)
):
    """🎥 [أدمن] توليد فيلم درس سينمائي — من المصنع"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_cinematic_core(grade, semester, unit, lesson, teacher_name, n_examples, n_challenges, ip)


@app.post("/api/teacher/cinematic_generate")
async def teacher_cinematic_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    teacher_name: str = Form(default=""),
    n_examples: int = Form(default=4),
    n_challenges: int = Form(default=1),
):
    """🎥 [معلم] توليد فيلm درس سينمائي — من لوحة المعلم"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_cinematic_core(grade, semester, unit, lesson, teacher_name, n_examples, n_challenges, ip)


# ════════════════════════════════════════════════════════════════════════════
# 👑 ROYAL LESSON — الدرس الإمبراطوري (قالب سينمائي متقن + AI للمحتوى)
# ════════════════════════════════════════════════════════════════════════════
async def _generate_royal_core(grade, semester, unit, lesson, teacher_name, ip):
    """👑 محتوى الدرس الإمبراطوري — مسرح + playlist + 6 أمثلة حركية"""
    if _is_rate_limited(ip, max_calls=12, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    ctx = []
    if grade: ctx.append(f"الصف: {grade}")
    if semester: ctx.append(f"الفصل: {semester}")
    if unit: ctx.append(f"الوحدة: {unit}")
    if lesson: ctx.append(f"الدرس: {lesson}")
    context = " | ".join(ctx) if ctx else "رياضيات عامة"
    
    prompt = f"""أنت خبير تعليم رياضيات في سلطنة عُمان وخبير في تصميم الدروس التفاعلية.

المهمة: أنشئ محتوى "درس إمبراطوري" متكامل سيُعرض في تطبيق ويب سينمائي بمسرح عرض وقائمة محطات.

السياق: {context}

الدرس يتكون من محطات بالترتيب التربوي:
1. تمهيد (hook): قصة أو لغز يربط الدرس بالواقع
2. مفهوم (concept): القانون/المفهوم الأساسي بشرح مركّز + 3 نقاط
3. ستة أمثلة محلولة (examples): متدرّجة الصعوبة، كل مثال: مسألة + خطوات حل واضحة + نتيجة نهائية
4. تقويم تكويني (formative): سؤال اختيار من متعدد + 4 خيارات + تغذية راجعة
5. تقويم ختامي (summative): تحدّي بإدخال رقمي + تلميحات + الإجابة

أخرج JSON فقط بهذا الشكل بالضبط:
{{
  "title": "عنوان الدرس",
  "code": "رمز مثل ٢٩-١",
  "hook": {{"story": "القصة التشويقية المرتبطة بالواقع", "highlight": "العبارة المفتاحية"}},
  "concept": {{"title": "عنوان المفهوم", "law": "القانون أو القاعدة الأساسية", "points": ["نقطة 1", "نقطة 2", "نقطة 3"]}},
  "examples": [
    {{"title": "مثال ١", "problem": "نص المسألة", "steps": ["خطوة 1", "خطوة 2"], "result": "النتيجة"}}
  ],
  "formative": {{"question": "السؤال", "options": ["أ", "ب", "ج", "د"], "correct_index": 0, "feedback": "شرح الإجابة الصحيحة"}},
  "summative": {{"question": "سؤال التحدّي", "hints": ["تلميح 1", "تلميح 2"], "answer": "الإجابة النهائية"}}
}}

قواعد صارمة:
- العربية الفصحى، أرقام عربية (٠-٩) في كل النصوص
- بالضبط ٦ أمثلة في examples، متدرّجة من السهل للصعب
- كل مثال خطوتان حل على الأقل + نتيجة واضحة
- 3 نقاط في المفهوم، 4 خيارات في formative
- correct_index من 0 إلى 3
- المحتوى دقيق رياضياً ومناسب لمستوى الصف"""
    
    import httpx, json as jl
    models = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_err = None
    for m in models:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{m}:generateContent?key={api_key}"
            with httpx.Client(timeout=80.0) as client:
                r = client.post(url, json={"contents":[{"parts":[{"text":prompt}]}],"generationConfig":{"temperature":0.82,"maxOutputTokens":9000,"responseMimeType":"application/json"}})
                if r.status_code == 200:
                    raw = r.json().get("candidates",[{}])[0].get("content",{}).get("parts",[{}])[0].get("text","").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n",1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```",1)[0] if "```" in raw else raw
                    parsed = jl.loads(raw)
                    if not parsed.get("examples"):
                        raise ValueError("لا أمثلة")
                    return {"status":"ok","royal":parsed,"model":m,"context":{"grade":grade,"semester":semester,"unit":unit,"lesson":lesson,"teacher_name":teacher_name}}
                else:
                    last_err = f"HTTP {r.status_code}"
        except jl.JSONDecodeError as e:
            last_err = f"JSON: {str(e)[:100]}"
        except Exception as e:
            last_err = f"{m}: {str(e)[:120]}"
            continue
    raise HTTPException(status_code=502, detail=f"❌ فشل توليد الدرس: {last_err}")


@app.post("/api/prep/royal_generate")
async def prep_royal_generate(
    request: Request,
    grade: str = Form(default=""), semester: str = Form(default=""),
    unit: str = Form(default=""), lesson: str = Form(default=""),
    teacher_name: str = Form(default="أ. رشدي سيد"),
    admin = Depends(get_current_admin)
):
    """👑 [أدمن] توليد درس إمبراطوري"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_royal_core(grade, semester, unit, lesson, teacher_name, ip)


@app.post("/api/teacher/royal_generate")
async def teacher_royal_generate(
    request: Request,
    grade: str = Form(default=""), semester: str = Form(default=""),
    unit: str = Form(default=""), lesson: str = Form(default=""),
    teacher_name: str = Form(default=""),
):
    """👑 [معلم] توليد درس إمبراطوري"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_royal_core(grade, semester, unit, lesson, teacher_name, ip)


# ════════════════════════════════════════════════════════════════════════════
# 🎮 INTERACTIVE LESSON — مولّد الدروس التفاعلية (مختبر + كويز + تحدّي)
# ════════════════════════════════════════════════════════════════════════════
async def _generate_interactive_core(grade, semester, unit, lesson, teacher_name, ip):
    """🎮 المنطق الأساسي للدرس التفاعلي — قوالب عامة قابلة لأي درس"""
    if _is_rate_limited(ip, max_calls=12, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise HTTPException(status_code=503, detail="❌ GEMINI_API_KEY مفقود")
    
    context_parts = []
    if grade: context_parts.append(f"الصف: {grade}")
    if semester: context_parts.append(f"الفصل: {semester}")
    if unit: context_parts.append(f"الوحدة: {unit}")
    if lesson: context_parts.append(f"الدرس: {lesson}")
    context = " | ".join(context_parts) if context_parts else "رياضيات عامة"
    
    prompt = f"""أنت خبير في تصميم الدروس التفاعلية للرياضيات في سلطنة عُمان.

المهمة: أنشئ محتوى درس تفاعلي كامل يتفاعل معه الطالب بنفسه.

السياق: {context}

الدرس يتكون من أقسام متتابعة:
1. قصة تشويقية (hook): فقرة جذابة تربط الدرس بالحياة
2. بطاقات مفاهيم (concepts): 3 مفاهيم أساسية لكل منها قانون/قاعدة
3. مختبر تفاعلي (lab): الطالب يحرّك منزلقاً (slider) ليرى نتيجة تتغير
4. مسرح أمثلة (examples): 3 أمثلة كل مثال عدة خطوات حل
5. نقطة تحقق (quiz): سؤال اختيار من متعدد + 4 خيارات
6. الفخ الشائع (pitfall): خطأ يقع فيه الطلاب + الصواب
7. نصائح ذهبية (tips): 3 نصائح قابلة للطي
8. تحدّي (challenge): سؤال صعب + 3 تلميحات + الإجابة
9. ملخص (summary): جملة ختامية + نسبة إتقان

أخرج JSON فقط بهذا الشكل بالضبط:
{{
  "lesson_title": "عنوان الدرس",
  "lesson_code": "رمز مثل ٢٩-١",
  "hook": {{"story": "القصة التشويقية", "highlight": "الكلمة المفتاحية"}},
  "concepts": [
    {{"title": "المفهوم", "desc": "الشرح", "law": "القانون أو القاعدة"}}
  ],
  "lab": {{
    "lab_type": "اختر النوع الأنسب للدرس: solids3d (مجسّمات ثلاثية الأبعاد) | fractions (كسور) | ratio (نسبة) | bars (إحصاء/أعمدة) | geometry (هندسة مستوية) | slider (عام)",
    "title": "عنوان المختبر",
    "description": "اشرح ماذا يفعل الطالب في المختبر",
    "note": "ملاحظة تعليمية",
    "config": {{
      "shape": "للـ solids3d فقط: cube | pyramid | cylinder | sphere | cone | prism",
      "fraction_parts": "للـ fractions فقط: عدد الأجزاء الكلي مثل 8",
      "fraction_shaded": "للـ fractions فقط: الأجزاء الملوّنة مثل 3",
      "ratio_a": "للـ ratio فقط: الجزء الأول مثل 3",
      "ratio_b": "للـ ratio فقط: الجزء الثاني مثل 1",
      "ratio_label_a": "للـ ratio: اسم الأول مثل ليمون",
      "ratio_label_b": "للـ ratio: اسم الثاني مثل نعناع",
      "bars_data": "للـ bars فقط: مصفوفة مثل [{{\\"label\\":\\"أ\\",\\"value\\":5}}]",
      "base_value": "للـ slider فقط: قيمة أساسية رقم 1-5",
      "unit_name": "للـ slider: اسم الوحدة",
      "multiplier_label": "للـ slider: وصف المنزلق"
    }}
  }},
  "examples": [
    {{"title": "مثال 1: العنوان", "steps": ["خطوة 1", "خطوة 2", "النتيجة"]}}
  ],
  "quiz": {{
    "question": "نص السؤال",
    "hint": "تلميح",
    "options": ["خيار 1", "خيار 2", "خيار 3", "خيار 4"],
    "correct_index": 1,
    "success_msg": "أحسنت! الشرح"
  }},
  "pitfall": {{"title": "الخطأ الشائع", "wrong": "الفهم الخاطئ", "right": "الفهم الصحيح"}},
  "tips": [
    {{"title": "عنوان النصيحة", "body": "تفاصيل النصيحة"}}
  ],
  "challenge": {{
    "question": "سؤال التحدّي",
    "hints": ["تلميح 1", "تلميح 2", "تلميح 3"],
    "answer": "الإجابة النهائية"
  }},
  "summary": {{"closing": "جملة ختامية ملهمة", "mastery": "ممتاز"}}
}}

قواعد صارمة:
- العربية الفصحى، أرقام عربية (٠-٩) في كل النصوص العربية (لكن أبقِ الأرقام في config بالإنجليزية)
- 3 مفاهيم، 3 أمثلة، 3 نصائح، 3 تلميحات بالضبط
- **مهم جداً — اختيار lab_type:** اختر النوع الذي يخدم الدرس فعلياً:
  * دروس المجسّمات/الحجم/المساحة السطحية → solids3d (حدد shape المناسب)
  * دروس الكسور → fractions
  * دروس النسبة/التناسب → ratio
  * دروس الإحصاء/البيانات → bars
  * دروس الهندسة المستوية/المضلّعات → geometry
  * أي درس آخر (أعداد، جبر، عمليات) → slider
- املأ config فقط بحقول النوع المختار (الباقي اتركه فارغاً أو احذفه)
- correct_index رقم من 0 إلى 3
- المحتوى دقيق رياضياً ومناسب لمستوى الصف"""
    
    import httpx
    import json as json_lib
    
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash", "gemini-2.5-flash-lite"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {"temperature": 0.85, "maxOutputTokens": 9000, "responseMimeType": "application/json"}
            }
            with httpx.Client(timeout=80.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```", 1)[0] if "```" in raw else raw
                    parsed = json_lib.loads(raw)
                    if not parsed.get("concepts"):
                        raise ValueError("محتوى ناقص")
                    return {
                        "status": "ok",
                        "lesson": parsed,
                        "model": model_name,
                        "context": {"grade": grade, "semester": semester, "unit": unit, "lesson": lesson, "teacher_name": teacher_name}
                    }
                else:
                    last_error = f"HTTP {resp.status_code}: {resp.text[:150]}"
        except json_lib.JSONDecodeError as e:
            last_error = f"JSON: {str(e)[:100]}"
        except Exception as e:
            last_error = f"{model_name}: {str(e)[:150]}"
            continue
    
    raise HTTPException(status_code=502, detail=f"❌ فشل توليد الدرس التفاعلي: {last_error}")


@app.post("/api/prep/interactive_generate")
async def prep_interactive_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    teacher_name: str = Form(default="أ. رشدي سيد"),
    admin = Depends(get_current_admin)
):
    """🎮 [أدمن] توليد درس تفاعلي — من المصنع"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_interactive_core(grade, semester, unit, lesson, teacher_name, ip)


@app.post("/api/teacher/interactive_generate")
async def teacher_interactive_generate(
    request: Request,
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    teacher_name: str = Form(default=""),
):
    """🎮 [معلم] توليد درس تفاعلي — من لوحة المعلم"""
    ip = request.client.host if request.client else "unknown"
    return await _generate_interactive_core(grade, semester, unit, lesson, teacher_name, ip)


@app.post("/api/prep/library_save")
async def prep_library_save(
    request: Request,
    title: str = Form(...),
    grade: str = Form(default=""),
    semester: str = Form(default=""),
    unit: str = Form(default=""),
    lesson: str = Form(default=""),
    lesson_type: str = Form(default="interactive"),
    html_content: str = Form(...),
):
    """📚 رفع درس (فيلم/تفاعلي) لمكتبة الدروس التفاعلية — يخزّن HTML في Supabase Storage"""
    if not html_content or len(html_content) < 100:
        raise HTTPException(status_code=400, detail="المحتوى فارغ أو قصير")
    if len(html_content) > 5 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="الملف أكبر من 5 ميجابايت")
    
    import time as _t
    safe_title = "".join(ch for ch in title if ch.isalnum() or ch in " _-").strip()[:50] or "lesson"
    file_name = f"lessons/{lesson_type}_{int(_t.time())}_{safe_title}.html"
    
    try:
        supabase.storage.from_("resources").upload(
            path=file_name,
            file=html_content.encode("utf-8"),
            file_options={"content-type": "text/html; charset=utf-8"}
        )
        file_url = supabase.storage.from_("resources").get_public_url(file_name)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل رفع الملف: {str(e)[:200]}")
    
    type_label = {"cinematic": "🎥 فيلم سينمائي", "interactive": "🎮 درس تفاعلي"}.get(lesson_type, "درس")
    row = {
        "title": title[:200],
        "grade": grade[:100],
        "semester": (semester or "")[:100],
        "unit": (unit or "")[:200],
        "lesson": (lesson or "")[:300],
        "description": type_label,
        "file_url": file_url,
        "file_size_kb": len(html_content.encode("utf-8")) // 1024,
        "sanitized": True,
    }
    try:
        supabase.table("html_lessons").insert(row).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل الحفظ: {str(e)[:200]}")
    
    return {"status": "success", "file_url": file_url, "message": "تم الرفع لمكتبة الدروس التفاعلية"}


# ════════════════════════════════════════════════════════════════════════════
# 📊 POWERPOINT EXPORT — تصدير العرض التقديمي كملف PPTX
# ════════════════════════════════════════════════════════════════════════════
async def _build_pptx(slides_json: str, lesson_title: str, grade: str, teacher_name: str):
    """يبني ملف PowerPoint من شرائح JSON ويُرجع مساره"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        raise HTTPException(status_code=503, detail="❌ مكتبة python-pptx غير مثبتة. شغّل: pip install python-pptx")
    
    import json as jl
    try:
        slides = jl.loads(slides_json)
    except Exception:
        raise HTTPException(status_code=400, detail="بيانات الشرائح غير صحيحة")
    if not isinstance(slides, list) or not slides:
        raise HTTPException(status_code=400, detail="لا توجد شرائح")
    
    # ألوان الإمبراطورية
    NAVY = RGBColor(0x0A, 0x1F, 0x4D)
    DARK = RGBColor(0x04, 0x0D, 0x21)
    GOLD = RGBColor(0xD4, 0xAF, 0x37)
    GOLD_L = RGBColor(0xFC, 0xD3, 0x4E)
    WHITE = RGBColor(0xF0, 0xEA, 0xD6)
    GREEN = RGBColor(0x2E, 0xCC, 0x71)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    
    def set_bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color
    
    def add_box(slide, x, y, w, h, text, size, color, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        lines = text.split("\n") if isinstance(text, str) else [str(text)]
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = ln
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = "Cairo"
        return tb
    
    def add_bar(slide, x, y, w, h, color):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        return shp
    
    for s in slides:
        stype = s.get("type", "content")
        slide = prs.slides.add_slide(blank)
        set_bg(slide, DARK if stype == "title" else NAVY)
        # شريط علوي ذهبي
        add_bar(slide, 0, 0, SW, Inches(0.15), GOLD)
        
        if stype == "title":
            add_box(slide, Inches(1), Inches(2.2), Inches(11.33), Inches(1.5), s.get("title", ""), 48, GOLD_L)
            if s.get("subtitle"):
                add_box(slide, Inches(1), Inches(3.8), Inches(11.33), Inches(1), s.get("subtitle", ""), 26, WHITE, bold=False)
            add_box(slide, Inches(1), Inches(6), Inches(11.33), Inches(0.8), "إمبراطورية الرياضيات 👑", 18, GOLD)
        elif stype == "question":
            add_box(slide, Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8), "❓ سؤال تفاعلي", 24, GOLD)
            add_box(slide, Inches(1), Inches(1.6), Inches(11.33), Inches(1), s.get("title", ""), 32, GOLD_L)
            add_box(slide, Inches(1), Inches(2.8), Inches(11.33), Inches(1.8), s.get("question", ""), 28, WHITE, bold=False)
            if s.get("answer"):
                add_box(slide, Inches(1.5), Inches(5), Inches(10.33), Inches(1.2), "✓ الإجابة: " + str(s.get("answer", "")), 24, GREEN)
        elif stype == "example":
            add_box(slide, Inches(0.5), Inches(0.5), Inches(12.33), Inches(1), s.get("title", ""), 32, GOLD_L)
            body = ""
            if s.get("problem"): body += str(s.get("problem")) + "\n\n"
            for i, st in enumerate(s.get("steps", []), 1):
                body += "• " + str(st) + "\n"
            add_box(slide, Inches(1), Inches(1.7), Inches(11.33), Inches(3.8), body.strip(), 22, WHITE, bold=False, anchor=MSO_ANCHOR.TOP)
            if s.get("result"):
                add_box(slide, Inches(1.5), Inches(5.7), Inches(10.33), Inches(1), "✓ " + str(s.get("result", "")), 26, GREEN)
        else:  # content / objectives / closing
            icon = {"objectives": "🎯 ", "closing": "🏁 ", "content": "📚 "}.get(stype, "")
            add_box(slide, Inches(0.5), Inches(0.5), Inches(12.33), Inches(1.2), icon + s.get("title", ""), 34, GOLD_L)
            pts = s.get("points", [])
            body = "\n".join("◆  " + str(p) for p in pts)
            add_box(slide, Inches(1.2), Inches(2), Inches(11), Inches(4.5), body, 24, WHITE, bold=False, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.RIGHT)
        
        # تذييل
        add_box(slide, Inches(0.5), Inches(6.95), Inches(12.33), Inches(0.4), (teacher_name or "") + "  ·  " + (grade or ""), 12, GOLD, bold=False)
    
    import tempfile, os as _os
    safe = "".join(ch for ch in (lesson_title or "presentation") if ch.isalnum() or ch in " _-").strip()[:50] or "presentation"
    out_path = _os.path.join(tempfile.gettempdir(), f"{safe}_{int(__import__('time').time())}.pptx")
    prs.save(out_path)
    return out_path


@app.post("/api/prep/slides_pptx")
async def prep_slides_pptx(
    slides_json: str = Form(...),
    lesson_title: str = Form(default="عرض تقديمي"),
    grade: str = Form(default=""),
    teacher_name: str = Form(default=""),
    admin = Depends(get_current_admin)
):
    """📊 [أدمن] تصدير العرض كملف PowerPoint"""
    path = await _build_pptx(slides_json, lesson_title, grade, teacher_name)
    safe = "".join(ch for ch in (lesson_title or "presentation") if ch.isalnum() or ch in " _-").strip()[:50] or "presentation"
    return FileResponse(path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=f"{safe}.pptx")


@app.post("/api/teacher/slides_pptx")
async def teacher_slides_pptx(
    slides_json: str = Form(...),
    lesson_title: str = Form(default="عرض تقديمي"),
    grade: str = Form(default=""),
    teacher_name: str = Form(default=""),
):
    """📊 [معلم] تصدير العرض كملف PowerPoint"""
    path = await _build_pptx(slides_json, lesson_title, grade, teacher_name)
    safe = "".join(ch for ch in (lesson_title or "presentation") if ch.isalnum() or ch in " _-").strip()[:50] or "presentation"
    return FileResponse(path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=f"{safe}.pptx")


# ════════════════════════════════════════════════════════════════════════════
# 📊 STUDENT REPORT — توليد تعليق تربوي ذكي لتقرير الطالب
# ════════════════════════════════════════════════════════════════════════════
@app.post("/api/teacher/student_report_comment")
async def teacher_student_report_comment(
    request: Request,
    student_name: str = Form(...),
    grade: str = Form(default=""),
    subject: str = Form(default="الرياضيات"),
    scores_json: str = Form(...),
    total_percent: float = Form(default=0),
    admin = Depends(get_current_admin)
):
    """
    📊 توليد تعليق تربوي مخصص لتقرير الطالب بناءً على درجاته
    يحلل نقاط القوة والضعف ويقترح توصيات
    """
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=30, window_seconds=120):
        raise HTTPException(status_code=429, detail="طلبات كثيرة، انتظر دقيقتين")
    
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        # بديل بدون AI - تعليق تلقائي حسب النسبة
        if total_percent >= 90:
            comment = f"أداء متميز ورائع! {student_name} يُظهر تفوقاً واضحاً. نشجعه على الاستمرار في هذا المستوى المشرّف."
        elif total_percent >= 80:
            comment = f"أداء جيد جداً. {student_name} يمتلك أساساً قوياً، ومع المزيد من التركيز يمكنه بلوغ التميز."
        elif total_percent >= 65:
            comment = f"أداء جيد. ننصح {student_name} بمراجعة النقاط الصعبة وزيادة التدريب لتحسين مستواه."
        elif total_percent >= 50:
            comment = f"أداء مقبول. يحتاج {student_name} إلى دعم إضافي ومتابعة مستمرة لتطوير مهاراته."
        else:
            comment = f"يحتاج {student_name} إلى دعم مكثّف وخطة علاجية. ننصح بالتواصل لوضع برنامج تحسين مناسب."
        return {"status": "ok", "comment": comment, "source": "auto"}
    
    import json as json_lib
    try:
        scores = json_lib.loads(scores_json)
    except json_lib.JSONDecodeError:
        scores = {}
    
    # بناء وصف الدرجات
    scores_desc = "\n".join([f"- {k}: {v}" for k, v in scores.items()]) if isinstance(scores, dict) else str(scores)
    
    prompt = f"""أنت معلم رياضيات خبير في سلطنة عُمان. اكتب تعليقاً تربوياً مختصراً واحترافياً لتقرير الطالب.

اسم الطالب: {student_name}
الصف: {grade}
المادة: {subject}
النسبة الإجمالية: {total_percent}%

تفاصيل الدرجات:
{scores_desc}

اكتب تعليقاً تربوياً:
- مختصر (2-3 جمل)
- إيجابي ومحفّز
- يذكر نقطة قوة محددة
- يقترح توصية عملية للتحسين
- بصيغة رسمية مناسبة لولي الأمر
- بالعربية الفصحى

أخرج JSON فقط:
{{"comment": "التعليق هنا"}}"""
    
    import httpx
    models_to_try = ["gemini-2.5-flash", "gemini-2.0-flash"]
    last_error = None
    
    for model_name in models_to_try:
        try:
            url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
            payload = {
                "contents": [{"parts": [{"text": prompt}]}],
                "generationConfig": {"temperature": 0.7, "maxOutputTokens": 500, "responseMimeType": "application/json"}
            }
            with httpx.Client(timeout=30.0) as client:
                resp = client.post(url, json=payload)
                if resp.status_code == 200:
                    data = resp.json()
                    raw = data.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "").strip()
                    if raw.startswith("```"):
                        raw = raw.split("\n", 1)[1] if "\n" in raw else raw
                        raw = raw.rsplit("```", 1)[0] if "```" in raw else raw
                    parsed = json_lib.loads(raw)
                    return {"status": "ok", "comment": parsed.get("comment", ""), "source": model_name}
                else:
                    last_error = f"HTTP {resp.status_code}"
        except Exception as e:
            last_error = str(e)[:150]
            continue
    
    # fallback لو فشل AI
    comment = f"{student_name} حقق نسبة {total_percent}%. نتمنى له مزيداً من التقدم والتفوق."
    return {"status": "ok", "comment": comment, "source": "fallback", "warning": last_error}


# ════════════════════════════════════════════════════════════════════════════
# 🎨 صفحة المصنع - تقديم HTML المخصص
# ════════════════════════════════════════════════════════════════════════════
@app.get("/prep", response_class=HTMLResponse)
async def prep_page():
    """
    🏭 صفحة مصنع الأسئلة الشخصي
    تتطلب تسجيل دخول بصلاحيات الأدمن (يتم التحقق client-side)
    """
    try:
        with open("templates/prep_factory.html", "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    except FileNotFoundError:
        return HTMLResponse(content="<h1>صفحة المصنع غير موجودة. أضف templates/prep_factory.html</h1>", status_code=404)


@app.post("/api/admin/update_password")
async def update_admin_password(
    new_password: str = Form(...),
    admin = Depends(get_current_admin)
):
    """
    تحديث كلمة مرور الأدمن — يُحفظ في system_state
    ملاحظة: لا يُحدّث env var ADMIN_PASSWORD مباشرة، لكن يُتيح override
    """
    if len(new_password) < 6:
        raise HTTPException(status_code=400, detail="كلمة المرور قصيرة جداً (6+ أحرف)")
    try:
        # نحفظ في system_state لـ override
        supabase.table("system_state").upsert({
            "key": "admin_password_override",
            "value": new_password.strip(),
            "updated_at": datetime.now(timezone.utc).isoformat()
        }).execute()
        return {"status": "success", "message": "تم تحديث كلمة المرور"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/admin/reports/full")
async def get_full_report(admin=Depends(get_current_admin)):
    """🛡️ تقرير شامل دقيق للإمبراطورية
    
    التحسينات:
    - XP من total_points (المصدر الموثوق)
    - إحصائيات زمنية (7/30 يوم)
    - تقسيم الطلاب (نشط/خامل/في خطر)
    - دقة الإحصائيات 100%
    """
    # جلب كل الطلاب — مع total_points للـ XP الحقيقي
    students = []
    offset = 0
    for _ in range(20):
        try:
            res_batch = supabase.table("students").select(
                "id, full_name, grade, username, created_at, total_points, is_active, is_elite"
            ).order("full_name").range(offset, offset + 999).execute()
            batch = res_batch.data or []
            if not batch:
                break
            students.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        except Exception as e:
            print(f"students pagination error: {e}")
            break

    # جلب كل النتائج (مع pagination لتجاوز حد 1000)
    results = []
    offset = 0
    page_size = 1000
    for _ in range(50):  # حد أقصى 50,000 نتيجة
        try:
            res_batch = supabase.table("results").select(
                "student_id, student_name, lesson, score, total, timestamp"
            ).order("timestamp", desc=True).range(offset, offset + page_size - 1).execute()
            batch = res_batch.data or []
            if not batch:
                break
            results.extend(batch)
            if len(batch) < page_size:
                break
            offset += page_size
        except Exception as e:
            print(f"results pagination error: {e}")
            break

    # جلب عدد الأسئلة الفعلي (count من Supabase وليس len(data) لأن data محدد بـ 1000)
    try:
        q_res = supabase.table("questions").select("id", count="exact").limit(1).execute()
        total_questions = q_res.count if hasattr(q_res, "count") and q_res.count is not None else 0
    except Exception:
        total_questions = 0

    # ════════════════════════════════════════════════════════════
    # 🧮 بناء إحصائيات لكل طالب (دقة 100%)
    # ════════════════════════════════════════════════════════════
    from collections import defaultdict
    from datetime import datetime, timezone, timedelta
    
    student_stats = defaultdict(lambda: {
        "tests": 0,
        "score_sum": 0,
        "total_sum": 0,
        "lessons": set(),
        "tests_7d": 0,
        "tests_30d": 0,
        "last_test": None,
    })
    
    now = datetime.now(timezone.utc)
    cutoff_7d = (now - timedelta(days=7)).isoformat()
    cutoff_30d = (now - timedelta(days=30)).isoformat()
    
    for r in results:
        sid = r.get("student_id")
        if not sid:
            continue
        ts = r.get("timestamp", "")
        student_stats[sid]["tests"]     += 1
        student_stats[sid]["score_sum"] += (r.get("score") or 0)
        student_stats[sid]["total_sum"] += (r.get("total") or 0)
        lesson = r.get("lesson")
        if lesson:
            student_stats[sid]["lessons"].add(lesson)
        if ts and ts >= cutoff_7d:
            student_stats[sid]["tests_7d"] += 1
        if ts and ts >= cutoff_30d:
            student_stats[sid]["tests_30d"] += 1
        if not student_stats[sid]["last_test"] or ts > student_stats[sid]["last_test"]:
            student_stats[sid]["last_test"] = ts

    # ════════════════════════════════════════════════════════════
    # 📊 دمج البيانات — XP من total_points (المصدر الموثوق)
    # ════════════════════════════════════════════════════════════
    students_report = []
    students_at_risk = []  # طلاب لم ينشطوا منذ 30 يوم
    students_inactive = []  # طلاب لم يبدأوا أبداً
    
    for s in students:
        sid = s["id"]
        stats = student_stats.get(sid, {})
        tests = stats.get("tests", 0)
        score_sum = stats.get("score_sum", 0)
        total_sum = stats.get("total_sum", 0)
        accuracy = round((score_sum / total_sum * 100), 1) if total_sum > 0 else 0
        
        # 🛡️ XP الحقيقي من DB
        real_xp = int(s.get("total_points", 0) or 0)
        last_test = stats.get("last_test")
        
        # تحديد حالة النشاط
        is_dormant = False
        if last_test:
            try:
                last_dt = datetime.fromisoformat(last_test.replace('Z', '+00:00'))
                days_inactive = (now - last_dt).days
                is_dormant = days_inactive >= 30
            except Exception:
                pass
        
        student_data = {
            "id":             sid,
            "full_name":      s.get("full_name", ""),
            "grade":          s.get("grade", ""),
            "username":       s.get("username", ""),
            "joined":         s.get("created_at", ""),
            "last_active":    s.get("last_active") or last_test or "",
            "xp":             real_xp,
            "total_points":   real_xp,
            "tests":          tests,
            "tests_7d":       stats.get("tests_7d", 0),
            "tests_30d":      stats.get("tests_30d", 0),
            "accuracy":       accuracy,
            "lessons_count":  len(stats.get("lessons", set())),
            "is_dormant":     is_dormant,
            "is_active":      bool(s.get("is_active", True)),
            "is_elite":       bool(s.get("is_elite", False)),
        }
        students_report.append(student_data)
        
        if tests == 0:
            students_inactive.append(student_data)
        elif is_dormant:
            students_at_risk.append(student_data)

    # ترتيب حسب XP
    students_report.sort(key=lambda x: x["xp"], reverse=True)

    # ════════════════════════════════════════════════════════════
    # 📈 إحصائيات الصفوف (مع XP إجمالي)
    # ════════════════════════════════════════════════════════════
    grade_stats = defaultdict(lambda: {"count": 0, "total_xp": 0, "total_tests": 0, "active": 0})
    for s in students_report:
        g = s.get("grade") or "غير محدد"
        grade_stats[g]["count"] += 1
        grade_stats[g]["total_xp"] += s["xp"]
        grade_stats[g]["total_tests"] += s["tests"]
        if s["tests"] > 0:
            grade_stats[g]["active"] += 1
    
    # حساب المتوسطات
    grade_distribution_detailed = {}
    for g, st in grade_stats.items():
        cnt = st["count"]
        grade_distribution_detailed[g] = {
            "count":      cnt,
            "total_xp":   st["total_xp"],
            "avg_xp":     round(st["total_xp"] / cnt) if cnt > 0 else 0,
            "active":     st["active"],
            "activity_pct": round((st["active"] / cnt) * 100, 1) if cnt > 0 else 0,
        }

    # ════════════════════════════════════════════════════════════
    # 🎯 إحصائيات شاملة
    # ════════════════════════════════════════════════════════════
    total_students = len(students)
    active_students = sum(1 for s in students_report if s["tests"] > 0)
    active_7d = sum(1 for s in students_report if s["tests_7d"] > 0)
    active_30d = sum(1 for s in students_report if s["tests_30d"] > 0)
    total_xp_sum = sum(s["xp"] for s in students_report)
    
    return {
        "summary": {
            "total_students":     total_students,
            "total_questions":    total_questions,
            "total_results":      len(results),
            "active_students":    active_students,
            "active_last_7days":  active_7d,
            "active_last_30days": active_30d,
            "inactive_count":     len(students_inactive),
            "at_risk_count":      len(students_at_risk),
            "total_xp_empire":    total_xp_sum,
            "avg_xp_per_student": round(total_xp_sum / total_students) if total_students > 0 else 0,
            "elite_count":        sum(1 for s in students_report if s.get("is_elite")),
            "engagement_rate":    round((active_30d / total_students) * 100, 1) if total_students > 0 else 0,
        },
        "grade_distribution":     {g: v["count"] for g, v in grade_distribution_detailed.items()},
        "grade_distribution_detailed": grade_distribution_detailed,
        "students":               students_report,
        "top10":                  students_report[:10],
        "students_at_risk":       students_at_risk[:50],     # أعلى 50 في خطر
        "students_inactive":      students_inactive[:50],    # أعلى 50 خاملين
    }





# ════════════════════════════════════════════════════════════
# 🧠 ENDPOINTS ذكية جديدة للتحليلات العميقة
# ════════════════════════════════════════════════════════════

@app.get("/api/admin/analytics/hardest_lessons")
async def get_hardest_lessons(limit: int = 10, admin=Depends(get_current_admin)):
    """🎯 الدروس الأكثر صعوبة على الطلاب
    يُرجع: الدروس + الطلاب المتعثرين في كل درس
    """
    from collections import defaultdict
    
    # 1. جلب كل النتائج مع اسم الطالب
    results = []
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("results").select(
                "student_id, student_name, grade, lesson, score, total"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            results.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # 2. تجميع حسب الدرس + اسم الطالب
    # lesson_stats[lesson] = {total_score, total_max, attempts, student_perfs: {sid: {name, score_sum, total_sum}}}
    lesson_stats = defaultdict(lambda: {
        "attempts": 0,
        "score_sum": 0,
        "total_sum": 0,
        "students": defaultdict(lambda: {
            "name": "", "grade": "", "score_sum": 0, "total_sum": 0, "attempts": 0
        })
    })
    
    for r in results:
        lesson = (r.get("lesson") or "").strip()
        if not lesson: continue
        lesson_stats[lesson]["attempts"]   += 1
        lesson_stats[lesson]["score_sum"]  += (r.get("score") or 0)
        lesson_stats[lesson]["total_sum"]  += (r.get("total") or 0)
        
        sid = r.get("student_id")
        sname = r.get("student_name") or "—"
        sgrade = r.get("grade") or ""
        if sid:
            sp = lesson_stats[lesson]["students"][sid]
            sp["name"] = sname
            sp["grade"] = sgrade
            sp["score_sum"] += (r.get("score") or 0)
            sp["total_sum"] += (r.get("total") or 0)
            sp["attempts"]  += 1
    
    # 3. حساب المتوسطات + استخراج الطلاب المتعثرين
    lessons_ranked = []
    for lesson, st in lesson_stats.items():
        if st["total_sum"] < 5: continue
        avg_pct = round((st["score_sum"] / st["total_sum"]) * 100, 1)
        
        # استخراج الطلاب المتعثرين في هذا الدرس (دقة < 60% مع عينة كافية)
        struggling = []
        for sid, sp in st["students"].items():
            if sp["total_sum"] < 3: continue
            sp_pct = round((sp["score_sum"] / sp["total_sum"]) * 100, 1)
            if sp_pct < 60:
                struggling.append({
                    "id":       sid,
                    "name":     sp["name"],
                    "grade":    sp["grade"],
                    "accuracy": sp_pct,
                    "attempts": sp["attempts"]
                })
        struggling.sort(key=lambda x: x["accuracy"])
        
        lessons_ranked.append({
            "lesson":              lesson,
            "attempts":            st["attempts"],
            "students_count":      len(st["students"]),
            "avg_score_pct":       avg_pct,
            "difficulty_rank":     "صعب جداً" if avg_pct < 50 else "صعب" if avg_pct < 70 else "متوسط" if avg_pct < 85 else "سهل",
            "struggling_students": struggling[:10]  # أصعب 10 طلاب فقط
        })
    
    lessons_ranked.sort(key=lambda x: x["avg_score_pct"])
    
    return {
        "hardest_lessons":  lessons_ranked[:limit],
        "easiest_lessons":  list(reversed(lessons_ranked[-limit:])) if len(lessons_ranked) >= limit else [],
        "total_lessons":    len(lessons_ranked)
    }


@app.get("/api/admin/analytics/activity_heatmap")
async def get_activity_heatmap(days: int = 30, admin=Depends(get_current_admin)):
    """🔥 خريطة حرارة النشاط (نشاط الطلاب حسب اليوم والساعة)"""
    from collections import defaultdict
    from datetime import datetime, timezone, timedelta
    
    cutoff = (datetime.now(timezone.utc) - timedelta(days=days)).isoformat()
    
    # جلب النتائج الحديثة
    results = []
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("results").select(
                "timestamp"
            ).gte("timestamp", cutoff).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            results.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # تجميع: heatmap[day_of_week][hour] = count
    # يوم: 0=الأحد, 6=السبت
    heatmap = [[0] * 24 for _ in range(7)]
    
    for r in results:
        ts = r.get("timestamp")
        if not ts: continue
        try:
            dt = datetime.fromisoformat(ts.replace('Z', '+00:00'))
            # تحويل لتوقيت عمان (UTC+4)
            local = dt + timedelta(hours=4)
            # weekday: Mon=0, Sun=6; نحوّل لـ الأحد=0
            day = (local.weekday() + 1) % 7
            hour = local.hour
            heatmap[day][hour] += 1
        except Exception:
            continue
    
    # أوقات الذروة
    max_val = 0
    peak_day, peak_hour = 0, 0
    for d in range(7):
        for h in range(24):
            if heatmap[d][h] > max_val:
                max_val = heatmap[d][h]
                peak_day, peak_hour = d, h
    
    days_ar = ["الأحد", "الإثنين", "الثلاثاء", "الأربعاء", "الخميس", "الجمعة", "السبت"]
    
    return {
        "heatmap": heatmap,
        "total_activity": sum(sum(row) for row in heatmap),
        "peak_time": {
            "day":      days_ar[peak_day],
            "hour":     peak_hour,
            "count":    max_val,
            "formatted": f"{days_ar[peak_day]} الساعة {peak_hour}:00"
        },
        "days_labels": days_ar,
        "period_days": days
    }


@app.get("/api/admin/analytics/grade_performance")
async def get_grade_performance(admin=Depends(get_current_admin)):
    """📚 أداء كل صف — متوسط XP، نسبة النشاط، عدد الإنجازات"""
    from collections import defaultdict
    
    # جلب الطلاب
    students = []
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("students").select(
                "id, grade, total_points, is_active, created_at"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            students.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # تجميع
    grade_perf = defaultdict(lambda: {
        "count": 0, "total_xp": 0, "active": 0, "elite_count": 0, "ids": []
    })
    
    for s in students:
        g = (s.get("grade") or "غير محدد").strip()
        sid = s.get("id")
        xp = int(s.get("total_points", 0) or 0)
        
        grade_perf[g]["count"] += 1
        grade_perf[g]["total_xp"] += xp
        grade_perf[g]["ids"].append(sid)
        if s.get("is_active", True):
            grade_perf[g]["active"] += 1
    
    # حساب المتوسطات + ترتيب
    grades_ranked = []
    for g, st in grade_perf.items():
        cnt = st["count"]
        grades_ranked.append({
            "grade":         g,
            "students_count": cnt,
            "total_xp":      st["total_xp"],
            "avg_xp":        round(st["total_xp"] / cnt) if cnt > 0 else 0,
            "active_count":  st["active"],
            "activity_pct":  round((st["active"] / cnt) * 100, 1) if cnt > 0 else 0,
        })
    
    grades_ranked.sort(key=lambda x: x["avg_xp"], reverse=True)
    
    return {
        "grades": grades_ranked,
        "best_grade": grades_ranked[0] if grades_ranked else None,
        "total_grades": len(grades_ranked)
    }


@app.get("/api/admin/analytics/student_rankings")
async def get_student_rankings(period: str = "all", limit: int = 20, admin=Depends(get_current_admin)):
    """🏆 ترتيب الطلاب — كل الأوقات / أسبوع / شهر
    
    period: 'all' | '7d' | '30d'
    """
    from collections import defaultdict
    from datetime import datetime, timezone, timedelta
    
    # جلب الطلاب
    students_map = {}
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("students").select(
                "id, full_name, grade, total_points, avatar_url"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            for s in batch:
                students_map[s["id"]] = s
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    if period == "all":
        # ترتيب من total_points مباشرة
        ranked = sorted(
            students_map.values(),
            key=lambda s: int(s.get("total_points", 0) or 0),
            reverse=True
        )
        result = [{
            "rank": i + 1,
            "id": s["id"],
            "full_name": s.get("full_name", ""),
            "grade": s.get("grade", ""),
            "xp": int(s.get("total_points", 0) or 0),
            "avatar_url": s.get("avatar_url", "")
        } for i, s in enumerate(ranked[:limit])]
        return {"rankings": result, "period": "all", "total": len(students_map)}
    
    # للفترة المحددة، نحتاج النتائج
    days = 7 if period == "7d" else 30
    cutoff = (datetime.now(timezone.utc) - timedelta(days=days)).isoformat()
    
    period_xp = defaultdict(int)
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("results").select(
                "student_id, score"
            ).gte("timestamp", cutoff).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            for r in batch:
                sid = r.get("student_id")
                if sid:
                    period_xp[sid] += (r.get("score") or 0)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # ترتيب
    ranked_ids = sorted(period_xp.items(), key=lambda x: x[1], reverse=True)[:limit]
    result = []
    for i, (sid, xp) in enumerate(ranked_ids):
        s = students_map.get(sid, {})
        result.append({
            "rank": i + 1,
            "id": sid,
            "full_name": s.get("full_name", ""),
            "grade": s.get("grade", ""),
            "xp_period": xp,
            "xp_total": int(s.get("total_points", 0) or 0),
            "avatar_url": s.get("avatar_url", "")
        })
    
    return {"rankings": result, "period": period, "total": len(period_xp)}


@app.get("/api/admin/analytics/struggling_students")
async def get_struggling_students(threshold: int = 50, admin=Depends(get_current_admin)):
    """⚠️ طلاب يحتاجون مساعدة (متوسط دقتهم أقل من threshold%)"""
    from collections import defaultdict
    
    # جلب النتائج
    results = []
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("results").select(
                "student_id, score, total"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            results.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # تجميع نسب الإجابة
    student_acc = defaultdict(lambda: {"score": 0, "total": 0, "tests": 0})
    for r in results:
        sid = r.get("student_id")
        if not sid: continue
        student_acc[sid]["score"] += (r.get("score") or 0)
        student_acc[sid]["total"] += (r.get("total") or 0)
        student_acc[sid]["tests"] += 1
    
    # جلب أسماء الطلاب
    students_map = {}
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("students").select(
                "id, full_name, grade, total_points"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            for s in batch:
                students_map[s["id"]] = s
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # تحديد الطلاب المتعثرين
    struggling = []
    for sid, acc in student_acc.items():
        if acc["total"] < 10 or acc["tests"] < 3:  # نحتاج عينة كافية
            continue
        pct = round((acc["score"] / acc["total"]) * 100, 1)
        if pct < threshold:
            s = students_map.get(sid, {})
            struggling.append({
                "id":         sid,
                "full_name":  s.get("full_name", ""),
                "grade":      s.get("grade", ""),
                "accuracy":   pct,
                "tests":      acc["tests"],
                "xp":         int(s.get("total_points", 0) or 0),
                "severity":   "حرج" if pct < 35 else "متعثر" if pct < 50 else "متوسط"
            })
    
    struggling.sort(key=lambda x: x["accuracy"])
    
    return {
        "struggling_students": struggling,
        "count":               len(struggling),
        "threshold":           threshold
    }




@app.get("/api/admin/export/students_csv")
async def export_students_csv(admin=Depends(get_current_admin)):
    """📥 تصدير قائمة الطلاب كـ CSV (يفتح في Excel مباشرة)"""
    from fastapi.responses import Response
    from collections import defaultdict
    from datetime import datetime, timezone
    import csv
    import io
    
    # جلب الطلاب
    students = []
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("students").select(
                "id, full_name, username, grade, school_name, total_points, is_active, is_elite, parent_phone, parent_email, created_at"
            ).order("total_points", desc=True).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            students.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # نتائج لكل طالب (عدد + متوسط)
    student_perf = defaultdict(lambda: {"tests": 0, "score": 0, "total": 0})
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("results").select(
                "student_id, score, total"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            for r in batch:
                sid = r.get("student_id")
                if sid:
                    student_perf[sid]["tests"] += 1
                    student_perf[sid]["score"] += (r.get("score") or 0)
                    student_perf[sid]["total"] += (r.get("total") or 0)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    # بناء CSV
    output = io.StringIO()
    # BOM لـ Excel ليفهم UTF-8
    output.write('﻿')
    writer = csv.writer(output)
    
    # رأس
    writer.writerow([
        "الترتيب", "الاسم الكامل", "اسم المستخدم", "الصف", "المدرسة",
        "XP", "عدد التحديات", "متوسط النسبة %", "نشط", "النخبة",
        "هاتف ولي الأمر", "بريد ولي الأمر", "آخر نشاط", "تاريخ التسجيل"
    ])
    
    # صفوف
    for i, s in enumerate(students):
        perf = student_perf.get(s["id"], {})
        avg = round((perf["score"] / perf["total"] * 100), 1) if perf.get("total", 0) > 0 else 0
        writer.writerow([
            i + 1,
            s.get("full_name", ""),
            s.get("username", ""),
            s.get("grade", ""),
            s.get("school_name", ""),
            s.get("total_points", 0) or 0,
            perf.get("tests", 0),
            avg,
            "نعم" if s.get("is_active", True) else "لا",
            "نعم" if s.get("is_elite") else "لا",
            s.get("parent_phone", ""),
            s.get("parent_email", ""),
            (s.get("last_active") or "")[:10],
            (s.get("created_at") or "")[:10]
        ])
    
    csv_content = output.getvalue()
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    
    return Response(
        content=csv_content.encode('utf-8'),
        media_type="text/csv; charset=utf-8",
        headers={
            "Content-Disposition": f'attachment; filename="students_report_{today}.csv"',
            "Content-Type": "text/csv; charset=utf-8"
        }
    )


@app.get("/api/admin/export/results_csv")
async def export_results_csv(
    days: int = 30,
    grade: str = "",
    admin=Depends(get_current_admin)
):
    """📥 تصدير نتائج التحديات كـ CSV"""
    from fastapi.responses import Response
    from datetime import datetime, timezone, timedelta
    import csv
    import io
    
    cutoff = (datetime.now(timezone.utc) - timedelta(days=days)).isoformat()
    
    results = []
    offset = 0
    for _ in range(50):
        try:
            q = supabase.table("results").select(
                "student_name, grade, lesson, score, total, timestamp"
            ).gte("timestamp", cutoff).order("timestamp", desc=True)
            if grade:
                q = q.eq("grade", grade)
            res = q.range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch: break
            results.extend(batch)
            if len(batch) < 1000: break
            offset += 1000
        except Exception:
            break
    
    output = io.StringIO()
    output.write('﻿')
    writer = csv.writer(output)
    writer.writerow(["اسم الطالب", "الصف", "الدرس", "الدرجة", "المجموع", "النسبة %", "التاريخ"])
    
    for r in results:
        score = r.get("score", 0) or 0
        total = r.get("total", 0) or 0
        pct = round((score / total * 100), 1) if total > 0 else 0
        writer.writerow([
            r.get("student_name", ""),
            r.get("grade", ""),
            r.get("lesson", ""),
            score,
            total,
            pct,
            (r.get("timestamp") or "")[:16].replace("T", " ")
        ])
    
    csv_content = output.getvalue()
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    
    return Response(
        content=csv_content.encode('utf-8'),
        media_type="text/csv; charset=utf-8",
        headers={
            "Content-Disposition": f'attachment; filename="results_last{days}days_{today}.csv"',
            "Content-Type": "text/csv; charset=utf-8"
        }
    )






# ════════════════════════════════════════════════════════════
# 🛡️ SUPERVISOR ENDPOINTS — لوحة المشرف
# ════════════════════════════════════════════════════════════

@app.post("/api/supervisor/login")
async def supervisor_login(
    request: Request,
    email: str = Form(...),
    password: str = Form(...)
):
    """🔐 تسجيل دخول المشرف مع حماية متدرجة"""
    email = (email or "").strip().lower()
    password = password or ""
    client_ip = request.client.host if request.client else "unknown"
    lockout_key = f"supervisor:{email}:{client_ip}"
    
    is_locked, seconds_left = check_progressive_lockout(lockout_key)
    if is_locked:
        minutes = max(1, seconds_left // 60)
        security_log("supervisor_login_blocked", client_ip, {"email": email})
        raise HTTPException(status_code=429, detail=f"🔒 الحساب مقفول مؤقتاً — حاول بعد {minutes} دقيقة")
    
    if _is_rate_limited(client_ip, max_calls=10, window_seconds=60):
        raise HTTPException(status_code=429, detail="⏳ تجاوزت عدد المحاولات")
    
    if not email or not password:
        raise HTTPException(status_code=400, detail="البريد وكلمة المرور مطلوبان")
    
    res = supabase.table("supervisors").select("*").eq("email", email).execute()
    if not res.data:
        record_login_failure(lockout_key)
        security_log("supervisor_login_failed", client_ip, {"email": email, "reason": "not_found"})
        raise HTTPException(status_code=401, detail="❌ بيانات خاطئة")
    
    sup = res.data[0]
    if not sup.get("is_active", True):
        raise HTTPException(status_code=403, detail="🚫 حسابك معطّل — تواصل مع الإدارة")
    
    if not verify_password(password, sup["password"]):
        count, duration = record_login_failure(lockout_key)
        security_log("supervisor_login_failed", client_ip, {"email": email, "attempts": count})
        if duration > 0:
            mins = max(1, duration // 60)
            raise HTTPException(status_code=401, detail=f"❌ بيانات خاطئة — قفل {mins} دقيقة")
        raise HTTPException(status_code=401, detail="❌ بيانات خاطئة")
    
    # نجاح
    record_login_success(lockout_key)
    # تحديث last_login
    try:
        from datetime import datetime, timezone
        supabase.table("supervisors").update({"last_login": datetime.now(timezone.utc).isoformat()}).eq("id", sup["id"]).execute()
    except Exception:
        pass
    token = create_access_token({"sub": str(sup["id"]), "role": "supervisor"})
    security_log("supervisor_login_success", client_ip, {"sup_id": sup["id"]})
    return {"access_token": token, "token_type": "bearer"}




@app.get("/api/supervisor/debug")
async def supervisor_debug(sup = Depends(get_current_supervisor)):
    """🔍 تشخيص للمشرف - يُظهر هويته وطلابه"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    
    # نجلب أسماء الطلاب
    student_names = []
    if student_ids:
        try:
            res = supabase.table("students").select("id, full_name").in_("id", student_ids).execute()
            student_names = [{"id": s["id"], "name": s.get("full_name", "—")} for s in (res.data or [])]
        except Exception as e:
            student_names = [{"error": str(e)[:200]}]
    
    return {
        "supervisor_id": sup["id"],
        "supervisor_name": sup.get("full_name", "—"),
        "supervisor_email": sup.get("email", "—"),
        "student_ids_count": len(student_ids),
        "student_ids": student_ids,
        "students": student_names,
        "diagnosis": "✅ كل شيء يعمل" if student_ids else "⚠️ لم يُربط أي طالب بهذا المشرف بعد"
    }

@app.get("/api/supervisor/me")
async def supervisor_me(sup = Depends(get_current_supervisor)):
    """👤 معلومات المشرف الحالي"""
    sup_copy = dict(sup)
    sup_copy.pop("password", None)
    return sup_copy


@app.get("/api/supervisor/students")
async def supervisor_students(sup = Depends(get_current_supervisor)):
    """👥 طلاب المشرف مع إحصائيات مبسطة"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    if not student_ids:
        return {"students": [], "total": 0}
    
    res = supabase.table("students").select("id, full_name, username, grade, total_points, created_at").in_("id", student_ids).execute()
    students = res.data or []
    
    # 🔧 normalize للتوافق مع الكود القديم
    for s in students:
        s["xp"] = s.get("total_points", 0) or 0
        s["points"] = s.get("total_points", 0) or 0
        s["total_xp"] = s.get("total_points", 0) or 0
        s["level"] = max(1, (s.get("total_points", 0) or 0) // 100)
        s["last_active"] = s.get("created_at", "") or ""
        s["avatar"] = ""
        s["curriculum"] = ""
    
    # نُضيف إحصائيات بسيطة (الدقة، عدد التحديات)
    for s in students:
        try:
            chal = supabase.table("student_challenges").select("is_correct").eq("student_id", s["id"]).execute()
            attempts = chal.data or []
            s["challenges_count"] = len(attempts)
            if attempts:
                correct = sum(1 for a in attempts if a.get("is_correct"))
                s["accuracy"] = round(correct / len(attempts) * 100, 1)
            else:
                s["accuracy"] = 0
        except Exception:
            s["challenges_count"] = 0
            s["accuracy"] = 0
    
    return {"students": students, "total": len(students)}


@app.get("/api/supervisor/students/{student_id}/stats")
async def supervisor_student_stats(student_id: int, sup = Depends(get_current_supervisor)):
    """📊 إحصائيات تفصيلية لطالب من طلاب المشرف"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    if student_id not in student_ids:
        raise HTTPException(status_code=403, detail="هذا الطالب ليس من طلابك")
    
    stu = supabase.table("students").select("*").eq("id", student_id).execute()
    if not stu.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    student = stu.data[0]
    student.pop("password", None)
    
    # التحديات
    chal = supabase.table("student_challenges").select("*").eq("student_id", student_id).order("created_at", desc=True).limit(50).execute()
    challenges = chal.data or []
    
    # حساب الإحصائيات
    total_attempts = len(challenges)
    correct = sum(1 for c in challenges if c.get("is_correct"))
    accuracy = round(correct / total_attempts * 100, 1) if total_attempts else 0
    
    # آخر 7 أيام
    from datetime import datetime, timezone, timedelta
    week_ago = datetime.now(timezone.utc) - timedelta(days=7)
    week_attempts = [c for c in challenges if c.get("created_at") and c["created_at"] > week_ago.isoformat()]
    
    return {
        "student": student,
        "stats": {
            "total_attempts": total_attempts,
            "correct": correct,
            "accuracy": accuracy,
            "weekly_attempts": len(week_attempts),
            "recent_challenges": challenges[:10]
        }
    }


@app.get("/api/supervisor/students/struggling")
async def supervisor_struggling_students(sup = Depends(get_current_supervisor)):
    """⚠️ الطلاب المتعثرون حسب معايير المشرف"""
    settings = sup.get("alert_settings") or {}
    min_accuracy = settings.get("min_accuracy", 50)
    lookback = settings.get("challenges_lookback", 3)
    absent_days = settings.get("absent_days", 3)
    
    student_ids = _get_supervisor_student_ids(sup["id"])
    if not student_ids:
        return {"students": [], "total": 0}
    
    students = supabase.table("students").select("id, full_name, username, grade, total_points, created_at").in_("id", student_ids).execute()
    # normalize للتوافق
    for s in (students.data or []):
        s["xp"] = s.get("total_points", 0) or 0
        s["total_xp"] = s.get("total_points", 0) or 0
        s["last_active"] = s.get("created_at", "") or ""
    
    from datetime import datetime, timezone, timedelta
    threshold_date = datetime.now(timezone.utc) - timedelta(days=absent_days)
    
    struggling = []
    for s in (students.data or []):
        reasons = []
        # 1. دقة منخفضة في آخر N تحديات
        try:
            chal = supabase.table("student_challenges").select("is_correct").eq("student_id", s["id"]).order("created_at", desc=True).limit(lookback).execute()
            attempts = chal.data or []
            if len(attempts) >= lookback:
                correct = sum(1 for a in attempts if a.get("is_correct"))
                acc = correct / len(attempts) * 100
                if acc < min_accuracy:
                    reasons.append(f"دقة {round(acc)}% < {min_accuracy}%")
                    s["recent_accuracy"] = round(acc, 1)
        except Exception:
            pass
        
        # 2. غياب
        last_active = s.get("last_active")
        if last_active and last_active < threshold_date.isoformat():
            days = (datetime.now(timezone.utc) - datetime.fromisoformat(last_active.replace("Z", "+00:00"))).days
            reasons.append(f"غاب {days} يوم")
            s["days_absent"] = days
        elif not last_active:
            reasons.append("لم يبدأ بعد")
        
        if reasons:
            s["reasons"] = reasons
            struggling.append(s)
    
    return {"students": struggling, "total": len(struggling)}


@app.get("/api/supervisor/questions")
async def supervisor_get_questions(
    limit: int = 0,
    grade: str = "",
    semester: str = "",
    unit: str = "",
    lesson: str = "",
    sup = Depends(get_current_supervisor)
):
    """📚 جلب أسئلة من البنك مع فلاتر (الصف/الفصل/الوحدة/الدرس)"""
    try:
        # نختار الأعمدة المتاحة فعلياً
        query = supabase.table("questions").select("*")
        if grade:    query = query.eq("grade", grade)
        if semester: query = query.eq("semester", semester)
        if unit:     query = query.eq("unit", unit)
        if lesson:   query = query.eq("lesson", lesson)
        
        # limit = 0 يعني كل شيء (max 2000)
        max_limit = 2000
        if limit and limit > 0:
            query = query.limit(min(limit, max_limit))
        else:
            query = query.limit(max_limit)
        
        q = query.execute()
        return {"questions": q.data or [], "total": len(q.data or [])}
    except Exception as e:
        return {"questions": [], "error": str(e)[:200]}


@app.post("/api/supervisor/exams")
async def supervisor_create_exam(
    payload: dict,
    sup = Depends(get_current_supervisor)
):
    """➕ إنشاء اختبار جديد"""
    title = (payload.get("title") or "").strip()
    questions = payload.get("questions_json") or payload.get("questions") or []
    if not title or not questions:
        raise HTTPException(status_code=400, detail="العنوان والأسئلة مطلوبة")
    new_exam = {
        "supervisor_id": sup["id"],
        "title": title,
        "description": payload.get("description", ""),
        "questions_json": questions,
        "total_marks": int(payload.get("total_marks", 20)),
        "duration_min": int(payload.get("duration_min", 30)),
        "is_published": bool(payload.get("is_published", False)),
        "target_student_ids": payload.get("target_student_ids", [])
    }
    res = supabase.table("supervisor_exams").insert(new_exam).execute()
    if not res.data:
        raise HTTPException(status_code=500, detail="فشل الإنشاء")
    return {"status": "success", "id": res.data[0]["id"], "exam": res.data[0]}


@app.get("/api/supervisor/exams")
async def supervisor_list_exams(sup = Depends(get_current_supervisor)):
    """📋 قائمة اختبارات المشرف"""
    res = supabase.table("supervisor_exams").select("*").eq("supervisor_id", sup["id"]).order("created_at", desc=True).execute()
    return {"exams": res.data or []}


@app.put("/api/supervisor/exams/{exam_id}")
async def supervisor_update_exam(
    exam_id: int,
    payload: dict,
    sup = Depends(get_current_supervisor)
):
    """✏️ تعديل اختبار"""
    # نتحقق من الملكية
    check = supabase.table("supervisor_exams").select("id").eq("id", exam_id).eq("supervisor_id", sup["id"]).execute()
    if not check.data:
        raise HTTPException(status_code=404, detail="الاختبار غير موجود")
    update_data = {}
    for k in ["title", "description", "total_marks", "duration_min", "is_published", "questions_json", "target_student_ids"]:
        if k in payload:
            update_data[k] = payload[k]
    if not update_data:
        raise HTTPException(status_code=400, detail="لا توجد بيانات للتحديث")
    from datetime import datetime, timezone
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    res = supabase.table("supervisor_exams").update(update_data).eq("id", exam_id).execute()
    return {"status": "success", "exam": res.data[0] if res.data else None}


@app.delete("/api/supervisor/exams/{exam_id}")
async def supervisor_delete_exam(exam_id: int, sup = Depends(get_current_supervisor)):
    """🗑️ حذف اختبار"""
    check = supabase.table("supervisor_exams").select("id").eq("id", exam_id).eq("supervisor_id", sup["id"]).execute()
    if not check.data:
        raise HTTPException(status_code=404, detail="الاختبار غير موجود")
    supabase.table("supervisor_exams").delete().eq("id", exam_id).execute()
    return {"status": "success"}


@app.get("/api/supervisor/exams/{exam_id}/results")
async def supervisor_exam_results(exam_id: int, sup = Depends(get_current_supervisor)):
    """📊 نتائج اختبار"""
    check = supabase.table("supervisor_exams").select("id,title,total_marks").eq("id", exam_id).eq("supervisor_id", sup["id"]).execute()
    if not check.data:
        raise HTTPException(status_code=404, detail="الاختبار غير موجود")
    exam = check.data[0]
    results = supabase.table("supervisor_exam_results").select("*").eq("exam_id", exam_id).execute()
    # نُضيف أسماء الطلاب
    if results.data:
        student_ids = list(set(r["student_id"] for r in results.data))
        stus = supabase.table("students").select("id, full_name, username").in_("id", student_ids).execute()
        stu_map = {s["id"]: s for s in (stus.data or [])}
        for r in results.data:
            r["student"] = stu_map.get(r["student_id"], {})
    return {"exam": exam, "results": results.data or []}




# ════════════════════════════════════════════════════════════
# 📨 رسائل المشرف للطالب
# ════════════════════════════════════════════════════════════
@app.get("/api/student/supervisor-messages")
async def get_supervisor_messages_for_student(request: Request):
    """📨 الطالب يجلب رسائل المشرفين الموجّهة له"""
    # نتحقق من توكن الطالب
    auth_header = request.headers.get("Authorization") or ""
    if not auth_header.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="غير مصرح")
    token = auth_header.split(" ")[1]
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("role") != "student":
            raise HTTPException(status_code=403, detail="هذا الـ endpoint للطلاب فقط")
        student_id = int(payload.get("sub", 0))
        if not student_id:
            raise HTTPException(status_code=401, detail="توكن غير صالح")
    except HTTPException:
        raise
    except Exception:
        raise HTTPException(status_code=401, detail="توكن غير صالح")
    
    try:
        # نجلب الرسائل + اسم المشرف
        res = supabase.table("supervisor_messages")\
            .select("*")\
            .eq("student_id", student_id)\
            .order("created_at", desc=True)\
            .limit(50).execute()
        msgs = res.data or []
        
        # نُضيف اسم المشرف لكل رسالة
        sup_ids = list(set(m["supervisor_id"] for m in msgs if m.get("supervisor_id")))
        if sup_ids:
            sups = supabase.table("supervisors").select("id,full_name").in_("id", sup_ids).execute()
            sup_map = {s["id"]: s["full_name"] for s in (sups.data or [])}
            for m in msgs:
                m["supervisor_name"] = sup_map.get(m.get("supervisor_id"), "المشرف")
        
        # عداد غير المقروء
        unread_count = sum(1 for m in msgs if not m.get("is_read"))
        return {"messages": msgs, "total": len(msgs), "unread": unread_count}
    except Exception as e:
        return {"messages": [], "total": 0, "unread": 0, "error": str(e)[:200]}


@app.put("/api/student/supervisor-messages/{msg_id}/read")
async def mark_supervisor_message_read(msg_id: int, request: Request):
    """✅ الطالب يُحدّد رسالة كمقروءة"""
    auth_header = request.headers.get("Authorization") or ""
    if not auth_header.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="غير مصرح")
    token = auth_header.split(" ")[1]
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("role") != "student":
            raise HTTPException(status_code=403, detail="ليس للطلاب")
        student_id = int(payload.get("sub", 0))
    except Exception:
        raise HTTPException(status_code=401, detail="توكن غير صالح")
    
    try:
        # نتحقق من ملكية الرسالة قبل التحديث
        check = supabase.table("supervisor_messages").select("id,student_id").eq("id", msg_id).execute()
        if not check.data or check.data[0].get("student_id") != student_id:
            raise HTTPException(status_code=404, detail="الرسالة غير موجودة")
        supabase.table("supervisor_messages").update({"is_read": True}).eq("id", msg_id).execute()
        return {"status": "success"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.get("/api/student/supervisor-messages/unread-count")
async def get_unread_supervisor_messages_count(request: Request):
    """🔔 عدد الرسائل غير المقروءة (للجرس)"""
    auth_header = request.headers.get("Authorization") or ""
    if not auth_header.startswith("Bearer "):
        return {"unread": 0}
    token = auth_header.split(" ")[1]
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        if payload.get("role") != "student":
            return {"unread": 0}
        student_id = int(payload.get("sub", 0))
        if not student_id:
            return {"unread": 0}
        res = supabase.table("supervisor_messages").select("id", count="exact")\
            .eq("student_id", student_id).eq("is_read", False).execute()
        return {"unread": res.count or 0}
    except Exception:
        return {"unread": 0}




# ════════════════════════════════════════════════════════════
# 🧭 نظام "بوصلة النجاح" - Student 360 + Smart Recommendations
# ════════════════════════════════════════════════════════════



@app.get("/api/supervisor/live-pulse")
async def supervisor_live_pulse(sup = Depends(get_current_supervisor)):
    """⚡ نبض اللحظة - بيانات لحظية عن نشاط الطلاب"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    if not student_ids:
        return {"active_now": 0, "today_attempts": 0, "today_accuracy": 0, "recent": []}
    
    from datetime import datetime, timezone, timedelta
    now = datetime.now(timezone.utc)
    today_start = now.replace(hour=0, minute=0, second=0, microsecond=0)
    fifteen_min_ago = now - timedelta(minutes=15)
    
    try:
        # الطلاب الذين أجابوا في آخر 15 دقيقة
        active_res = supabase.table("student_challenges").select(
            "student_id"
        ).in_("student_id", student_ids).gte("created_at", fifteen_min_ago.isoformat()).execute()
        
        active_ids = set([c["student_id"] for c in (active_res.data or [])])
        active_now = len(active_ids)
        
        # محاولات اليوم
        today_res = supabase.table("student_challenges").select(
            "is_correct, student_id, created_at"
        ).in_("student_id", student_ids).gte("created_at", today_start.isoformat()).order(
            "created_at", desc=True
        ).limit(500).execute()
        
        today_data = today_res.data or []
        today_attempts = len(today_data)
        today_correct = sum(1 for c in today_data if c.get("is_correct"))
        today_accuracy = round((today_correct / today_attempts * 100) if today_attempts else 0, 1)
        
        # آخر 5 نشاطات
        recent = []
        if today_data:
            # نجلب أسماء الطلاب
            recent_ids = list(set([c["student_id"] for c in today_data[:10]]))
            names_res = supabase.table("students").select("id, full_name").in_("id", recent_ids).execute()
            names_map = {s["id"]: s.get("full_name", "—") for s in (names_res.data or [])}
            
            for c in today_data[:5]:
                try:
                    t = datetime.fromisoformat(c["created_at"].replace("Z", "+00:00"))
                    diff = (now - t).total_seconds()
                    if diff < 60:
                        ago = "الآن"
                    elif diff < 3600:
                        ago = f"منذ {int(diff//60)} د"
                    elif diff < 86400:
                        ago = f"منذ {int(diff//3600)} س"
                    else:
                        ago = "اليوم"
                    
                    recent.append({
                        "student_name": names_map.get(c["student_id"], "—"),
                        "correct": c.get("is_correct", False),
                        "time_ago": ago
                    })
                except Exception:
                    pass
        
        return {
            "active_now": active_now,
            "today_attempts": today_attempts,
            "today_accuracy": today_accuracy,
            "recent": recent,
            "total_students": len(student_ids)
        }
    except Exception as e:
        print(f"[live-pulse] error: {e}")
        return {"active_now": 0, "today_attempts": 0, "today_accuracy": 0, "recent": [], "error": str(e)[:100]}



# ════════════════════════════════════════════════════════════
# 🏆 نظام شهادات التقدير من المشرف
# ════════════════════════════════════════════════════════════

# قوالب الشهادات المتاحة
CERTIFICATE_TEMPLATES = {
    "excellence": {
        "name": "👑 شهادة التميّز الإمبراطوري",
        "icon": "👑",
        "color": "#d4af37",
        "category": "academic",
        "description": "للطلاب المتميّزين أكاديمياً"
    },
    "champion": {
        "name": "🏆 شهادة بطل التحديات",
        "icon": "🏆",
        "color": "#e74c3c",
        "category": "sport",
        "description": "للنشطين في التحديات"
    },
    "elite": {
        "name": "💎 شهادة النخبة الذهبية",
        "icon": "💎",
        "color": "#9b59b6",
        "category": "elite",
        "description": "لأعضاء النادي الإمبراطوري"
    },
    "progress": {
        "name": "🌟 شهادة التحسّن المستمر",
        "icon": "🌟",
        "color": "#f39c12",
        "category": "motivation",
        "description": "للطلاب الذين تحسّنوا بشكل ملحوظ"
    },
    "mastery": {
        "name": "🎯 شهادة الإتقان الكامل",
        "icon": "🎯",
        "color": "#3498db",
        "category": "achievement",
        "description": "لمن حقق 100% في الاختبارات"
    },
    "dedication": {
        "name": "📚 شهادة المثابر الذهبي",
        "icon": "📚",
        "color": "#27ae60",
        "category": "attendance",
        "description": "للمواظبة والمثابرة"
    }
}

def _generate_qr_code() -> str:
    """يُولّد رمز فريد للشهادة"""
    import secrets, string
    chars = string.ascii_uppercase + string.digits
    return "ME-" + "".join(secrets.choice(chars) for _ in range(12))


@app.get("/api/supervisor/certificates/templates")
async def get_certificate_templates(sup = Depends(get_current_supervisor)):
    """📋 قائمة قوالب الشهادات المتاحة"""
    return {"templates": CERTIFICATE_TEMPLATES}


@app.post("/api/supervisor/certificates/issue")
async def issue_certificate(
    student_id: int = Form(...),
    template_type: str = Form("excellence"),
    title: str = Form(...),
    message: str = Form(""),
    sup = Depends(get_current_supervisor)
):
    """🏆 إصدار شهادة تقدير لطالب"""
    
    # تحقق: التيمبليت موجود
    if template_type not in CERTIFICATE_TEMPLATES:
        raise HTTPException(status_code=400, detail=f"قالب غير صالح: {template_type}")
    
    # تحقق: الطالب من طلابي
    sup_students = _get_supervisor_student_ids(sup["id"])
    if student_id not in sup_students:
        raise HTTPException(status_code=403, detail="هذا الطالب ليس من طلابك")
    
    # جلب اسم الطالب
    try:
        st_res = supabase.table("students").select("id, full_name").eq("id", student_id).limit(1).execute()
        if not st_res.data:
            raise HTTPException(status_code=404, detail="الطالب غير موجود")
        student_name = st_res.data[0].get("full_name", "—")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ جلب الطالب: {str(e)[:150]}")
    
    # توليد رمز QR فريد
    qr_code = _generate_qr_code()
    
    # الإدراج
    try:
        new_cert = {
            "student_id": student_id,
            "supervisor_id": sup["id"],
            "supervisor_name": sup.get("full_name", "—"),
            "template_type": template_type,
            "title": title.strip()[:200],
            "message": message.strip()[:1000],
            "student_name": student_name,
            "qr_code": qr_code,
        }
        res = supabase.table("certificates").insert(new_cert).execute()
        if not res.data:
            raise HTTPException(status_code=500, detail="فشل الإدراج")
        
        cert = res.data[0]
        print(f"[cert] ✅ issued #{cert['id']} for student {student_id} by sup {sup['id']}")
        return {"status": "success", "certificate": cert}
    except HTTPException:
        raise
    except Exception as e:
        err_str = str(e)
        if "does not exist" in err_str.lower() or "relation" in err_str.lower():
            raise HTTPException(status_code=500, detail="جدول certificates غير موجود في Supabase!")
        raise HTTPException(status_code=500, detail=f"خطأ الإصدار: {err_str[:200]}")


@app.get("/api/supervisor/certificates/list")
async def list_supervisor_certificates(sup = Depends(get_current_supervisor)):
    """📋 قائمة الشهادات التي منحها المشرف"""
    try:
        res = supabase.table("certificates").select(
            "id, student_id, student_name, template_type, title, message, qr_code, issued_at, is_revoked"
        ).eq("supervisor_id", sup["id"]).order("issued_at", desc=True).limit(200).execute()
        
        certs = res.data or []
        # إضافة معلومات القالب
        for c in certs:
            tpl = CERTIFICATE_TEMPLATES.get(c.get("template_type", "excellence"), {})
            c["template_name"] = tpl.get("name", "—")
            c["template_icon"] = tpl.get("icon", "🏆")
            c["template_color"] = tpl.get("color", "#d4af37")
        
        return {"certificates": certs, "total": len(certs)}
    except Exception as e:
        print(f"[cert list] {e}")
        return {"certificates": [], "total": 0, "error": str(e)[:200]}


@app.get("/api/student/certificates")
async def get_my_certificates(student = Depends(get_current_student)):
    """🏆 شهادات الطالب"""
    try:
        res = supabase.table("certificates").select(
            "id, supervisor_name, template_type, title, message, qr_code, issued_at"
        ).eq("student_id", student["id"]).eq("is_revoked", False).order("issued_at", desc=True).limit(100).execute()
        
        certs = res.data or []
        for c in certs:
            tpl = CERTIFICATE_TEMPLATES.get(c.get("template_type", "excellence"), {})
            c["template_name"] = tpl.get("name", "—")
            c["template_icon"] = tpl.get("icon", "🏆")
            c["template_color"] = tpl.get("color", "#d4af37")
        
        return {"certificates": certs, "total": len(certs)}
    except Exception as e:
        print(f"[student certs] {e}")
        return {"certificates": [], "total": 0}


@app.get("/api/certificate/verify/{qr_code}")
async def verify_certificate(qr_code: str):
    """🔍 التحقق من صحة شهادة بـ QR Code (عام - لا يحتاج تسجيل دخول)"""
    try:
        res = supabase.table("certificates").select(
            "id, student_name, supervisor_name, template_type, title, message, issued_at, is_revoked"
        ).eq("qr_code", qr_code.upper().strip()).limit(1).execute()
        
        if not res.data:
            return {"valid": False, "reason": "الشهادة غير موجودة"}
        
        cert = res.data[0]
        if cert.get("is_revoked"):
            return {"valid": False, "reason": "تم سحب الشهادة", "certificate": cert}
        
        tpl = CERTIFICATE_TEMPLATES.get(cert.get("template_type", "excellence"), {})
        cert["template_name"] = tpl.get("name", "—")
        cert["template_icon"] = tpl.get("icon", "🏆")
        cert["template_color"] = tpl.get("color", "#d4af37")
        
        return {"valid": True, "certificate": cert}
    except Exception as e:
        return {"valid": False, "reason": f"خطأ في التحقق: {str(e)[:100]}"}


@app.delete("/api/supervisor/certificates/{cert_id}")
async def revoke_certificate(cert_id: int, sup = Depends(get_current_supervisor)):
    """❌ سحب شهادة (لا تُحذف، تُعلَّم كملغاة)"""
    try:
        # نتحقق أنها شهادة المشرف
        check = supabase.table("certificates").select("id, supervisor_id").eq("id", cert_id).limit(1).execute()
        if not check.data:
            raise HTTPException(status_code=404, detail="الشهادة غير موجودة")
        if check.data[0]["supervisor_id"] != sup["id"]:
            raise HTTPException(status_code=403, detail="لا يمكنك سحب شهادة لم تمنحها")
        
        # نُعلّمها كملغاة
        from datetime import datetime, timezone
        supabase.table("certificates").update({
            "is_revoked": True,
            "revoked_at": datetime.now(timezone.utc).isoformat()
        }).eq("id", cert_id).execute()
        
        return {"status": "success", "message": "تم سحب الشهادة"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ السحب: {str(e)[:200]}")


@app.get("/api/supervisor/certificates/suggestions/{student_id}")
async def suggest_certificate(student_id: int, sup = Depends(get_current_supervisor)):
    """💡 اقتراح شهادة مناسبة لطالب بناءً على بياناته"""
    
    # تحقق: الطالب من طلابي
    sup_students = _get_supervisor_student_ids(sup["id"])
    if student_id not in sup_students:
        raise HTTPException(status_code=403, detail="هذا الطالب ليس من طلابك")
    
    suggestions = []
    
    try:
        # جلب بيانات الطالب
        st_res = supabase.table("students").select("id, full_name, total_points, is_elite").eq("id", student_id).limit(1).execute()
        if not st_res.data:
            return {"suggestions": []}
        student = st_res.data[0]
        total_points = student.get("total_points", 0) or 0
        is_elite = student.get("is_elite", False)
        
        # جلب التحدّيات
        chal_res = supabase.table("student_challenges").select("is_correct, created_at").eq("student_id", student_id).execute()
        chals = chal_res.data or []
        total_chals = len(chals)
        correct_chals = sum(1 for c in chals if c.get("is_correct"))
        accuracy = (correct_chals / total_chals * 100) if total_chals else 0
        
        # القواعد الذكية
        
        # 1. للنخبة
        if is_elite:
            suggestions.append({
                "template": "elite",
                "title": "عضو النخبة الذهبية",
                "message": f"تكريماً للطالب/ـة {student['full_name']} لعضويته في النادي الإمبراطوري وتميّزه المستمر.",
                "reason": "💎 عضو في النادي الإمبراطوري",
                "priority": 10
            })
        
        # 2. للإتقان (دقة 95%+)
        if total_chals >= 10 and accuracy >= 95:
            suggestions.append({
                "template": "mastery",
                "title": "إتقان كامل في التعلّم",
                "message": f"تكريماً للطالب/ـة {student['full_name']} على تحقيق دقة استثنائية {accuracy:.1f}% في {total_chals} تحدي.",
                "reason": f"🎯 دقة {accuracy:.1f}% في {total_chals} تحدي",
                "priority": 9
            })
        
        # 3. لبطل التحديات (50+ تحدي)
        if total_chals >= 50:
            suggestions.append({
                "template": "champion",
                "title": "بطل التحديات",
                "message": f"تكريماً للطالب/ـة {student['full_name']} على إنجاز {total_chals} تحدي بمعدل {accuracy:.0f}% دقة.",
                "reason": f"🏆 {total_chals} تحدي مُكتمل",
                "priority": 8
            })
        
        # 4. للمتفوقين (XP عالي)
        if total_points >= 1000:
            suggestions.append({
                "template": "excellence",
                "title": "التميّز الأكاديمي",
                "message": f"تكريماً للطالب/ـة {student['full_name']} على جمع {total_points} نقطة وتميّزه في الأداء الأكاديمي.",
                "reason": f"👑 {total_points} نقطة XP",
                "priority": 7
            })
        
        # 5. للتحسّن (إن كان من الضعاف وتحسّن)
        if total_chals >= 20 and 50 <= accuracy < 80:
            # نقارن الأسبوع الأخير بالأول
            from datetime import datetime, timezone, timedelta
            week_ago = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
            recent = [c for c in chals if c.get("created_at", "") > week_ago]
            if recent:
                recent_acc = sum(1 for c in recent if c.get("is_correct")) / len(recent) * 100
                if recent_acc > accuracy + 10:
                    suggestions.append({
                        "template": "progress",
                        "title": "تحسّن ملحوظ ومستمر",
                        "message": f"تكريماً للطالب/ـة {student['full_name']} على التحسّن الكبير من {accuracy:.0f}% إلى {recent_acc:.0f}% في أداء التحديات.",
                        "reason": f"🌟 تحسّن +{recent_acc-accuracy:.0f}% هذا الأسبوع",
                        "priority": 9
                    })
        
        # 6. للمواظبة (نشط مؤخراً)
        if total_chals >= 30:
            suggestions.append({
                "template": "dedication",
                "title": "المثابرة والمواظبة",
                "message": f"تكريماً للطالب/ـة {student['full_name']} على المواظبة والمثابرة في حل {total_chals} تحدي.",
                "reason": f"📚 {total_chals} تحدي بمواظبة",
                "priority": 5
            })
        
        # ترتيب حسب الأولوية
        suggestions.sort(key=lambda x: x.get("priority", 0), reverse=True)
        
        return {
            "suggestions": suggestions[:3],  # أفضل 3
            "student_stats": {
                "total_points": total_points,
                "is_elite": is_elite,
                "total_challenges": total_chals,
                "accuracy": round(accuracy, 1)
            }
        }
    except Exception as e:
        print(f"[cert suggest] {e}")
        return {"suggestions": [], "error": str(e)[:200]}


@app.get("/verify-certificate/{qr_code}")
async def verify_certificate_page(request: Request, qr_code: str):
    """صفحة التحقق العامة من الشهادة"""
    return templates.TemplateResponse(request=request, name="verify_certificate.html")






# ════════════════════════════════════════════════════════════
# 🎁 الميزات المتقدمة للشهادات
# ════════════════════════════════════════════════════════════

# قواعد الإصدار التلقائي
AUTO_CERT_RULES = {
    "first_100_correct": {
        "name": "🎊 أول 100 إجابة صحيحة",
        "template": "champion",
        "title": "إنجاز أول 100 إجابة صحيحة",
        "message_template": "تكريماً للطالب/ـة {name} على إنجاز أول 100 إجابة صحيحة في رحلته التعليمية. هذا إنجاز يستحق الفخر!",
        "check": lambda stats: stats.get("correct_total", 0) >= 100 and stats.get("correct_total", 0) < 110
    },
    "first_500_correct": {
        "name": "🏆 خمسمئة إجابة صحيحة",
        "template": "champion",
        "title": "بطل الـ 500 إجابة",
        "message_template": "تكريماً للطالب/ـة {name} على بلوغه/ا 500 إجابة صحيحة. تميّز رائع يستحق التكريم!",
        "check": lambda stats: 500 <= stats.get("correct_total", 0) < 510
    },
    "perfect_week": {
        "name": "💎 أسبوع مثالي",
        "template": "mastery",
        "title": "أسبوع الإتقان الكامل",
        "message_template": "تكريماً للطالب/ـة {name} على تحقيق دقة 100% لمدة أسبوع كامل. إنجاز استثنائي!",
        "check": lambda stats: stats.get("week_accuracy", 0) == 100 and stats.get("week_attempts", 0) >= 10
    },
    "high_xp": {
        "name": "👑 صاحب الـ 2000 نقطة",
        "template": "excellence",
        "title": "صاحب الـ 2000 نقطة الذهبية",
        "message_template": "تكريماً للطالب/ـة {name} على جمع 2000 نقطة. تميّز أكاديمي بارز يستحق الإشادة!",
        "check": lambda stats: 2000 <= stats.get("total_points", 0) < 2100
    },
    "improved_50": {
        "name": "🌟 تحسّن بنسبة 50%",
        "template": "progress",
        "title": "صاحب التحسّن الاستثنائي",
        "message_template": "تكريماً للطالب/ـة {name} على التحسّن الملحوظ بنسبة 50% في الدقة خلال الأسابيع الماضية.",
        "check": lambda stats: stats.get("improvement", 0) >= 50
    }
}


@app.get("/api/public/wall-of-fame")
async def get_wall_of_fame(limit: int = 20):
    """🏅 جدار الفخر العام - آخر الشهادات المُمنوحة (للجميع)"""
    try:
        limit = min(max(limit, 1), 50)
        res = supabase.table("certificates").select(
            "id, student_name, supervisor_name, template_type, title, issued_at, qr_code"
        ).eq("is_revoked", False).order("issued_at", desc=True).limit(limit).execute()
        
        certs = res.data or []
        for c in certs:
            tpl = CERTIFICATE_TEMPLATES.get(c.get("template_type", "excellence"), {})
            c["template_icon"] = tpl.get("icon", "🏆")
            c["template_color"] = tpl.get("color", "#d4af37")
            c["template_name"] = tpl.get("name", "—")
        
        # إخفاء جزئي من اسم الطالب لحماية الخصوصية
        for c in certs:
            name = c.get("student_name", "")
            if len(name) > 3:
                parts = name.split()
                if len(parts) >= 2:
                    c["student_name_display"] = parts[0] + " " + parts[-1][0] + "."
                else:
                    c["student_name_display"] = name[:3] + "..."
            else:
                c["student_name_display"] = name
        
        return {"certificates": certs, "total": len(certs)}
    except Exception as e:
        print(f"[wall-of-fame] {e}")
        return {"certificates": [], "total": 0}


@app.get("/api/supervisor/certificates/auto-check")
async def check_auto_certificates(sup = Depends(get_current_supervisor)):
    """🤖 فحص أتمتة الشهادات - يقترح من يستحق شهادة تلقائية"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    if not student_ids:
        return {"candidates": []}
    
    candidates = []
    
    try:
        # جلب بيانات كل الطلاب
        st_res = supabase.table("students").select("id, full_name, total_points, is_elite").in_("id", student_ids).execute()
        students = {s["id"]: s for s in (st_res.data or [])}
        
        # جلب التحديات لكل الطلاب
        from datetime import datetime, timezone, timedelta
        week_ago = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
        
        chal_res = supabase.table("student_challenges").select(
            "student_id, is_correct, created_at"
        ).in_("student_id", student_ids).execute()
        all_chals = chal_res.data or []
        
        # تجميع الإحصائيات لكل طالب
        student_stats = {}
        for sid in student_ids:
            chals = [c for c in all_chals if c.get("student_id") == sid]
            week_chals = [c for c in chals if c.get("created_at", "") > week_ago]
            
            correct_total = sum(1 for c in chals if c.get("is_correct"))
            week_correct = sum(1 for c in week_chals if c.get("is_correct"))
            week_accuracy = (week_correct / len(week_chals) * 100) if week_chals else 0
            
            student = students.get(sid, {})
            
            # حساب التحسّن (مقارنة آخر أسبوع بما قبله)
            two_weeks_ago = (datetime.now(timezone.utc) - timedelta(days=14)).isoformat()
            prev_week_chals = [c for c in chals if two_weeks_ago < c.get("created_at", "") <= week_ago]
            prev_correct = sum(1 for c in prev_week_chals if c.get("is_correct"))
            prev_acc = (prev_correct / len(prev_week_chals) * 100) if prev_week_chals else 0
            improvement = week_accuracy - prev_acc
            
            student_stats[sid] = {
                "name": student.get("full_name", "—"),
                "total_points": student.get("total_points", 0) or 0,
                "is_elite": student.get("is_elite", False),
                "correct_total": correct_total,
                "total_attempts": len(chals),
                "week_attempts": len(week_chals),
                "week_accuracy": round(week_accuracy, 1),
                "improvement": round(improvement, 1)
            }
        
        # جلب الشهادات الحالية لتجنب التكرار
        existing_res = supabase.table("certificates").select("student_id, metadata").eq(
            "supervisor_id", sup["id"]
        ).execute()
        existing_keys = set()
        for c in (existing_res.data or []):
            meta = c.get("metadata", {}) or {}
            if isinstance(meta, dict) and "auto_rule" in meta:
                existing_keys.add(f"{c['student_id']}:{meta['auto_rule']}")
        
        # تطبيق القواعد
        for sid, stats in student_stats.items():
            for rule_key, rule in AUTO_CERT_RULES.items():
                if f"{sid}:{rule_key}" in existing_keys:
                    continue  # سبق منحه
                try:
                    if rule["check"](stats):
                        candidates.append({
                            "student_id": sid,
                            "student_name": stats["name"],
                            "rule_key": rule_key,
                            "rule_name": rule["name"],
                            "template": rule["template"],
                            "title": rule["title"],
                            "message": rule["message_template"].format(name=stats["name"]),
                            "stats": stats
                        })
                except Exception:
                    pass
        
        return {"candidates": candidates[:20], "total": len(candidates)}
    except Exception as e:
        print(f"[auto-check] {e}")
        return {"candidates": [], "error": str(e)[:200]}


@app.post("/api/supervisor/certificates/auto-issue")
async def auto_issue_certificates(
    candidates_json: str = Form(...),
    sup = Depends(get_current_supervisor)
):
    """🎁 إصدار جماعي لشهادات تلقائية"""
    import json as _json
    
    try:
        candidates = _json.loads(candidates_json)
        if not isinstance(candidates, list):
            raise ValueError("candidates يجب أن يكون list")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"صيغة JSON خاطئة: {str(e)[:100]}")
    
    sup_students = _get_supervisor_student_ids(sup["id"])
    sup_set = set(sup_students)
    
    issued = []
    failed = []
    
    for cand in candidates:
        try:
            sid = int(cand.get("student_id", 0))
            if sid not in sup_set:
                failed.append({"student_id": sid, "reason": "ليس من طلابك"})
                continue
            
            qr = _generate_qr_code()
            new_cert = {
                "student_id": sid,
                "supervisor_id": sup["id"],
                "supervisor_name": sup.get("full_name", "—"),
                "template_type": cand.get("template", "excellence"),
                "title": cand.get("title", "")[:200],
                "message": cand.get("message", "")[:1000],
                "student_name": cand.get("student_name", ""),
                "qr_code": qr,
                "metadata": {"auto_rule": cand.get("rule_key", ""), "auto_issued": True}
            }
            res = supabase.table("certificates").insert(new_cert).execute()
            if res.data:
                issued.append(res.data[0])
        except Exception as e:
            failed.append({"student_id": cand.get("student_id"), "reason": str(e)[:100]})
    
    return {
        "status": "success",
        "issued_count": len(issued),
        "failed_count": len(failed),
        "issued": issued,
        "failed": failed
    }


@app.post("/api/supervisor/certificates/bulk-issue")
async def bulk_issue_certificates(
    student_ids_json: str = Form(...),
    template_type: str = Form("excellence"),
    title: str = Form(...),
    message: str = Form(""),
    sup = Depends(get_current_supervisor)
):
    """🎓 إصدار جماعي لمجموعة طلاب (مثل نهاية الفصل)"""
    import json as _json
    
    try:
        student_ids = _json.loads(student_ids_json)
        student_ids = [int(x) for x in student_ids]
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"صيغة IDs خاطئة: {str(e)[:100]}")
    
    if template_type not in CERTIFICATE_TEMPLATES:
        raise HTTPException(status_code=400, detail=f"قالب غير صالح: {template_type}")
    
    sup_students = _get_supervisor_student_ids(sup["id"])
    sup_set = set(sup_students)
    
    issued = []
    failed = []
    
    # جلب أسماء الطلاب
    valid_ids = [i for i in student_ids if i in sup_set]
    if not valid_ids:
        return {"status": "error", "message": "لا يوجد طلاب صالحين"}
    
    try:
        names_res = supabase.table("students").select("id, full_name").in_("id", valid_ids).execute()
        names_map = {s["id"]: s.get("full_name", "—") for s in (names_res.data or [])}
    except Exception:
        names_map = {}
    
    for sid in valid_ids:
        try:
            student_name = names_map.get(sid, "—")
            # نُشخصن الرسالة
            personalized_msg = message.replace("{name}", student_name).replace("{student}", student_name)
            
            qr = _generate_qr_code()
            new_cert = {
                "student_id": sid,
                "supervisor_id": sup["id"],
                "supervisor_name": sup.get("full_name", "—"),
                "template_type": template_type,
                "title": title.strip()[:200],
                "message": personalized_msg[:1000],
                "student_name": student_name,
                "qr_code": qr,
                "metadata": {"bulk_issued": True}
            }
            res = supabase.table("certificates").insert(new_cert).execute()
            if res.data:
                issued.append(res.data[0])
        except Exception as e:
            failed.append({"student_id": sid, "reason": str(e)[:100]})
    
    return {
        "status": "success",
        "issued_count": len(issued),
        "failed_count": len(failed),
        "skipped_count": len(student_ids) - len(valid_ids),
        "issued": issued[:50],  # نُرجع أول 50 فقط
        "failed": failed[:10]
    }


@app.get("/wall-of-fame")
async def wall_of_fame_page(request: Request):
    """🏅 صفحة جدار الفخر العامة"""
    return templates.TemplateResponse(request=request, name="wall_of_fame.html")


@app.get("/api/student/certificates/share-link/{cert_id}")
async def get_share_link(cert_id: int, student = Depends(get_current_student)):
    """📤 الحصول على رابط مشاركة الشهادة"""
    try:
        res = supabase.table("certificates").select(
            "id, qr_code, title, student_name, supervisor_name, template_type"
        ).eq("id", cert_id).eq("student_id", student["id"]).limit(1).execute()
        
        if not res.data:
            raise HTTPException(status_code=404, detail="الشهادة غير موجودة")
        
        cert = res.data[0]
        tpl = CERTIFICATE_TEMPLATES.get(cert.get("template_type", "excellence"), {})
        
        # إنشاء رسائل جاهزة للمشاركة
        verify_url = f"/verify-certificate/{cert['qr_code']}"
        
        whatsapp_msg = f"🏆 حصلت على شهادة تقدير: {cert['title']}\n\n👨‍🏫 من المشرف: {cert['supervisor_name']}\n\n✅ للتحقق من صحة الشهادة:\n{verify_url}"
        
        return {
            "qr_code": cert["qr_code"],
            "verify_url": verify_url,
            "whatsapp_text": whatsapp_msg,
            "whatsapp_link": f"https://wa.me/?text={_url_encode(whatsapp_msg)}",
            "title": cert["title"],
            "template_icon": tpl.get("icon", "🏆")
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


def _url_encode(s: str) -> str:
    """تشفير URL آمن"""
    from urllib.parse import quote
    return quote(s, safe="")






# ════════════════════════════════════════════════════════════
# 💬 منتدى المعلمين
# ════════════════════════════════════════════════════════════

FORUM_CATEGORIES = {
    "announcements": {"name": "📢 إعلانات الإدارة", "icon": "📢", "color": "#e74c3c", "admin_only": True},
    "ideas": {"name": "💡 اقتراحات وأفكار", "icon": "💡", "color": "#f1c40f"},
    "teaching": {"name": "🎓 أساليب تدريس", "icon": "🎓", "color": "#3498db"},
    "resources": {"name": "📚 موارد ومصادر", "icon": "📚", "color": "#2ecc71"},
    "questions": {"name": "❓ أسئلة واستفسارات", "icon": "❓", "color": "#9b59b6"},
    "collaboration": {"name": "🤝 تعاون ومشاركة", "icon": "🤝", "color": "#e67e22"},
    "general": {"name": "☕ نقاش عام", "icon": "☕", "color": "#7f8c8d"}
}


async def _verify_teacher(teacher_id: int) -> dict:
    """التحقق من المعلم وجلب بياناته"""
    if not teacher_id or teacher_id <= 0:
        raise HTTPException(status_code=401, detail="معرف المعلم مطلوب")
    try:
        res = supabase.table("teachers").select("id, full_name, username").eq("id", teacher_id).limit(1).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="المعلم غير موجود")
        return res.data[0]
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ التحقق: {str(e)[:150]}")


@app.get("/api/teachers/forum/categories")
async def forum_get_categories():
    """📋 قائمة فئات المنتدى"""
    return {"categories": FORUM_CATEGORIES}


@app.get("/api/teachers/forum/topics")
async def forum_list_topics(
    category: str = "",
    search: str = "",
    sort: str = "recent",  # recent | popular | replies
    page: int = 1,
    page_size: int = 20
):
    """📋 قائمة المواضيع - مع فلاتر وبحث"""
    try:
        page = max(1, page)
        page_size = min(max(page_size, 5), 50)
        offset = (page - 1) * page_size
        
        query = supabase.table("teacher_forum_topics").select(
            "id, teacher_id, teacher_name, category, title, content, is_pinned, is_locked, views, likes_count, replies_count, last_reply_at, last_reply_by, created_at"
        )
        
        if category and category != "all":
            query = query.eq("category", category)
        
        if search:
            # بحث في العنوان والمحتوى
            search = search.strip()[:100]
            query = query.or_(f"title.ilike.%{search}%,content.ilike.%{search}%")
        
        # ترتيب
        if sort == "popular":
            query = query.order("is_pinned", desc=True).order("likes_count", desc=True)
        elif sort == "replies":
            query = query.order("is_pinned", desc=True).order("replies_count", desc=True)
        else:
            query = query.order("is_pinned", desc=True).order("created_at", desc=True)
        
        res = query.range(offset, offset + page_size - 1).execute()
        topics = res.data or []
        
        # إضافة معلومات الفئة
        for t in topics:
            cat_info = FORUM_CATEGORIES.get(t.get("category", "general"), FORUM_CATEGORIES["general"])
            t["category_name"] = cat_info["name"]
            t["category_icon"] = cat_info["icon"]
            t["category_color"] = cat_info["color"]
            # نُختصر المحتوى للمعاينة
            if t.get("content") and len(t["content"]) > 200:
                t["content_preview"] = t["content"][:200] + "..."
            else:
                t["content_preview"] = t.get("content", "")
        
        return {"topics": topics, "total": len(topics), "page": page}
    except Exception as e:
        print(f"[forum list] {e}")
        return {"topics": [], "total": 0, "error": str(e)[:200]}


@app.get("/api/teachers/forum/topics/{topic_id}")
async def forum_get_topic(topic_id: int, teacher_id: int = 0):
    """📖 موضوع كامل + ردوده"""
    try:
        # نجلب الموضوع
        t_res = supabase.table("teacher_forum_topics").select("*").eq("id", topic_id).limit(1).execute()
        if not t_res.data:
            raise HTTPException(status_code=404, detail="الموضوع غير موجود")
        
        topic = t_res.data[0]
        cat_info = FORUM_CATEGORIES.get(topic.get("category", "general"), FORUM_CATEGORIES["general"])
        topic["category_name"] = cat_info["name"]
        topic["category_icon"] = cat_info["icon"]
        topic["category_color"] = cat_info["color"]
        
        # زيادة المشاهدات (لو ليس المؤلف)
        if teacher_id and teacher_id != topic.get("teacher_id"):
            try:
                supabase.table("teacher_forum_topics").update({
                    "views": (topic.get("views", 0) or 0) + 1
                }).eq("id", topic_id).execute()
                topic["views"] = (topic.get("views", 0) or 0) + 1
            except Exception:
                pass
        
        # نجلب الردود
        r_res = supabase.table("teacher_forum_replies").select("*").eq(
            "topic_id", topic_id
        ).order("created_at", desc=False).limit(200).execute()
        replies = r_res.data or []
        
        # نتحقق من إعجابات المعلم الحالي
        liked_topic = False
        liked_replies = set()
        if teacher_id:
            try:
                likes_res = supabase.table("teacher_forum_likes").select("topic_id, reply_id").eq(
                    "teacher_id", teacher_id
                ).or_(f"topic_id.eq.{topic_id},reply_id.in.({','.join(str(r['id']) for r in replies) or '0'})").execute()
                for like in (likes_res.data or []):
                    if like.get("topic_id") == topic_id:
                        liked_topic = True
                    if like.get("reply_id"):
                        liked_replies.add(like["reply_id"])
            except Exception:
                pass
        
        topic["liked_by_me"] = liked_topic
        for r in replies:
            r["liked_by_me"] = r["id"] in liked_replies
        
        return {"topic": topic, "replies": replies}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/teachers/forum/topics")
async def forum_create_topic(
    teacher_id: int = Form(...),
    category: str = Form("general"),
    title: str = Form(...),
    content: str = Form(...)
):
    """✍️ إنشاء موضوع جديد"""
    teacher = await _verify_teacher(teacher_id)
    
    if category not in FORUM_CATEGORIES:
        raise HTTPException(status_code=400, detail=f"فئة غير صالحة: {category}")
    
    title = title.strip()[:200]
    content = content.strip()[:5000]
    
    if len(title) < 5:
        raise HTTPException(status_code=400, detail="العنوان قصير جداً (5+ أحرف)")
    if len(content) < 10:
        raise HTTPException(status_code=400, detail="المحتوى قصير جداً (10+ أحرف)")
    
    try:
        new_topic = {
            "teacher_id": teacher["id"],
            "teacher_name": teacher.get("full_name", "—"),
            "category": category,
            "title": title,
            "content": content,
        }
        res = supabase.table("teacher_forum_topics").insert(new_topic).execute()
        return {"status": "success", "topic": res.data[0] if res.data else None}
    except Exception as e:
        err = str(e)
        if "does not exist" in err.lower():
            raise HTTPException(status_code=500, detail="جدول teacher_forum_topics غير موجود!")
        raise HTTPException(status_code=500, detail=f"خطأ الإنشاء: {err[:200]}")


@app.post("/api/teachers/forum/topics/{topic_id}/reply")
async def forum_add_reply(
    topic_id: int,
    teacher_id: int = Form(...),
    content: str = Form(...),
    parent_reply_id: int = Form(0)
):
    """💬 إضافة رد على موضوع"""
    teacher = await _verify_teacher(teacher_id)
    
    content = content.strip()[:3000]
    if len(content) < 2:
        raise HTTPException(status_code=400, detail="الرد قصير جداً")
    
    try:
        # تحقق من الموضوع
        t_res = supabase.table("teacher_forum_topics").select("id, teacher_id, title, is_locked").eq("id", topic_id).limit(1).execute()
        if not t_res.data:
            raise HTTPException(status_code=404, detail="الموضوع غير موجود")
        
        topic = t_res.data[0]
        if topic.get("is_locked"):
            raise HTTPException(status_code=403, detail="الموضوع مغلق للردود")
        
        new_reply = {
            "topic_id": topic_id,
            "teacher_id": teacher["id"],
            "teacher_name": teacher.get("full_name", "—"),
            "content": content,
            "parent_reply_id": parent_reply_id if parent_reply_id > 0 else None
        }
        res = supabase.table("teacher_forum_replies").insert(new_reply).execute()
        reply = res.data[0] if res.data else None
        
        # تحديث الموضوع (عدد الردود + آخر رد)
        try:
            from datetime import datetime, timezone
            supabase.table("teacher_forum_topics").update({
                "replies_count": _count_topic_replies(topic_id),
                "last_reply_at": datetime.now(timezone.utc).isoformat(),
                "last_reply_by": teacher.get("full_name", "—")
            }).eq("id", topic_id).execute()
        except Exception:
            pass
        
        # إشعار صاحب الموضوع (إن لم يكن هو نفسه)
        if topic.get("teacher_id") and topic["teacher_id"] != teacher["id"]:
            try:
                supabase.table("teacher_forum_notifications").insert({
                    "teacher_id": topic["teacher_id"],
                    "type": "reply",
                    "topic_id": topic_id,
                    "reply_id": reply.get("id") if reply else None,
                    "from_teacher_id": teacher["id"],
                    "from_teacher_name": teacher.get("full_name", "—"),
                    "message": f"رد {teacher.get('full_name', '—')} على موضوعك: {topic.get('title', '')[:50]}"
                }).execute()
            except Exception:
                pass
        
        return {"status": "success", "reply": reply}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ الرد: {str(e)[:200]}")


def _count_topic_replies(topic_id: int) -> int:
    try:
        res = supabase.table("teacher_forum_replies").select("id", count="exact").eq("topic_id", topic_id).limit(1).execute()
        return res.count or 0
    except Exception:
        return 0


@app.post("/api/teachers/forum/like")
async def forum_toggle_like(
    teacher_id: int = Form(...),
    topic_id: int = Form(0),
    reply_id: int = Form(0)
):
    """❤️ إعجاب/إلغاء على موضوع أو رد"""
    await _verify_teacher(teacher_id)
    
    if topic_id <= 0 and reply_id <= 0:
        raise HTTPException(status_code=400, detail="يجب تحديد موضوع أو رد")
    
    try:
        # نتحقق إذا كان معجباً مسبقاً
        q = supabase.table("teacher_forum_likes").select("id").eq("teacher_id", teacher_id)
        if topic_id > 0:
            q = q.eq("topic_id", topic_id).is_("reply_id", "null")
        else:
            q = q.eq("reply_id", reply_id).is_("topic_id", "null")
        
        existing = q.limit(1).execute()
        
        if existing.data:
            # إلغاء الإعجاب
            supabase.table("teacher_forum_likes").delete().eq("id", existing.data[0]["id"]).execute()
            action = "unliked"
        else:
            # إضافة إعجاب
            new_like = {"teacher_id": teacher_id}
            if topic_id > 0:
                new_like["topic_id"] = topic_id
            else:
                new_like["reply_id"] = reply_id
            supabase.table("teacher_forum_likes").insert(new_like).execute()
            action = "liked"
        
        # تحديث العداد
        if topic_id > 0:
            count_res = supabase.table("teacher_forum_likes").select("id", count="exact").eq(
                "topic_id", topic_id
            ).is_("reply_id", "null").limit(1).execute()
            new_count = count_res.count or 0
            supabase.table("teacher_forum_topics").update({"likes_count": new_count}).eq("id", topic_id).execute()
        else:
            count_res = supabase.table("teacher_forum_likes").select("id", count="exact").eq(
                "reply_id", reply_id
            ).is_("topic_id", "null").limit(1).execute()
            new_count = count_res.count or 0
            supabase.table("teacher_forum_replies").update({"likes_count": new_count}).eq("id", reply_id).execute()
        
        return {"status": "success", "action": action, "likes_count": new_count}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/teachers/forum/topics/{topic_id}")
async def forum_delete_topic(topic_id: int, teacher_id: int):
    """🗑️ حذف موضوع (المؤلف فقط)"""
    teacher = await _verify_teacher(teacher_id)
    
    try:
        # نتحقق من الملكية
        t_res = supabase.table("teacher_forum_topics").select("teacher_id").eq("id", topic_id).limit(1).execute()
        if not t_res.data:
            raise HTTPException(status_code=404, detail="الموضوع غير موجود")
        if t_res.data[0]["teacher_id"] != teacher["id"]:
            raise HTTPException(status_code=403, detail="لا يمكنك حذف موضوع لم تكتبه")
        
        # حذف الردود أولاً
        supabase.table("teacher_forum_replies").delete().eq("topic_id", topic_id).execute()
        # حذف الإعجابات
        supabase.table("teacher_forum_likes").delete().eq("topic_id", topic_id).execute()
        # حذف الموضوع
        supabase.table("teacher_forum_topics").delete().eq("id", topic_id).execute()
        
        return {"status": "success", "message": "تم الحذف"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/teachers/forum/replies/{reply_id}")
async def forum_delete_reply(reply_id: int, teacher_id: int):
    """🗑️ حذف رد (الكاتب فقط)"""
    teacher = await _verify_teacher(teacher_id)
    
    try:
        r_res = supabase.table("teacher_forum_replies").select("teacher_id, topic_id").eq("id", reply_id).limit(1).execute()
        if not r_res.data:
            raise HTTPException(status_code=404, detail="الرد غير موجود")
        if r_res.data[0]["teacher_id"] != teacher["id"]:
            raise HTTPException(status_code=403, detail="لا يمكنك حذف رد ليس لك")
        
        topic_id = r_res.data[0].get("topic_id")
        supabase.table("teacher_forum_likes").delete().eq("reply_id", reply_id).execute()
        supabase.table("teacher_forum_replies").delete().eq("id", reply_id).execute()
        
        # تحديث عداد الردود
        if topic_id:
            try:
                supabase.table("teacher_forum_topics").update({
                    "replies_count": _count_topic_replies(topic_id)
                }).eq("id", topic_id).execute()
            except Exception:
                pass
        
        return {"status": "success"}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teachers/forum/leaderboard")
async def forum_leaderboard():
    """🏆 الأكثر نشاطاً في المنتدى"""
    try:
        # نُحسب نشاط كل معلم
        topics_res = supabase.table("teacher_forum_topics").select(
            "teacher_id, teacher_name, likes_count, replies_count"
        ).limit(500).execute()
        replies_res = supabase.table("teacher_forum_replies").select(
            "teacher_id, teacher_name, likes_count"
        ).limit(1000).execute()
        
        stats = {}
        for t in (topics_res.data or []):
            tid = t.get("teacher_id")
            if not tid:
                continue
            if tid not in stats:
                stats[tid] = {"teacher_id": tid, "teacher_name": t.get("teacher_name", "—"), "topics": 0, "replies": 0, "likes_received": 0, "score": 0}
            stats[tid]["topics"] += 1
            stats[tid]["likes_received"] += t.get("likes_count", 0) or 0
        
        for r in (replies_res.data or []):
            tid = r.get("teacher_id")
            if not tid:
                continue
            if tid not in stats:
                stats[tid] = {"teacher_id": tid, "teacher_name": r.get("teacher_name", "—"), "topics": 0, "replies": 0, "likes_received": 0, "score": 0}
            stats[tid]["replies"] += 1
            stats[tid]["likes_received"] += r.get("likes_count", 0) or 0
        
        # حساب النقاط: 10 لكل موضوع، 3 لكل رد، 2 لكل إعجاب مُستقبل
        for s in stats.values():
            s["score"] = s["topics"] * 10 + s["replies"] * 3 + s["likes_received"] * 2
        
        # ترتيب
        ranked = sorted(stats.values(), key=lambda x: x["score"], reverse=True)[:20]
        
        # رتب
        for i, t in enumerate(ranked):
            t["rank"] = i + 1
            if i == 0: t["badge"] = "👑"
            elif i == 1: t["badge"] = "🥈"
            elif i == 2: t["badge"] = "🥉"
            else: t["badge"] = ""
        
        return {"leaderboard": ranked, "total": len(stats)}
    except Exception as e:
        print(f"[forum leaderboard] {e}")
        return {"leaderboard": [], "error": str(e)[:200]}


@app.get("/api/teachers/forum/notifications")
async def forum_get_notifications(teacher_id: int, unread_only: bool = False):
    """🔔 إشعارات المعلم"""
    await _verify_teacher(teacher_id)
    try:
        q = supabase.table("teacher_forum_notifications").select("*").eq("teacher_id", teacher_id)
        if unread_only:
            q = q.eq("is_read", False)
        res = q.order("created_at", desc=True).limit(50).execute()
        notifs = res.data or []
        unread_count = sum(1 for n in notifs if not n.get("is_read"))
        return {"notifications": notifs, "unread_count": unread_count, "total": len(notifs)}
    except Exception as e:
        return {"notifications": [], "unread_count": 0, "error": str(e)[:200]}


@app.post("/api/teachers/forum/notifications/mark-read")
async def forum_mark_read(teacher_id: int = Form(...), notification_id: int = Form(0)):
    """✅ تعليم إشعار كمقروء"""
    await _verify_teacher(teacher_id)
    try:
        q = supabase.table("teacher_forum_notifications").update({"is_read": True}).eq("teacher_id", teacher_id)
        if notification_id > 0:
            q = q.eq("id", notification_id)
        q.execute()
        return {"status": "success"}
    except Exception as e:
        return {"status": "error", "message": str(e)[:200]}


@app.get("/api/teachers/forum/my-stats")
async def forum_my_stats(teacher_id: int):
    """📊 إحصاءات المعلم في المنتدى"""
    teacher = await _verify_teacher(teacher_id)
    try:
        topics_res = supabase.table("teacher_forum_topics").select("id, likes_count, replies_count, views").eq("teacher_id", teacher_id).execute()
        replies_res = supabase.table("teacher_forum_replies").select("id, likes_count").eq("teacher_id", teacher_id).execute()
        
        topics = topics_res.data or []
        replies = replies_res.data or []
        
        total_views = sum(t.get("views", 0) or 0 for t in topics)
        total_likes_received = sum(t.get("likes_count", 0) or 0 for t in topics) + sum(r.get("likes_count", 0) or 0 for r in replies)
        score = len(topics) * 10 + len(replies) * 3 + total_likes_received * 2
        
        return {
            "teacher_name": teacher.get("full_name", "—"),
            "topics_count": len(topics),
            "replies_count": len(replies),
            "likes_received": total_likes_received,
            "total_views": total_views,
            "score": score
        }
    except Exception as e:
        return {"error": str(e)[:200]}






# ════════════════════════════════════════════════════════════
# 📎 المرفقات + 🔝 التثبيت + 📊 إدارة الأدمن + 🔔 الاشتراكات
# ════════════════════════════════════════════════════════════

@app.post("/api/teachers/forum/upload-attachment")
async def forum_upload_attachment(
    teacher_id: int = Form(...),
    topic_id: int = Form(0),
    reply_id: int = Form(0),
    file: UploadFile = File(...)
):
    """📎 رفع مرفق لموضوع أو رد"""
    teacher = await _verify_teacher(teacher_id)
    
    if topic_id <= 0 and reply_id <= 0:
        raise HTTPException(status_code=400, detail="يجب تحديد موضوع أو رد")
    
    # تحقق من نوع وحجم الملف
    ALLOWED_TYPES = {
        "image/jpeg", "image/jpg", "image/png", "image/gif", "image/webp",
        "application/pdf",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # docx
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # xlsx
        "application/msword", "application/vnd.ms-excel"
    }
    
    if file.content_type not in ALLOWED_TYPES:
        raise HTTPException(status_code=400, detail=f"نوع الملف غير مسموح: {file.content_type}")
    
    # قراءة الملف وفحص الحجم (5MB max)
    contents = await file.read()
    size_kb = len(contents) // 1024
    if size_kb > 5120:  # 5 MB
        raise HTTPException(status_code=400, detail=f"الملف كبير جداً ({size_kb} KB). الحد الأقصى 5 MB")
    
    try:
        # رفع الملف لـ Supabase Storage
        import secrets
        ext = file.filename.split(".")[-1] if "." in file.filename else "bin"
        safe_name = f"forum/{secrets.token_urlsafe(12)}.{ext}"
        
        try:
            upload_res = supabase.storage.from_("uploads").upload(
                safe_name, contents, {"content-type": file.content_type or "application/octet-stream"}
            )
            file_url = supabase.storage.from_("uploads").get_public_url(safe_name)
        except Exception as storage_err:
            # محاولة بـ bucket آخر
            try:
                upload_res = supabase.storage.from_("public").upload(
                    safe_name, contents, {"content-type": file.content_type or "application/octet-stream"}
                )
                file_url = supabase.storage.from_("public").get_public_url(safe_name)
            except Exception:
                raise HTTPException(status_code=500, detail=f"فشل رفع الملف: {str(storage_err)[:150]}")
        
        # نُسجّل في جدول المرفقات
        attach_data = {
            "teacher_id": teacher["id"],
            "file_url": file_url,
            "file_name": file.filename[:200] if file.filename else "ملف",
            "file_type": file.content_type or "unknown",
            "file_size_kb": size_kb
        }
        if topic_id > 0:
            attach_data["topic_id"] = topic_id
        if reply_id > 0:
            attach_data["reply_id"] = reply_id
        
        res = supabase.table("teacher_forum_attachments").insert(attach_data).execute()
        
        return {
            "status": "success",
            "attachment": res.data[0] if res.data else None,
            "file_url": file_url
        }
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teachers/forum/attachments")
async def forum_get_attachments(topic_id: int = 0, reply_id: int = 0):
    """📋 جلب مرفقات موضوع أو رد"""
    if topic_id <= 0 and reply_id <= 0:
        return {"attachments": []}
    
    try:
        q = supabase.table("teacher_forum_attachments").select("*")
        if topic_id > 0:
            q = q.eq("topic_id", topic_id)
        else:
            q = q.eq("reply_id", reply_id)
        res = q.order("created_at", desc=False).execute()
        return {"attachments": res.data or []}
    except Exception as e:
        return {"attachments": [], "error": str(e)[:200]}


# ──────── 🔝 التثبيت والقفل (الأدمن) ────────

@app.post("/api/admin/forum/pin-topic")
async def admin_pin_topic(
    topic_id: int = Form(...),
    pin: bool = Form(True),
    admin = Depends(get_current_admin)
):
    """📌 تثبيت/إلغاء تثبيت موضوع (الأدمن فقط)"""
    try:
        supabase.table("teacher_forum_topics").update({"is_pinned": pin}).eq("id", topic_id).execute()
        return {"status": "success", "is_pinned": pin}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/admin/forum/lock-topic")
async def admin_lock_topic(
    topic_id: int = Form(...),
    lock: bool = Form(True),
    admin = Depends(get_current_admin)
):
    """🔒 قفل/فتح موضوع (الأدمن فقط)"""
    try:
        supabase.table("teacher_forum_topics").update({"is_locked": lock}).eq("id", topic_id).execute()
        return {"status": "success", "is_locked": lock}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.delete("/api/admin/forum/topics/{topic_id}")
async def admin_delete_topic(topic_id: int, admin = Depends(get_current_admin)):
    """🗑️ حذف موضوع (الأدمن - بدون قيود)"""
    try:
        supabase.table("teacher_forum_replies").delete().eq("topic_id", topic_id).execute()
        supabase.table("teacher_forum_likes").delete().eq("topic_id", topic_id).execute()
        supabase.table("teacher_forum_attachments").delete().eq("topic_id", topic_id).execute()
        supabase.table("teacher_forum_subscriptions").delete().eq("topic_id", topic_id).execute()
        supabase.table("teacher_forum_topics").delete().eq("id", topic_id).execute()
        return {"status": "success", "message": "تم الحذف بواسطة الأدمن"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


# ──────── 📊 لوحة تحكم الأدمن للمنتدى ────────

@app.get("/api/admin/forum/stats")
async def admin_forum_stats(admin = Depends(get_current_admin)):
    """📊 إحصاءات شاملة عن المنتدى"""
    try:
        # العدد الإجمالي
        topics_res = supabase.table("teacher_forum_topics").select("id", count="exact").limit(1).execute()
        replies_res = supabase.table("teacher_forum_replies").select("id", count="exact").limit(1).execute()
        likes_res = supabase.table("teacher_forum_likes").select("id", count="exact").limit(1).execute()
        
        total_topics = topics_res.count or 0
        total_replies = replies_res.count or 0
        total_likes = likes_res.count or 0
        
        # عدد المعلمين النشطين
        active_res = supabase.table("teacher_forum_topics").select("teacher_id").limit(1000).execute()
        active_in_topics = set(t.get("teacher_id") for t in (active_res.data or []))
        replies_active_res = supabase.table("teacher_forum_replies").select("teacher_id").limit(1000).execute()
        active_in_replies = set(r.get("teacher_id") for r in (replies_active_res.data or []))
        active_teachers = len(active_in_topics | active_in_replies)
        
        # هذا الأسبوع
        from datetime import datetime, timezone, timedelta
        week_ago = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
        
        week_topics_res = supabase.table("teacher_forum_topics").select("id", count="exact").gte("created_at", week_ago).limit(1).execute()
        week_replies_res = supabase.table("teacher_forum_replies").select("id", count="exact").gte("created_at", week_ago).limit(1).execute()
        
        # توزيع الفئات
        cat_res = supabase.table("teacher_forum_topics").select("category").limit(1000).execute()
        cat_counts = {}
        for t in (cat_res.data or []):
            c = t.get("category", "general")
            cat_counts[c] = cat_counts.get(c, 0) + 1
        
        # أكثر المواضيع تفاعلاً
        top_topics_res = supabase.table("teacher_forum_topics").select(
            "id, title, teacher_name, likes_count, replies_count, views, created_at"
        ).order("likes_count", desc=True).limit(10).execute()
        
        # أكثر المعلمين نشاطاً
        sup_stats = {}
        for t in (active_res.data or []):
            tid = t.get("teacher_id")
            if tid:
                sup_stats[tid] = sup_stats.get(tid, 0) + 1
        top_teachers = sorted(sup_stats.items(), key=lambda x: x[1], reverse=True)[:5]
        
        # نمو أسبوعي (آخر 7 أيام)
        daily_growth = []
        for i in range(7):
            day_start = (datetime.now(timezone.utc) - timedelta(days=6-i)).replace(hour=0, minute=0, second=0, microsecond=0)
            day_end = day_start + timedelta(days=1)
            day_topics_res = supabase.table("teacher_forum_topics").select("id", count="exact").gte(
                "created_at", day_start.isoformat()
            ).lt("created_at", day_end.isoformat()).limit(1).execute()
            daily_growth.append({
                "date": day_start.strftime("%m-%d"),
                "topics": day_topics_res.count or 0
            })
        
        return {
            "totals": {
                "topics": total_topics,
                "replies": total_replies,
                "likes": total_likes,
                "active_teachers": active_teachers
            },
            "this_week": {
                "topics": week_topics_res.count or 0,
                "replies": week_replies_res.count or 0
            },
            "categories": cat_counts,
            "top_topics": top_topics_res.data or [],
            "daily_growth": daily_growth
        }
    except Exception as e:
        return {"error": str(e)[:200]}


@app.get("/api/admin/forum/topics")
async def admin_list_topics(
    category: str = "",
    page: int = 1,
    page_size: int = 30,
    admin = Depends(get_current_admin)
):
    """📋 قائمة المواضيع للأدمن (مع تفاصيل إدارية)"""
    try:
        page = max(1, page)
        page_size = min(max(page_size, 10), 100)
        offset = (page - 1) * page_size
        
        q = supabase.table("teacher_forum_topics").select("*")
        if category:
            q = q.eq("category", category)
        
        res = q.order("created_at", desc=True).range(offset, offset + page_size - 1).execute()
        return {"topics": res.data or [], "page": page}
    except Exception as e:
        return {"topics": [], "error": str(e)[:200]}


# ──────── 🔔 الاشتراكات والكتم ────────

@app.post("/api/teachers/forum/subscribe")
async def forum_subscribe(
    teacher_id: int = Form(...),
    topic_id: int = Form(...),
    action: str = Form("subscribe")  # subscribe | unsubscribe | mute | unmute
):
    """🔔 اشتراك/كتم/إلغاء موضوع"""
    teacher = await _verify_teacher(teacher_id)
    
    try:
        # نتحقق من الموضوع
        t_res = supabase.table("teacher_forum_topics").select("id").eq("id", topic_id).limit(1).execute()
        if not t_res.data:
            raise HTTPException(status_code=404, detail="الموضوع غير موجود")
        
        # نتحقق من وجود اشتراك حالي
        existing = supabase.table("teacher_forum_subscriptions").select("id, is_muted").eq(
            "teacher_id", teacher_id
        ).eq("topic_id", topic_id).limit(1).execute()
        
        if action == "subscribe":
            if existing.data:
                supabase.table("teacher_forum_subscriptions").update({"is_muted": False}).eq("id", existing.data[0]["id"]).execute()
            else:
                supabase.table("teacher_forum_subscriptions").insert({
                    "teacher_id": teacher_id,
                    "topic_id": topic_id,
                    "is_muted": False
                }).execute()
            return {"status": "success", "subscribed": True, "muted": False}
        
        elif action == "unsubscribe":
            if existing.data:
                supabase.table("teacher_forum_subscriptions").delete().eq("id", existing.data[0]["id"]).execute()
            return {"status": "success", "subscribed": False, "muted": False}
        
        elif action == "mute":
            if existing.data:
                supabase.table("teacher_forum_subscriptions").update({"is_muted": True}).eq("id", existing.data[0]["id"]).execute()
            else:
                supabase.table("teacher_forum_subscriptions").insert({
                    "teacher_id": teacher_id,
                    "topic_id": topic_id,
                    "is_muted": True
                }).execute()
            return {"status": "success", "subscribed": True, "muted": True}
        
        elif action == "unmute":
            if existing.data:
                supabase.table("teacher_forum_subscriptions").update({"is_muted": False}).eq("id", existing.data[0]["id"]).execute()
            return {"status": "success", "subscribed": True, "muted": False}
        
        else:
            raise HTTPException(status_code=400, detail=f"إجراء غير صالح: {action}")
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.get("/api/teachers/forum/my-subscription/{topic_id}")
async def forum_my_subscription(topic_id: int, teacher_id: int):
    """🔍 حالة اشتراكي في موضوع"""
    try:
        res = supabase.table("teacher_forum_subscriptions").select("is_muted").eq(
            "teacher_id", teacher_id
        ).eq("topic_id", topic_id).limit(1).execute()
        
        if not res.data:
            return {"subscribed": False, "muted": False}
        
        return {"subscribed": True, "muted": res.data[0].get("is_muted", False)}
    except Exception:
        return {"subscribed": False, "muted": False}


@app.get("/api/teachers/forum/my-subscriptions")
async def forum_my_subscriptions(teacher_id: int):
    """📋 قائمة اشتراكاتي"""
    await _verify_teacher(teacher_id)
    try:
        subs_res = supabase.table("teacher_forum_subscriptions").select("topic_id, is_muted, created_at").eq(
            "teacher_id", teacher_id
        ).order("created_at", desc=True).execute()
        
        subs = subs_res.data or []
        if not subs:
            return {"subscriptions": []}
        
        topic_ids = [s["topic_id"] for s in subs]
        topics_res = supabase.table("teacher_forum_topics").select(
            "id, title, category, replies_count, likes_count, last_reply_at, teacher_name"
        ).in_("id", topic_ids).execute()
        
        topics_map = {t["id"]: t for t in (topics_res.data or [])}
        result = []
        for s in subs:
            t = topics_map.get(s["topic_id"])
            if t:
                t["is_muted"] = s["is_muted"]
                t["subscribed_at"] = s["created_at"]
                cat_info = FORUM_CATEGORIES.get(t.get("category", "general"), FORUM_CATEGORIES["general"])
                t["category_icon"] = cat_info["icon"]
                t["category_color"] = cat_info["color"]
                result.append(t)
        
        return {"subscriptions": result, "total": len(result)}
    except Exception as e:
        return {"subscriptions": [], "error": str(e)[:200]}




@app.get("/api/supervisor/student-360/{student_id}")
async def supervisor_student_360(student_id: int, sup = Depends(get_current_supervisor)):
    """📊 بطاقة الطالب التفصيلية - Student 360 view"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    if student_id not in student_ids:
        raise HTTPException(status_code=403, detail="هذا الطالب ليس من طلابك")
    
    # 1. بيانات الطالب الأساسية
    stu_res = supabase.table("students").select("*").eq("id", student_id).execute()
    if not stu_res.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    student = stu_res.data[0]
    student.pop("password", None)
    
    # 2. آخر 100 محاولة
    chal_res = supabase.table("student_challenges").select("*").eq("student_id", student_id).order("created_at", desc=True).limit(100).execute()
    challenges = chal_res.data or []
    
    # 3. حساب الإحصائيات
    from datetime import datetime, timezone, timedelta
    now = datetime.now(timezone.utc)
    week_ago = now - timedelta(days=7)
    month_ago = now - timedelta(days=30)
    
    total = len(challenges)
    correct = sum(1 for c in challenges if c.get("is_correct"))
    accuracy = round(correct / total * 100, 1) if total else 0
    
    week_chal = [c for c in challenges if c.get("created_at") and c["created_at"] > week_ago.isoformat()]
    week_correct = sum(1 for c in week_chal if c.get("is_correct"))
    week_acc = round(week_correct / len(week_chal) * 100, 1) if week_chal else 0
    
    month_chal = [c for c in challenges if c.get("created_at") and c["created_at"] > month_ago.isoformat()]
    
    # 4. تحليل نقاط القوة والضعف حسب الدرس
    lessons_stats = {}
    for c in challenges:
        lesson = c.get("lesson") or c.get("topic") or "غير محدد"
        if lesson not in lessons_stats:
            lessons_stats[lesson] = {"total": 0, "correct": 0, "wrong": 0}
        lessons_stats[lesson]["total"] += 1
        if c.get("is_correct"):
            lessons_stats[lesson]["correct"] += 1
        else:
            lessons_stats[lesson]["wrong"] += 1
    
    # حساب النسب لكل درس
    strengths = []  # 85%+ مع 3+ محاولات
    weaknesses = []  # <60% مع 3+ محاولات
    for lesson, st in lessons_stats.items():
        if st["total"] < 3:
            continue
        pct = round(st["correct"] / st["total"] * 100, 1)
        item = {"lesson": lesson, "accuracy": pct, "attempts": st["total"], "correct": st["correct"], "wrong": st["wrong"]}
        if pct >= 85:
            strengths.append(item)
        elif pct < 60:
            weaknesses.append(item)
    
    strengths.sort(key=lambda x: -x["accuracy"])
    weaknesses.sort(key=lambda x: x["accuracy"])
    
    # 5. Heatmap - أيام الأسبوع الأكثر نشاطاً
    day_counts = {0: 0, 1: 0, 2: 0, 3: 0, 4: 0, 5: 0, 6: 0}  # Sat-Fri
    hour_counts = {h: 0 for h in range(24)}
    for c in month_chal:
        try:
            d = datetime.fromisoformat(c["created_at"].replace("Z", "+00:00"))
            day_counts[d.weekday()] += 1
            hour_counts[d.hour] += 1
        except Exception:
            pass
    
    # 6. تحليل سلوكي بسيط
    behavior_insights = []
    
    # سرعة الإجابة (لو متاحة)
    fast_wrong = sum(1 for c in challenges if c.get("time_taken", 999) < 10 and not c.get("is_correct"))
    if fast_wrong >= 5:
        behavior_insights.append({
            "icon": "⚡",
            "type": "warning",
            "title": "يميل للإجابة بسرعة",
            "desc": f"{fast_wrong} إجابة خاطئة في أقل من 10 ثوانٍ — يُنصح بتشجيعه على التفكير قبل الإجابة"
        })
    
    # محاولات فاشلة متتالية
    consec_fail = 0
    max_consec_fail = 0
    for c in challenges:
        if not c.get("is_correct"):
            consec_fail += 1
            max_consec_fail = max(max_consec_fail, consec_fail)
        else:
            consec_fail = 0
    
    if max_consec_fail >= 5:
        behavior_insights.append({
            "icon": "🔁",
            "type": "alert",
            "title": "إحباط محتمل",
            "desc": f"{max_consec_fail} محاولات فاشلة متتالية — قد يحتاج تشجيعاً ومراجعة المفاهيم الأساسية"
        })
    
    # نشاط متواصل
    if week_acc > accuracy + 10:
        behavior_insights.append({
            "icon": "📈",
            "type": "success",
            "title": "تحسن ملحوظ",
            "desc": f"دقة هذا الأسبوع {week_acc}% مقابل المتوسط {accuracy}% — قدّم له تقديراً!"
        })
    elif week_acc < accuracy - 10 and len(week_chal) >= 3:
        behavior_insights.append({
            "icon": "📉",
            "type": "warning",
            "title": "تراجع في الأداء",
            "desc": f"دقة هذا الأسبوع {week_acc}% أقل من المتوسط {accuracy}% — يحتاج متابعة"
        })
    
    # غياب
    if student.get("last_active"):
        try:
            la = datetime.fromisoformat(student["last_active"].replace("Z", "+00:00"))
            days_absent = (now - la).days
            if days_absent >= 5:
                behavior_insights.append({
                    "icon": "🕰️",
                    "type": "alert",
                    "title": f"غياب {days_absent} يوم",
                    "desc": "أرسل له رسالة تحفيز للعودة"
                })
        except Exception:
            pass
    
    # نشاط ممتاز
    if total >= 50 and accuracy >= 80:
        behavior_insights.append({
            "icon": "⭐",
            "type": "success",
            "title": "أداء متميز",
            "desc": f"{total} محاولة بدقة {accuracy}% — مرشح لنادي النخبة!"
        })
    
    # 7. سجل التدخلات السابقة
    msgs_res = supabase.table("supervisor_messages").select("*")\
        .eq("supervisor_id", sup["id"]).eq("student_id", student_id)\
        .order("created_at", desc=True).limit(10).execute()
    interventions = msgs_res.data or []
    
    return {
        "student": student,
        "stats": {
            "total_attempts": total,
            "correct": correct,
            "accuracy": accuracy,
            "weekly_attempts": len(week_chal),
            "weekly_accuracy": week_acc,
            "monthly_attempts": len(month_chal),
            "xp": student.get("total_xp", 0) or 0,
            "total_points": student.get("total_points", 0) or 0,
            "level": max(1, (student.get("total_xp", 0) or 0) // 100)
        },
        "strengths": strengths[:5],
        "weaknesses": weaknesses[:5],
        "heatmap": {
            "by_day": day_counts,  # 0=Mon, 6=Sun
            "by_hour": hour_counts
        },
        "behavior_insights": behavior_insights,
        "interventions": interventions,
        "recent_challenges": challenges[:15]
    }


@app.get("/api/supervisor/recommendations/{student_id}")
async def supervisor_get_recommendations(student_id: int, sup = Depends(get_current_supervisor)):
    """🎯 محرك التوصيات الذكي - يقترح إجراءات بناءً على بيانات الطالب"""
    student_ids = _get_supervisor_student_ids(sup["id"])
    if student_id not in student_ids:
        raise HTTPException(status_code=403, detail="هذا الطالب ليس من طلابك")
    
    # نجلب البيانات
    stu_res = supabase.table("students").select("*").eq("id", student_id).execute()
    if not stu_res.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    student = stu_res.data[0]
    
    chal_res = supabase.table("student_challenges").select("*").eq("student_id", student_id).order("created_at", desc=True).limit(50).execute()
    challenges = chal_res.data or []
    
    from datetime import datetime, timezone, timedelta
    now = datetime.now(timezone.utc)
    
    recommendations = []
    
    # القاعدة 1: غياب 3+ أيام
    if student.get("last_active"):
        try:
            la = datetime.fromisoformat(student["last_active"].replace("Z", "+00:00"))
            days_absent = (now - la).days
            if days_absent >= 3:
                recommendations.append({
                    "id": "absence",
                    "priority": "high" if days_absent >= 7 else "medium",
                    "icon": "🕰️",
                    "title": f"غياب لمدة {days_absent} يوم",
                    "action": "أرسل رسالة تحفيزية",
                    "template": f"مرحباً {student.get('full_name','')}! 👋\n\nنفتقدك في إمبراطورية الرياضيات! 👑\nلقد مرّت {days_absent} أيام منذ آخر زيارة لك. تعال وأكمل رحلتك نحو النجاح!\n\nبانتظارك 💪",
                    "type": "note"
                })
        except Exception:
            pass
    
    # القاعدة 2: تراجع دقة في درس معين
    lessons_stats = {}
    for c in challenges:
        lesson = c.get("lesson") or c.get("topic")
        if not lesson:
            continue
        if lesson not in lessons_stats:
            lessons_stats[lesson] = {"total": 0, "correct": 0}
        lessons_stats[lesson]["total"] += 1
        if c.get("is_correct"):
            lessons_stats[lesson]["correct"] += 1
    
    for lesson, st in lessons_stats.items():
        if st["total"] >= 5:
            acc = st["correct"] / st["total"] * 100
            if acc < 50:
                recommendations.append({
                    "id": f"weak_lesson_{lesson}",
                    "priority": "high",
                    "icon": "📉",
                    "title": f"ضعف في درس: {lesson} (دقة {round(acc)}%)",
                    "action": "اقترح مراجعة الدرس + تمارين إضافية",
                    "template": f"عزيزي {student.get('full_name','')}، 📚\n\nلاحظت أن درس \"{lesson}\" يحتاج لمزيد من التركيز.\nأقترح عليك:\n1. مراجعة شرح الدرس من جديد\n2. حل 5 تمارين إضافية\n3. سؤالي عن أي نقطة غير واضحة\n\nأنا واثق من قدرتك على التحسن! 💪",
                    "type": "note"
                })
                break  # نأخذ واحداً فقط
    
    # القاعدة 3: سرعة + أخطاء
    fast_wrong = sum(1 for c in challenges if c.get("time_taken", 999) < 10 and not c.get("is_correct"))
    if fast_wrong >= 5:
        recommendations.append({
            "id": "fast_wrong",
            "priority": "medium",
            "icon": "⚡",
            "title": "إجابات سريعة بأخطاء كثيرة",
            "action": "نصيحة بتأنّي القراءة",
            "template": f"نصيحة ذهبية لك {student.get('full_name','')} 🌟\n\nلاحظت أنك تُجيب بسرعة كبيرة. جرّب:\n• اقرأ السؤال مرتين قبل الإجابة\n• فكّر 5 ثوانٍ على الأقل\n• راجع إجابتك قبل التأكيد\n\nالأناة تأتي بنتائج رائعة! 🎯",
            "type": "note"
        })
    
    # القاعدة 4: محاولات فاشلة متتالية
    consec_fail = 0
    for c in challenges:
        if not c.get("is_correct"):
            consec_fail += 1
        else:
            break
    if consec_fail >= 5:
        recommendations.append({
            "id": "consecutive_fail",
            "priority": "high",
            "icon": "🔁",
            "title": f"{consec_fail} محاولات فاشلة متتالية",
            "action": "تشجيع ومراجعة المفاهيم",
            "template": f"مرحباً {student.get('full_name','')} 🤗\n\nلا تيأس! كل بطل واجه صعوبات.\nخذ استراحة قصيرة، ثم:\n• راجع المفاهيم الأساسية للدرس\n• ابدأ بأسئلة أسهل لتستعيد ثقتك\n• تذكّر: الخطأ خطوة نحو التعلم!\n\nأنا معك في هذه الرحلة 💎",
            "type": "alert"
        })
    
    # القاعدة 5: أداء ممتاز (5+ صحيحات متتالية)
    consec_correct = 0
    for c in challenges:
        if c.get("is_correct"):
            consec_correct += 1
        else:
            break
    if consec_correct >= 5:
        recommendations.append({
            "id": "excellent_streak",
            "priority": "celebrate",
            "icon": "🏆",
            "title": f"{consec_correct} إجابات صحيحة متتالية!",
            "action": "تقدير وشهادة",
            "template": f"🎉 مبارك {student.get('full_name','')}! 🎉\n\nأداء مذهل! {consec_correct} إجابات صحيحة متتالية.\n\nأنت تستحق التقدير والاحترام. استمر بهذا التميّز!\n\n👑 من فخر إمبراطورية الرياضيات",
            "type": "note"
        })
    
    # القاعدة 6: XP عالي (للترقية)
    total_pts = student.get("total_points", 0) or student.get("points", 0)
    if total_pts >= 500:
        recommendations.append({
            "id": "elite_candidate",
            "priority": "celebrate",
            "icon": "💎",
            "title": f"مرشح لنادي النخبة ({total_pts} XP)",
            "action": "رسالة دعوة",
            "template": f"🌟 دعوة خاصة لـ {student.get('full_name','')} 🌟\n\nبفضل اجتهادك ومثابرتك، حقّقت {total_pts} XP!\nأنت الآن مرشّح للانضمام إلى نادي النخبة 💎\n\nتواصل معي للتفاصيل والشروط.\n\nبفخر، مشرفك",
            "type": "note"
        })
    
    # ترتيب حسب الأولوية
    priority_order = {"high": 0, "medium": 1, "celebrate": 2, "low": 3}
    recommendations.sort(key=lambda x: priority_order.get(x.get("priority", "low"), 9))
    
    return {
        "student_id": student_id,
        "student_name": student.get("full_name", ""),
        "recommendations": recommendations,
        "total": len(recommendations)
    }


@app.post("/api/supervisor/intervention-log")
async def supervisor_log_intervention(payload: dict, sup = Depends(get_current_supervisor)):
    """📝 تسجيل تدخل في سجل الطالب (بعد إرسال رسالة من التوصيات)"""
    student_id = payload.get("student_id")
    rec_id = payload.get("recommendation_id", "")
    msg_id = payload.get("message_id")
    
    if not student_id:
        raise HTTPException(status_code=400, detail="student_id مطلوب")
    
    student_ids = _get_supervisor_student_ids(sup["id"])
    if int(student_id) not in student_ids:
        raise HTTPException(status_code=403, detail="ليس من طلابك")
    
    # حالياً نُسجّل في supervisor_messages مع metadata
    # في المستقبل: جدول منفصل interventions
    return {"status": "logged", "student_id": student_id, "rec_id": rec_id}


@app.get("/api/supervisor/messages")
async def supervisor_list_messages(sup = Depends(get_current_supervisor)):
    """📨 قائمة الرسائل التي أرسلها المشرف"""
    res = supabase.table("supervisor_messages").select("*").eq("supervisor_id", sup["id"]).order("created_at", desc=True).limit(100).execute()
    msgs = res.data or []
    # نُضيف أسماء الطلاب
    student_ids = list(set(m["student_id"] for m in msgs if m.get("student_id")))
    if student_ids:
        stus = supabase.table("students").select("id, full_name").in_("id", student_ids).execute()
        stu_map = {s["id"]: s["full_name"] for s in (stus.data or [])}
        for m in msgs:
            if m.get("student_id"):
                m["student_name"] = stu_map.get(m["student_id"], "—")
    return {"messages": msgs}


@app.post("/api/supervisor/messages/send")
async def supervisor_send_message(
    payload: dict,
    sup = Depends(get_current_supervisor)
):
    """📤 إرسال رسالة لطالب واحد أو لكل الطلاب"""
    message = (payload.get("message") or "").strip()
    if not message:
        raise HTTPException(status_code=400, detail="الرسالة فارغة")
    msg_type = payload.get("type", "note")
    link = payload.get("link")
    student_id = payload.get("student_id")
    
    student_ids = _get_supervisor_student_ids(sup["id"])
    if not student_ids:
        raise HTTPException(status_code=400, detail="لا يوجد طلاب لإرسال الرسالة")
    
    rows = []
    if student_id is None:
        # رسالة جماعية لكل الطلاب
        for sid in student_ids:
            rows.append({
                "supervisor_id": sup["id"],
                "student_id": sid,
                "message": message,
                "type": msg_type,
                "link": link,
                "is_broadcast": True
            })
    else:
        sid = int(student_id)
        if sid not in student_ids:
            raise HTTPException(status_code=403, detail="هذا الطالب ليس من طلابك")
        rows.append({
            "supervisor_id": sup["id"],
            "student_id": sid,
            "message": message,
            "type": msg_type,
            "link": link,
            "is_broadcast": False
        })
    res = supabase.table("supervisor_messages").insert(rows).execute()
    return {"status": "success", "sent": len(rows)}


@app.put("/api/supervisor/settings/alerts")
async def supervisor_update_alerts(
    payload: dict,
    sup = Depends(get_current_supervisor)
):
    """⚙️ تحديث معايير الإنذار"""
    allowed = ["min_accuracy", "challenges_lookback", "absent_days", "xp_drop_percent"]
    new_settings = {k: payload[k] for k in allowed if k in payload}
    if not new_settings:
        raise HTTPException(status_code=400, detail="لا توجد بيانات")
    # ندمج مع الإعدادات الحالية
    current = sup.get("alert_settings") or {}
    current.update(new_settings)
    supabase.table("supervisors").update({"alert_settings": current}).eq("id", sup["id"]).execute()
    return {"status": "success", "alert_settings": current}


@app.put("/api/supervisor/settings/password")
async def supervisor_change_password(
    payload: dict,
    sup = Depends(get_current_supervisor)
):
    """🔑 تغيير كلمة المرور"""
    current_pw = payload.get("current_password", "")
    new_pw = payload.get("new_password", "")
    if not verify_password(current_pw, sup["password"]):
        raise HTTPException(status_code=401, detail="كلمة المرور الحالية خاطئة")
    if len(new_pw) < 6:
        raise HTTPException(status_code=400, detail="الكلمة الجديدة قصيرة جداً (6+ أحرف)")
    supabase.table("supervisors").update({"password": hash_password(new_pw)}).eq("id", sup["id"]).execute()
    security_log("supervisor_password_changed", "self", {"sup_id": sup["id"]})
    return {"status": "success"}


@app.put("/api/supervisor/settings/theme")
async def supervisor_update_theme(
    theme: str = Form(...),
    sup = Depends(get_current_supervisor)
):
    """🎨 تحديث الثيم"""
    allowed = ["royal-gold", "ocean-blue", "emerald-forest", "rose-modern", "cosmic-purple", "sunset-warm", "turquoise-crystal", "luxury-black", "bright-day"]
    if theme not in allowed:
        raise HTTPException(status_code=400, detail="ثيم غير صالح")
    supabase.table("supervisors").update({"theme": theme}).eq("id", sup["id"]).execute()
    return {"status": "success", "theme": theme}



# ════════════════════════════════════════════════════════════
# 🛡️ ADMIN SECURITY ENDPOINTS
# ════════════════════════════════════════════════════════════

@app.get("/api/admin/security/locked_accounts")
async def get_locked_accounts(admin=Depends(get_current_admin)):
    """🔒 يُرجع قائمة الحسابات المقفولة حالياً"""
    import time as _time
    now = _time.time()
    
    locked = []
    for key, rec in list(_login_attempts.items()):
        locked_until = rec.get("locked_until", 0)
        if locked_until > now:
            seconds_left = int(locked_until - now)
            locked.append({
                "key": key,
                "attempts": rec.get("count", 0),
                "first_attempt": rec.get("first", 0),
                "last_attempt": rec.get("last", 0),
                "locked_until": locked_until,
                "seconds_left": seconds_left,
                "minutes_left": max(1, seconds_left // 60)
            })
    
    locked.sort(key=lambda x: x["seconds_left"], reverse=True)
    return {
        "locked_count": len(locked),
        "locked_accounts": locked,
        "total_tracked": len(_login_attempts)
    }


@app.post("/api/admin/security/unlock")
async def unlock_account(
    key: str = Form(...),
    admin=Depends(get_current_admin)
):
    """🔓 فك قفل حساب يدوياً (للأدمن)"""
    if key in _login_attempts:
        del _login_attempts[key]
        security_log("manual_unlock", "admin", {"key": key})
        return {"status": "success", "message": f"تم فك قفل {key}"}
    return {"status": "not_found", "message": "الحساب غير موجود في قائمة القفل"}


@app.post("/api/admin/security/cleanup")
async def cleanup_security_records(admin=Depends(get_current_admin)):
    """🧹 تنظيف السجلات القديمة (أكثر من 24 ساعة)"""
    before = len(_login_attempts)
    cleanup_old_attempts()
    after = len(_login_attempts)
    cleaned = before - after
    return {
        "status": "success",
        "cleaned": cleaned,
        "remaining": after
    }


@app.get("/api/admin/security/status")
async def security_status(admin=Depends(get_current_admin)):
    """📊 لوحة معلومات أمنية"""
    import time as _time
    now = _time.time()
    
    total = len(_login_attempts)
    locked = sum(1 for r in _login_attempts.values() if r.get("locked_until", 0) > now)
    
    # تجميع حسب النوع
    by_type = {"admin": 0, "student": 0, "teacher": 0, "parent": 0, "other": 0}
    for key in _login_attempts.keys():
        if key.startswith("admin:"):
            by_type["admin"] += 1
        elif key.startswith("student:"):
            by_type["student"] += 1
        elif key.startswith("teacher:"):
            by_type["teacher"] += 1
        elif key.startswith("parent:"):
            by_type["parent"] += 1
        else:
            by_type["other"] += 1
    
    return {
        "total_tracked": total,
        "currently_locked": locked,
        "by_type": by_type,
        "https_enforced": os.getenv("ENV", "production").lower() == "production",
        "trusted_hosts": _TRUSTED_HOSTS
    }




# ════════════════════════════════════════════════════════════
# 👑 ADMIN: إدارة المشرفين
# ════════════════════════════════════════════════════════════

@app.post("/api/admin/supervisors")
async def admin_create_supervisor(
    full_name: str = Form(...),
    email: str = Form(...),
    password: str = Form(...),
    phone: str = Form(""),
    admin = Depends(get_current_admin)
):
    """➕ إنشاء حساب مشرف جديد"""
    full_name = (full_name or "").strip()
    email = (email or "").strip().lower()
    password = password or ""
    if not full_name or not email or len(password) < 6:
        raise HTTPException(status_code=400, detail="الاسم والبريد مطلوبان + كلمة مرور 6 أحرف+")
    # نتحقق من البريد فريد
    existing = supabase.table("supervisors").select("id").eq("email", email).execute()
    if existing.data:
        raise HTTPException(status_code=400, detail="البريد مستخدم مسبقاً")
    # نُنشئ
    new_sup = {
        "full_name": full_name,
        "email": email,
        "phone": (phone or "").strip(),
        "password": hash_password(password),
        "is_active": True
    }
    res = supabase.table("supervisors").insert(new_sup).execute()
    if not res.data:
        raise HTTPException(status_code=500, detail="فشل الإنشاء")
    sup = res.data[0]
    sup.pop("password", None)
    security_log("supervisor_created", "admin", {"sup_id": sup["id"], "email": email})
    return {"status": "success", "supervisor": sup}


@app.get("/api/admin/supervisors")
async def admin_list_supervisors(admin = Depends(get_current_admin)):
    """📋 قائمة كل المشرفين مع عدد الطلاب لكل واحد"""
    try:
        sups = supabase.table("supervisors").select("id,full_name,email,phone,is_active,created_at,last_login").order("created_at", desc=True).execute()
        result = []
        for s in (sups.data or []):
            # نعدّ الطلاب
            cnt = supabase.table("supervisor_students").select("id", count="exact").eq("supervisor_id", s["id"]).execute()
            s["students_count"] = cnt.count or 0
            # نعدّ الاختبارات
            ex = supabase.table("supervisor_exams").select("id", count="exact").eq("supervisor_id", s["id"]).execute()
            s["exams_count"] = ex.count or 0
            result.append(s)
        return {"supervisors": result, "total": len(result)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.put("/api/admin/supervisors/{sup_id}")
async def admin_update_supervisor(
    sup_id: int,
    full_name: str = Form(None),
    email: str = Form(None),
    phone: str = Form(None),
    password: str = Form(None),
    is_active: str = Form(None),
    admin = Depends(get_current_admin)
):
    """✏️ تحديث بيانات مشرف"""
    update_data = {}
    if full_name and full_name.strip():
        update_data["full_name"] = full_name.strip()
    if email and email.strip():
        update_data["email"] = email.strip().lower()
    if phone is not None:
        update_data["phone"] = phone.strip()
    if password and len(password) >= 6:
        update_data["password"] = hash_password(password)
    if is_active is not None:
        update_data["is_active"] = str(is_active).lower() in ("true", "1", "yes", "نعم")
    if not update_data:
        raise HTTPException(status_code=400, detail="لا توجد بيانات للتحديث")
    res = supabase.table("supervisors").update(update_data).eq("id", sup_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="المشرف غير موجود")
    sup = res.data[0]
    sup.pop("password", None)
    security_log("supervisor_updated", "admin", {"sup_id": sup_id, "fields": list(update_data.keys())})
    return {"status": "success", "supervisor": sup}


@app.delete("/api/admin/supervisors/{sup_id}")
async def admin_delete_supervisor(sup_id: int, admin = Depends(get_current_admin)):
    """🗑️ حذف مشرف نهائياً"""
    res = supabase.table("supervisors").delete().eq("id", sup_id).execute()
    security_log("supervisor_deleted", "admin", {"sup_id": sup_id})
    return {"status": "success", "deleted": sup_id}


@app.get("/api/admin/supervisors/{sup_id}/students")
async def admin_get_supervisor_students(sup_id: int, admin = Depends(get_current_admin)):
    """👥 طلاب مشرف معيّن"""
    links = supabase.table("supervisor_students").select("student_id,assigned_at,notes").eq("supervisor_id", sup_id).execute()
    student_ids = [l["student_id"] for l in (links.data or [])]
    if not student_ids:
        return {"students": []}
    students_data = supabase.table("students").select("id, full_name, username, grade, total_points, created_at").in_("id", student_ids).execute()
    result = students_data.data or []
    # normalize للتوافق
    for s in result:
        s["xp"] = s.get("total_points", 0) or 0
        s["points"] = s.get("total_points", 0) or 0
        s["total_xp"] = s.get("total_points", 0) or 0
        s["last_active"] = s.get("created_at", "") or ""
    return {"students": result}


@app.post("/api/admin/supervisors/{sup_id}/assign-student")
async def admin_assign_student(
    sup_id: int,
    student_id: int = Form(...),
    admin = Depends(get_current_admin)
):
    """🔗 ربط طالب بمشرف"""
    # نتحقق من وجود المشرف والطالب
    sup_check = supabase.table("supervisors").select("id").eq("id", sup_id).execute()
    if not sup_check.data:
        raise HTTPException(status_code=404, detail="المشرف غير موجود")
    stu_check = supabase.table("students").select("id").eq("id", student_id).execute()
    if not stu_check.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    # نتحقق من عدم التكرار
    existing = supabase.table("supervisor_students").select("id").eq("supervisor_id", sup_id).eq("student_id", student_id).execute()
    if existing.data:
        return {"status": "exists", "message": "الطالب مرتبط مسبقاً"}
    res = supabase.table("supervisor_students").insert({
        "supervisor_id": sup_id,
        "student_id": student_id
    }).execute()
    return {"status": "success", "link": res.data[0] if res.data else None}


@app.delete("/api/admin/supervisors/{sup_id}/unassign-student/{student_id}")
async def admin_unassign_student(
    sup_id: int,
    student_id: int,
    admin = Depends(get_current_admin)
):
    """🔓 فك ربط طالب من مشرف"""
    supabase.table("supervisor_students").delete().eq("supervisor_id", sup_id).eq("student_id", student_id).execute()
    return {"status": "success"}




@app.get("/api/admin/check-supervisor-table")
async def admin_check_supervisor_table(admin = Depends(get_current_admin)):
    """🔍 فحص حالة جدول supervisor_students"""
    result = {
        "table_exists": False,
        "row_count": 0,
        "sample": [],
        "error": None,
        "create_sql": ""
    }
    
    try:
        res = supabase.table("supervisor_students").select("*").limit(5).execute()
        result["table_exists"] = True
        result["sample"] = res.data or []
        
        # عدد الصفوف
        count_res = supabase.table("supervisor_students").select("supervisor_id", count="exact").limit(1).execute()
        result["row_count"] = count_res.count or 0
    except Exception as e:
        result["error"] = str(e)[:300]
    
    result["create_sql"] = """-- SQL لإنشاء جدول supervisor_students في Supabase:

CREATE TABLE IF NOT EXISTS supervisor_students (
    id SERIAL PRIMARY KEY,
    supervisor_id INTEGER NOT NULL REFERENCES supervisors(id) ON DELETE CASCADE,
    student_id INTEGER NOT NULL REFERENCES students(id) ON DELETE CASCADE,
    assigned_at TIMESTAMPTZ DEFAULT NOW(),
    notes TEXT,
    UNIQUE(supervisor_id, student_id)
);

CREATE INDEX IF NOT EXISTS idx_supstu_sup ON supervisor_students(supervisor_id);
CREATE INDEX IF NOT EXISTS idx_supstu_stu ON supervisor_students(student_id);
"""
    return result

@app.post("/api/admin/supervisors/{sup_id}/assign-bulk")
async def admin_assign_bulk(
    sup_id: int,
    student_ids: str = Form(...),
    admin = Depends(get_current_admin)
):
    """🔗 ربط مجموعة طلاب دفعة واحدة. student_ids: '1,2,3' أو JSON array"""
    import json as _json
    
    # 1. parse IDs
    try:
        if student_ids.strip().startswith("["):
            ids = _json.loads(student_ids)
        else:
            ids = [int(x.strip()) for x in student_ids.split(",") if x.strip()]
        ids = [int(i) for i in ids]
    except Exception as e:
        print(f"[assign-bulk] parse error: {e}, input: {student_ids[:200]}")
        raise HTTPException(status_code=400, detail=f"صيغة IDs خاطئة: {str(e)[:100]}")
    
    if not ids:
        return {"status": "success", "added": 0, "message": "لا توجد IDs"}
    
    print(f"[assign-bulk] sup={sup_id}, ids={ids}")
    
    # 2. التحقق من وجود المشرف
    try:
        sup_check = supabase.table("supervisors").select("id, full_name").eq("id", sup_id).execute()
        if not sup_check.data:
            raise HTTPException(status_code=404, detail=f"المشرف #{sup_id} غير موجود")
    except HTTPException:
        raise
    except Exception as e:
        print(f"[assign-bulk] supervisor check error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ التحقق من المشرف: {str(e)[:100]}")
    
    # 3. التحقق من وجود الطلاب
    try:
        students_check = supabase.table("students").select("id").in_("id", ids).execute()
        valid_ids = [s["id"] for s in (students_check.data or [])]
        if len(valid_ids) != len(ids):
            missing = set(ids) - set(valid_ids)
            print(f"[assign-bulk] missing students: {missing}")
            ids = valid_ids  # نستخدم فقط الموجودين
            if not ids:
                raise HTTPException(status_code=400, detail=f"لم يُعثر على الطلاب: {list(missing)[:5]}")
    except HTTPException:
        raise
    except Exception as e:
        print(f"[assign-bulk] students check error: {e}")
        raise HTTPException(status_code=500, detail=f"خطأ التحقق من الطلاب: {str(e)[:100]}")
    
    # 4. جلب الموجودين لتجنب duplicates
    try:
        existing = supabase.table("supervisor_students").select("student_id").eq("supervisor_id", sup_id).execute()
        existing_ids = set(r["student_id"] for r in (existing.data or []))
    except Exception as e:
        # ربما الجدول غير موجود!
        print(f"[assign-bulk] table check error: {e}")
        raise HTTPException(status_code=500, detail=f"جدول supervisor_students قد لا يكون موجوداً في Supabase: {str(e)[:150]}")
    
    new_ids = [i for i in ids if i not in existing_ids]
    if not new_ids:
        return {"status": "success", "added": 0, "message": f"كلهم مرتبطون مسبقاً ({len(ids)} طالب)"}
    
    # 5. الـ insert
    try:
        rows = [{"supervisor_id": sup_id, "student_id": sid} for sid in new_ids]
        res = supabase.table("supervisor_students").insert(rows).execute()
        print(f"[assign-bulk] ✅ added {len(new_ids)} students to sup {sup_id}")
        
        # إبطال الـ cache (لو موجود)
        try:
            _cache.invalidate_prefix("supervisor:")
            _cache.invalidate_prefix("students:")
        except Exception:
            pass
        
        return {
            "status": "success", 
            "added": len(new_ids),
            "skipped": len(ids) - len(new_ids),
            "total_attempted": len(ids)
        }
    except Exception as e:
        err_str = str(e)
        print(f"[assign-bulk] insert error: {err_str}")
        # رسالة واضحة حسب نوع الخطأ
        if "duplicate" in err_str.lower() or "unique" in err_str.lower():
            raise HTTPException(status_code=409, detail="بعض الطلاب مرتبطون مسبقاً")
        elif "foreign key" in err_str.lower():
            raise HTTPException(status_code=400, detail="مرجع غير صالح (طالب أو مشرف)")
        elif "does not exist" in err_str.lower() or "relation" in err_str.lower():
            raise HTTPException(status_code=500, detail="جدول supervisor_students غير موجود في Supabase!")
        else:
            raise HTTPException(status_code=500, detail=f"خطأ الإدراج: {err_str[:200]}")


@app.get("/api/admin/supervisors/{sup_id}/report")
async def admin_supervisor_report(sup_id: int, admin = Depends(get_current_admin)):
    """📊 تقرير أداء مشرف"""
    sup = supabase.table("supervisors").select("*").eq("id", sup_id).execute()
    if not sup.data:
        raise HTTPException(status_code=404, detail="المشرف غير موجود")
    sup_data = sup.data[0]
    sup_data.pop("password", None)
    # عدد الطلاب
    students = supabase.table("supervisor_students").select("student_id").eq("supervisor_id", sup_id).execute()
    student_ids = [r["student_id"] for r in (students.data or [])]
    # الاختبارات
    exams = supabase.table("supervisor_exams").select("id,title,is_published,created_at").eq("supervisor_id", sup_id).execute()
    # متوسط نتائج الاختبارات
    exam_ids = [e["id"] for e in (exams.data or [])]
    avg_score = 0
    if exam_ids:
        results = supabase.table("supervisor_exam_results").select("score,max_score").in_("exam_id", exam_ids).execute()
        if results.data:
            scores = [(float(r["score"]) / float(r["max_score"]) * 100) for r in results.data if r.get("max_score")]
            if scores:
                avg_score = round(sum(scores) / len(scores), 1)
    # عدد الرسائل
    msgs = supabase.table("supervisor_messages").select("id", count="exact").eq("supervisor_id", sup_id).execute()
    return {
        "supervisor": sup_data,
        "students_count": len(student_ids),
        "exams_count": len(exams.data or []),
        "exams_published": sum(1 for e in (exams.data or []) if e.get("is_published")),
        "avg_score_percent": avg_score,
        "messages_sent": msgs.count or 0
    }


# ==========================================
# --- 14. نظام النخبة (Elite System) ---
# ==========================================

class EliteArenaManager:
    """مدير ساحة مبارزة النخبة — WebSocket مستقل"""
    def __init__(self):
        self.waiting_players: dict = {}
        self.active_rooms:    dict = {}

    async def connect(self, websocket: WebSocket, student_name: str, grade: str):
        await websocket.accept()
        if grade not in self.waiting_players:
            self.waiting_players[grade] = []
        self.waiting_players[grade].append({"ws": websocket, "name": student_name})
        await self.elite_matchmake(grade)

    async def elite_matchmake(self, grade: str):
        queue = self.waiting_players[grade]
        if len(queue) < 2:
            try:
                await queue[-1]["ws"].send_json({"type": "waiting", "msg": "⏳ بحث عن منافس من النخبة..."})
            except: pass
            return
        p1 = queue.pop(0)
        p2 = queue.pop(0)
        room_id = f"elite_{id(p1['ws'])}_{id(p2['ws'])}"
        self.active_rooms[room_id] = {"p1": p1, "p2": p2, "scores": {p1["name"]: 0, p2["name"]: 0}}
        for p, opp in [(p1, p2), (p2, p1)]:
            try:
                await p["ws"].send_json({
                    "type": "matched",
                    "room_id": room_id,
                    "opponent": opp["name"],
                    "msg": f"⚔️ تم إيجاد منافس: {opp['name']}"
                })
            except: pass

    async def broadcast_score(self, room_id: str, sender_name: str, new_score: int):
        room = self.active_rooms.get(room_id)
        if not room: return
        room["scores"][sender_name] = new_score
        for key in ["p1", "p2"]:
            try:
                await room[key]["ws"].send_json({
                    "type": "score_update",
                    "scores": room["scores"]
                })
            except: pass

    async def disconnect(self, websocket: WebSocket, grade: str):
        if grade in self.waiting_players:
            self.waiting_players[grade] = [
                p for p in self.waiting_players[grade] if p["ws"] != websocket
            ]

elite_arena_manager = EliteArenaManager()


@app.websocket("/api/elite/arena/ws/{student_name}/{grade}")
async def elite_arena_websocket(websocket: WebSocket, student_name: str, grade: str):
    """WebSocket لساحة مبارزة النخبة"""
    from urllib.parse import unquote
    clean_name  = unquote(student_name)
    clean_grade = unquote(grade)
    await elite_arena_manager.connect(websocket, clean_name, clean_grade)
    try:
        while True:
            data = await websocket.receive_json()
            if data.get("type") == "score_update":
                await elite_arena_manager.broadcast_score(
                    data["room_id"], clean_name, data["score"]
                )
    except WebSocketDisconnect:
        await elite_arena_manager.disconnect(websocket, clean_grade)


@app.post("/api/elite/request")
async def submit_elite_request(
    request: Request,
    student_id: int = Form(...),
    username:   str = Form(...),
    full_name:  str = Form(...),
    grade:      str = Form(...),
    xp:         int = Form(default=0),
    lessons_85: int = Form(default=0),
):
    """طلب انضمام تلقائي لنادي النخبة"""
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=5, window_seconds=3600):
        raise HTTPException(status_code=429, detail="تم الإرسال مسبقاً")
    # تحقق: هل قدّم طلباً سابقاً؟
    existing = supabase.table("elite_requests").select("id,status")         .eq("student_id", student_id).execute()
    if existing.data:
        st = existing.data[0]["status"]
        if st == "approved": return {"status": "already_elite"}
        if st == "pending":  return {"status": "pending"}
    supabase.table("elite_requests").insert({
        "student_id": student_id, "username": username,
        "full_name":  full_name,  "grade":    grade,
        "xp":         xp,         "lessons_85": lessons_85,
        "status":     "pending"
    }).execute()
    return {"status": "submitted"}


@app.get("/api/elite/check/{student_id}")
async def check_elite_status(student_id: int):
    """هل الطالب معتمد كفائق؟"""
    res = supabase.table("students").select("is_elite").eq("id", student_id).execute()
    if res.data:
        return {"is_elite": bool(res.data[0].get("is_elite", False))}
    return {"is_elite": False}


@app.get("/api/elite/questions")
async def get_elite_questions(grade: str = "", lesson: str = ""):
    """أسئلة النخبة — مصنّفة"""
    query = supabase.table("questions").select(
        "id,grade,lesson,subject,q_type,question,options,answer,image_url,difficulty"
    ).eq("is_elite", True)
    if grade:   query = query.eq("grade", grade.strip())
    if lesson:  query = query.ilike("lesson", lesson.strip())
    res = query.execute()
    return res.data or []


@app.get("/api/elite/leaderboard")
async def elite_leaderboard():
    """ترتيب الفائقين — استعلام واحد بدلاً من N+1"""
    # ═══ 1) جلب كل طلاب النخبة (استعلام واحد) ═══
    res = supabase.table("students").select(
        "id, full_name, grade, school_name, avatar_url"
    ).eq("is_elite", True).execute()
    students = res.data or []

    if not students:
        return []

    # ═══ 2) جلب كل النتائج لهم دفعة واحدة (استعلام واحد باستخدام in_) ═══
    student_ids = [st["id"] for st in students]
    all_results = supabase.table("results").select("student_id,score,total")\
        .in_("student_id", student_ids).execute().data or []

    # ═══ 3) تجميع النتائج حسب الطالب في الذاكرة ═══
    results_by_student = {}
    for r in all_results:
        sid = r.get("student_id")
        if sid is None:
            continue
        results_by_student.setdefault(sid, []).append(r)

    # ═══ 4) بناء لوحة الترتيب ═══
    board = []
    for st in students:
        results = results_by_student.get(st["id"], [])
        total_correct = sum(x.get("score", 0) or 0 for x in results)
        total_q = sum(x.get("total", 0) or 0 for x in results if (x.get("total") or 0) > 0)
        accuracy = round((total_correct / total_q * 100)) if total_q > 0 else 0
        board.append({**st, "xp": total_correct, "accuracy": accuracy, "tests": len(results)})

    board.sort(key=lambda x: (-x["xp"], -x["accuracy"]))
    return board[:50]


# ── ADMIN: إدارة النخبة ──
@app.get("/api/admin/elite/requests")
async def get_elite_requests(admin=Depends(get_current_admin)):
    res = supabase.table("elite_requests").select("*").order("created_at", desc=True).execute()
    return res.data or []


@app.post("/api/admin/elite/approve/{request_id}")
async def approve_elite(request_id: int, admin=Depends(get_current_admin)):
    req = supabase.table("elite_requests").select("*").eq("id", request_id).execute()
    if not req.data: raise HTTPException(status_code=404, detail="الطلب غير موجود")
    r = req.data[0]
    supabase.table("students").update({
        "is_elite": True,
        "elite_approved_at": datetime.now(timezone.utc).isoformat()
    }).eq("id", r["student_id"]).execute()
    supabase.table("elite_requests").update({"status": "approved"}).eq("id", request_id).execute()
    return {"status": "approved"}


@app.post("/api/admin/elite/reject/{request_id}")
async def reject_elite(request_id: int, admin=Depends(get_current_admin)):
    supabase.table("elite_requests").update({"status": "rejected"}).eq("id", request_id).execute()
    return {"status": "rejected"}


@app.post("/api/admin/elite/revoke/{student_id}")
async def revoke_elite(student_id: int, admin=Depends(get_current_admin)):
    supabase.table("students").update({"is_elite": False}).eq("id", student_id).execute()
    supabase.table("elite_requests").update({"status": "rejected"})         .eq("student_id", student_id).execute()
    return {"status": "revoked"}


@app.post("/api/admin/elite/grant/{student_id}")
async def grant_elite_manually(student_id: int, admin=Depends(get_current_admin)):
    """منح لقب الفائق يدوياً"""
    supabase.table("students").update({
        "is_elite": True,
        "elite_approved_at": datetime.now(timezone.utc).isoformat()
    }).eq("id", student_id).execute()
    return {"status": "granted"}


@app.get("/api/admin/elite/members")
async def get_elite_members(admin=Depends(get_current_admin)):
    res = supabase.table("students").select(
        "id, full_name, username, grade, school_name, is_elite, elite_approved_at"
    ).eq("is_elite", True).execute()
    return res.data or []


# ==========================================
# --- الإشعارات العامة ---
# ==========================================
@app.post("/api/admin/notifications")
async def send_notification(
    request: Request,
    title:    str = Form(...),
    body:     str = Form(...),
    grade:    str = Form(default="all"),
    priority: str = Form(default="normal"),
    type:     str = Form(default="announcement"),
    admin=Depends(get_current_admin)
):
    """إرسال إشعار للطلاب — يُخزَّن في Supabase"""
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=20, window_seconds=60):
        raise HTTPException(status_code=429, detail="طلبات كثيرة جداً")
    row = {
        "title":    title.strip()[:200],
        "body":     body.strip()[:5000],
        "grade":    grade,
        "priority": priority,
        "type":     type,
        "is_active": True,
    }
    try:
        supabase.table("notifications").insert(row).execute()
    except Exception as e:
        print(f"notifications insert error: {e}")
    return {"status": "success", "message": "تم إرسال الإشعار"}


@app.get("/api/notifications")
async def get_notifications(grade: str = ""):
    """جلب الإشعارات النشطة للطلاب"""
    try:
        if grade and grade != "all":
            res1 = supabase.table("notifications").select("*").eq("is_active", True).eq("grade", grade).order("created_at", desc=True).limit(5).execute()
            res2 = supabase.table("notifications").select("*").eq("is_active", True).eq("grade", "all").order("created_at", desc=True).limit(5).execute()
            data = (res1.data or []) + (res2.data or [])
            data.sort(key=lambda x: x.get("created_at", ""), reverse=True)
            return data[:10]
        res = supabase.table("notifications").select("*").eq("is_active", True).order("created_at", desc=True).limit(10).execute()
        return res.data or []
    except Exception:
        return []


@app.delete("/api/admin/notifications/{notif_id}")
async def delete_notification(notif_id: int, admin=Depends(get_current_admin)):
    """تعطيل إشعار"""
    try:
        supabase.table("notifications").update({"is_active": False}).eq("id", notif_id).execute()
    except Exception:
        pass
    return {"status": "success"}


# ==========================================
# --- endpoints إدارة الحسابات ---
# ==========================================

@app.get("/api/admin/students/full")
async def get_all_students_full(admin=Depends(get_current_admin)):
    """جلب قائمة الطلاب الكاملة + إحصائيات XP والتحديات والوقت الفعلي
    
    🛡️ مُحدّث: يستخدم total_points من جدول students (المصدر الموثوق للـ XP)
    بدلاً من جمع الدرجات من results.
    """
    # 1. جلب كل الطلاب (pagination) — مع total_points للـ XP الحقيقي
    students = []
    offset = 0
    for _ in range(20):
        try:
            res = supabase.table("students").select(
                "id, full_name, username, grade, school_name, avatar_url, is_active, is_elite, created_at, parent_code, parent_phone, parent_name, parent_email, total_points"
            ).order("full_name").range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            students.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        except Exception:
            break

    # 2. جلب كل النتائج (pagination) لحساب XP والتحديات
    all_results = []
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("results").select(
                "student_id, student_name, score, total, timestamp"
            ).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            all_results.extend(batch)
            if len(batch) < 1000:
                break
            offset += 1000
        except Exception:
            break

    # 3. تجميع إحصائيات لكل طالب (عدد التحديات + متوسط النسبة)
    # ⚠️ ملاحظة: XP يأتي من students.total_points (المصدر الموثوق)
    # حساب XP من results مُلغى لأنه يُعطي قيم خاطئة
    from collections import defaultdict
    stats = defaultdict(lambda: {"tests": 0, "score_sum": 0, "total_sum": 0})
    for r in all_results:
        sid = r.get("student_id")
        if not sid:
            continue
        stats[sid]["tests"]     += 1
        stats[sid]["score_sum"] += r.get("score", 0) or 0
        stats[sid]["total_sum"] += r.get("total", 0) or 0

    # 4. جلب وقت الجلسات (آخر 30 يوم) لكل طالب
    from datetime import datetime, timezone, timedelta
    range_30 = (datetime.now(timezone.utc) - timedelta(days=30)).isoformat()
    sessions_by_student = defaultdict(int)  # student_id -> bucket count
    offset = 0
    for _ in range(50):
        try:
            res = supabase.table("student_sessions").select(
                "student_id"
            ).gte("session_bucket", range_30).range(offset, offset + 999).execute()
            batch = res.data or []
            if not batch:
                break
            for s in batch:
                sid = s.get("student_id")
                if sid:
                    sessions_by_student[sid] += 1
            if len(batch) < 1000:
                break
            offset += 1000
        except Exception:
            break

    # 5. دمج البيانات — XP من total_points الموثوق
    enriched = []
    for s in students:
        sid = s.get("id")
        sst = stats.get(sid, {"tests": 0, "score_sum": 0, "total_sum": 0})
        avg_pct = round((sst["score_sum"] / sst["total_sum"]) * 100, 1) if sst["total_sum"] > 0 else 0
        # كل bucket يمثّل 5 دقائق نشاط حقيقي
        minutes_30d = sessions_by_student.get(sid, 0) * 5
        # 🛡️ XP الموثوق من DB (وليس مجموع الدرجات)
        real_xp = int(s.get("total_points", 0) or 0)
        enriched.append({
            **s,
            "xp":            real_xp,
            "total_points":  real_xp,  # alias للتوافق
            "tests":         sst["tests"],
            "avg_score_pct": avg_pct,
            "minutes_30d":   minutes_30d,
            "hours_30d":     round(minutes_30d / 60, 1),
        })

    return enriched


@app.post("/api/admin/students/{student_id}/toggle")
async def toggle_student(student_id: int, is_active: str = Form(...), admin=Depends(get_current_admin)):
    """تعطيل أو تفعيل حساب طالب"""
    active = is_active.lower() not in ('false', '0', 'no')
    supabase.table("students").update({"is_active": active}).eq("id", student_id).execute()
    return {"status": "success", "is_active": active}


@app.delete("/api/admin/students/{student_id}/delete")
async def delete_student(student_id: int, request: Request, admin=Depends(get_current_admin)):
    """حذف طالب نهائياً — يتطلب كلمة مرور الأدمن"""
    body = await request.form()
    admin_pass = body.get("admin_password", "")
    if not admin_pass or not (admin_pass == ADMIN_PASSWORD or verify_password(admin_pass, ADMIN_PASSWORD)):
        raise HTTPException(status_code=403, detail="كلمة مرور الأدمن خاطئة")
    supabase.table("results").delete().eq("student_id", student_id).execute()
    supabase.table("elite_requests").delete().eq("student_id", student_id).execute()
    res = supabase.table("students").delete().eq("id", student_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    return {"status": "deleted"}


@app.get("/api/admin/teachers")
async def get_all_teachers(admin=Depends(get_current_admin)):
    """جلب قائمة المعلمين"""
    try:
        res = supabase.table("teachers").select(
            "id, full_name, username, subject, is_active, created_at, last_login"
        ).order("full_name").execute()
    except Exception:
        # last_login قد لا يوجد بعد
        res = supabase.table("teachers").select(
            "id, full_name, username, subject, is_active, created_at"
        ).order("full_name").execute()
    return res.data or []


@app.get("/api/admin/teachers/activity")
async def teachers_activity(admin=Depends(get_current_admin)):
    """📊 إحصاء نشاط المعلمين: نشطون / غائبون / لم يدخلوا قط"""
    try:
        res = supabase.table("teachers").select(
            "id, full_name, username, subject, is_active, created_at, last_login"
        ).order("last_login", desc=True).execute()
        teachers = res.data or []
    except Exception:
        # العمود غير موجود — نُرجع القائمة بلا تتبّع
        res = supabase.table("teachers").select(
            "id, full_name, username, subject, is_active, created_at"
        ).execute()
        return {
            "has_tracking": False,
            "message": "عمود last_login غير موجود — نفّذ migration أولاً",
            "total": len(res.data or []),
            "teachers": res.data or []
        }
    now = datetime.now(timezone.utc)
    active_7d = absent = never = 0
    enriched = []
    for t in teachers:
        ll = t.get("last_login")
        days = None
        status = "never"
        if ll:
            try:
                dt = datetime.fromisoformat(str(ll).replace("Z", "+00:00"))
                days = (now - dt).days
                status = "active" if days <= 7 else "absent"
            except Exception:
                pass
        if status == "active": active_7d += 1
        elif status == "absent": absent += 1
        else: never += 1
        t["_days_since"] = days
        t["_status"] = status
        enriched.append(t)
    return {
        "has_tracking": True,
        "total": len(teachers),
        "active_7d": active_7d,
        "absent": absent,
        "never_logged": never,
        "teachers": enriched
    }


@app.post("/api/admin/teachers/{teacher_id}/toggle")
async def toggle_teacher(teacher_id: int, is_active: str = Form(...), admin=Depends(get_current_admin)):
    """تعطيل أو تفعيل حساب معلم"""
    active = is_active.lower() not in ('false', '0', 'no')
    supabase.table("teachers").update({"is_active": active}).eq("id", teacher_id).execute()
    return {"status": "success", "is_active": active}


@app.delete("/api/admin/teachers/{teacher_id}/delete")
async def delete_teacher(teacher_id: int, request: Request, admin=Depends(get_current_admin)):
    """حذف معلم نهائياً — يتطلب كلمة مرور الأدمن"""
    body = await request.form()
    admin_pass = body.get("admin_password", "")
    if not admin_pass or not (admin_pass == ADMIN_PASSWORD or verify_password(admin_pass, ADMIN_PASSWORD)):
        raise HTTPException(status_code=403, detail="كلمة مرور الأدمن خاطئة")
    res = supabase.table("teachers").delete().eq("id", teacher_id).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="المعلم غير موجود")
    return {"status": "deleted"}


@app.post("/api/admin/teachers/add")
async def add_teacher_admin(
    full_name: str = Form(...),
    username:  str = Form(...),
    password:  str = Form(...),
    subject:   str = Form(default="رياضيات"),
    admin=Depends(get_current_admin)
):
    """إضافة معلم جديد من لوحة الأدمن"""
    if len(password) < 6:
        raise HTTPException(status_code=400, detail="كلمة المرور أقل من 6 أحرف")
    existing = supabase.table("teachers").select("username").eq("username", username).execute()
    if existing.data:
        raise HTTPException(status_code=400, detail="اسم المستخدم موجود مسبقاً")
    supabase.table("teachers").insert({
        "full_name": full_name.strip(),
        "username":  username.strip().lower(),
        "password":  hash_password(password),
        "subject":   subject.strip(),
        "is_active": True,
    }).execute()
    return {"status": "success", "message": f"تم إضافة المعلم {full_name}"}

# ==========================================
# --- 14. نظام التحضيرات الملكية ---
# ==========================================
# SQL لإنشاء الجدول في Supabase (نفّذه مرة واحدة):
# CREATE TABLE lesson_preparations (
#   id              BIGSERIAL PRIMARY KEY,
#   grade           TEXT NOT NULL,
#   semester        TEXT NOT NULL,
#   unit            TEXT NOT NULL,
#   lesson          TEXT NOT NULL,
#   concepts        TEXT DEFAULT '',
#   warm_up         TEXT DEFAULT '',
#   activities      TEXT DEFAULT '',
#   formative_eval  TEXT DEFAULT '',
#   summative_eval  TEXT DEFAULT '',
#   attachments     JSONB DEFAULT '[]',
#   created_at      TIMESTAMPTZ DEFAULT NOW(),
#   updated_at      TIMESTAMPTZ DEFAULT NOW(),
#   UNIQUE(grade, semester, unit, lesson)
# );


@app.get("/api/preparations")
async def get_preparation(grade: str, semester: str, unit: str, lesson: str):
    """جلب تحضير درس محدد — متاح بدون مصادقة للعرض"""
    try:
        res = supabase.table("lesson_preparations").select("*") \
            .eq("grade",    grade.strip()) \
            .eq("semester", semester.strip()) \
            .eq("unit",     unit.strip()) \
            .eq("lesson",   lesson.strip()) \
            .execute()
        if res.data:
            return res.data[0]
        return {
            "id": None, "concepts": "", "warm_up": "", "activities": "",
            "formative_eval": "", "summative_eval": "", "attachments": []
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/admin/preparations")
async def save_preparation(request: Request, admin=Depends(get_current_admin)):
    """حفظ أو تحديث تحضير درس — يتطلب صلاحيات الأدمن"""
    ip = request.client.host if request.client else "unknown"
    if _is_rate_limited(ip, max_calls=30, window_seconds=60):
        raise HTTPException(status_code=429, detail="طلبات كثيرة جداً")
    try:
        body = await request.json()
    except Exception:
        raise HTTPException(status_code=400, detail="صيغة الطلب غير صحيحة")

    grade    = (body.get("grade",    "") or "").strip()
    semester = (body.get("semester", "") or "").strip()
    unit     = (body.get("unit",     "") or "").strip()
    lesson   = (body.get("lesson",   "") or "").strip()

    if not all([grade, semester, unit, lesson]):
        raise HTTPException(status_code=400, detail="الصف والفصل والوحدة والدرس مطلوبة")

    row = {
        "grade":          grade,
        "semester":       semester,
        "unit":           unit,
        "lesson":         lesson,
        "concepts":       (body.get("concepts")       or "")[:5000],
        "warm_up":        (body.get("warm_up")        or "")[:5000],
        "activities":     (body.get("activities")     or "")[:10000],
        "formative_eval": (body.get("formative_eval") or "")[:5000],
        "summative_eval": (body.get("summative_eval") or "")[:5000],
        "attachments":    body.get("attachments") or [],
        "updated_at":     datetime.now(timezone.utc).isoformat(),
    }

    existing = supabase.table("lesson_preparations").select("id") \
        .eq("grade", grade).eq("semester", semester) \
        .eq("unit",  unit).eq("lesson",   lesson).execute()

    if existing.data:
        supabase.table("lesson_preparations") \
            .update(row).eq("id", existing.data[0]["id"]).execute()
        return {"status": "updated"}
    else:
        supabase.table("lesson_preparations").insert(row).execute()
        return {"status": "created"}


@app.post("/api/admin/preparations/attachment")
async def upload_preparation_attachment(
    file: UploadFile = File(...),
    admin=Depends(get_current_admin)
):
    """رفع مرفق لتحضير الدرس (PDF، صورة، HTML، فيديو، SCORM)"""
    if not file.filename:
        raise HTTPException(status_code=400, detail="لم يُرسل ملف")

    ext  = os.path.splitext(file.filename)[1].lower()
    name = f"prep_{uuid.uuid4().hex}{ext}"
    content_bytes = await file.read()

    mime = file.content_type or "application/octet-stream"
    if ext in [".zip", ".scorm"]:  mime = "application/zip"
    elif ext in [".html", ".htm"]: mime = "text/html"

    supabase.storage.from_("resources").upload(
        path=name, file=content_bytes,
        file_options={"content-type": mime}
    )
    url = supabase.storage.from_("resources").get_public_url(name)

    file_type = "file"
    if   ext in [".pdf"]:                                   file_type = "pdf"
    elif ext in [".html", ".htm"]:                          file_type = "html"
    elif ext in [".zip", ".scorm"]:                         file_type = "scorm"
    elif ext in [".mp4", ".webm", ".ogg"]:                  file_type = "video"
    elif ext in [".jpg", ".jpeg", ".png", ".gif", ".webp"]: file_type = "image"

    return {"url": url, "name": file.filename, "type": file_type, "mime": mime}


@app.delete("/api/admin/preparations/{prep_id}")
async def delete_preparation(prep_id: int, admin=Depends(get_current_admin)):
    """حذف تحضير درس كامل"""
    supabase.table("lesson_preparations").delete().eq("id", prep_id).execute()
    return {"status": "deleted"}



# ==========================================
# --- 15. اشتراك المعلمين ---
# ==========================================

@app.post("/api/teacher/subscription/activate")
async def activate_teacher_subscription(
    code:       str = Form(...),
    teacher_id: Optional[int] = Form(default=None)
):
    """تفعيل كود اشتراك من قِبَل المعلم"""
    code_upper = code.strip().upper()

    res = supabase.table("subscription_codes").select("*").eq("code", code_upper).execute()
    if not res.data:
        raise HTTPException(status_code=404, detail="الكود غير موجود")

    entry = res.data[0]

    if entry.get("is_used"):
        raise HTTPException(status_code=400, detail="هذا الكود مستخدَم مسبقاً")

    # التحقق من نوع المستخدم — الكود يجب أن يكون للمعلم أو عاماً
    user_type = entry.get("user_type", "student")
    if user_type == "student":
        raise HTTPException(status_code=403, detail="هذا الكود مخصص للطلاب فقط")

    # حساب تاريخ الانتهاء
    months = entry.get("months", 1)
    if months == -1:
        expiry = None
    else:
        now    = datetime.now(timezone.utc)
        expiry = (now.replace(
            month=(((now.month - 1) + months) % 12) + 1,
            year=now.year + (((now.month - 1) + months) // 12)
        )).isoformat()

    update_data = {
        "is_used":              True,
        "used_at":              datetime.now(timezone.utc).isoformat(),
        "activated_by_student": teacher_id,
    }
    supabase.table("subscription_codes").update(update_data).eq("id", entry["id"]).execute()

    return {
        "status": "success",
        "months": months,
        "expiry": expiry,
        "note":   entry.get("note", "")
    }


@app.get("/api/teacher/subscription/check")
async def check_teacher_subscription(teacher_id: int):
    """التحقق من حالة اشتراك المعلم"""
    try:
        res = supabase.table("subscription_codes").select("*") \
            .eq("activated_by_student", teacher_id) \
            .eq("is_used", True) \
            .eq("user_type", "teacher") \
            .order("used_at", desc=True) \
            .limit(1).execute()

        if not res.data:
            return {"active": False, "expiry": None, "months": 0}

        entry  = res.data[0]
        months = entry.get("months", 1)
        expiry = None

        if months == -1:
            return {"active": True, "expiry": None, "months": -1, "label": "👑 دائم"}

        used_at = entry.get("used_at")
        if used_at:
            from datetime import timezone as _tz
            activated = datetime.fromisoformat(used_at.replace("Z", "+00:00"))
            expiry_dt = activated.replace(
                month=(((activated.month - 1) + months) % 12) + 1,
                year=activated.year + (((activated.month - 1) + months) // 12)
            )
            now_utc = datetime.now(timezone.utc)
            expiry  = expiry_dt.isoformat()
            active  = expiry_dt > now_utc
            return {"active": active, "expiry": expiry, "months": months}

        return {"active": True, "expiry": None, "months": months}
    except Exception as e:
        return {"active": False, "expiry": None, "months": 0}


# ==========================================
# --- 12.5 🎓 نظام السنة الدراسية (Academic Year) ---
# ==========================================
"""
آلية العمل:
- السنة الدراسية: سبتمبر → يوليو (2025-2026 ← مثال)
- 1 أغسطس: أرشفة تلقائية + تصفير XP (لكن الحفاظ على الأوسمة والأفاتار)
- 1 سبتمبر: ترقية تلقائية للصف التالي + لقب "خريج"
- كل عمليات التصفير/الأرشفة idempotent (لن تتكرر في نفس السنة)
"""

# ═══ مخطط ترقية الصفوف ═══
GRADE_PROGRESSION = {
    "الصف الأول الابتدائي":  "الصف الثاني الابتدائي",
    "الصف الثاني الابتدائي": "الصف الثالث الابتدائي",
    "الصف الثالث الابتدائي": "الصف الرابع الابتدائي",
    "الصف الرابع الابتدائي": "الصف الخامس الابتدائي",
    "الصف الخامس الابتدائي": "الصف السادس الابتدائي",
    "الصف السادس الابتدائي": "الصف السابع",
    "الصف السابع":           "الصف الثامن",
    "الصف الثامن":           "الصف التاسع",
    "الصف التاسع":           "الصف العاشر",
    "الصف العاشر":           "الصف الحادي عشر",
    "الصف الحادي عشر":       "الصف الثاني عشر",
    "الصف الثاني عشر":       "خريج الثانوية",  # نهاية المسار
}


def _get_current_academic_year() -> str:
    """يحسب السنة الدراسية الحالية — من سبتمبر لأغسطس"""
    now = datetime.now(timezone.utc)
    if now.month >= 9:
        # من سبتمبر = بداية سنة جديدة
        return f"{now.year}-{now.year + 1}"
    else:
        # من يناير إلى أغسطس = السنة ما زالت الحالية
        return f"{now.year - 1}-{now.year}"


def _get_system_state(key: str, default: str = "") -> str:
    """قراءة قيمة من system_state"""
    try:
        res = supabase.table("system_state").select("value").eq("key", key).execute()
        if res.data and len(res.data) > 0:
            return res.data[0].get("value", default)
    except Exception:
        pass
    return default


def _set_system_state(key: str, value: str):
    """حفظ قيمة في system_state (upsert)"""
    try:
        supabase.table("system_state").upsert({
            "key": key,
            "value": value,
            "updated_at": datetime.now(timezone.utc).isoformat()
        }).execute()
    except Exception as e:
        print(f"⚠️ _set_system_state error: {e}")


async def _archive_student_year(student: dict, academic_year: str) -> dict:
    """أرشفة سنة واحدة لطالب واحد — يُرجع الـ archive record"""
    student_id   = student["id"]
    student_name = student.get("full_name", "")
    grade        = student.get("grade", "")

    # جلب كل نتائج الطالب
    results = supabase.table("results").select("*").eq("student_id", student_id).execute().data or []

    # حساب الإحصائيات
    total_xp        = sum(r.get("score", 0) for r in results)
    challenges_done = len(results)
    avg_score       = 0
    if challenges_done > 0:
        total_max = sum(r.get("total", 0) for r in results)
        if total_max > 0:
            avg_score = round((total_xp / total_max) * 100, 1)

    # أفضل درس
    lesson_scores = {}
    for r in results:
        l = r.get("lesson", "")
        if l and l not in lesson_scores:
            lesson_scores[l] = 0
        if l:
            lesson_scores[l] += r.get("score", 0)
    best_lesson = max(lesson_scores.items(), key=lambda x: x[1])[0] if lesson_scores else ""

    # حساب الترتيب في الصف
    all_in_grade = supabase.table("students").select("id").eq("grade", grade).execute().data or []
    rank_data = []
    for s in all_in_grade:
        s_results = supabase.table("results").select("score").eq("student_id", s["id"]).execute().data or []
        s_xp = sum(r.get("score", 0) for r in s_results)
        rank_data.append((s["id"], s_xp))
    rank_data.sort(key=lambda x: x[1], reverse=True)
    rank_in_grade = next((i + 1 for i, (sid, _) in enumerate(rank_data) if sid == student_id), 0)

    # اللقب
    graduated_title = f"🎓 خريج السنة {academic_year} — {grade}"

    # إدراج في الأرشيف
    archive_record = {
        "student_id":      student_id,
        "student_name":    student_name,
        "academic_year":   academic_year,
        "grade":           grade,
        "total_xp":        total_xp,
        "challenges_done": challenges_done,
        "avg_score":       avg_score,
        "best_lesson":     best_lesson,
        "rank_in_grade":   rank_in_grade,
        "graduated_title": graduated_title,
        "full_history":    results,  # JSONB
        "badges_earned":   [],       # placeholder — client-side badges
    }

    try:
        supabase.table("results_archive").insert(archive_record).execute()
    except Exception as e:
        print(f"⚠️ archive insert failed for {student_name}: {e}")
        return {}

    return archive_record


async def _run_annual_archive(academic_year: str, dry_run: bool = False) -> dict:
    """
    ينفّذ الأرشفة السنوية لكل الطلاب:
    1. يُؤرشف كل نتائج السنة في results_archive
    2. يحذف النتائج القديمة من results (يبدأون من الصفر)
    3. يُضيف اللقب لـ graduation_titles في students
    4. يحفظ تاريخ الأرشفة لمنع التكرار
    
    dry_run=True: يُحصي بدون تنفيذ فعلي
    """
    # تحقق من عدم التكرار
    last = _get_system_state("last_archive_year")
    if last == academic_year and not dry_run:
        return {"status": "already_archived", "academic_year": academic_year, "students_count": 0}

    # جلب كل الطلاب النشطين
    students = supabase.table("students").select(
        "id, full_name, grade, graduation_titles"
    ).execute().data or []

    archived_count = 0
    failed_count   = 0
    archives       = []

    for student in students:
        try:
            archive = await _archive_student_year(student, academic_year)
            if archive:
                archives.append({
                    "student_id":   student["id"],
                    "student_name": student.get("full_name"),
                    "total_xp":     archive.get("total_xp"),
                    "title":        archive.get("graduated_title"),
                })
                
                if not dry_run:
                    # تحديث ألقاب الطالب + مسح نتائجه
                    old_titles = student.get("graduation_titles") or []
                    if isinstance(old_titles, str):
                        import json as _json
                        try: old_titles = _json.loads(old_titles)
                        except: old_titles = []
                    old_titles.append(archive.get("graduated_title"))
                    
                    supabase.table("students").update({
                        "graduation_titles": old_titles,
                        "last_archived_at":  datetime.now(timezone.utc).isoformat(),
                        "current_academic_year": "",  # مؤقت — حتى سبتمبر تأتي
                    }).eq("id", student["id"]).execute()
                    
                    # حذف نتائج السنة المُنقضية
                    supabase.table("results").delete().eq("student_id", student["id"]).execute()
                
                archived_count += 1
        except Exception as e:
            print(f"❌ archive error for {student.get('full_name')}: {e}")
            failed_count += 1

    if not dry_run:
        _set_system_state("last_archive_year", academic_year)

    return {
        "status":          "success" if not dry_run else "dry_run",
        "academic_year":   academic_year,
        "students_count":  archived_count,
        "failed_count":    failed_count,
        "archives":        archives[:20],  # أول 20 فقط في الرد
    }


async def _run_grade_promotion(target_year: str, dry_run: bool = False) -> dict:
    """
    ترقية كل الطلاب للصف التالي:
    — تُنفّذ في سبتمبر (بداية السنة الجديدة)
    — تُحدّث grade + current_academic_year
    """
    last_promo = _get_system_state("last_promotion_year")
    if last_promo == target_year and not dry_run:
        return {"status": "already_promoted", "academic_year": target_year, "students_count": 0}

    students = supabase.table("students").select("id, full_name, grade").execute().data or []

    promoted   = 0
    graduated  = 0
    promotions = []

    for s in students:
        current_grade = s.get("grade", "").strip()
        next_grade    = GRADE_PROGRESSION.get(current_grade)
        
        if not next_grade:
            continue
        
        if next_grade == "خريج الثانوية":
            graduated += 1
            if not dry_run:
                supabase.table("students").update({
                    "is_active": False,
                    "grade": next_grade,
                    "grade_promoted_at": datetime.now(timezone.utc).isoformat()
                }).eq("id", s["id"]).execute()
        else:
            promoted += 1
            if not dry_run:
                supabase.table("students").update({
                    "grade": next_grade,
                    "current_academic_year": target_year,
                    "grade_promoted_at": datetime.now(timezone.utc).isoformat()
                }).eq("id", s["id"]).execute()
        
        promotions.append({
            "student_id":    s["id"],
            "student_name":  s.get("full_name"),
            "from_grade":    current_grade,
            "to_grade":      next_grade,
        })

    if not dry_run:
        _set_system_state("last_promotion_year", target_year)

    return {
        "status":         "success" if not dry_run else "dry_run",
        "academic_year":  target_year,
        "promoted":       promoted,
        "graduated":      graduated,
        "details":        promotions[:20],
    }


# ═══ Endpoints عامة (متاحة للطلاب والأدمن) ═══

@app.get("/api/academic/current_year")
async def get_current_year():
    """السنة الدراسية الحالية + حالة الأرشفة"""
    now = datetime.now(timezone.utc)
    current_year = _get_current_academic_year()
    last_archive = _get_system_state("last_archive_year")
    
    # هل يجب عرض تنبيه للطالب بنهاية السنة؟
    show_end_warning = now.month == 7  # يوليو = شهر التحذير
    
    return {
        "current_year":      current_year,
        "last_archived":     last_archive,
        "month":             now.month,
        "show_end_warning":  show_end_warning,
        "year_ends_on":      f"{now.year if now.month >= 9 else now.year}-07-31",
    }


@app.get("/api/student/archive/{student_id}")
async def get_student_archive(student_id: int):
    """جلب أرشيف الطالب (كل سنواته السابقة)"""
    res = supabase.table("results_archive").select(
        "academic_year, grade, total_xp, challenges_done, avg_score, "
        "best_lesson, rank_in_grade, graduated_title, archived_at"
    ).eq("student_id", student_id).order("archived_at", desc=True).execute()
    
    return {
        "archives": res.data if res.data else [],
        "count":    len(res.data) if res.data else 0,
    }


@app.get("/api/student/{student_id}/year_summary")
async def get_year_summary(student_id: int):
    """ملخص السنة الحالية للطالب — يُعرض في يوليو قبل الأرشفة"""
    st = supabase.table("students").select("full_name, grade").eq("id", student_id).execute()
    if not st.data:
        raise HTTPException(status_code=404, detail="الطالب غير موجود")
    
    results = supabase.table("results").select("*").eq("student_id", student_id).execute().data or []
    
    total_xp = sum(r.get("score", 0) for r in results)
    challenges = len(results)
    avg_pct = 0
    if challenges > 0:
        total_max = sum(r.get("total", 0) for r in results)
        if total_max > 0:
            avg_pct = round((total_xp / total_max) * 100, 1)
    
    # أفضل 3 دروس
    lesson_scores = {}
    for r in results:
        l = r.get("lesson", "")
        if l:
            lesson_scores[l] = lesson_scores.get(l, 0) + r.get("score", 0)
    top_lessons = sorted(lesson_scores.items(), key=lambda x: x[1], reverse=True)[:3]
    
    return {
        "student":          st.data[0],
        "academic_year":    _get_current_academic_year(),
        "total_xp":         total_xp,
        "challenges_done":  challenges,
        "avg_score":        avg_pct,
        "top_lessons":      [{"lesson": l, "xp": x} for l, x in top_lessons],
    }


# ═══ Endpoints للأدمن فقط (تحكم يدوي) ═══

@app.post("/api/admin/academic/archive_year")
async def trigger_archive(
    request: Request,
    admin=Depends(get_current_admin)
):
    """أرشفة سنوية يدوية (يستخدمها الأدمن أو Cron)"""
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    academic_year = body.get("academic_year", _get_current_academic_year())
    dry_run       = bool(body.get("dry_run", False))
    
    result = await _run_annual_archive(academic_year, dry_run=dry_run)
    return result


@app.post("/api/admin/academic/promote_grades")
async def trigger_promotion(
    request: Request,
    admin=Depends(get_current_admin)
):
    """ترقية جماعية للصف التالي"""
    body = await request.json() if request.headers.get("content-type", "").startswith("application/json") else {}
    target_year = body.get("academic_year", _get_current_academic_year())
    dry_run     = bool(body.get("dry_run", False))
    
    result = await _run_grade_promotion(target_year, dry_run=dry_run)
    return result


@app.get("/api/admin/academic/archives")
async def list_all_archives(admin=Depends(get_current_admin)):
    """عرض كل الأرشيفات السنوية (مجمعة حسب السنة)"""
    res = supabase.table("results_archive").select(
        "id, student_id, student_name, academic_year, grade, total_xp, "
        "challenges_done, avg_score, rank_in_grade, archived_at"
    ).order("archived_at", desc=True).execute()
    
    archives = res.data or []
    
    # تجميع حسب السنة
    by_year = {}
    for a in archives:
        year = a.get("academic_year", "—")
        if year not in by_year:
            by_year[year] = []
        by_year[year].append(a)
    
    return {
        "total_archives": len(archives),
        "by_year":        by_year,
    }


# ═══ Auto-check في كل request (خفيف جداً) ═══
_LAST_AUTO_CHECK = {"date": None}

@app.middleware("http")
async def auto_academic_tasks(request: Request, call_next):
    """
    Middleware يتحقق مرة في اليوم من:
    - 1 أغسطس: تشغيل الأرشفة التلقائية
    - 1 سبتمبر: تشغيل ترقية الصفوف
    """
    try:
        today = datetime.now(timezone.utc).date()
        
        # مرة في اليوم فقط
        if _LAST_AUTO_CHECK["date"] != today:
            _LAST_AUTO_CHECK["date"] = today
            
            # 1 أغسطس → أرشفة
            if today.month == 8 and today.day == 1:
                academic_year = f"{today.year - 1}-{today.year}"
                last = _get_system_state("last_archive_year")
                if last != academic_year:
                    print(f"🎓 [AUTO] بدء الأرشفة السنوية: {academic_year}")
                    try:
                        result = await _run_annual_archive(academic_year, dry_run=False)
                        print(f"✅ [AUTO] أرشفة: {result}")
                    except Exception as e:
                        print(f"❌ [AUTO] خطأ في الأرشفة: {e}")
            
            # 1 سبتمبر → ترقية
            if today.month == 9 and today.day == 1:
                academic_year = f"{today.year}-{today.year + 1}"
                last_promo = _get_system_state("last_promotion_year")
                if last_promo != academic_year:
                    print(f"🎓 [AUTO] بدء ترقية الصفوف: {academic_year}")
                    try:
                        result = await _run_grade_promotion(academic_year, dry_run=False)
                        print(f"✅ [AUTO] ترقية: {result}")
                    except Exception as e:
                        print(f"❌ [AUTO] خطأ في الترقية: {e}")
    except Exception as e:
        print(f"⚠️ auto_academic_tasks error (non-fatal): {e}")
    
    response = await call_next(request)
    return response


# ==========================================
# --- 13. تشغيل المحرك المركزي ---
# ==========================================


# ==========================================
# --- 12.6 📘 الدروس التفاعلية (HTML Lessons) ---
# ==========================================

# CDNs المسموح بها داخل ملفات HTML
ALLOWED_CDN_HOSTS = {
    "fonts.googleapis.com", "fonts.gstatic.com",
    "cdnjs.cloudflare.com", "cdn.jsdelivr.net",
    "unpkg.com", "cdn.tailwindcss.com",
}

DANGEROUS_HTML_PATTERNS = [
    r'fetch\s*\(', r'XMLHttpRequest', r'navigator\.sendBeacon',
    r'WebSocket\s*\(', r'EventSource\s*\(',
    r'localStorage', r'sessionStorage', r'document\.cookie', r'indexedDB',
    r'\beval\s*\(', r'new\s+Function\s*\(',
    r'setTimeout\s*\(\s*[\'"`]', r'setInterval\s*\(\s*[\'"`]',
    r'<iframe', r'<frame', r'<embed', r'<object',
    r'import\s*\([\'"]', r'importScripts\s*\(',
    r'<meta\s+[^>]*http-equiv\s*=\s*[\'"]?refresh',
    r'<form\s+[^>]*action\s*=\s*[\'"]https?://',
    r'navigator\.serviceWorker', r'navigator\.clipboard',
    r'window\.location\s*=', r'location\.href\s*=', r'location\.replace',
    r'window\.open', r'window\.top', r'window\.parent',
]

MAX_HTML_SIZE_MB = 2
MAX_HTML_SIZE_BYTES = MAX_HTML_SIZE_MB * 1024 * 1024


def _scan_html_threats(html_text: str) -> list:
    import re as _re
    warnings = []
    for pattern in DANGEROUS_HTML_PATTERNS:
        try:
            matches = _re.findall(pattern, html_text, flags=_re.IGNORECASE)
            if matches:
                sev = "high" if any(x in pattern for x in ["eval", "Function", "localStorage", "cookie", "fetch"]) else "medium"
                warnings.append({"pattern": pattern, "count": len(matches), "severity": sev})
        except Exception:
            pass
    return warnings


def _sanitize_html_lesson(html_text: str):
    """ينظّف HTML — يُرجع (cleaned, report)"""
    import re as _re
    original_size = len(html_text)
    removed = {"iframes": 0, "bad_scripts": 0, "meta_refresh": 0}
    cleaned = html_text

    # 1. حذف iframe/frame/embed/object
    for tag in ["iframe", "frame", "embed", "object"]:
        pattern = rf"<{tag}\b[^>]*>.*?</{tag}>"
        count = len(_re.findall(pattern, cleaned, flags=_re.IGNORECASE | _re.DOTALL))
        cleaned = _re.sub(pattern, "", cleaned, flags=_re.IGNORECASE | _re.DOTALL)
        self_close = rf"<{tag}\b[^>]*/?>"
        count += len(_re.findall(self_close, cleaned, flags=_re.IGNORECASE))
        cleaned = _re.sub(self_close, "", cleaned, flags=_re.IGNORECASE)
        if tag == "iframe":
            removed["iframes"] = count

    # 2. meta refresh
    meta_pattern = r'<meta\s+[^>]*http-equiv\s*=\s*[\'"]?refresh[^>]*>'
    count = len(_re.findall(meta_pattern, cleaned, flags=_re.IGNORECASE))
    cleaned = _re.sub(meta_pattern, "", cleaned, flags=_re.IGNORECASE)
    removed["meta_refresh"] = count

    # 3. scripts خارجية من نطاقات غير مسموحة
    def _check_script(match):
        src = match.group(1).lower()
        if src.startswith("//"): src = "https:" + src
        if not src.startswith(("http://", "https://")):
            return match.group(0)
        try:
            from urllib.parse import urlparse
            host = urlparse(src).netloc.lower()
            if host.startswith("www."): host = host[4:]
            allowed = any(host == a or host.endswith("." + a) for a in ALLOWED_CDN_HOSTS)
            if not allowed:
                removed["bad_scripts"] += 1
                return ""
        except Exception:
            removed["bad_scripts"] += 1
            return ""
        return match.group(0)

    script_src_pattern = r'<script\b[^>]*\bsrc\s*=\s*[\'"]([^\'"]+)[\'"][^>]*>\s*</script>'
    cleaned = _re.sub(script_src_pattern, _check_script, cleaned, flags=_re.IGNORECASE)

    report = {
        "original_size": original_size,
        "cleaned_size": len(cleaned),
        "removed": removed,
        "warnings": _scan_html_threats(cleaned),
    }
    return cleaned, report


def _wrap_html_lesson_safe(html_text: str) -> str:
    """يضيف CSP + secure links"""
    import re as _re
    csp = (
        "default-src 'self' data: blob:; "
        "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com https://cdnjs.cloudflare.com https://cdn.jsdelivr.net https://unpkg.com; "
        "font-src 'self' https://fonts.gstatic.com https://cdnjs.cloudflare.com data:; "
        "script-src 'self' 'unsafe-inline' https://cdnjs.cloudflare.com https://cdn.jsdelivr.net https://unpkg.com https://cdn.tailwindcss.com; "
        "img-src 'self' data: https:; "
        "connect-src 'none'; "
        "frame-ancestors 'none'; "
        "base-uri 'none'; "
        "form-action 'none';"
    )
    csp_tag = f'<meta http-equiv="Content-Security-Policy" content="{csp}">'
    if "<head>" in html_text.lower():
        result = _re.sub(r"(<head[^>]*>)", r"\\1\n" + csp_tag, html_text, count=1, flags=_re.IGNORECASE)
    elif "<html" in html_text.lower():
        result = _re.sub(r"(<html[^>]*>)", r"\\1\n<head>\n" + csp_tag + "\n</head>", html_text, count=1, flags=_re.IGNORECASE)
    else:
        result = csp_tag + "\n" + html_text
    return result


@app.post("/api/admin/html_lessons")
async def upload_html_lesson(
    title: str        = Form(...),
    grade: str        = Form(...),
    semester: str     = Form(default=""),
    unit: str         = Form(default=""),
    lesson: str       = Form(...),
    description: str  = Form(default=""),
    file: UploadFile  = File(...),
    admin = Depends(get_current_admin)
):
    """رفع درس HTML تفاعلي — مع تنظيف أمني"""
    content_bytes = await file.read()
    if not content_bytes:
        raise HTTPException(status_code=400, detail="الملف فارغ")
    if len(content_bytes) > MAX_HTML_SIZE_BYTES:
        size_mb = len(content_bytes) / (1024 * 1024)
        raise HTTPException(status_code=413, detail=f"الملف كبير جداً ({size_mb:.1f} MB). الحد الأقصى {MAX_HTML_SIZE_MB} MB")
    filename = (file.filename or "").strip()
    if not filename:
        raise HTTPException(status_code=400, detail="اسم الملف مفقود")
    ext = os.path.splitext(filename)[1].lower()
    if ext not in {".html", ".htm"}:
        raise HTTPException(status_code=415, detail="الامتداد يجب .html أو .htm")
    if ".." in filename or "/" in filename or "\\" in filename:
        raise HTTPException(status_code=400, detail="اسم الملف يحتوي رموز ممنوعة")

    try:
        html_text = content_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            html_text = content_bytes.decode("windows-1256")
        except Exception:
            raise HTTPException(status_code=400, detail="الملف ليس UTF-8")

    cleaned, report = _sanitize_html_lesson(html_text)
    final_html = _wrap_html_lesson_safe(cleaned)

    file_name = f"html_lesson_{uuid.uuid4().hex}.html"
    try:
        supabase.storage.from_("resources").upload(
            path=file_name,
            file=final_html.encode("utf-8"),
            file_options={"content-type": "text/html; charset=utf-8"}
        )
        file_url = supabase.storage.from_("resources").get_public_url(file_name)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل رفع الملف: {str(e)[:200]}")

    row = {
        "title": title[:200],
        "grade": grade[:100],
        "semester": (semester or "")[:100],
        "unit": (unit or "")[:200],
        "lesson": lesson[:300],
        "description": (description or "")[:500],
        "file_url": file_url,
        "file_size_kb": len(final_html.encode("utf-8")) // 1024,
        "sanitized": True,
    }
    try:
        supabase.table("html_lessons").insert(row).execute()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"فشل الحفظ: {str(e)[:200]}")

    return {
        "status": "success",
        "file_url": file_url,
        "sanitization": report,
        "message": f"حُذف: {report['removed']['iframes']} iframe، {report['removed']['bad_scripts']} script خارجي"
    }



@app.get("/api/admin/html_lessons/diagnose")
async def diagnose_html_lessons(admin = Depends(get_current_admin)):
    """🩺 endpoint تشخيصي — يعرض حالة جدول html_lessons"""
    try:
        res = supabase.table("html_lessons").select(
            "id, title, grade, semester, unit, lesson, file_url"
        ).execute()
        rows = res.data or []
        
        # عيّنة من القيم الفعلية
        unique_grades = list(set(r.get("grade", "") for r in rows if r.get("grade")))
        unique_lessons = list(set(r.get("lesson", "") for r in rows if r.get("lesson")))[:20]
        
        # فحص الروابط
        broken_files = []
        for r in rows:
            if not r.get("file_url"):
                broken_files.append({"id": r["id"], "title": r.get("title"), "issue": "no file_url"})
        
        return {
            "total_lessons": len(rows),
            "unique_grades_in_db": unique_grades,
            "sample_lessons": unique_lessons,
            "broken_files": broken_files,
            "all_lessons": rows[:30],
        }
    except Exception as e:
        return {"error": str(e)[:300]}

@app.get("/api/html_lessons")
async def list_html_lessons(grade: str = "", lesson: str = ""):
    """قائمة الدروس التفاعلية — مع فلتر مرن للصف"""
    # نجلب كل الدروس
    query = supabase.table("html_lessons").select(
        "id, title, grade, semester, unit, lesson, description, file_url, file_size_kb, uploaded_at, view_count"
    )
    if lesson:
        query = query.eq("lesson", lesson)
    res = query.order("uploaded_at", desc=True).execute()
    all_data = res.data or []
    
    # فلترة مرنة على الصف (بـ Python لمعالجة اختلاف الصياغة)
    if grade:
        def _norm(s):
            if not s: return ""
            return str(s).strip().lower().replace("الصف", "").replace("الصّف", "").replace(" ", "")
        norm_grade = _norm(grade)
        filtered = []
        for row in all_data:
            row_grade = _norm(row.get("grade", ""))
            if row_grade == norm_grade or norm_grade in row_grade or row_grade in norm_grade:
                filtered.append(row)
        return filtered
    
    return all_data




# ═══════════════════════════════════════════════════════════════
# 📚 STUDENT LESSONS — دروس الطالب (تفاعلية + مخططة من المنهج)
# ═══════════════════════════════════════════════════════════════

@app.get("/api/student/lessons_overview")
async def student_lessons_overview(grade: str = ""):
    """
    📚 يُرجع نظرة شاملة لطالب صف معين:
    - الدروس التفاعلية المتوفرة (HTML uploaded)
    - الدروس المخططة من المنهج (لم تُرفع بعد)
    
    يُساعد الطالب على رؤية كل دروس صفه + معرفة المتوفر والقادم
    """
    if not grade.strip():
        raise HTTPException(status_code=400, detail="grade مطلوب")
    
    # دالة تنظيف
    def _norm(s):
        if not s: return ""
        return str(s).strip().lower().replace("الصف", "").replace("الصّف", "").replace(" ", "")
    
    norm_grade = _norm(grade)
    
    result = {
        "grade": grade,
        "interactive_lessons": [],   # دروس HTML مرفوعة
        "planned_lessons": [],        # دروس من المنهج (لم تُرفع)
        "stats": {
            "total_curriculum_lessons": 0,
            "uploaded_count": 0,
            "planned_count": 0,
        }
    }
    
    # 1️⃣ جلب الدروس التفاعلية المرفوعة
    try:
        all_html = supabase.table("html_lessons").select(
            "id, title, grade, semester, unit, lesson, description, file_url, uploaded_at, view_count"
        ).order("uploaded_at", desc=True).execute()
        
        for r in (all_html.data or []):
            row_grade = _norm(r.get("grade", ""))
            if row_grade == norm_grade or norm_grade in row_grade or row_grade in norm_grade:
                result["interactive_lessons"].append(r)
        
        result["stats"]["uploaded_count"] = len(result["interactive_lessons"])
    except Exception as e:
        print(f"[lessons_overview] html_lessons error: {e}")
    
    # 2️⃣ جلب المنهج المخطط لصف الطالب
    try:
        # نُحاول جلب hierarchy: grades → semesters → units → lessons
        struct = supabase.table("grades").select(
            "id, name, semesters(id, name, units(id, name, lessons(id, name)))"
        ).execute()
        grades_data = struct.data or []
        
        # ابحث عن صف الطالب
        student_grade_node = None
        for g in grades_data:
            if _norm(g.get("name", "")) == norm_grade:
                student_grade_node = g
                break
            # محاولة partial match
            g_norm = _norm(g.get("name", ""))
            if g_norm and (g_norm in norm_grade or norm_grade in g_norm):
                student_grade_node = g
                break
        
        if not student_grade_node:
            print(f"[lessons_overview] لم يُعثر على grade node لـ '{grade}'")
            return result
        
        # نبني قائمة بكل دروس المنهج لهذا الصف
        all_curriculum_lessons = []
        for s in (student_grade_node.get("semesters") or []):
            sem_name = s.get("name", "")
            for u in (s.get("units") or []):
                unit_name = u.get("name", "")
                for L in (u.get("lessons") or []):
                    lesson_name = L.get("name", "")
                    if not lesson_name: continue
                    all_curriculum_lessons.append({
                        "id": L.get("id"),
                        "name": lesson_name,
                        "unit": unit_name,
                        "semester": sem_name,
                    })
        
        result["stats"]["total_curriculum_lessons"] = len(all_curriculum_lessons)
        
        # 3️⃣ نحدّد أي درس من المنهج "تم رفعه" بالفعل
        uploaded_lesson_names = set()
        for L in result["interactive_lessons"]:
            ln = (L.get("lesson") or "").strip().lower()
            if ln:
                uploaded_lesson_names.add(ln)
        
        # 4️⃣ الدروس المخططة = دروس المنهج التي لم تُرفع
        for L in all_curriculum_lessons:
            ln = L["name"].strip().lower()
            if ln not in uploaded_lesson_names:
                # نتحقق من partial match أيضاً
                is_uploaded = False
                for uploaded_name in uploaded_lesson_names:
                    if uploaded_name in ln or ln in uploaded_name:
                        is_uploaded = True
                        break
                if not is_uploaded:
                    result["planned_lessons"].append(L)
        
        result["stats"]["planned_count"] = len(result["planned_lessons"])
        
    except Exception as e:
        print(f"[lessons_overview] curriculum error: {e}")
        import traceback
        traceback.print_exc()
    
    return result


@app.get("/api/html_lessons/{lesson_id}/render")
async def render_html_lesson(lesson_id: int, request: Request):
    """
    🎯 Proxy لتقديم محتوى HTML lesson مع headers صحيحة لعرضه في iframe.
    
    لماذا؟ Supabase Storage يُعيد الملفات مع X-Frame-Options: DENY
    → لا يمكن عرضها في iframe مباشرةً.
    → نجلبها من الخادم ونُعيدها مع CSP خفيف يسمح بـ iframe embedding.
    """
    try:
        # جلب الـ file_url من DB
        res = supabase.table("html_lessons").select("file_url, title").eq("id", lesson_id).limit(1).execute()
        if not res.data:
            raise HTTPException(status_code=404, detail="الدرس غير موجود")
        
        file_url = res.data[0].get("file_url", "")
        if not file_url:
            raise HTTPException(status_code=404, detail="الملف غير موجود")
        
        # جلب المحتوى من Supabase Storage
        import httpx
        async with httpx.AsyncClient(timeout=15.0) as client:
            r = await client.get(file_url)
            if r.status_code != 200:
                raise HTTPException(status_code=502, detail="فشل تحميل المحتوى")
            html_content = r.text
        
        # تنظيف بسيط — منع inline events دون كسر الدرس
        # نُضيف base target للروابط لتفتح في تبويب جديد
        if "<head>" in html_content.lower():
            base_inject = '<base target="_blank">'
            # ابحث عن <head> case-insensitive
            import re
            html_content = re.sub(
                r'(<head[^>]*>)', 
                r'\1' + base_inject, 
                html_content, 
                count=1, 
                flags=re.IGNORECASE
            )
        
        # أعد المحتوى مع headers تسمح بـ iframe
        from fastapi.responses import HTMLResponse
        response = HTMLResponse(content=html_content)
        response.headers["X-Frame-Options"] = "SAMEORIGIN"
        response.headers["Content-Security-Policy"] = "default-src * data: blob: 'unsafe-inline' 'unsafe-eval'; img-src * data: blob:; style-src * 'unsafe-inline'; script-src * 'unsafe-inline' 'unsafe-eval'; frame-ancestors *"
        # عداد المشاهدات
        try:
            view_res = supabase.table("html_lessons").select("view_count").eq("id", lesson_id).limit(1).execute()
            current = view_res.data[0].get("view_count", 0) if view_res.data else 0
            supabase.table("html_lessons").update({"view_count": current + 1}).eq("id", lesson_id).execute()
        except Exception:
            pass
        
        return response
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/html_lessons/{lesson_id}/view")
async def increment_view_count(lesson_id: int):
    """زيادة عداد المشاهدات"""
    try:
        res = supabase.table("html_lessons").select("view_count").eq("id", lesson_id).execute()
        if res.data:
            current = res.data[0].get("view_count", 0) or 0
            supabase.table("html_lessons").update({"view_count": current + 1}).eq("id", lesson_id).execute()
        return {"status": "ok"}
    except Exception:
        return {"status": "error"}


@app.delete("/api/admin/html_lessons/{lesson_id}")
async def delete_html_lesson(lesson_id: int, admin = Depends(get_current_admin)):
    """حذف درس تفاعلي"""
    res = supabase.table("html_lessons").select("file_url").eq("id", lesson_id).execute()
    if res.data:
        file_url = res.data[0].get("file_url", "")
        file_name = file_url.rsplit("/", 1)[-1] if "/" in file_url else ""
        if file_name and file_name.startswith("html_lesson_"):
            try:
                supabase.storage.from_("resources").remove([file_name])
            except Exception:
                pass
    supabase.table("html_lessons").delete().eq("id", lesson_id).execute()
    return {"status": "deleted"}




# ═══════════════════════════════════════════════════════════════
# 📡 PUSH NOTIFICATIONS — تسجيل + إرسال
# ═══════════════════════════════════════════════════════════════
@app.get("/api/push/vapid_public_key")
async def get_vapid_public_key():
    """يُرجع المفتاح العام للعميل ليُسجّل اشتراك push"""
    return {"key": VAPID_PUBLIC_KEY, "enabled": PUSH_ENABLED}


@app.post("/api/push/subscribe")
async def push_subscribe(
    student_id: int = Form(...),
    endpoint: str   = Form(...),
    p256dh: str     = Form(...),
    auth: str       = Form(...),
    user_agent: str = Form(default=""),
):
    """تسجيل اشتراك push من جهاز الطالب"""
    if not endpoint or not p256dh or not auth:
        raise HTTPException(status_code=400, detail="بيانات الاشتراك ناقصة")
    try:
        # حذف اشتراك سابق بنفس endpoint إن وُجد
        supabase.table("push_subscriptions").delete().eq("endpoint", endpoint).execute()
        # إدراج الجديد
        supabase.table("push_subscriptions").insert({
            "student_id": student_id,
            "endpoint":   endpoint,
            "p256dh":     p256dh,
            "auth":       auth,
            "user_agent": user_agent[:500],
        }).execute()
        return {"status": "ok"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"خطأ: {str(e)[:200]}")


@app.post("/api/push/unsubscribe")
async def push_unsubscribe(endpoint: str = Form(...)):
    """إلغاء اشتراك push"""
    try:
        supabase.table("push_subscriptions").delete().eq("endpoint", endpoint).execute()
    except Exception:
        pass
    return {"status": "ok"}


@app.post("/api/admin/push/send_to_student/{student_id}")
async def admin_push_to_student(
    student_id: int,
    title: str = Form(...),
    body: str  = Form(...),
    url: str   = Form(default="/student"),
    admin = Depends(get_current_admin),
):
    """إرسال push يدوي لطالب"""
    sent = _push_to_student(student_id, title, body, url=url, tag="admin_push")
    return {"sent": sent, "enabled": PUSH_ENABLED}


# ═══════════════════════════════════════════════════════════════
# 📓 ADMIN TASKS — دفتر الأعمال
# ═══════════════════════════════════════════════════════════════
@app.get("/api/admin/tasks")
async def list_tasks(status: str = "", admin = Depends(get_current_admin)):
    """قائمة المهام (يمكن فلترتها بالحالة)"""
    try:
        q = supabase.table("admin_tasks").select("*")
        if status:
            q = q.eq("status", status)
        res = q.order("priority").order("due_date", nullsfirst=False).order("created_at", desc=True).execute()
        return res.data or []
    except Exception as e:
        return {"error": str(e)[:200]}


@app.post("/api/admin/tasks")
async def create_task(
    title: str       = Form(...),
    description: str = Form(default=""),
    priority: str    = Form(default="normal"),
    category: str    = Form(default="general"),
    due_date: str    = Form(default=""),
    admin = Depends(get_current_admin)
):
    """إضافة مهمة جديدة"""
    if priority not in ("low", "normal", "high", "urgent"):
        priority = "normal"
    row = {
        "title": title.strip()[:200],
        "description": description.strip()[:2000],
        "priority": priority,
        "category": category[:50],
        "status": "pending",
    }
    if due_date:
        row["due_date"] = due_date
    try:
        res = supabase.table("admin_tasks").insert(row).execute()
        return {"status": "created", "id": res.data[0]["id"] if res.data else None}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.put("/api/admin/tasks/{task_id}")
async def update_task(
    task_id: int,
    title: str       = Form(default=None),
    description: str = Form(default=None),
    priority: str    = Form(default=None),
    status: str      = Form(default=None),
    due_date: str    = Form(default=None),
    admin = Depends(get_current_admin)
):
    """تحديث مهمة"""
    update = {"updated_at": datetime.now(timezone.utc).isoformat()}
    if title is not None:       update["title"]       = title.strip()[:200]
    if description is not None: update["description"] = description.strip()[:2000]
    if priority is not None:    update["priority"]    = priority
    if status is not None:
        update["status"] = status
        if status == "done":
            update["completed_at"] = datetime.now(timezone.utc).isoformat()
    if due_date is not None:    update["due_date"]    = due_date or None
    try:
        supabase.table("admin_tasks").update(update).eq("id", task_id).execute()
        return {"status": "updated"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.delete("/api/admin/tasks/{task_id}")
async def delete_task(task_id: int, admin = Depends(get_current_admin)):
    try:
        supabase.table("admin_tasks").delete().eq("id", task_id).execute()
        return {"status": "deleted"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


# ═══════════════════════════════════════════════════════════════
# 🎯 MOTIVATIONAL NOTIFICATIONS — تحفيز الطلاب تلقائياً
# ═══════════════════════════════════════════════════════════════
def _was_motivation_sent_today(student_id: int, notif_type: str) -> bool:
    """فحص هل أُرسل هذا النوع لهذا الطالب اليوم"""
    try:
        from datetime import date
        today = date.today().isoformat()
        res = supabase.table("motivation_log").select("id").eq(
            "student_id", student_id
        ).eq("notif_type", notif_type).eq("notif_date", today).limit(1).execute()
        return bool(res.data)
    except Exception:
        return False


def _log_motivation(student_id: int, notif_type: str):
    """سجّل أن الإشعار أُرسل"""
    try:
        supabase.table("motivation_log").insert({
            "student_id": student_id,
            "notif_type": notif_type,
        }).execute()
    except Exception:
        pass


# قوالب الرسائل التحفيزية (يمكن للأدمن تخصيصها لاحقاً)
MOTIVATION_TEMPLATES = {
    "inactive_3days": {
        "titles": [
            "🏰 إمبراطوريتك تشتاق إليك!",
            "⚔️ أبطالك ينتظرون عودتك",
            "👑 العرش يحتاج بطله",
        ],
        "bodies": [
            "غبت عن المنصة 3 أيام — تعال خض تحدياً جديداً وارفع نقاطك! ⚡",
            "لا تترك أصدقاءك يسبقونك في الترتيب — عد للمعركة! 🏆",
            "تحدٍ جديد بانتظارك في ساحة المبارزة — اضغط وابدأ! 🎯",
        ],
    },
    "inactive_7days": {
        "titles": [
            "🚨 أسبوع كامل بدون تحديات!",
            "💔 افتقدناك في إمبراطوريتنا",
        ],
        "bodies": [
            "مرّ أسبوع — استعد عرشك بحلّ تحدٍ سريع الآن! 5 دقائق فقط ⏱️",
            "زملاؤك حصلوا على 200+ XP هذا الأسبوع — لا تتأخر عنهم! 🎖️",
        ],
    },
    "streak_break": {
        "titles": ["🔥 لا تكسر سلسلة إنجازاتك!"],
        "bodies": ["كنت في طريقك لرقم قياسي — حلّ تحدٍ واحد فقط لتحافظ على السلسلة! 💪"],
    },
    "morning_motivation": {
        "titles": [
            "☀️ صباح المبارزات!",
            "🌅 يوم جديد لتحديات جديدة",
        ],
        "bodies": [
            "ابدأ يومك بحلّ تحدٍ سريع — 10 دقائق تعطيك طاقة لليوم كله! ⚡",
            "أبطال اليوم يبدؤون باكراً — كن منهم! 🏆",
        ],
    },
    "evening_reminder": {
        "titles": ["🌙 لم تتحدّ اليوم بعد!"],
        "bodies": ["لا تنهي يومك بدون تحدٍ واحد على الأقل — اكسب نقاطك! ⭐"],
    },
}


import random as _random_mod


@app.post("/api/admin/motivation/send_inactive")
async def send_motivation_inactive(
    days: int = Form(default=3),
    admin = Depends(get_current_admin)
):
    """
    إرسال إشعار تحفيزي للطلاب غير النشطين منذ N أيام
    يفحص آخر heartbeat ويُرسل push للطلاب الغائبين
    """
    from datetime import datetime, timezone, timedelta
    cutoff = datetime.now(timezone.utc) - timedelta(days=days)
    
    notif_type = f"inactive_{days}days" if days in (3, 7) else "inactive_3days"
    
    # اجلب كل الطلاب
    try:
        students_res = supabase.table("students").select(
            "id, full_name"
        ).execute()
        students = students_res.data or []
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])
    
    template = MOTIVATION_TEMPLATES.get(notif_type, MOTIVATION_TEMPLATES["inactive_3days"])
    sent_count = 0
    skipped_count = 0
    
    for st in students:
        sid = st["id"]
        last_active = st.get("last_active")
        
        # تخطّى الطلاب النشطين
        if last_active:
            try:
                la_dt = datetime.fromisoformat(last_active.replace("Z", "+00:00"))
                if la_dt.tzinfo is None:
                    la_dt = la_dt.replace(tzinfo=timezone.utc)
                if la_dt > cutoff:
                    continue  # نشط — تخطَّ
            except Exception:
                pass
        
        # تخطّى لو أُرسل اليوم
        if _was_motivation_sent_today(sid, notif_type):
            skipped_count += 1
            continue
        
        # اختر رسالة عشوائية
        title = _random_mod.choice(template["titles"])
        body  = _random_mod.choice(template["bodies"])
        
        sent = _push_to_student(sid, title, body, url="/student", tag=notif_type)
        if sent > 0:
            _log_motivation(sid, notif_type)
            sent_count += 1
    
    return {
        "sent": sent_count,
        "skipped_today": skipped_count,
        "total_inactive_students": len(students),
        "push_enabled": PUSH_ENABLED,
    }


@app.post("/api/admin/motivation/send_custom")
async def send_motivation_custom(
    title: str        = Form(...),
    body: str         = Form(...),
    target: str       = Form(default="all"),  # all / grade:X / student:N
    url: str          = Form(default="/student"),
    admin = Depends(get_current_admin)
):
    """إرسال رسالة تحفيزية مخصصة"""
    if not title or not body:
        raise HTTPException(status_code=400, detail="title و body مطلوبان")
    
    sent = 0
    if target.startswith("student:"):
        try:
            sid = int(target.split(":")[1])
            sent = _push_to_student(sid, title, body, url=url, tag="custom_motivation")
        except Exception:
            pass
    elif target.startswith("grade:"):
        grade = target.split(":", 1)[1]
        sent = _push_to_grade(grade, title, body, url=url, tag="custom_motivation")
    else:
        # all — لكل الطلاب
        try:
            res = supabase.table("students").select("id").execute()
            for s in (res.data or []):
                sent += _push_to_student(s["id"], title, body, url=url, tag="custom_motivation")
        except Exception:
            pass
    
    return {"sent": sent, "push_enabled": PUSH_ENABLED}


@app.get("/api/admin/motivation/templates")
async def get_motivation_templates(admin = Depends(get_current_admin)):
    """قوالب الرسائل التحفيزية الجاهزة"""
    return MOTIVATION_TEMPLATES



if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)